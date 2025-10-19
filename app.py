"""Streamlit dashboard for くらしいきいき社の計数管理アプリ."""
from __future__ import annotations

# TODO: Streamlit UIコンポーネントを使ってダッシュボードを構築
import base64
import calendar
import html
import hashlib
import importlib
import io
import json
import logging
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
from contextlib import contextmanager
from datetime import date, datetime, timedelta
from pathlib import Path
import traceback
from typing import Any, Callable, Dict, List, Optional, Sequence, Tuple, TypeVar, Union
from urllib.parse import parse_qsl, quote

import numpy as np
import pandas as pd
import altair as alt
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
import streamlit.components.v1 as components
from streamlit_plotly_events import plotly_events


logger = logging.getLogger(__name__)


REPORTLAB_AVAILABLE = True
try:  # pragma: no cover - optional dependency check
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
except ImportError:  # pragma: no cover - gracefully handle missing dependency
    REPORTLAB_AVAILABLE = False


PPTX_AVAILABLE = True
try:  # pragma: no cover - optional dependency check
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - gracefully handle missing dependency
    PPTX_AVAILABLE = False


DATA_PROCESSING_IMPORT_ERROR: Optional[BaseException] = None
try:
    from data_processing import (
        DEFAULT_FIXED_COST,
        annotate_customer_segments,
        build_alerts,
        calculate_kpis,
        create_current_pl,
        create_default_cashflow_plan,
        fetch_sales_from_endpoint,
        forecast_cashflow,
        train_forecast_model,
        predict_sales_forecast,
        apply_forecast_to_cashflow,
        estimate_forecast_savings,
        generate_sample_cost_data,
        generate_sample_sales_data,
        generate_sample_subscription_data,
        detect_duplicate_rows,
        validate_channel_fees,
        ValidationReport,
        load_cost_workbook,
        load_sales_files,
        load_sales_workbook,
        load_subscription_workbook,
        merge_sales_and_costs,
        monthly_sales_summary,
        simulate_pl,
        compute_channel_share,
        compute_category_share,
        compute_kpi_breakdown,
    )
except ImportError as import_error:
    DATA_PROCESSING_IMPORT_ERROR = import_error


_PAGE_CONFIG = {
    "page_title": "経営ダッシュボード",
    "page_icon": ":bar_chart:",
    "layout": "wide",
    "initial_sidebar_state": "expanded",
}

_THEME_CONFIG = {
    "primaryColor": "#2563EB",
    "backgroundColor": "#F8FAFC",
    "secondaryBackgroundColor": "#FFFFFF",
    "textColor": "#0F172A",
    "font": "Inter",
}


def _apply_legacy_theme_css(theme: Dict[str, str]) -> None:
    """Inject CSS to emulate the desired theme on Streamlit versions without native support."""

    primary = theme.get("primaryColor", "#2563EB")
    background = theme.get("backgroundColor", "#F8FAFC")
    secondary = theme.get("secondaryBackgroundColor", "#FFFFFF")
    text_color = theme.get("textColor", "#0F172A")
    font = theme.get("font", "Inter")

    st.markdown(
        f"""
        <style>
        :root {{
            --app-primary-color: {primary};
            --app-text-color: {text_color};
        }}

        html, body, [data-testid="stAppViewContainer"] > .main {{
            background-color: {background};
            color: {text_color};
            font-family: '{font}', sans-serif;
        }}

        [data-testid="stSidebar"] > div:first-child {{
            background-color: {secondary};
        }}

        .stButton>button, .stDownloadButton>button {{
            background-color: {primary};
            color: white;
            border-radius: 0.5rem;
            border: none;
        }}

        .stButton>button:hover, .stDownloadButton>button:hover {{
            filter: brightness(0.9);
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )


def _configure_page() -> None:
    """Configure the Streamlit page, supporting older versions without theme parameter."""

    try:
        st.set_page_config(**_PAGE_CONFIG, theme=_THEME_CONFIG)
    except TypeError as error:
        logger.warning(
            "Streamlit theme configuration failed; falling back to CSS override. Error: %s",
            error,
        )
        st.set_page_config(**_PAGE_CONFIG)
        _apply_legacy_theme_css(_THEME_CONFIG)


_configure_page()


def _render_data_processing_import_error(error: BaseException) -> None:
    """Display a user-friendly message when data_processing import fails."""

    logger.exception("Failed to import data_processing module", exc_info=error)
    st.title("初期化エラー")
    st.error("データ処理モジュールの読み込みに失敗しました。依存関係を確認してください。")

    missing_package: Optional[str] = None
    if isinstance(error, ModuleNotFoundError):
        missing_package = getattr(error, "name", None)
    if missing_package:
        st.markdown(f"- 不足しているモジュール: `{missing_package}`")

    st.markdown(
        "- 仮想環境で `pip install -r requirements.txt` を実行し、依存関係を再インストールしてください。"
    )
    st.markdown(
        "- Streamlit Cloud を利用している場合は、アプリの `Manage app` からログを確認し不足モジュールを追加してください。"
    )

    with st.expander("エラーログの詳細"):
        formatted = "".join(
            traceback.format_exception(type(error), error, error.__traceback__)
        )
        st.code(formatted, language="python")


if DATA_PROCESSING_IMPORT_ERROR is not None:
    _render_data_processing_import_error(DATA_PROCESSING_IMPORT_ERROR)
    st.stop()


def trigger_rerun() -> None:
    """Streamlitの再実行を互換性を保ちながら呼び出す。"""

    rerun_callable = getattr(st, "rerun", None)
    if rerun_callable is None:
        rerun_callable = getattr(st, "experimental_rerun", None)
    if rerun_callable is None:
        raise RuntimeError("Streamlit rerun function is unavailable.")
    rerun_callable()


HOTKEY_ACTION_STATE = "hotkey_action"
HOTKEY_REFETCH_DATA = "refetch_data"


def _queue_hotkey_action(action: str) -> None:
    """Store a hotkey action request with a timestamp for later processing."""

    st.session_state[HOTKEY_ACTION_STATE] = {
        "action": action,
        "timestamp": datetime.utcnow().timestamp(),
    }


def _parse_api_params(params_raw: Optional[str]) -> Optional[Dict[str, str]]:
    """Convert a query-string style parameter value into a dictionary."""

    if not params_raw:
        return None
    parsed_pairs = parse_qsl(params_raw, keep_blank_values=False)
    if not parsed_pairs:
        return None
    return {key: value for key, value in parsed_pairs}


def refresh_configured_automations(channel_files: Dict[str, Any]) -> Dict[str, Any]:
    """Fetch latest sales data for channels with configured API endpoints."""

    feedback: Dict[str, Any] = {
        "success": [],
        "warnings": [],
        "errors": [],
        "info": None,
    }

    if not channel_files:
        feedback["info"] = "売上データのチャネルがまだ設定されていません。"
        return feedback

    configured_channels = [
        channel for channel in channel_files.keys() if st.session_state.get(f"api_endpoint_{channel}")
    ]
    if not configured_channels:
        feedback["info"] = "APIエンドポイントが設定されていません。サイドバーで接続先を登録してください。"
        return feedback

    for channel in configured_channels:
        endpoint = st.session_state.get(f"api_endpoint_{channel}")
        token = st.session_state.get(f"api_token_{channel}") or None
        params_raw = st.session_state.get(f"api_params_{channel}")
        params_dict = _parse_api_params(params_raw)

        try:
            with st.spinner(f"{channel}のデータを再取得中..."):
                fetched_df, fetch_report = fetch_sales_from_endpoint(
                    endpoint,
                    token=token,
                    params=params_dict,
                    channel_hint=channel,
                )
        except Exception as exc:  # pragma: no cover - safeguards remote failures
            logger.exception("Quick action refresh failed for %s", channel)
            feedback["errors"].append(f"{channel}: APIリクエストでエラーが発生しました ({exc}).")
            continue

        st.session_state.setdefault("api_sales_data", {})[channel] = fetched_df
        st.session_state.setdefault("api_sales_validation", {})[channel] = fetch_report
        st.session_state.setdefault("api_last_fetched", {})[channel] = datetime.now()

        error_messages = [
            msg.message
            for msg in getattr(fetch_report, "messages", [])
            if getattr(msg, "level", "").lower() == "error"
        ]
        warning_messages = [
            msg.message
            for msg in getattr(fetch_report, "messages", [])
            if getattr(msg, "level", "").lower() == "warning"
        ]

        if error_messages:
            feedback["errors"].append(
                f"{channel}: " + " / ".join(error_messages[:3])
                + (" ..." if len(error_messages) > 3 else "")
            )
            continue

        feedback["success"].append(channel)
        if warning_messages:
            feedback["warnings"].append(
                f"{channel}: " + " / ".join(warning_messages[:3])
                + (" ..." if len(warning_messages) > 3 else "")
            )

    if not any((feedback["success"], feedback["warnings"], feedback["errors"])):
        feedback["info"] = "再取得の対象となるチャネルが見つかりませんでした。"

    return feedback

PERIOD_FREQ_OPTIONS: List[Tuple[str, str]] = [
    ("月次", "M"),
    ("週次", "W-MON"),
    ("四半期", "Q"),
    ("年次", "Y"),
]

PERIOD_YOY_LAG: Dict[str, int] = {
    "M": 12,
    "W-MON": 52,
    "Q": 4,
    "Y": 1,
}


PLAN_WIZARD_STEPS: List[Dict[str, str]] = [
    {
        "title": "基本情報入力",
        "description": "会社名や計画期間を設定し、計画の前提条件を整理します。",
    },
    {
        "title": "売上予測",
        "description": "チャネル別の売上計画をCSV取り込みやテンプレートで作成します。",
    },
    {
        "title": "経費入力",
        "description": "固定費・変動費のテンプレートや自動補完を使ってコスト計画を整えます。",
    },
    {
        "title": "財務指標計算",
        "description": "売上と経費から利益率などの主要指標を自動計算します。",
    },
    {
        "title": "結果確認",
        "description": "入力内容を確認し、計画サマリーを共有用に出力します。",
    },
]


SALES_PLAN_COLUMNS = ["項目", "月次売上", "チャネル"]
EXPENSE_PLAN_COLUMNS = ["費目", "月次金額", "区分"]

COMMON_SALES_ITEMS = [
    "自社サイト売上",
    "楽天市場売上",
    "Amazon売上",
    "Yahoo!ショッピング売上",
    "サブスク売上",
    "卸売売上",
    "定期便アップセル",
    "店頭販売",
]

COMMON_EXPENSE_ITEMS = [
    "人件費",
    "家賃",
    "広告宣伝費",
    "配送費",
    "外注費",
    "システム利用料",
    "水道光熱費",
    "雑費",
]

PLAN_CHANNEL_OPTIONS_BASE = [
    "自社サイト",
    "楽天市場",
    "Amazon",
    "Yahoo!ショッピング",
    "卸売",
    "サブスク",
    "広告流入",
    "その他",
]

PLAN_EXPENSE_CLASSIFICATIONS = ["固定費", "変動費", "投資", "その他"]

CATEGORY_SUGGESTION_LIMIT = 30

SALES_PLAN_TEMPLATES: Dict[str, List[Dict[str, Any]]] = {
    "EC標準チャネル構成": [
        {"項目": "自社サイト売上", "月次売上": 1_200_000, "チャネル": "自社サイト"},
        {"項目": "楽天市場売上", "月次売上": 950_000, "チャネル": "楽天市場"},
        {"項目": "Amazon売上", "月次売上": 780_000, "チャネル": "Amazon"},
        {"項目": "Yahoo!ショッピング売上", "月次売上": 320_000, "チャネル": "Yahoo!ショッピング"},
    ],
    "サブスク強化モデル": [
        {"項目": "サブスク売上", "月次売上": 850_000, "チャネル": "サブスク"},
        {"項目": "定期便アップセル", "月次売上": 420_000, "チャネル": "サブスク"},
        {"項目": "新規顧客向け単品", "月次売上": 380_000, "チャネル": "広告流入"},
    ],
}

EXPENSE_PLAN_TEMPLATES: Dict[str, List[Dict[str, Any]]] = {
    "スリム型コスト構成": [
        {"費目": "人件費", "月次金額": 600_000, "区分": "固定費"},
        {"費目": "家賃", "月次金額": 200_000, "区分": "固定費"},
        {"費目": "広告宣伝費", "月次金額": 180_000, "区分": "変動費"},
        {"費目": "システム利用料", "月次金額": 90_000, "区分": "固定費"},
    ],
    "成長投資モデル": [
        {"費目": "人件費", "月次金額": 850_000, "区分": "固定費"},
        {"費目": "広告宣伝費", "月次金額": 320_000, "区分": "変動費"},
        {"費目": "外注費", "月次金額": 160_000, "区分": "変動費"},
        {"費目": "研究開発費", "月次金額": 120_000, "区分": "投資"},
    ],
}

DEFAULT_STORE_OPTIONS = ["全社", "本店", "那覇本店", "浦添物流センター", "EC本部"]

T = TypeVar("T")

FILTER_STATE_KEYS = {
    "store": "filter_store",
    "channels": "filter_channels",
    "categories": "filter_categories",
    "period": "filter_period",
    "freq": "filter_frequency",
    "signature": "filter_signature",
}

UPLOAD_STORAGE_ROOT = Path("uploaded_data")
MAX_UPLOAD_HISTORY = 50

MANUAL_INPUT_FORMATS: Dict[str, Dict[str, Any]] = {
    "manual_active_customers": {"is_integer": True, "text_format": ",d"},
    "manual_new_customers": {"is_integer": True, "text_format": ",d"},
    "manual_repeat_customers": {"is_integer": True, "text_format": ",d"},
    "manual_cancel_customers": {"is_integer": True, "text_format": ",d"},
    "manual_previous_customers": {"is_integer": True, "text_format": ",d"},
    "manual_marketing_cost": {"is_integer": False, "text_format": ",.0f"},
    "manual_ltv_value": {"is_integer": False, "text_format": ",.0f"},
    "manual_inventory_days": {"is_integer": False, "text_format": ",.0f"},
    "manual_stockout_pct": {"is_integer": False, "text_format": ",.1f"},
    "manual_training_sessions": {"is_integer": True, "text_format": ",d"},
    "manual_new_products": {"is_integer": True, "text_format": ",d"},
}


def widget_key_for(state_key: str) -> str:
    """アプリ状態に対応するウィジェットkeyを生成する。"""

    return f"_{state_key}_widget"


VALID_ICON_TONES = {"accent", "muted", "success", "warning", "error"}


NUMERIC_TEXT_TRANSLATION = str.maketrans(
    {
        "０": "0",
        "１": "1",
        "２": "2",
        "３": "3",
        "４": "4",
        "５": "5",
        "６": "6",
        "７": "7",
        "８": "8",
        "９": "9",
        "．": ".",
        "，": ",",
        "－": "-",
        "ー": "-",
        "＋": "+",
        "％": "%",
    }
)


def build_ui_icon(code: str, *, tone: str = "accent", circle: bool = True) -> str:
    """指定したコードから単色スタイルのUIアイコンHTMLを生成する。"""

    safe_code = html.escape(code or "")
    icon_tone = tone if tone in VALID_ICON_TONES else "accent"
    classes = ["ui-icon"]
    if circle:
        classes.append("ui-icon--circle")
    classes.append(f"ui-icon--{icon_tone}")
    class_attr = " ".join(classes)
    return (
        f"<span class=\"{class_attr}\" data-icon=\"{safe_code}\" "
        "aria-hidden=\"true\"></span>"
    )


def _clone_state_value(value: Any) -> Any:
    """リストなどのミュータブル値をコピーし、副作用を防ぐ。"""

    if isinstance(value, list):
        return list(value)
    if isinstance(value, dict):
        return dict(value)
    if isinstance(value, tuple):
        return tuple(value)
    return value


def set_state_and_widget(state_key: str, value: Any) -> None:
    """状態と対応するウィジェット値を同期して更新する。"""

    cloned_value = _clone_state_value(value)
    st.session_state[state_key] = cloned_value
    st.session_state[widget_key_for(state_key)] = _clone_state_value(cloned_value)


def ensure_widget_mirror(state_key: str) -> str:
    """状態の値を反映したウィジェットkeyを初期化する。"""

    widget_key = widget_key_for(state_key)
    if widget_key not in st.session_state:
        st.session_state[widget_key] = _clone_state_value(st.session_state.get(state_key))
    return widget_key


def update_state_from_widget(state_key: str) -> None:
    """ウィジェット値を読み取りアプリ状態に反映する。"""

    widget_key = widget_key_for(state_key)
    if widget_key in st.session_state:
        st.session_state[state_key] = _clone_state_value(st.session_state[widget_key])


def add_manual_filter_value(
    filter_state_key: str,
    manual_state_key: str,
    input_widget_key: str,
) -> None:
    """テキスト入力されたフィルタ値を永続化し、選択中の値に追加する。"""

    raw_value = st.session_state.get(input_widget_key)
    if raw_value is None:
        return

    value = str(raw_value).strip()
    st.session_state[input_widget_key] = ""
    if not value:
        return

    manual_values = list(st.session_state.get(manual_state_key, []))
    if value not in manual_values:
        manual_values.append(value)
        st.session_state[manual_state_key] = manual_values

    current_selection = list(st.session_state.get(filter_state_key, []))
    if value not in current_selection:
        current_selection.append(value)
    set_state_and_widget(filter_state_key, current_selection)
    trigger_rerun()


STATE_MESSAGES: Dict[str, Dict[str, Any]] = {
    "data_unloaded": {
        "type": "warning",
        "text": "データが読み込まれていません。サイドバー上部の『はじめに』からサンプルデータを読み込むか、CSVをアップロードしてください。",
        "action_label": "サンプルデータを読み込む",
        "secondary_action_label": "ファイルを選択",
    },
    "loading": {
        "type": "info",
        "text": "データを準備しています。完了まで数秒お待ちください…",
    },
    "filter_no_result": {
        "type": "warning",
        "text": "該当するデータが見つかりません。期間や店舗フィルタを変更して再検索してください。",
        "action_label": "フィルタをリセット",
    },
    "upload_failed": {
        "type": "error",
        "text": "CSVファイルの形式が正しくありません。サンプルテンプレートをダウンロードしてフォーマットを確認してください。",
        "action_label": "再アップロード",
        "secondary_action_label": "テンプレートをダウンロード",
    },
    "server_error": {
        "type": "error",
        "text": "予期しないエラーが発生しました。ページを再読み込みして再試行してください。",
        "action_label": "再読み込み",
    },
    "success": {
        "type": "success",
        "text": "フィルタ設定を更新しました。",
    },
    "warning_gross_margin": {
        "type": "warning",
        "text": "粗利率が目標を下回っています。商品構成を見直しましょう。",
    },
    "csv_done": {
        "type": "info",
        "text": "CSVをダウンロードしました。",
    },
    "unauthorized": {
        "type": "error",
        "text": "この操作を行う権限がありません。管理者にお問い合わせください。",
    },
}

SALES_IMPORT_CANDIDATES: Dict[str, List[str]] = {
    "項目": ["項目", "科目", "勘定科目", "売上科目", "部門"],
    "月次売上": ["月次売上", "金額", "売上高", "予測額"],
    "チャネル": ["チャネル", "分類", "モール", "部門", "経路"],
}

EXPENSE_IMPORT_CANDIDATES: Dict[str, List[str]] = {
    "費目": ["費目", "科目", "勘定科目", "費用科目"],
    "月次金額": ["月次金額", "金額", "予算額", "支出額"],
    "区分": ["区分", "分類", "タイプ", "費用区分"],
}


UPLOAD_META_MULTIPLE = "対応形式: CSV, Excel（最大10MB・複数ファイル対応）"
UPLOAD_META_SINGLE = "対応形式: CSV, Excel（最大10MB・1ファイル）"
UPLOAD_HELP_MULTIPLE = "CSVまたはExcelファイルをドラッグ＆ドロップで追加できます。複数ファイルをまとめてアップロードできます。"
UPLOAD_HELP_SINGLE = "CSVまたはExcelファイルをドラッグ＆ドロップでアップロードしてください。1ファイルのみアップロードできます。"

SALES_UPLOAD_CONFIGS: List[Dict[str, str]] = [
    {
        "channel": "自社サイト",
        "label": "自社サイト売上データ",
        "description": "公式ECサイトの受注・売上明細ファイルをアップロードしてください。",
    },
    {
        "channel": "楽天市場",
        "label": "楽天市場売上データ",
        "description": "楽天RMSなどからダウンロードした売上CSV/Excelを読み込みます。",
    },
    {
        "channel": "Amazon",
        "label": "Amazon売上データ",
        "description": "Amazonセラーセントラルのレポートをアップロードします。",
    },
    {
        "channel": "Yahoo!ショッピング",
        "label": "Yahoo!ショッピング売上データ",
        "description": "ストアクリエイターProから出力した受注データを取り込みます。",
    },
]

CHANNEL_ASSIGNMENT_PLACEHOLDER = "チャネルを選択"

ANCILLARY_UPLOAD_CONFIGS: List[Dict[str, Any]] = [
    {
        "key": "cost",
        "label": "商品原価率一覧",
        "description": "商品別の売価・原価・原価率がまとまったファイルをアップロードします。",
        "meta_text": UPLOAD_META_SINGLE,
        "help_text": "商品原価率表のCSVまたはExcelを1ファイルだけアップロードできます。",
        "multiple": False,
    },
    {
        "key": "subscription",
        "label": "定期購買/KPIデータ",
        "description": "サブスク会員数・解約数などの月次KPIを含むファイルを読み込みます。",
        "meta_text": UPLOAD_META_SINGLE,
        "help_text": "サブスクリプションのKPIを記載したCSVまたはExcelを1ファイルアップロードしてください。",
        "multiple": False,
    },
]


STATUS_PILL_DETAILS: Dict[str, Tuple[str, str, str]] = {
    "ok": ("OK", "success", "正常"),
    "warning": ("!", "warning", "警告"),
    "error": ("×", "error", "エラー"),
}


PRIMARY_NAV_ITEMS: List[Dict[str, str]] = [
    {"key": "dashboard", "label": "ダッシュボード", "icon": "DB"},
    {"key": "sales", "label": "売上", "icon": "SL"},
    {"key": "gross", "label": "粗利", "icon": "GR"},
    {"key": "inventory", "label": "在庫", "icon": "IV"},
    {"key": "cash", "label": "資金", "icon": "CS"},
    {"key": "kpi", "label": "KPI", "icon": "KP"},
    {"key": "scenario", "label": "シナリオ分析", "icon": "SC"},
    {"key": "data", "label": "データ管理", "icon": "DT"},
]

NAV_LABEL_LOOKUP: Dict[str, str] = {item["key"]: item["label"] for item in PRIMARY_NAV_ITEMS}
NAV_OPTION_LOOKUP: Dict[str, str] = {
    item["key"]: f"[{item['icon']}] {item['label']}" for item in PRIMARY_NAV_ITEMS
}

TUTORIAL_INDEX: List[Dict[str, Any]] = [
    {
        "title": "KPIの読み解き方と活用ガイド",
        "keywords": ["kpi", "活用", "レポート", "ダッシュボード"],
        "path": "docs/01_user_research_and_kpi.md",
    }
]


PRIMARY_COLOR = "#0B3C8C"
SECONDARY_COLOR = "#1E293B"
ACCENT_COLOR = "#2563EB"
BACKGROUND_COLOR = "#F8FAFC"
SURFACE_COLOR = "#FFFFFF"
SUCCESS_COLOR = "#0F766E"
WARNING_COLOR = "#B45309"
ERROR_COLOR = "#B91C1C"
TEXT_COLOR = "#0F172A"
CAPTION_TEXT_COLOR = "#475569"
MUTED_TEXT_COLOR = SECONDARY_COLOR
LIGHT_THEME_TOKENS: Dict[str, str] = {
    "background": BACKGROUND_COLOR,
    "surface": SURFACE_COLOR,
    "text": TEXT_COLOR,
    "caption": CAPTION_TEXT_COLOR,
    "muted": MUTED_TEXT_COLOR,
    "border_subtle": "rgba(15,23,42,0.08)",
    "border_strong": "rgba(15,23,42,0.22)",
    "grid": "rgba(15,23,42,0.08)",
    "domain": "rgba(15,23,42,0.2)",
    "surface_tint": "rgba(15,23,42,0.05)",
}
MCKINSEY_FONT_STACK = (
    "'Inter', 'Inter var', 'Source Sans 3', '-apple-system', 'BlinkMacSystemFont', "
    "'Segoe UI', 'Helvetica Neue', 'Arial', 'Noto Sans JP', sans-serif"
)
ALT_FONT_FAMILY = (
    "'Source Sans 3', '-apple-system', 'BlinkMacSystemFont', 'Segoe UI', "
    "'Helvetica Neue', 'Arial', 'Noto Sans JP', sans-serif"
)
NUMERIC_FONT_STACK = (
    "'Inter Tight', 'Inter', 'Inter var', 'Source Sans 3', '-apple-system', "
    "'BlinkMacSystemFont', 'Segoe UI', 'Helvetica Neue', 'Arial', 'Noto Sans JP', sans-serif"
)
MONO_FONT_STACK = "'Roboto Mono', 'Source Code Pro', monospace"

COMPANY_LOGO_URL = "https://raw.githubusercontent.com/streamlit/brand/main/logos/mark/streamlit-mark-color.png"
COLOR_TOKENS: Dict[str, str] = {
    "primary": PRIMARY_COLOR,
    "secondary": SECONDARY_COLOR,
    "accent": ACCENT_COLOR,
    "background": BACKGROUND_COLOR,
    "surface": SURFACE_COLOR,
    "text": TEXT_COLOR,
    "caption": CAPTION_TEXT_COLOR,
    "muted": MUTED_TEXT_COLOR,
    "success": SUCCESS_COLOR,
    "warning": WARNING_COLOR,
    "error": ERROR_COLOR,
}

DARK_THEME_VARIANTS: Dict[str, Dict[str, str]] = {
    "deep": {
        "background": "#050B18",
        "surface": "#111E2E",
        "text": "#EEF3FF",
        "caption": "#A8B5CB",
        "muted": "#8FA5C6",
        "border_subtle": "rgba(167,189,219,0.32)",
        "border_strong": "rgba(199,214,238,0.6)",
        "grid": "rgba(162,189,227,0.28)",
        "domain": "rgba(176,204,240,0.45)",
        "surface_tint": "rgba(90,126,173,0.22)",
    },
    "high_contrast": {
        "background": "#010409",
        "surface": "#0C1A2A",
        "text": "#F8FBFF",
        "caption": "#D0DCF2",
        "muted": "#9FB7DD",
        "border_subtle": "rgba(120,178,255,0.45)",
        "border_strong": "rgba(127,196,255,0.85)",
        "grid": "rgba(147,200,255,0.45)",
        "domain": "rgba(160,210,255,0.7)",
        "surface_tint": "rgba(116,170,250,0.32)",
    },
}

DEFAULT_DARK_THEME_VARIANT = "deep"
DARK_THEME_VARIANT_LABELS: Dict[str, str] = {
    "deep": "ディープブルー",
    "high_contrast": "ハイコントラスト",
}
COLOR_PALETTE_PRESETS: Dict[str, Dict[str, Union[str, List[str]]]] = {
    "brand": {
        "label": "ブランド",
        "colors": [ACCENT_COLOR, SUCCESS_COLOR],
    },
    "colorblind": {
        "label": "ユニバーサル",
        "colors": ["#0072B2", "#E69F00"],
    },
}
DEFAULT_CHART_PALETTE_KEY = "brand"

LANGUAGE_OPTIONS: List[Tuple[str, str]] = [
    ("ja", "日本語"),
    ("en", "English"),
]
DEFAULT_LANGUAGE = "ja"

I18N_STRINGS: Dict[str, Dict[str, str]] = {
    "theme_language_label": {"ja": "表示言語", "en": "Display language"},
    "theme_language_help": {
        "ja": "日本語/英語の表記とヘルプを切り替えます。",
        "en": "Switch between Japanese and English labels and help text.",
    },
    "theme_toggle_label": {"ja": "ダークテーマ", "en": "Dark theme"},
    "theme_toggle_help": {
        "ja": "ライトテーマに切り替えると背景が明るい配色になります。",
        "en": "Toggle to switch between light and dark color schemes.",
    },
    "theme_caption": {
        "ja": "配色やフォントサイズを調整して見やすいダッシュボードにカスタマイズできます。",
        "en": "Tweak colors and typography for better readability.",
    },
    "font_size_label": {"ja": "本文フォントサイズ", "en": "Body font size"},
    "font_size_help": {
        "ja": "本文や表の文字サイズを調整します (基準値=100)。",
        "en": "Adjust type scale for body text and tables (base = 100).",
    },
    "dark_variant_label": {
        "ja": "暗色テーマのコントラスト",
        "en": "Dark theme contrast",
    },
    "dark_variant_help": {
        "ja": "暗色テーマ時の背景/境界のコントラストを調整します。",
        "en": "Fine-tune contrast tokens when dark mode is enabled.",
    },
    "palette_label": {
        "ja": "チャートカラーパレット",
        "en": "Chart color palette",
    },
    "palette_help": {
        "ja": "色覚多様性に配慮した配色に切り替えられます。",
        "en": "Choose palettes with improved color-vision accessibility.",
    },
    "quick_search_title": {"ja": "クイック検索", "en": "Quick search"},
    "quick_search_placeholder": {
        "ja": "商品名、チャネル、チュートリアルを検索",
        "en": "Search products, channels, or tutorials",
    },
    "quick_search_heading": {
        "ja": "クイック検索結果",
        "en": "Quick search results",
    },
    "quick_search_no_data": {
        "ja": "売上データが読み込まれていないため検索できません。",
        "en": "Sales data has not been loaded yet.",
    },
    "quick_search_no_match": {
        "ja": "売上データに一致する項目は見つかりませんでした。",
        "en": "No matching rows were found in sales data.",
    },
    "quick_search_tab_sales": {
        "ja": "売上データ",
        "en": "Sales data",
    },
    "quick_search_tab_tutorials": {
        "ja": "チュートリアル",
        "en": "Tutorials",
    },
    "quick_search_tab_help": {
        "ja": "検索結果の種類を切り替えます。",
        "en": "Switch between sales matches and documentation.",
    },
    "quick_search_tutorial_heading": {
        "ja": "関連チュートリアル",
        "en": "Related tutorials",
    },
    "quick_search_caption": {
        "ja": "直近の一致10件を表示しています。",
        "en": "Showing the 10 most recent matches.",
    },
    "quick_search_no_tutorial": {
        "ja": "該当するチュートリアルは見つかりませんでした。",
        "en": "No tutorials matched the query.",
    },
    "back_to_top_label": {"ja": "トップへ戻る", "en": "Back to top"},
    "back_to_top_aria": {
        "ja": "ページの先頭へ戻る",
        "en": "Return to the top of the page",
    },
    "streamlit_badge_aria": {
        "ja": "Streamlit公式サイト (外部リンク)",
        "en": "Streamlit official site (opens in new tab)",
    },
    "kpi_help_language_label": {"ja": "KPIヘルプ言語", "en": "KPI help language"},
    "kpi_help_button_label": {"ja": "KPIヘルプ", "en": "KPI help"},
    "kpi_help_button_help": {
        "ja": "指標の定義一覧を表示します。",
        "en": "Open the glossary of KPI definitions.",
    },
    "kpi_help_panel_title": {"ja": "KPI定義ヘルプ", "en": "KPI definitions"},
    "kpi_help_close": {"ja": "ヘルプを閉じる", "en": "Close help"},
}

PHASE2_SESSION_DEFAULTS: Dict[str, Any] = {
    "scenario_inputs": [],
    "scenario_uploaded_df": None,
    "scenario_results": None,
    "phase2_summary_df": None,
    "phase2_swot": None,
    "phase2_benchmark": None,
    "phase2_report_summary": "",
    "ui_theme_mode": "light",
    "ui_language": DEFAULT_LANGUAGE,
}


def init_phase2_session_state() -> None:
    """初回アクセス時にPhase2で利用するセッション状態を初期化する。"""

    for key, default in PHASE2_SESSION_DEFAULTS.items():
        if key not in st.session_state:
            if isinstance(default, list):
                st.session_state[key] = list(default)
            elif isinstance(default, dict):
                st.session_state[key] = dict(default)
            else:
                st.session_state[key] = default


def ensure_language_state_defaults() -> None:
    """UI言語設定のデフォルト値をセッションに登録する。"""

    st.session_state.setdefault("ui_language", DEFAULT_LANGUAGE)


def get_ui_language() -> str:
    """現在選択されているUI言語コードを返す。"""

    return str(st.session_state.get("ui_language", DEFAULT_LANGUAGE))


def translate(key: str, *, default: Optional[str] = None, language: Optional[str] = None) -> str:
    """I18N辞書から翻訳文字列を取得する。"""

    lang = language or get_ui_language()
    if key in I18N_STRINGS:
        entry = I18N_STRINGS[key]
        if lang in entry:
            return entry[lang]
        if DEFAULT_LANGUAGE in entry:
            return entry[DEFAULT_LANGUAGE]
    return default if default is not None else key

TYPOGRAPHY_TOKENS: Dict[str, Dict[str, Union[str, int, float]]] = {
    "h1": {"size": "1.75rem", "weight": 700, "line_height": 1.35},
    "h2": {"size": "1.35rem", "weight": 600, "line_height": 1.4},
    "body": {"size": "0.95rem", "weight": 400, "line_height": 1.55},
    "body_small": {"size": "0.85rem", "weight": 400, "line_height": 1.5},
    "numeric": {"size": "1.1rem", "weight": 600, "line_height": 1.35},
    "caption": {"size": "0.78rem", "weight": 400, "line_height": 1.45},
}

SPACING_UNIT = 8
SPACING_SCALE: Dict[str, str] = {
    "xs": f"{SPACING_UNIT * 1:g}px",
    "sm": f"{SPACING_UNIT * 1.5:g}px",
    "md": f"{SPACING_UNIT * 2:g}px",
    "lg": f"{SPACING_UNIT * 3:g}px",
    "xl": f"{SPACING_UNIT * 4:g}px",
}

RADIUS_TOKENS: Dict[str, str] = {
    "card": "12px",
    "panel": "16px",
    "chip": "999px",
    "input": "10px",
}

SHADOW_TOKENS: Dict[str, str] = {
    "sm": "0 8px 16px rgba(15,23,42,0.08)",
    "md": "0 12px 24px rgba(15,23,42,0.08)",
    "lg": "0 24px 48px rgba(15,23,42,0.16)",
}


def _hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """16進カラーコードをRGBタプルに変換する。"""

    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def _rgba_from_hex(hex_color: str, alpha: float) -> str:
    """16進カラーコードと透過率からCSSのrgba文字列を生成する。"""

    r, g, b = _hex_to_rgb(hex_color)
    alpha_str = f"{alpha:.3f}".rstrip("0").rstrip(".")
    return f"rgba({r},{g},{b},{alpha_str})"


SUCCESS_RGB = _hex_to_rgb(SUCCESS_COLOR)
WARNING_RGB = _hex_to_rgb(WARNING_COLOR)
ERROR_RGB = _hex_to_rgb(ERROR_COLOR)

SUCCESS_SURFACE_COLOR = _rgba_from_hex(SUCCESS_COLOR, 0.12)
WARNING_SURFACE_COLOR = _rgba_from_hex(WARNING_COLOR, 0.128)
ERROR_SURFACE_COLOR = _rgba_from_hex(ERROR_COLOR, 0.128)
SUCCESS_GAUGE_RANGE_COLOR = _rgba_from_hex(SUCCESS_COLOR, 0.144)
ERROR_GAUGE_RANGE_COLOR = _rgba_from_hex(ERROR_COLOR, 0.144)

SALES_SERIES_COLOR = PRIMARY_COLOR
GROSS_SERIES_COLOR = ACCENT_COLOR
OPERATING_SERIES_COLOR = "#38BDF8"
INVENTORY_SERIES_COLOR = "#1D4ED8"
CASH_SERIES_COLOR = ACCENT_COLOR
YOY_SERIES_COLOR = SECONDARY_COLOR
BASELINE_SERIES_COLOR = "#94A3B8"
FORECAST_SERIES_COLOR = "#0EA5E9"
FORECAST_BAND_COLOR = _rgba_from_hex(FORECAST_SERIES_COLOR, 0.18)

CF_COLOR_MAPPING = {
    "営業CF": SALES_SERIES_COLOR,
    "投資CF": YOY_SERIES_COLOR,
    "財務CF": GROSS_SERIES_COLOR,
    "返済": WARNING_COLOR,
}

def _parse_rem(value: Union[str, float, int]) -> float:
    """rem単位の文字列から数値部分を抽出し、小数として返す。"""

    if isinstance(value, (int, float)):
        return float(value)
    stripped = value.strip().lower().replace("rem", "")
    try:
        return float(stripped)
    except ValueError as exc:  # pragma: no cover - defensive
        raise ValueError(f"Invalid rem value: {value!r}") from exc


def _scaled_rem(value: Union[str, float, int], scale: float) -> str:
    """rem表記の値にスケールを掛け合わせ、rem文字列として返す。"""

    return f"{_parse_rem(value) * scale:.2f}rem"


def _rem_to_px(value: Union[str, float, int], scale: float, *, base: int = 16) -> int:
    """rem表記をpxに換算する。Streamlit/ブラウザ標準の16pxを基準とする。"""

    return int(round(_parse_rem(value) * scale * base))


def get_font_scale() -> float:
    """現在のフォントスケールをセッション状態から取得する。"""

    return float(st.session_state.get("ui_font_scale", 1.0))


def get_active_chart_colorway() -> List[str]:
    """ユーザー設定に基づいたチャートの配色セットを返す。"""

    palette_key = st.session_state.get("ui_color_palette", DEFAULT_CHART_PALETTE_KEY)
    palette = COLOR_PALETTE_PRESETS.get(palette_key)
    if not palette:
        palette = COLOR_PALETTE_PRESETS[DEFAULT_CHART_PALETTE_KEY]
    colors = list(palette["colors"])  # type: ignore[index]
    if len(colors) < 2:
        colors.append(SECONDARY_COLOR)
    return colors


def get_theme_tokens(*, dark_mode: bool, variant: Optional[str] = None) -> Dict[str, str]:
    """現在のテーマに応じたトークンセットを返す。"""

    if not dark_mode:
        return LIGHT_THEME_TOKENS
    chosen = variant or st.session_state.get("ui_dark_variant", DEFAULT_DARK_THEME_VARIANT)
    return DARK_THEME_VARIANTS.get(chosen, DARK_THEME_VARIANTS[DEFAULT_DARK_THEME_VARIANT])


def ensure_theme_state_defaults() -> None:
    """テーマ関連のセッションデフォルト値を設定する。"""

    st.session_state.setdefault("ui_color_palette", DEFAULT_CHART_PALETTE_KEY)
    st.session_state.setdefault("ui_dark_variant", DEFAULT_DARK_THEME_VARIANT)
    st.session_state.setdefault("ui_font_scale", 1.0)
    st.session_state.setdefault("ui_dark_palette_saved", DEFAULT_CHART_PALETTE_KEY)
    st.session_state.setdefault("ui_dark_variant_saved", DEFAULT_DARK_THEME_VARIANT)

HEATMAP_BLUE_SCALE = [[0.0, "#E2E8F0"], [0.5, "#60A5FA"], [1.0, ACCENT_COLOR]]


KGI_TARGETS = {
    "sales": 7_000_000,
    "gross_margin_rate": 0.62,
    "cash_balance": 5_000_000,
}


BSC_TARGETS = {
    "ltv": 60_000,
    "repeat_rate": 0.45,
    "inventory_turnover_days": 45,
    "training_sessions": 6,
}


KPI_HELP_CONTENT: Dict[str, Dict[str, Dict[str, str]]] = {
    "active_customers": {
        "ja": {
            "title": "月次顧客数",
            "formula": "対象期間中に購入・契約したユニーク顧客数。",
            "interpretation": "前月比や目標との差からリテンション施策の成果を把握します。",
        },
        "en": {
            "title": "Monthly Active Customers",
            "formula": "Number of unique customers who purchased or subscribed in the period.",
            "interpretation": "Track changes versus the previous month to understand retention performance.",
        },
    },
    "ltv": {
        "ja": {
            "title": "LTV",
            "formula": "平均購入単価 × リピート回数 × 粗利率。",
            "interpretation": "顧客1人から得られる長期的な利益規模を示し、マーケ投資の上限設定に活用します。",
        },
        "en": {
            "title": "LTV",
            "formula": "Average order value × Repeat frequency × Gross margin rate.",
            "interpretation": "Represents expected lifetime profit per customer and guides marketing spend limits.",
        },
    },
    "arpu": {
        "ja": {
            "title": "ARPU",
            "formula": "売上高 ÷ アクティブ顧客数。",
            "interpretation": "顧客1人あたりの売上規模を示し、単価施策やアップセルの効果を確認します。",
        },
        "en": {
            "title": "ARPU",
            "formula": "Revenue ÷ Active customers.",
            "interpretation": "Highlights monetisation per customer and the effectiveness of upsell initiatives.",
        },
    },
    "churn_rate": {
        "ja": {
            "title": "解約率",
            "formula": "当月の解約件数 ÷ 前月の契約者数。",
            "interpretation": "値が高いほど解約が増加しているため、CSやプロダクト改善の緊急度が高まります。",
        },
        "en": {
            "title": "Churn Rate",
            "formula": "Cancelled subscriptions ÷ Prior month subscribers.",
            "interpretation": "A higher value indicates accelerated churn and calls for urgent retention actions.",
        },
    },
    "gross_margin_rate": {
        "ja": {
            "title": "粗利率",
            "formula": "粗利 ÷ 売上高。",
            "interpretation": "仕入・原価管理の効率を測る指標で、価格改定や値引きの影響を確認します。",
        },
        "en": {
            "title": "Gross Margin Rate",
            "formula": "Gross profit ÷ Revenue.",
            "interpretation": "Measures purchasing efficiency and the impact of pricing or discounting decisions.",
        },
    },
    "repeat_rate": {
        "ja": {
            "title": "リピート率",
            "formula": "リピート顧客 ÷ アクティブ顧客数。",
            "interpretation": "ファン化の度合いを示し、CRM施策の成果を把握します。",
        },
        "en": {
            "title": "Repeat Rate",
            "formula": "Repeat customers ÷ Active customers.",
            "interpretation": "Shows how effectively customers are retained through CRM initiatives.",
        },
    },
    "inventory_turnover_days": {
        "ja": {
            "title": "在庫回転日数",
            "formula": "平均在庫金額 ÷ 売上原価 × 30日換算。",
            "interpretation": "在庫の滞留状況や安全在庫の適正さを判断します。値が小さいほど現金化が早いです。",
        },
        "en": {
            "title": "Inventory Turnover Days",
            "formula": "Average inventory cost ÷ Cost of goods sold × 30 days.",
            "interpretation": "Evaluates inventory velocity and safety stock; lower values indicate faster cash conversion.",
        },
    },
}


def apply_chart_theme(fig):
    """デザイン・トークンに基づいたPlotly共通スタイルを適用する。"""

    ensure_theme_state_defaults()
    font_scale = get_font_scale()
    tokens = st.session_state.get("ui_active_tokens", LIGHT_THEME_TOKENS)
    colorway = get_active_chart_colorway()
    body_px = _rem_to_px(TYPOGRAPHY_TOKENS["body"]["size"], font_scale)
    caption_px = _rem_to_px(TYPOGRAPHY_TOKENS["caption"]["size"], font_scale)
    title_px = _rem_to_px(TYPOGRAPHY_TOKENS["h2"]["size"], font_scale)

    fig.update_layout(
        font=dict(family=MCKINSEY_FONT_STACK, color=tokens["text"], size=body_px),
        title=dict(
            font=dict(size=title_px, color=tokens["text"], family=MCKINSEY_FONT_STACK)
        ),
        legend=dict(
            bgcolor="rgba(0,0,0,0)",
            font=dict(size=caption_px, color=tokens["text"]),
        ),
        plot_bgcolor="rgba(0,0,0,0)",
        paper_bgcolor="rgba(0,0,0,0)",
        margin=dict(l=48, r=36, t=60, b=48),
        hoverlabel=dict(
            font=dict(
                family=MCKINSEY_FONT_STACK,
                color=tokens["text"],
                size=body_px,
            )
        ),
        colorway=colorway,
    )
    axis_tick_font = dict(color=tokens["muted"], size=caption_px)
    axis_title_font = dict(color=tokens["muted"], size=caption_px)
    fig.update_xaxes(
        showgrid=True,
        gridcolor=tokens["grid"],
        linecolor=tokens["domain"],
        tickfont=axis_tick_font,
        title_font=axis_title_font,
    )
    fig.update_yaxes(
        showgrid=True,
        gridcolor=tokens["grid"],
        linecolor=tokens["domain"],
        tickfont=axis_tick_font,
        title_font=axis_title_font,
    )
    return fig


def apply_altair_theme(chart: alt.Chart) -> alt.Chart:
    """Altairグラフに共通のスタイル・タイポグラフィを適用する。"""

    ensure_theme_state_defaults()
    font_scale = get_font_scale()
    tokens = st.session_state.get("ui_active_tokens", LIGHT_THEME_TOKENS)
    palette = get_active_chart_colorway()
    axis_label_size = _rem_to_px(TYPOGRAPHY_TOKENS["body_small"]["size"], font_scale)
    axis_title_size = _rem_to_px(TYPOGRAPHY_TOKENS["body"]["size"], font_scale)
    legend_size = _rem_to_px(TYPOGRAPHY_TOKENS["body_small"]["size"], font_scale)
    title_size = _rem_to_px(TYPOGRAPHY_TOKENS["h2"]["size"], font_scale)

    return (
        chart.configure_axis(
            labelFont=MCKINSEY_FONT_STACK,
            titleFont=MCKINSEY_FONT_STACK,
            labelColor=tokens["muted"],
            titleColor=tokens["text"],
            labelFontSize=axis_label_size,
            titleFontSize=axis_title_size,
            gridColor=tokens["grid"],
            domainColor=tokens["domain"],
        )
        .configure_legend(
            titleFont=MCKINSEY_FONT_STACK,
            labelFont=MCKINSEY_FONT_STACK,
            labelColor=tokens["text"],
            titleColor=tokens["muted"],
            titleFontSize=legend_size,
            labelFontSize=legend_size,
            orient="top",
            direction="horizontal",
            symbolSize=120,
        )
        .configure_range(category=palette, ordinal=palette)
        .configure_view(strokeOpacity=0)
        .configure_title(
            font=MCKINSEY_FONT_STACK,
            color=tokens["text"],
            fontSize=title_size,
        )
        .configure_mark(
            font=MCKINSEY_FONT_STACK,
            color=palette[0],
            fill=palette[0],
            stroke=palette[0],
        )
    )


def inject_mckinsey_style(
    *, dark_mode: bool = False, theme_variant: Optional[str] = None, font_scale: Optional[float] = None
) -> None:
    """デザイン・トークンとマッキンゼー風スタイルをアプリに適用する。"""

    ensure_theme_state_defaults()
    tokens = get_theme_tokens(dark_mode=dark_mode, variant=theme_variant)
    st.session_state["ui_active_tokens"] = tokens

    resolved_font_scale = font_scale if font_scale is not None else get_font_scale()
    if font_scale is not None:
        st.session_state["ui_font_scale"] = resolved_font_scale

    chart_colors = get_active_chart_colorway()
    chart_primary = chart_colors[0]
    chart_secondary = chart_colors[1] if len(chart_colors) > 1 else chart_colors[0]

    typography_scaled = {
        "h1": _scaled_rem(TYPOGRAPHY_TOKENS["h1"]["size"], resolved_font_scale),
        "h2": _scaled_rem(TYPOGRAPHY_TOKENS["h2"]["size"], resolved_font_scale),
        "body": _scaled_rem(TYPOGRAPHY_TOKENS["body"]["size"], resolved_font_scale),
        "caption": _scaled_rem(TYPOGRAPHY_TOKENS["caption"]["size"], resolved_font_scale),
    }

    st.markdown(
        f"""
        <style>
        :root {{
            --primary-color: {PRIMARY_COLOR};
            --secondary-color: {SECONDARY_COLOR};
            --accent-color: {ACCENT_COLOR};
            --success-color: {SUCCESS_COLOR};
            --warning-color: {WARNING_COLOR};
            --error-color: {ERROR_COLOR};
            --success-rgb: {",".join(map(str, SUCCESS_RGB))};
            --warning-rgb: {",".join(map(str, WARNING_RGB))};
            --error-rgb: {",".join(map(str, ERROR_RGB))};
            --success-surface: {SUCCESS_SURFACE_COLOR};
            --warning-surface: {WARNING_SURFACE_COLOR};
            --error-surface: {ERROR_SURFACE_COLOR};
            --surface-color: {tokens['surface']};
            --background-color: {tokens['background']};
            --text-color: {tokens['text']};
            --muted-text-color: {tokens['muted']};
            --caption-text-color: {tokens['caption']};
            --border-subtle-color: {tokens['border_subtle']};
            --border-strong-color: {tokens['border_strong']};
            --surface-tint-color: {tokens['surface_tint']};
            --grid-color: {tokens['grid']};
            --domain-color: {tokens['domain']};
            --font-family: {MCKINSEY_FONT_STACK};
            --alt-font-family: {ALT_FONT_FAMILY};
            --numeric-font-family: {NUMERIC_FONT_STACK};
            --h1-size: {typography_scaled['h1']};
            --h1-line-height: {TYPOGRAPHY_TOKENS['h1']['line_height']};
            --h2-size: {typography_scaled['h2']};
            --h2-line-height: {TYPOGRAPHY_TOKENS['h2']['line_height']};
            --body-size: {typography_scaled['body']};
            --body-line-height: {TYPOGRAPHY_TOKENS['body']['line_height']};
            --caption-size: {typography_scaled['caption']};
            --caption-line-height: {TYPOGRAPHY_TOKENS['caption']['line_height']};
            --font-scale: {resolved_font_scale:.2f};
            --chart-primary: {chart_primary};
            --chart-secondary: {chart_secondary};
            --spacing-xs: {SPACING_SCALE['xs']};
            --spacing-sm: {SPACING_SCALE['sm']};
            --spacing-md: {SPACING_SCALE['md']};
            --spacing-lg: {SPACING_SCALE['lg']};
            --spacing-xl: {SPACING_SCALE['xl']};
            --radius-card: {RADIUS_TOKENS['card']};
            --radius-panel: {RADIUS_TOKENS['panel']};
            --radius-chip: {RADIUS_TOKENS['chip']};
            --radius-input: {RADIUS_TOKENS['input']};
            --shadow-sm: {SHADOW_TOKENS['sm']};
            --shadow-md: {SHADOW_TOKENS['md']};
            --shadow-lg: {SHADOW_TOKENS['lg']};
            --grid-max-width: 1200px;
        }}
        html, body {{
            background: var(--background-color);
            background-color: var(--background-color);
            background-image: none;
            color: var(--text-color);
            font-family: var(--font-family);
            font-size: var(--body-size);
            line-height: var(--body-line-height);
        }}
        .stApp, [data-testid="stAppViewContainer"] {{
            background: var(--background-color);
            background-color: var(--background-color);
            background-image: none;
            color: var(--text-color);
        }}
        main .block-container {{
            max-width: var(--grid-max-width);
            padding: 3rem 2.5rem 3.5rem;
            color: var(--text-color);
        }}
        main .block-container p {{
            color: var(--text-color);
            font-size: var(--body-size);
            line-height: var(--body-line-height);
        }}
        .with-icon {{
            display: inline-flex;
            align-items: center;
            gap: 0.4rem;
        }}
        .with-icon--tight {{
            gap: 0.3rem;
        }}
        .ui-icon {{
            display: inline-flex;
            align-items: center;
            justify-content: center;
            min-width: 1.3rem;
            width: 1.3rem;
            height: 1.3rem;
            border-radius: 0.4rem;
            font-size: 0.68rem;
            font-weight: 700;
            letter-spacing: 0.05em;
            text-transform: uppercase;
            font-family: var(--alt-font-family);
        }}
        .ui-icon::before {{
            content: attr(data-icon);
            line-height: 1;
        }}
        .ui-icon--circle {{
            border-radius: 999px;
        }}
        .ui-icon--accent {{
            background: var(--accent-color);
            color: #ffffff;
        }}
        .ui-icon--muted {{
            background: var(--secondary-color);
            color: #ffffff;
        }}
        .ui-icon--success {{
            background: var(--success-color);
            color: #ffffff;
        }}
        .ui-icon--warning {{
            background: var(--warning-color);
            color: #ffffff;
        }}
        .ui-icon--error {{
            background: var(--error-color);
            color: #ffffff;
        }}
        .caption-with-icon {{
            display: inline-flex;
            align-items: center;
            gap: 0.3rem;
            color: var(--caption-text-color);
            font-size: var(--caption-size);
            line-height: var(--caption-line-height);
        }}
        .caption-with-icon .ui-icon {{
            min-width: 1.1rem;
            width: 1.1rem;
            height: 1.1rem;
            font-size: 0.6rem;
        }}
        main .block-container h1,
        main .block-container h2,
        main .block-container h3,
        main .block-container h4 {{
            color: var(--text-color);
            font-family: var(--font-family);
        }}
        main .block-container h1 {{
            font-size: var(--h1-size);
            font-weight: 700;
            line-height: var(--h1-line-height);
            margin-bottom: var(--spacing-sm);
        }}
        main .block-container h2 {{
            font-size: var(--h2-size);
            font-weight: 600;
            line-height: var(--h2-line-height);
            margin-bottom: {SPACING_SCALE['sm']};
        }}
        main .block-container h3 {{
            font-size: calc(1.15rem * var(--font-scale));
            font-weight: 600;
            margin-bottom: var(--spacing-xs);
        }}
        main .block-container h4 {{
            font-size: calc(1rem * var(--font-scale));
            font-weight: 600;
        }}
        .stMarkdown li,
        label,
        span[data-testid="stMarkdownContainer"] p {{
            font-size: var(--body-size);
            line-height: var(--body-line-height);
        }}
        small,
        .stCaption,
        .stMarkdown small,
        .caption {{
            color: var(--caption-text-color);
            font-size: var(--caption-size);
            line-height: var(--caption-line-height);
        }}
        .surface-card {{
            background: var(--surface-color);
            border-radius: var(--radius-card);
            padding: 2rem;
            border: 1px solid var(--border-subtle-color);
            box-shadow: var(--shadow-md);
            overflow: visible;
        }}
        .surface-card.main-nav-block {{
            position: sticky;
            top: 1.5rem;
            z-index: 40;
            padding: 1.5rem 2rem;
            margin-bottom: var(--spacing-md);
            backdrop-filter: blur(8px);
        }}
        .surface-card.quick-actions-block {{
            margin-bottom: var(--spacing-md);
        }}
        .quick-action-info {{
            background: rgba(11, 60, 140, 0.08);
            border-radius: var(--radius-panel);
            padding: 0.75rem 1rem;
            font-size: calc(0.95rem * var(--font-scale));
            color: var(--text-color);
            box-shadow: var(--shadow-sm);
        }}
        .quick-action-info strong {{
            color: var(--primary-color);
            font-weight: 700;
        }}
        .stButton button, .stDownloadButton button {{
            min-height: 3rem;
            font-weight: 600;
            border-radius: var(--radius-input);
            padding: 0.75rem 1.5rem;
        }}
        .main-nav-tabs {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(140px, 1fr));
            gap: 0.6rem;
        }}
        .main-nav-tabs .stButton > button {{
            border-radius: 999px;
            font-weight: 700;
            letter-spacing: 0.02em;
            padding: 0.65rem 1.25rem;
            transition: transform 0.2s ease, box-shadow 0.2s ease;
        }}
        .main-nav-tabs .stButton > button[data-testid="baseButton-secondary"] {{
            background: rgba(255, 255, 255, 0.05);
            border: 1px solid var(--border-subtle-color);
            color: var(--muted-text-color);
        }}
        .main-nav-tabs .stButton > button[data-testid="baseButton-secondary"]:hover {{
            transform: translateY(-1px);
            box-shadow: var(--shadow-sm);
        }}
        .main-nav-tabs .stButton > button[data-testid="baseButton-primary"] {{
            box-shadow: var(--shadow-md);
        }}
        .stButton button:focus-visible, .stDownloadButton button:focus-visible {{
            outline: 3px solid var(--accent-color);
            outline-offset: 2px;
        }}
        .back-to-top {{
            position: fixed;
            bottom: 2.5rem;
            right: 2.5rem;
            z-index: 120;
        }}
        .back-to-top a {{
            display: inline-flex;
            align-items: center;
            gap: 0.4rem;
            background: var(--accent-color);
            color: #ffffff !important;
            padding: 0.65rem 1rem;
            border-radius: 999px;
            font-weight: 600;
            text-decoration: none;
            box-shadow: var(--shadow-md);
        }}
        .back-to-top a:hover {{
            box-shadow: var(--shadow-lg);
        }}
        footer [data-testid="stFooter"] a[href*="streamlit.io"] {{
            position: fixed;
            bottom: 1.5rem;
            left: 1.5rem;
            display: inline-flex;
            align-items: center;
            gap: 0.35rem;
            background: rgba(15, 23, 42, 0.85);
            border: 1px solid rgba(255, 255, 255, 0.2);
            border-radius: var(--radius-chip);
            padding: 0.35rem 0.85rem;
            color: #ffffff;
            box-shadow: var(--shadow-md);
            opacity: 0.45;
            z-index: 140;
            transition: opacity 0.2s ease, transform 0.2s ease;
        }}
        footer [data-testid="stFooter"] a[href*="streamlit.io"]:hover {{
            opacity: 0.9;
            transform: translateY(-2px);
        }}
        footer [data-testid="stFooter"] a[href*="streamlit.io"]::after {{
            content: attr(data-external-label);
            font-size: 0.75rem;
            color: rgba(255, 255, 255, 0.75);
            margin-left: 0.25rem;
        }}
        .hero-panel {{
            background: linear-gradient(135deg, rgba(11,31,59,0.9), rgba(30,136,229,0.75));
            border-radius: var(--radius-panel);
            padding: 2.5rem 3rem;
            box-shadow: var(--shadow-lg);
            color: #ffffff;
            margin-bottom: 1.75rem;
        }}
        .hero-title {{
            font-size: 1.65rem;
            font-weight: 700;
            margin-bottom: 0.6rem;
        }}
        .hero-subtitle {{
            font-size: 1rem;
            margin-bottom: 1.2rem;
            color: rgba(255,255,255,0.82);
        }}
        .hero-meta {{
            display: flex;
            flex-wrap: wrap;
            gap: 0.6rem;
            margin-bottom: 1.2rem;
        }}
        .hero-badge {{
            display: inline-flex;
            align-items: center;
            gap: 0.35rem;
            padding: 0.5rem 1rem;
            border-radius: var(--radius-chip);
            background: rgba(255,255,255,0.18);
            color: #ffffff;
            font-size: 0.85rem;
            font-weight: 600;
        }}
        .hero-badge--alert {{
            background: rgba(208,135,0,0.22);
        }}
        .hero-badge--accent {{
            background: rgba(30,136,229,0.28);
        }}
        .hero-badge .ui-icon {{
            background: rgba(255,255,255,0.35);
        }}
        .hero-badge--alert .ui-icon {{
            background: var(--warning-color);
        }}
        .hero-badge--accent .ui-icon {{
            background: var(--accent-color);
        }}
        .hero-badge__label {{
            display: inline-flex;
            align-items: center;
        }}
        .hero-persona {{
            display: flex;
            flex-wrap: wrap;
            gap: 0.45rem;
        }}
        .hero-chip {{
            background: rgba(255,255,255,0.15);
            color: rgba(255,255,255,0.9);
            padding: 0.35rem 0.85rem;
            border-radius: 999px;
            font-size: 0.8rem;
            font-weight: 600;
            display: inline-flex;
            align-items: center;
            gap: 0.4rem;
        }}
        .hero-chip .ui-icon {{
            background: rgba(255,255,255,0.35);
            color: #0b1f3b;
        }}
        .chart-section {{
            background: var(--surface-color);
            border-radius: var(--radius-card);
            padding: 1.5rem 1.75rem;
            border: 1px solid var(--border-subtle-color);
            box-shadow: var(--shadow-md);
            margin-bottom: var(--spacing-lg);
        }}
        .bsc-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(260px, 1fr));
            gap: var(--spacing-md);
            margin: var(--spacing-md) 0;
        }}
        .bsc-card {{
            background: var(--surface-color);
            border-radius: var(--radius-card);
            border: 1px solid var(--border-subtle-color);
            padding: 1.25rem 1.5rem;
            box-shadow: var(--shadow-sm);
            transition: border-color 0.2s ease, box-shadow 0.2s ease;
        }}
        .bsc-card--success {{
            border-color: var(--success-color);
            box-shadow: 0 0 0 1px var(--success-color), var(--shadow-md);
            background: var(--success-surface);
        }}
        .bsc-card--warning {{
            border-color: var(--warning-color);
            box-shadow: 0 0 0 1px var(--warning-color), var(--shadow-md);
            background: var(--warning-surface);
        }}
        .bsc-card__title {{
            font-size: calc(1.05rem * var(--font-scale));
            font-weight: 600;
            margin-bottom: 0.35rem;
        }}
        .bsc-card__subtitle {{
            color: var(--muted-text-color);
            font-size: calc(0.85rem * var(--font-scale));
            margin-bottom: 0.75rem;
        }}
        .bsc-card .stMetric {{
            padding: 0.5rem 0;
        }}
        .bsc-card .stMetric label {{
            font-size: calc(0.9rem * var(--font-scale));
        }}
        .bsc-card .stMetric div[data-testid="stMetricValue"] {{
            font-family: var(--numeric-font-family);
            font-size: calc(1.25rem * var(--font-scale));
        }}
        .bsc-card .stMetric div[data-testid="stMetricDeltaValue"] {{
            font-size: calc(0.85rem * var(--font-scale));
        }}
        .chart-section__header {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: var(--spacing-sm);
        }}
        .chart-section__title {{
            font-size: 1.05rem;
            font-weight: 700;
            color: var(--text-color);
        }}
        .kpi-strip {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: var(--spacing-sm);
            margin-bottom: var(--spacing-lg);
        }}
        .kpi-strip__card {{
            background: var(--surface-color);
            border-radius: var(--radius-card);
            padding: 1.25rem 1.5rem;
            border: 1px solid var(--border-subtle-color);
            box-shadow: var(--shadow-sm);
        }}
        .kpi-strip__label {{
            font-size: 0.78rem;
            letter-spacing: 0.08em;
            text-transform: uppercase;
            color: var(--muted-text-color);
            margin-bottom: var(--spacing-xs);
        }}
        .kpi-strip__value {{
            font-size: 1.4rem;
            font-weight: 700;
            color: var(--primary-color);
            font-family: var(--numeric-font-family);
            font-variant-numeric: tabular-nums;
            font-feature-settings: 'tnum';
        }}
        .kpi-strip__delta {{
            font-size: 0.85rem;
            font-weight: 600;
            color: var(--muted-text-color);
        }}
        .dashboard-meta {{
            display: flex;
            flex-wrap: wrap;
            gap: var(--spacing-xs);
            margin-bottom: var(--spacing-md);
            background: var(--surface-color);
            border-radius: var(--radius-card);
            border: 1px solid var(--border-subtle-color);
            padding: 0.75rem 1rem;
            position: sticky;
            top: 4.75rem;
            z-index: 60;
            box-shadow: var(--shadow-sm);
        }}
        .dashboard-meta__chip {{
            padding: 0.4rem 0.85rem;
            border-radius: var(--radius-chip);
            background: var(--surface-tint-color);
            color: var(--text-color);
            font-size: 0.8rem;
            font-weight: 600;
            display: inline-flex;
            align-items: center;
            gap: 0.45rem;
        }}
        .dashboard-meta__chip .ui-icon {{
            background: var(--accent-color);
            color: #ffffff;
        }}
        .dashboard-meta--empty {{
            padding: 0.4rem 0.85rem;
            border-radius: var(--radius-chip);
            background: var(--surface-tint-color);
            color: var(--muted-text-color);
            font-size: 0.8rem;
            font-weight: 500;
            display: inline-flex;
            align-items: center;
            gap: 0.35rem;
        }}
        .dashboard-filter-chips-anchor + div[data-testid="stHorizontalBlock"] {{
            display: flex;
            flex-wrap: wrap;
            gap: var(--spacing-xs);
            margin-bottom: var(--spacing-md);
        }}
        .dashboard-filter-chips-anchor + div[data-testid="stHorizontalBlock"] > div[data-testid="column"] {{
            width: auto !important;
            flex: 0 0 auto;
        }}
        .dashboard-filter-chips-anchor + div[data-testid="stHorizontalBlock"] div[data-testid="stButton"] {{
            margin: 0;
        }}
        .dashboard-filter-chips-anchor + div[data-testid="stHorizontalBlock"] button {{
            padding: 0.35rem 0.85rem;
            border-radius: var(--radius-chip);
            background: var(--surface-tint-color);
            color: var(--text-color);
            font-size: 0.8rem;
            font-weight: 600;
            border: 1px solid var(--border-strong-color);
            cursor: pointer;
            transition: background-color 0.15s ease, box-shadow 0.15s ease, border-color 0.15s ease;
        }}
        .dashboard-filter-chips-anchor + div[data-testid="stHorizontalBlock"] button:hover {{
            background: var(--surface-tint-color);
            border-color: var(--border-strong-color);
            box-shadow: 0 2px 6px rgba(15,23,42,0.16);
        }}
        .dashboard-filter-chips-anchor + div[data-testid="stHorizontalBlock"] button:active {{
            background: var(--surface-tint-color);
            border-color: var(--border-strong-color);
        }}
        div[data-testid="stMetric"] {{
            background: var(--surface-color);
            border-radius: var(--radius-card);
            padding: 1.25rem 1.5rem;
            border: 1px solid var(--border-subtle-color);
            box-shadow: var(--shadow-md);
        }}
        div[data-testid="stMetricLabel"] {{
            color: var(--muted-text-color);
            font-size: 0.8rem;
            letter-spacing: 0.02em;
        }}
        div[data-testid="stMetricValue"] {{
            font-family: var(--numeric-font-family);
            font-variant-numeric: tabular-nums;
            font-feature-settings: 'tnum';
            color: var(--text-color);
            font-size: calc(1.4rem * var(--font-scale));
        }}
        div[data-testid="stMetricDelta"] {{
            font-weight: 600;
            font-size: calc(0.9rem * var(--font-scale));
            display: flex;
            align-items: center;
        }}
        div[data-testid="stTabs"] > div[role="tablist"] {{
            position: sticky;
            top: 6.6rem;
            z-index: 50;
            background: var(--background-color);
            padding: 0.4rem 0.25rem;
            margin-bottom: var(--spacing-md);
            gap: var(--spacing-xs);
            border-bottom: 1px solid var(--border-subtle-color);
        }}
        div[data-testid="stTabs"] button[role="tab"] {{
            border-radius: 999px;
            padding: 0.55rem 1.15rem;
            font-weight: 600;
            border: 1px solid transparent;
            background: transparent;
            color: var(--text-color);
            transition: background-color 0.15s ease, color 0.15s ease, box-shadow 0.15s ease, border-color 0.15s ease;
        }}
        div[data-testid="stTabs"] button[role="tab"][aria-selected="true"] {{
            background: var(--accent-color);
            color: #ffffff;
            border-color: var(--accent-color);
            box-shadow: var(--shadow-md);
        }}
        div[data-testid="stTabs"] button[role="tab"][aria-selected="false"] {{
            background: var(--surface-tint-color);
            border-color: var(--border-subtle-color);
            color: var(--muted-text-color);
        }}
        div[data-testid="stTabs"] button[role="tab"]:focus-visible {{
            outline: 3px solid var(--accent-color);
            outline-offset: 2px;
        }}
        div[data-testid="stMetricDelta"] svg {{
            fill: currentColor !important;
        }}
        .status-pill {{
            display: inline-flex;
            align-items: center;
            gap: 0.35rem;
            padding: 0.35rem 0.75rem;
            border-radius: var(--radius-chip);
            font-size: 0.8rem;
            font-weight: 600;
        }}
        .status-pill .ui-icon {{
            background: rgba(255,255,255,0.25);
            color: inherit;
        }}
        .status-pill--ok {{
            background: var(--success-surface);
            color: var(--success-color);
        }}
        .status-pill--warning {{
            background: var(--warning-surface);
            color: var(--warning-color);
        }}
        .status-pill--error {{
            background: var(--error-surface);
            color: var(--error-color);
        }}
        .kgi-card__delta--up,
        .kpi-strip__delta--up {{
            color: var(--success-color);
        }}
        .kgi-card__delta--down,
        .kpi-strip__delta--down {{
            color: var(--error-color);
        }}
        .alert-banner {{
            border-radius: var(--radius-card);
            padding: 1rem 1.25rem;
            margin-bottom: var(--spacing-md);
            border: 1px solid transparent;
            background: var(--surface-color);
            color: var(--text-color);
        }}
        .alert-banner--warning {{
            background: var(--warning-surface);
            border-color: rgba(var(--warning-rgb), 0.35);
            color: var(--warning-color);
        }}
        .alert-banner__title {{
            display: inline-flex;
            align-items: center;
            gap: 0.45rem;
            font-weight: 700;
        }}
        .alert-banner__title .ui-icon {{
            color: #ffffff;
        }}
        .alert-banner--ok {{
            background: var(--success-surface);
            border-color: rgba(var(--success-rgb), 0.35);
            color: var(--success-color);
        }}
        .data-status-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(240px, 1fr));
            gap: var(--spacing-sm);
            margin-top: var(--spacing-sm);
        }}
        .data-status-card {{
            background: var(--surface-color);
            border-radius: var(--radius-card);
            padding: 1.25rem 1.4rem;
            border: 1px solid var(--border-subtle-color);
            box-shadow: var(--shadow-md);
        }}
        .data-status-card__status {{
            margin-top: 0.8rem;
            font-weight: 700;
            font-size: 0.85rem;
            display: inline-flex;
            align-items: center;
            gap: 0.4rem;
            padding: 0.35rem 0.75rem;
            border-radius: var(--radius-chip);
        }}
        .data-status-card__status .ui-icon {{
            background: currentColor;
            color: #ffffff;
        }}
        .data-status-card__status--ok {{
            background: var(--success-surface);
            color: var(--success-color);
        }}
        .data-status-card__status--warning {{
            background: var(--warning-surface);
            color: var(--warning-color);
        }}
        .data-status-card__status--error {{
            background: var(--error-surface);
            color: var(--error-color);
        }}
        section[data-testid="stSidebar"] {{
            background: var(--surface-color);
            color: var(--text-color);
            border-right: 1px solid var(--border-subtle-color);
        }}
        section[data-testid="stSidebar"] .sidebar-section {{
            background: var(--surface-color);
            border-radius: var(--radius-card);
            border: 1px solid var(--border-subtle-color);
            padding: 1.25rem 1.35rem;
            box-shadow: var(--shadow-sm);
            margin-bottom: var(--spacing-sm);
        }}
        section[data-testid="stSidebar"] .sidebar-section__status {{
            display: inline-flex;
            align-items: center;
            gap: 0.3rem;
            padding: 0.35rem 0.7rem;
            border-radius: var(--radius-chip);
            background: rgba(30,136,229,0.12);
            color: var(--accent-color);
            font-size: 0.8rem;
            font-weight: 600;
        }}
        .main-nav-block div[role="radiogroup"] {{
            display: flex;
            flex-wrap: wrap;
            gap: 0.6rem;
        }}
        .main-nav-block div[role="radiogroup"] label {{
            padding: 0.5rem 1.2rem;
            border-radius: var(--radius-chip);
            border: 1px solid var(--border-strong-color);
            background: var(--surface-color);
            font-weight: 600;
            color: var(--text-color);
        }}
        .main-nav-block div[role="radiogroup"] label[aria-checked="true"] {{
            background: var(--primary-color);
            color: #ffffff;
            border-color: var(--border-strong-color);
        }}
        .search-card input {{
            border-radius: var(--radius-input);
            border: 1px solid var(--border-strong-color);
            padding: 0.6rem 0.9rem;
        }}
        .stApp main .stButton>button,
        .stApp main .stDownloadButton>button {{
            border-radius: var(--radius-input);
            padding: 0.65rem 1.4rem;
            font-weight: 600;
            background: var(--accent-color);
            color: #ffffff;
            border: none;
            box-shadow: 0 12px 24px rgba(30,136,229,0.2);
        }}
        .stApp main .stButton>button:hover,
        .stApp main .stDownloadButton>button:hover {{
            filter: brightness(0.95);
        }}
        div[data-baseweb="tab-list"] {{
            gap: var(--spacing-xs);
            border-bottom: 1px solid var(--border-subtle-color);
            margin-bottom: var(--spacing-sm);
        }}
        div[data-baseweb="tab"] {{
            padding: 0.75rem 1.4rem;
            border-radius: var(--radius-chip) var(--radius-chip) 0 0;
            font-weight: 600;
            color: var(--muted-text-color);
        }}
        div[data-baseweb="tab"][aria-selected="true"] {{
            background: var(--surface-color);
            color: var(--primary-color);
            box-shadow: inset 0 -3px 0 var(--primary-color);
        }}
        input, select, textarea {{
            border-radius: var(--radius-input) !important;
            border: 1px solid var(--border-strong-color) !important;
            padding: 0.6rem 0.9rem !important;
            font-size: var(--body-size) !important;
        }}
        @media (max-width: 900px) {{
            main .block-container {{
                padding: 2rem 1.5rem 2.75rem;
            }}
            .kpi-strip {{
                grid-template-columns: 1fr;
            }}
            .hero-panel {{
                padding: 2rem 1.8rem;
            }}
        }}
        .onboarding-wizard {{
            border: 1px solid var(--border-strong-color);
            border-radius: var(--radius-card);
            padding: 1.25rem 1.1rem 1.4rem;
            margin-bottom: var(--spacing-sm);
            background: var(--surface-tint-color);
            box-shadow: var(--shadow-sm);
        }}
        .onboarding-wizard--sidebar {{
            margin-bottom: 0;
        }}
        .onboarding-wizard__summary {{
            list-style: none;
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 0.5rem;
            cursor: pointer;
            margin: 0 0 var(--spacing-xs);
        }}
        .onboarding-wizard__summary::-webkit-details-marker {{
            display: none;
        }}
        .onboarding-wizard__title {{
            font-size: 0.95rem;
            font-weight: 700;
        }}
        .onboarding-wizard__chevron {{
            display: inline-flex;
            align-items: center;
            justify-content: center;
            width: 1.4rem;
            height: 1.4rem;
            border-radius: 50%;
            background: rgba(var(--primary-rgb),0.12);
            color: var(--primary-color);
            font-size: 0.75rem;
            transition: transform 0.2s ease;
        }}
        .onboarding-wizard:not([open]) .onboarding-wizard__chevron {{
            transform: rotate(-90deg);
        }}
        .onboarding-wizard[open] .onboarding-wizard__chevron {{
            transform: rotate(0deg);
        }}
        section[data-testid="stSidebar"] button[data-testid="stBaseButton-primary"] {{
            width: 100%;
            justify-content: flex-start;
            gap: 0.6rem;
            background: var(--surface-tint-color);
            border: 1px solid var(--border-strong-color);
            color: var(--text-color);
            box-shadow: none;
            padding: 0.65rem 0.85rem;
            border-radius: var(--radius-card);
            font-weight: 600;
        }}
        section[data-testid="stSidebar"] button[data-testid="stBaseButton-primary"] p {{
            margin: 0;
            font-weight: 600;
        }}
        section[data-testid="stSidebar"] button[data-testid="stBaseButton-primary"]:hover {{
            background: var(--surface-tint-color);
            border-color: var(--border-strong-color);
        }}
        .onboarding-step {{
            border-radius: var(--radius-card);
            padding: 0.75rem 0.85rem;
            background: var(--surface-tint-color);
            border: 1px dashed var(--border-subtle-color);
            margin-bottom: var(--spacing-xs);
        }}
        .onboarding-step--done {{
            border-style: solid;
            background: rgba(var(--success-rgb),0.08);
            border-color: rgba(var(--success-rgb),0.4);
        }}
        .onboarding-step__header {{
            display: flex;
            align-items: center;
            gap: 0.5rem;
            font-weight: 600;
            margin-bottom: 0.35rem;
        }}
        .onboarding-step__badge {{
            display: inline-flex;
            align-items: center;
            justify-content: center;
            width: 1.6rem;
            height: 1.6rem;
            border-radius: 50%;
            background: var(--surface-tint-color);
            color: var(--text-color);
            font-size: 0.8rem;
            font-weight: 700;
        }}
        .onboarding-step--done .onboarding-step__badge {{
            background: rgba(var(--success-rgb),0.25);
            color: var(--success-color);
        }}
        .onboarding-step__title {{
            font-size: 0.85rem;
        }}
        .onboarding-step__desc {{
            font-size: 0.8rem;
            color: var(--muted-text-color);
            line-height: 1.5;
        }}
        .sidebar-disabled {{
            border-radius: var(--radius-card);
            border: 1px dashed var(--border-subtle-color);
            padding: 0.85rem 1rem;
            margin-bottom: var(--spacing-sm);
            background: var(--surface-tint-color);
            color: var(--muted-text-color);
            font-size: 0.82rem;
        }}
        .empty-dashboard {{
            border-radius: var(--radius-card);
            border: 1px dashed var(--border-subtle-color);
            padding: 2rem 2.25rem;
            text-align: center;
            background: var(--surface-tint-color);
            color: var(--muted-text-color);
            margin-bottom: var(--spacing-lg);
        }}
        .empty-dashboard__hint {{
            display: block;
            margin-top: 0.75rem;
            font-size: 0.85rem;
        }}
        .quick-tutorial {{
            border-radius: var(--radius-card);
            border: 1px solid var(--border-subtle-color);
            padding: 1.1rem 1.2rem;
            background: var(--surface-tint-color);
            box-shadow: var(--shadow-sm);
        }}
        .quick-tutorial__title {{
            font-weight: 700;
            margin-bottom: 0.6rem;
            font-size: 0.9rem;
        }}
        .quick-tutorial ol {{
            padding-left: 1.2rem;
            margin: 0;
        }}
        .quick-tutorial li {{
            font-size: 0.85rem;
            margin-bottom: 0.45rem;
        }}
        .sidebar-wizard-title {{
            font-weight: 700;
            font-size: 0.95rem;
            margin-bottom: 0.4rem;
        }}
        .wizard-file-item {{
            display: flex;
            align-items: center;
            justify-content: space-between;
            padding: 0.35rem 0;
            border-bottom: 1px dashed var(--border-subtle-color);
            font-size: 0.82rem;
        }}
        .wizard-file-item:last-of-type {{
            border-bottom: none;
        }}
        .wizard-file-item__name {{
            flex: 1;
            margin-right: 0.5rem;
            word-break: break-all;
        }}
        .wizard-file-item__size {{
            color: var(--muted-text-color);
            font-size: 0.75rem;
        }}
        .viewerBadge_container__1QSob,
        .viewerBadge_link__1S137,
        .styles_viewerBadge__1yB5S,
        [data-testid="stHeaderActionButton"],
        [data-testid="stAppDeployButton"] {{
            display: none !important;
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )


def apply_accessibility_enhancements() -> None:
    """主要コンポーネントにARIA属性を付与し支援技術対応を強化する。"""

    aria_back_to_top = json.dumps(
        translate("back_to_top_aria", default="ページの先頭へ戻る")
    )
    aria_badge = json.dumps(
        translate("streamlit_badge_aria", default="Streamlit公式サイト (外部リンク)")
    )

    components.html(
        f"""
        <script>
        (function() {{
            const doc = window.parent.document;
            const backToTop = doc.querySelector('.back-to-top a');
            const backToTopLabel = {aria_back_to_top};
            if (backToTop) {{
                backToTop.setAttribute('aria-label', backToTopLabel);
                backToTop.setAttribute('title', backToTopLabel);
            }}
            const badge = doc.querySelector('footer [data-testid="stFooter"] a[href*="streamlit.io"]');
            const badgeLabel = {aria_badge};
            if (badge) {{
                badge.setAttribute('aria-label', badgeLabel);
                badge.setAttribute('title', badgeLabel);
                badge.setAttribute('rel', 'noopener noreferrer');
                badge.setAttribute('target', '_blank');
                badge.setAttribute('data-external-label', badgeLabel);
                if (!badge.dataset.confirmBound) {{
                    badge.dataset.confirmBound = 'true';
                    badge.addEventListener('click', (event) => {{
                        const message = badgeLabel + ' を新しいタブで開きます。続行しますか？';
                        if (!window.confirm(message)) {{
                            event.preventDefault();
                            event.stopPropagation();
                        }}
                    }});
                }}
            }}
            doc.querySelectorAll('button[data-testid^="baseButton"]').forEach((button) => {{
                if (!button.getAttribute('aria-label')) {{
                    const text = button.textContent.trim();
                    if (text) {{
                        button.setAttribute('aria-label', text);
                    }}
                }}
            }});
        }})();
        </script>
        """,
        height=0,
    )


def remember_last_uploaded_files(
    uploaded_sales: Dict[str, Any],
    cost_file: Any,
    subscription_file: Any,
) -> None:
    """最新のアップロードファイル名をセッションに保存する。"""

    file_names: List[str] = []

    for files in uploaded_sales.values():
        if isinstance(files, list):
            for file in files:
                if file is not None and hasattr(file, "name"):
                    file_names.append(file.name)
        elif files is not None and hasattr(files, "name"):
            file_names.append(files.name)

    for extra in (cost_file, subscription_file):
        if isinstance(extra, list):
            for file in extra:
                if file is not None and hasattr(file, "name"):
                    file_names.append(file.name)
        elif extra is not None and hasattr(extra, "name"):
            file_names.append(extra.name)

    if file_names:
        unique_names = list(dict.fromkeys(file_names))
        st.session_state["last_uploaded"] = unique_names


def _ensure_upload_directory(category: str) -> Path:
    """Ensure that the storage directory for a given upload category exists."""

    root = UPLOAD_STORAGE_ROOT
    root.mkdir(parents=True, exist_ok=True)
    target = root / category
    target.mkdir(parents=True, exist_ok=True)
    return target


def _build_upload_signature(uploaded_file: Any, category: str) -> Optional[str]:
    """Generate a stable signature for an uploaded file."""

    if uploaded_file is None:
        return None
    name = getattr(uploaded_file, "name", None)
    if not name:
        return None
    safe_name = Path(str(name)).name
    size = getattr(uploaded_file, "size", None)
    if size is None:
        try:
            uploaded_file.seek(0, io.SEEK_END)
            size = uploaded_file.tell()
        except Exception:
            size = None
        finally:
            try:
                uploaded_file.seek(0)
            except Exception:
                pass
    return f"{category}:{safe_name}:{size}"


def persist_uploaded_artifact(
    uploaded_file: Any,
    *,
    category: str,
    label: Optional[str] = None,
) -> Tuple[Optional[Path], bool, Optional[str]]:
    """Persist an uploaded file to disk with a progress indicator."""

    signature = _build_upload_signature(uploaded_file, category)
    if signature is None:
        return None, False, None

    index = st.session_state.setdefault("_persisted_upload_index", {})
    existing_path = index.get(signature)
    if existing_path and Path(existing_path).exists():
        metadata_index = st.session_state.setdefault("_persisted_upload_metadata", {})
        if signature not in metadata_index:
            metadata_index[signature] = {
                "category": category,
                "name": Path(existing_path).name,
                "label": label or Path(existing_path).name,
                "path": existing_path,
                "saved_at": datetime.utcnow().isoformat(timespec="seconds"),
                "size": getattr(uploaded_file, "size", None),
            }
        return Path(existing_path), False, signature

    destination_dir = _ensure_upload_directory(category)
    safe_name = Path(getattr(uploaded_file, "name", "uploaded"))
    timestamp = datetime.utcnow().strftime("%Y%m%d-%H%M%S")
    digest = hashlib.sha256(signature.encode("utf-8")).hexdigest()[:10]
    destination_path = destination_dir / f"{timestamp}_{digest}_{safe_name.name}"

    total = getattr(uploaded_file, "size", None)
    progress_placeholder = st.empty()
    display_label = label or safe_name.name
    progress = progress_placeholder.progress(0.0, text=f"「{display_label}」を保存しています…")
    bytes_written = 0

    try:
        with destination_path.open("wb") as output:
            while True:
                chunk = uploaded_file.read(1024 * 1024)
                if not chunk:
                    break
                output.write(chunk)
                bytes_written += len(chunk)
                if total:
                    progress_value = min(bytes_written / total, 0.995)
                    progress.progress(progress_value, text=f"「{display_label}」を保存しています…")
        progress.progress(1.0, text=f"「{display_label}」の保存が完了しました")
    except Exception as exc:  # pragma: no cover - runtime safeguard
        progress_placeholder.empty()
        try:
            destination_path.unlink(missing_ok=True)
        except Exception:
            pass
        try:
            uploaded_file.seek(0)
        except Exception:
            pass
        st.error(f"ファイルの保存に失敗しました: {exc}")
        return None, False, signature
    finally:
        progress_placeholder.empty()

    try:
        uploaded_file.seek(0)
    except Exception:
        pass

    index[signature] = str(destination_path)
    entry = {
        "category": category,
        "label": display_label,
        "name": safe_name.name,
        "path": str(destination_path),
        "saved_at": datetime.utcnow().isoformat(timespec="seconds"),
        "size": getattr(uploaded_file, "size", None),
    }
    metadata_index = st.session_state.setdefault("_persisted_upload_metadata", {})
    metadata_index[signature] = entry

    history = st.session_state.setdefault("upload_history", [])
    history.insert(0, entry.copy())
    st.session_state["upload_history"] = history[:MAX_UPLOAD_HISTORY]

    st.success(f"「{display_label}」を保存しました。保存先: {destination_path}")
    return destination_path, True, signature


def update_upload_metadata(signature: Optional[str], **updates: Any) -> None:
    """Update stored metadata for a persisted upload."""

    if not signature:
        return
    metadata_index = st.session_state.setdefault("_persisted_upload_metadata", {})
    entry = metadata_index.get(signature)
    if not entry:
        return
    for key, value in updates.items():
        if value is not None:
            entry[key] = value
    metadata_index[signature] = entry

    for record in st.session_state.get("upload_history", []):
        if record.get("path") == entry.get("path"):
            for key, value in updates.items():
                if value is not None:
                    record[key] = value


def list_saved_uploads(category: str) -> List[Dict[str, Any]]:
    """Return saved upload entries filtered by category."""

    history = st.session_state.get("upload_history", [])
    return [entry for entry in history if entry.get("category") == category]


def load_saved_upload(entry: Dict[str, Any]) -> Optional[pd.DataFrame]:
    """Read a saved upload from disk into a DataFrame."""

    path = Path(entry.get("path", ""))
    if not path.exists():
        st.warning(f"保存済みファイル『{entry.get('name', path.name)}』が見つかりません。")
        return None

    suffix = path.suffix.lower()
    try:
        if suffix == ".csv":
            return pd.read_csv(path)
        if suffix in {".xlsx", ".xls"}:
            return pd.read_excel(path)
        if suffix == ".json":
            return pd.read_json(path)
    except Exception as exc:  # pragma: no cover - runtime safeguard
        st.error(f"保存済みファイルの読み込みに失敗しました: {exc}")
        return None

    st.warning("現在のところ、このファイル形式の再利用には対応していません。")
    return None


def _make_json_serializable(value: Any) -> Any:
    """Streamlitのセッション状態から取り出した値をJSONシリアライズ可能な型に変換する。"""

    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, np.generic):
        return value.item()
    if isinstance(value, (date, datetime)):
        return value.isoformat()
    if isinstance(value, Path):
        return str(value)
    if isinstance(value, dict):
        return {str(key): _make_json_serializable(val) for key, val in value.items()}
    if isinstance(value, (list, tuple, set)):
        return [_make_json_serializable(item) for item in value]
    tolist = getattr(value, "tolist", None)
    if callable(tolist):  # numpy配列などをPythonのリストに変換
        return _make_json_serializable(tolist())
    return str(value)


def collect_dashboard_settings() -> Dict[str, Any]:
    """現在のフィルタやシナリオ設定をJSON化して返す。"""

    filters = {
        key: st.session_state.get(FILTER_STATE_KEYS[key])
        for key in ("store", "channels", "categories", "period", "freq")
    }
    manual_inputs = {
        key: st.session_state.get(key)
        for key in MANUAL_INPUT_FORMATS.keys()
    }
    scenarios = st.session_state.get("scenario_inputs", [])
    scenario_defaults = st.session_state.get("scenario_form_defaults", {})
    saved_sales = st.session_state.get("active_saved_sales", [])
    settings = {
        "filters": filters,
        "manual_inputs": manual_inputs,
        "scenarios": scenarios,
        "scenario_defaults": scenario_defaults,
        "saved_sales_paths": saved_sales,
    }
    return {key: _make_json_serializable(value) for key, value in settings.items()}


def apply_dashboard_settings(settings: Dict[str, Any]) -> None:
    """保存された設定をセッションに反映する。"""

    filters = settings.get("filters", {})
    for name in ("store", "channels", "categories", "period", "freq"):
        if name in filters:
            set_state_and_widget(FILTER_STATE_KEYS[name], filters[name])

    manual_inputs = settings.get("manual_inputs", {})
    for key, value in manual_inputs.items():
        if value is None or key not in MANUAL_INPUT_FORMATS:
            continue
        config = MANUAL_INPUT_FORMATS[key]
        st.session_state[key] = value
        spinner_key = f"{key}_spinner"
        text_key = f"{key}_text"
        st.session_state[spinner_key] = value
        st.session_state[text_key] = _format_numeric_text(
            value,
            is_integer=config.get("is_integer", False),
            text_format=config.get("text_format"),
        )

    scenarios = settings.get("scenarios")
    if isinstance(scenarios, list):
        st.session_state["scenario_inputs"] = scenarios

    defaults = settings.get("scenario_defaults")
    if isinstance(defaults, dict):
        st.session_state["scenario_form_defaults"] = defaults

    saved_sales_paths = settings.get("saved_sales_paths")
    if isinstance(saved_sales_paths, (list, tuple, set)):
        st.session_state["active_saved_sales"] = list(saved_sales_paths)


@st.cache_data(ttl=24 * 60 * 60)
def load_sample_data() -> Tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    """サンプルデータをキャッシュして高速に提供する。"""

    sales = generate_sample_sales_data()
    if len(sales) > 3000:
        sales = sales.head(3000).copy()
    return (
        sales,
        generate_sample_cost_data(),
        generate_sample_subscription_data(),
    )


def ensure_sample_data_cached() -> None:
    """初回アクセス時にサンプルデータのキャッシュを温める。"""

    if st.session_state.get("sample_data_warmed"):
        return

    with st.spinner("サンプルデータを初期化しています…"):
        sales, _, _ = load_sample_data()
    st.session_state["sample_data_warmed"] = True
    st.session_state["sample_data_rows"] = int(len(sales))


def _build_sample_filename(prefix: str, key: str) -> str:
    """アップロードカードで使うサンプルCSVのファイル名を生成する。"""

    digest = hashlib.md5(key.encode("utf-8")).hexdigest()[:8]
    return f"{prefix}_{digest}.csv"


def get_sample_sales_template(channel: Optional[str] = None, limit: int = 200) -> pd.DataFrame:
    """チャネル別に整形したサンプル売上フォーマットを返す。"""

    sales_df, _, _ = load_sample_data()
    sample = sales_df.copy()
    if channel:
        filtered = sample[sample["channel"] == channel].copy()
        sample = filtered if not filtered.empty else sample
    if limit:
        sample = sample.head(limit)
    columns = [
        "order_date",
        "channel",
        "store",
        "product_code",
        "product_name",
        "category",
        "quantity",
        "sales_amount",
        "customer_id",
        "campaign",
    ]
    existing = [col for col in columns if col in sample.columns]
    return sample[existing] if existing else sample


def get_sample_cost_template() -> pd.DataFrame:
    """原価率データのサンプルフォーマットを返す。"""

    _, cost_df, _ = load_sample_data()
    return cost_df.copy()


def get_sample_subscription_template() -> pd.DataFrame:
    """サブスク/KPIデータのサンプルフォーマットを返す。"""

    _, _, subscription_df = load_sample_data()
    sample = subscription_df.copy()
    if "month" in sample.columns:
        sample["month"] = sample["month"].astype(str)
    return sample


def get_plan_sales_template() -> pd.DataFrame:
    """事業計画ウィザード向けの売上CSVテンプレートを返す。"""

    sample_rows = [
        {"項目": "自社サイト売上", "月次売上": 1_200_000, "チャネル": "自社サイト"},
        {"項目": "楽天市場売上", "月次売上": 950_000, "チャネル": "楽天市場"},
        {"項目": "Amazon売上", "月次売上": 780_000, "チャネル": "Amazon"},
    ]
    return pd.DataFrame(sample_rows, columns=SALES_PLAN_COLUMNS)


def get_plan_expense_template() -> pd.DataFrame:
    """事業計画ウィザード向けの経費CSVテンプレートを返す。"""

    sample_rows = [
        {"費目": "人件費", "月次金額": 600_000, "区分": "固定費"},
        {"費目": "家賃", "月次金額": 200_000, "区分": "固定費"},
        {"費目": "広告宣伝費", "月次金額": 180_000, "区分": "変動費"},
    ]
    return pd.DataFrame(sample_rows, columns=EXPENSE_PLAN_COLUMNS)


def render_onboarding_wizard(
    container: Any,
    *,
    data_loaded: bool,
    filters_ready: bool,
    analysis_ready: bool,
    sample_checked: bool,
    visible: bool = True,
) -> None:
    """サイドバー上部にオンボーディングウィザードを描画する。"""

    desired_sample_state = bool(sample_checked)
    if st.session_state.get("use_sample_data") != desired_sample_state:
        set_state_and_widget("use_sample_data", desired_sample_state)

    if not visible:
        container.empty()
        return

    container.empty()
    wizard_box = container.container()

    step1_class = "onboarding-step onboarding-step--done" if data_loaded else "onboarding-step"
    step2_class = (
        "onboarding-step onboarding-step--done"
        if data_loaded and filters_ready
        else "onboarding-step"
    )
    step3_class = (
        "onboarding-step onboarding-step--done"
        if data_loaded and analysis_ready
        else "onboarding-step"
    )

    step1_badge = "✓" if data_loaded else "1"
    step2_badge = "✓" if data_loaded and filters_ready else "2"
    step3_badge = "✓" if data_loaded and analysis_ready else "3"

    wizard_box.markdown(
        f"""
        <details class="onboarding-wizard onboarding-wizard--sidebar" open>
            <summary class="onboarding-wizard__summary">
                <span class="onboarding-wizard__title">セットアップ手順</span>
                <span class="onboarding-wizard__chevron">⌄</span>
            </summary>
            <div class="{step1_class}">
                <div class="onboarding-step__header">
                    <span class="onboarding-step__badge">{step1_badge}</span>
                    <span class="onboarding-step__title">ステップ1: データを読み込む</span>
                </div>
                <div class="onboarding-step__desc">
                    サンプルデータを読み込むか、自社のCSV/Excelをアップロードしてダッシュボードを起動します。
                </div>
            </div>
            <div class="{step2_class}">
                <div class="onboarding-step__header">
                    <span class="onboarding-step__badge">{step2_badge}</span>
                    <span class="onboarding-step__title">ステップ2: フィルタを設定</span>
                </div>
                <div class="onboarding-step__desc">
                    店舗・期間・チャネルのフィルタを選択すると、分析対象が絞り込まれます。
                </div>
            </div>
            <div class="{step3_class}">
                <div class="onboarding-step__header">
                    <span class="onboarding-step__badge">{step3_badge}</span>
                    <span class="onboarding-step__title">ステップ3: ダッシュボード閲覧</span>
                </div>
                <div class="onboarding-step__desc">
                    KPIや資金繰り、シナリオ分析タブで意思決定に必要な示唆を確認しましょう。
                </div>
            </div>
        </details>
        """,
        unsafe_allow_html=True,
    )

    sample_widget_key = ensure_widget_mirror("use_sample_data")

    def _toggle_sample_checkbox() -> None:
        update_state_from_widget("use_sample_data")
        trigger_rerun()

    wizard_box.checkbox(
        "サンプルデータを使用して試す",
        value=bool(st.session_state.get("use_sample_data", True)),
        key=sample_widget_key,
        help="チェックを外すとアップロードした実データのみでダッシュボードを構成します。",
        on_change=_toggle_sample_checkbox,
    )

    use_sample = bool(st.session_state.get("use_sample_data", True))

    if not data_loaded:
        if wizard_box.button("サンプルデータを読み込む", key="wizard_load_sample_button"):
            set_state_and_widget("use_sample_data", True)
            st.session_state.pop("sample_data_warmed", None)
            st.session_state.pop("sample_data_rows", None)
            trigger_rerun()
    else:
        warmed_rows = st.session_state.get("sample_data_rows")
        if use_sample and warmed_rows:
            wizard_box.success(f"サンプルデータ {warmed_rows:,} 行を読み込み済みです。")
        elif use_sample:
            wizard_box.success("サンプルデータを読み込み済みです。")
        else:
            wizard_box.info("アップロードしたデータを表示しています。")

    wizard_box.caption("自社データはサイドバー下部のアップロードセクションから追加できます。")


def render_sidebar_disabled_placeholder() -> None:
    """データ未投入時にフィルタのプレースホルダを表示する。"""

    st.sidebar.markdown(
        "<div class='sidebar-disabled' title='データを読み込んでください'>データを読み込むと店舗やチャネルのフィルタが利用できます。</div>",
        unsafe_allow_html=True,
    )


def render_empty_dashboard_placeholder() -> None:
    """データがない場合のメイン画面プレースホルダを表示する。"""

    st.markdown(
        """
        <div class="empty-dashboard">
            データがまだ読み込まれていません。
            <span class="empty-dashboard__hint">サイドバーの「はじめに」でサンプルデータを読み込むか、売上ファイルをアップロードしてください。</span>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_intro_section() -> None:
    """アプリの目的と操作ガイドをまとめた導入セクションを表示する。"""

    st.title("意思決定支援ダッシュボード")
    st.caption(
        "くらしいきいき社の経営データを統合し、現場の意思決定を素早く支援するためのダッシュボードです。"
    )

    lead_col, media_col = st.columns([2, 1])
    with lead_col:
        st.markdown(
            """
            **初めての方へ**: サイドバーからサンプルデータを読み込むか、自社の売上・コストデータをアップロードすると、
            売上推移や粗利、資金繰りまで一気に可視化できます。経営会議で使える示唆を最小限の操作で得られるように設計されています。
            """
        )
        st.markdown(
            """
            1. 左側のサイドバーで対象期間とチャネルを選択します。
            2. 売上や固定費などを設定すると、主要KPI・キャッシュフロー・アラートが自動で更新されます。
            3. 必要に応じてタブで詳細分析（売上 / 粗利 / 在庫 / 資金 / KPI）を切り替えてください。
            """
        )
        with st.expander("このアプリで実現できること", expanded=False):
            st.markdown(
                """
                - **意思決定の高速化**: 店舗やチャネル別の売上と粗利をリアルタイムに把握できます。
                - **ボトルネックの特定**: KPIアラートで粗利率や在庫回転日数の悪化を早期に察知します。
                - **経営シナリオの比較**: 固定費や投資額を調整してPLシミュレーションを即座に確認できます。
                """
            )

    with media_col:
        st.markdown(
            """
            <div class="quick-tutorial">
                <div class="quick-tutorial__title">60秒クイックガイド</div>
                <ol>
                    <li>サイドバーの「はじめに」でサンプルデータを読み込むか、自社データをアップロードします。</li>
                    <li>対象となる店舗・期間・チャネルを選択して指標を絞り込みます。</li>
                    <li>ダッシュボードの各タブで売上/KPI/資金繰りを確認し、アラートに基づいて次のアクションを検討します。</li>
                </ol>
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.caption("動画の代わりに主要な操作手順をテキストで確認できます。")

    st.markdown("---")


def load_data(
    use_sample: bool,
    uploaded_sales: Dict[str, List],
    cost_file,
    subscription_file,
    *,
    automated_sales: Optional[Dict[str, pd.DataFrame]] = None,
    automated_reports: Optional[List[ValidationReport]] = None,
) -> Dict[str, Any]:
    """アップロード状況に応じてデータを読み込む。"""
    # TODO: アップロードされたExcelファイルを読み込んでデータフレームに統合
    sales_frames: List[pd.DataFrame] = []
    cost_frames: List[pd.DataFrame] = []
    subscription_frames: List[pd.DataFrame] = []
    validation_report = ValidationReport()

    if use_sample:
        sample_sales, sample_cost, sample_subscription = load_sample_data()
        sales_frames.append(sample_sales)
        cost_frames.append(sample_cost)
        subscription_frames.append(sample_subscription)

    loaded_sales, uploaded_validation = load_sales_files(uploaded_sales)
    validation_report.extend(uploaded_validation)
    if not loaded_sales.empty:
        sales_frames.append(loaded_sales)

    if cost_file:
        cost_frames.append(load_cost_workbook(cost_file))
    if subscription_file:
        subscription_frames.append(load_subscription_workbook(subscription_file))

    if automated_sales:
        for df in automated_sales.values():
            if df is not None and not df.empty:
                sales_frames.append(df)

    if automated_reports:
        for report in automated_reports:
            validation_report.extend(report)

    active_saved_sales = st.session_state.get("active_saved_sales", set())
    if isinstance(active_saved_sales, list):
        active_saved_sales = set(active_saved_sales)
    saved_sales_cache: Dict[str, Dict[str, Any]] = st.session_state.setdefault(
        "_saved_sales_cache", {}
    )
    cleaned_paths: set = set()
    for path_str in list(active_saved_sales):
        path = Path(path_str)
        try:
            stat_info = path.stat()
        except OSError:
            validation_report.add_message(
                "warning",
                f"保存済みファイル {path_str} が見つからなかったため読み込みをスキップしました。",
            )
            continue

        current_mtime = stat_info.st_mtime
        cached_entry = saved_sales_cache.get(path_str)
        if (
            cached_entry
            and cached_entry.get("mtime") == current_mtime
            and cached_entry.get("data") is not None
        ):
            saved_df = cached_entry.get("data")
            saved_report = cached_entry.get("report")
        else:
            with path.open("rb") as handle:
                saved_df, saved_report = load_sales_workbook(handle)
            saved_sales_cache[path_str] = {
                "data": saved_df,
                "report": saved_report,
                "mtime": current_mtime,
                "size": stat_info.st_size,
            }

        if isinstance(saved_report, ValidationReport):
            validation_report.extend(saved_report)
        if isinstance(saved_df, pd.DataFrame) and not saved_df.empty:
            sales_frames.append(saved_df)
            cleaned_paths.add(path_str)

    stale_cache_keys = set(saved_sales_cache.keys()) - cleaned_paths
    for stale_path in stale_cache_keys:
        saved_sales_cache.pop(stale_path, None)
    st.session_state["active_saved_sales"] = list(cleaned_paths)

    sales_df = pd.concat(sales_frames, ignore_index=True) if sales_frames else pd.DataFrame()
    cost_df = pd.concat(cost_frames, ignore_index=True) if cost_frames else pd.DataFrame()
    subscription_df = pd.concat(subscription_frames, ignore_index=True) if subscription_frames else pd.DataFrame()

    if not sales_df.empty:
        combined_duplicates = detect_duplicate_rows(sales_df)
        if not combined_duplicates.empty:
            before = len(validation_report.duplicate_rows)
            validation_report.add_duplicates(combined_duplicates)
            if len(validation_report.duplicate_rows) > before:
                validation_report.add_message(
                    "warning",
                    f"全チャネルの売上データで重複しているレコードが{len(combined_duplicates)}件検出されました。",
                    count=int(combined_duplicates.shape[0]),
                )

    return {
        "sales": sales_df,
        "cost": cost_df,
        "subscription": subscription_df,
        "sales_validation": validation_report,
    }


def apply_filters(
    sales_df: pd.DataFrame,
    channels: List[str],
    date_range: List[date],
    categories: Optional[List[str]] = None,
    stores: Optional[List[str]] = None,
) -> pd.DataFrame:
    """サイドバーで選択した条件をもとに売上データを抽出する。"""
    if sales_df.empty:
        return sales_df

    filtered = sales_df.copy()
    if stores and "store" in filtered.columns:
        if isinstance(stores, (str, bytes)):
            stores = [stores]
        filtered = filtered[filtered["store"].isin(stores)]
    if channels:
        filtered = filtered[filtered["channel"].isin(channels)]
    if categories:
        filtered = filtered[filtered["category"].isin(categories)]
    if date_range:
        start_date = pd.to_datetime(date_range[0]) if date_range[0] else filtered["order_date"].min()
        end_date = pd.to_datetime(date_range[1]) if date_range[1] else filtered["order_date"].max()
        filtered = filtered[(filtered["order_date"] >= start_date) & (filtered["order_date"] <= end_date)]
    return filtered


def download_button_from_df(label: str, df: pd.DataFrame, filename: str) -> None:
    """データフレームをCSVとしてダウンロードするボタンを配置。"""
    if df is None or df.empty:
        return
    buffer = io.StringIO()
    df.to_csv(buffer, index=False)
    clicked = st.download_button(label, buffer.getvalue(), file_name=filename, mime="text/csv")
    if clicked:
        display_state_message("csv_done", action_key=f"csv_done_{filename}")


def download_excel_from_df(label: str, df: pd.DataFrame, filename: str) -> None:
    """Excel形式でデータフレームをダウンロードするボタンを表示する。"""

    if df is None or df.empty:
        return

    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name="data", index=False)
    st.download_button(
        label,
        buffer.getvalue(),
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


def _dataframe_to_pdf(df: pd.DataFrame, title: str, subtitle: Optional[str] = None) -> bytes:
    """Convert DataFrame to PDF bytes using ReportLab."""

    if not REPORTLAB_AVAILABLE:
        raise RuntimeError("reportlab library is not available")

    working_df = df.copy()
    working_df = working_df.fillna("-")
    str_df = working_df.astype(str)
    data = [list(str_df.columns)] + str_df.values.tolist()

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=landscape(A4))
    styles = getSampleStyleSheet()
    elements: List[Any] = []
    elements.append(Paragraph(title, styles["Title"]))
    if subtitle:
        elements.append(Paragraph(subtitle, styles["Normal"]))
    elements.append(Spacer(1, 12))
    table = Table(data, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0b1f33")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#f5f7fa")),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#c3c8d4")),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
            ]
        )
    )
    elements.append(table)
    doc.build(elements)
    buffer.seek(0)
    return buffer.getvalue()


def download_pdf_from_df(
    label: str,
    df: pd.DataFrame,
    filename: str,
    *,
    title: str,
    subtitle: Optional[str] = None,
) -> None:
    """PDF形式でダウンロードできるボタンを表示する。"""

    if df is None or df.empty:
        return

    if not REPORTLAB_AVAILABLE:
        st.button(label, disabled=True, help="reportlabのインストールが必要です。")
        return

    try:
        pdf_bytes = _dataframe_to_pdf(df, title, subtitle)
    except Exception as exc:  # pragma: no cover - PDF生成失敗時の保護
        st.button(label, disabled=True, help=f"PDF生成でエラーが発生しました: {exc}")
        return

    st.download_button(
        label,
        pdf_bytes,
        file_name=filename,
        mime="application/pdf",
    )


def _dataframe_to_ppt(df: pd.DataFrame, title: str, subtitle: Optional[str] = None) -> bytes:
    """Create a simple PowerPoint slide containing a table representation of the DataFrame."""

    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx library is not available")

    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(14)

    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.8), Inches(9.0), Inches(5.0)).table

    for col_idx, column in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(column)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)

    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = "-" if pd.isna(value) else str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def download_ppt_from_df(
    label: str,
    df: pd.DataFrame,
    filename: str,
    *,
    title: str,
    subtitle: Optional[str] = None,
) -> None:
    """PowerPoint形式でダウンロードできるボタンを表示する。"""

    if df is None or df.empty:
        return

    if not PPTX_AVAILABLE:
        st.button(label, disabled=True, help="python-pptxのインストールが必要です。")
        return

    try:
        ppt_bytes = _dataframe_to_ppt(df, title, subtitle)
    except Exception as exc:  # pragma: no cover - PPT生成失敗時の保護
        st.button(label, disabled=True, help=f"PowerPoint生成でエラーが発生しました: {exc}")
        return

    st.download_button(
        label,
        ppt_bytes,
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


def encode_share_payload(payload: Dict[str, Any]) -> str:
    """Convert a payload to a URL-safe base64 token."""

    json_bytes = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    token = base64.urlsafe_b64encode(json_bytes).decode("utf-8")
    return token.rstrip("=")


def decode_share_payload(token: str) -> Optional[Dict[str, Any]]:
    """Decode a URL-safe base64 token back into a dictionary."""

    if not token:
        return None
    padding = "=" * (-len(token) % 4)
    try:
        decoded = base64.urlsafe_b64decode(token + padding)
        return json.loads(decoded.decode("utf-8"))
    except Exception:  # pragma: no cover - invalid share token recovery
        return None


def load_share_payload_from_query() -> None:
    """URLパラメータの共有トークンを読み込みセッションに格納する。"""

    params = st.experimental_get_query_params()
    token_list = params.get("share")
    if not token_list:
        return
    token = token_list[0]
    if not token or st.session_state.get("scenario_share_token_loaded") == token:
        return

    payload = decode_share_payload(token)
    if not payload:
        st.warning("共有リンクの読み込みに失敗しました。URLをご確認ください。")
        return

    st.session_state["scenario_share_token"] = token
    st.session_state["scenario_share_token_loaded"] = token
    st.session_state["scenario_share_payload"] = payload

    summary_records = payload.get("summary")
    if summary_records:
        try:
            st.session_state["shared_scenario_summary_df"] = pd.DataFrame(summary_records)
        except Exception:  # pragma: no cover - invalid payload failsafe
            st.session_state["shared_scenario_summary_df"] = None
    range_records = payload.get("range_summary")
    if range_records:
        try:
            st.session_state["shared_scenario_range_df"] = pd.DataFrame(range_records)
        except Exception:
            st.session_state["shared_scenario_range_df"] = None


def display_state_message(
    state: str,
    *,
    format_kwargs: Optional[Dict[str, Any]] = None,
    action: Optional[Callable[[], None]] = None,
    action_label: Optional[str] = None,
    action_key: Optional[str] = None,
    secondary_action: Optional[Callable[[], None]] = None,
    secondary_action_label: Optional[str] = None,
    secondary_action_key: Optional[str] = None,
    container: Optional[Any] = None,
) -> None:
    """状態に応じたフィードバックメッセージを表示する。"""

    config = STATE_MESSAGES.get(state)
    if not config:
        return

    target = container or st
    format_kwargs = format_kwargs or {}
    message_text = config["text"].format(**format_kwargs)
    message_type = config.get("type", "info")
    display_fn = getattr(target, message_type, target.info)
    display_fn(message_text)

    label = action_label or config.get("action_label")
    if label:
        button_key = action_key or config.get("action_key") or f"{state}_primary_action"
        primary_action = action or config.get("action")
        if target.button(label, key=button_key):
            if callable(primary_action):
                primary_action()

    secondary_label = secondary_action_label or config.get("secondary_action_label")
    if secondary_label:
        secondary_key = (
            secondary_action_key
            or config.get("secondary_action_key")
            or f"{state}_secondary_action"
        )
        secondary_callable = secondary_action or config.get("secondary_action")
        if target.button(secondary_label, key=secondary_key):
            if callable(secondary_callable):
                secondary_callable()


def _subtract_calendar_months(reference: date, months: int) -> date:
    """月単位で遡りつつ可能な限り日付を保つ。"""

    year = reference.year
    month = reference.month
    day = reference.day
    for _ in range(months):
        month -= 1
        if month == 0:
            month = 12
            year -= 1
        month_days = calendar.monthrange(year, month)[1]
        if day > month_days:
            day = month_days
    return date(year, month, day)


def suggest_default_period(min_date: date, max_date: date) -> Tuple[date, date]:
    """利用可能なデータから今日を上限に直近12か月（不足時は365日分）の推奨期間を返す。"""

    if min_date > max_date:
        min_date, max_date = max_date, min_date

    today = date.today()
    end_day = min(max_date, today)
    if end_day < min_date:
        end_day = min_date

    available_span_days = (end_day - min_date).days

    start_candidate: date
    if available_span_days >= 365:
        relativedelta_fn: Optional[Callable[..., Any]] = None
        relativedelta_spec = importlib.util.find_spec("dateutil.relativedelta")
        if relativedelta_spec is not None:
            relativedelta_module = importlib.import_module("dateutil.relativedelta")
            relativedelta_fn = getattr(relativedelta_module, "relativedelta", None)
        if callable(relativedelta_fn):
            start_candidate = end_day - relativedelta_fn(months=11)
        else:
            start_candidate = _subtract_calendar_months(end_day, 11)
    else:
        start_candidate = end_day - timedelta(days=365)

    start_day = max(min_date, start_candidate)
    if start_day > end_day:
        start_day = min_date

    return start_day, end_day


def reset_filters(defaults: Dict[str, Any]) -> None:
    """フィルタ関連のセッション状態を初期値に戻す。"""

    for key, value in defaults.items():
        set_state_and_widget(key, value)
    trigger_rerun()


def jump_to_section(section_key: str) -> None:
    """ナビゲーションの選択を強制的に切り替えてリロードする。"""

    if section_key not in NAV_OPTION_LOOKUP:
        return
    st.session_state["main_nav"] = section_key
    st.session_state["main_nav_display"] = NAV_OPTION_LOOKUP[section_key]
    trigger_rerun()


def build_filter_signature(
    store: Optional[str],
    channels: Optional[List[str]],
    categories: Optional[List[str]],
    date_range: Any,
    freq_label: str,
) -> Tuple[Any, ...]:
    """フィルタの状態を比較可能なタプルに変換する。"""

    if isinstance(date_range, (list, tuple)) and len(date_range) == 2:
        start_value, end_value = date_range
    else:
        start_value = end_value = date_range

    def _normalize_date(value: Any) -> Optional[str]:
        if value is None:
            return None
        if hasattr(value, "isoformat"):
            return value.isoformat()
        return str(value)

    return (
        store or "all",
        tuple(channels or []),
        tuple(categories or []),
        _normalize_date(start_value),
        _normalize_date(end_value),
        freq_label,
    )


def normalize_date_input(value: Any) -> Optional[date]:
    """様々な入力値をdate型に揃える。"""

    if value is None:
        return None
    if isinstance(value, date) and not isinstance(value, datetime):
        return value
    if isinstance(value, datetime):
        return value.date()
    if hasattr(value, "to_pydatetime"):
        try:
            return value.to_pydatetime().date()
        except Exception:
            return None
    try:
        return pd.to_datetime(value).date()
    except Exception:
        return None


def normalize_period_state_value(
    value: Any,
    min_date: date,
    max_date: date,
    default_period: Tuple[date, date],
) -> Tuple[date, date]:
    """セッション状態に保存されている期間情報を安全なタプルに正規化する。"""

    start_default, end_default = default_period
    start_default = max(min_date, min(max_date, start_default))
    end_default = max(min_date, min(max_date, end_default))

    start_candidate: Optional[date]
    end_candidate: Optional[date]

    if isinstance(value, (list, tuple)) and len(value) == 2:
        start_candidate = normalize_date_input(value[0])
        end_candidate = normalize_date_input(value[1])
    else:
        single_value = normalize_date_input(value)
        start_candidate = single_value
        end_candidate = single_value

    start = start_candidate or start_default
    end = end_candidate or end_default

    start = max(min_date, min(max_date, start))
    end = max(min_date, min(max_date, end))
    if start > end:
        start, end = end, start
        start = max(min_date, min(max_date, start))
        end = max(min_date, min(max_date, end))

    return start, end


def prepare_plan_table(
    data: Any,
    columns: List[str],
    numeric_columns: Optional[List[str]] = None,
) -> pd.DataFrame:
    """ウィザード用の表を指定の列構成と数値型に整形する。"""

    numeric_columns = numeric_columns or []
    if isinstance(data, pd.DataFrame):
        df = data.copy()
    elif data is None or (hasattr(data, "__len__") and len(data) == 0):
        df = pd.DataFrame(columns=columns)
    else:
        df = pd.DataFrame(data)

    if df.empty and not list(df.columns):
        df = pd.DataFrame(columns=columns)

    df.columns = [str(col).strip() for col in df.columns]
    for column in columns:
        if column not in df.columns:
            df[column] = 0.0 if column in numeric_columns else ""
    df = df[columns]
    for column in numeric_columns:
        df[column] = pd.to_numeric(df[column], errors="coerce").fillna(0.0)
    if columns:
        label_column = columns[0]
        df[label_column] = df[label_column].astype(str).str.strip()
    return df


def reset_plan_wizard_state() -> None:
    """経営計画ウィザードのセッション状態を初期化する。"""

    default_start = date.today().replace(day=1)
    st.session_state["plan_wizard"] = {
        "current_step": 0,
        "completed": False,
        "basic_info": {
            "company_name": "",
            "preparer": "",
            "fiscal_year_start": default_start,
            "plan_period_months": 12,
            "target_margin": 15.0,
            "strategic_focus": "",
        },
        "sales_table": pd.DataFrame(columns=SALES_PLAN_COLUMNS),
        "expense_table": pd.DataFrame(columns=EXPENSE_PLAN_COLUMNS),
        "sales_import_hash": None,
        "expense_import_hash": None,
        "sales_import_feedback": None,
        "expense_import_feedback": None,
        "metrics": {},
    }
    for key in [
        "plan_sales_editor",
        "plan_expense_editor",
        "plan_sales_common_select",
        "plan_expense_common_select",
    ]:
        st.session_state.pop(key, None)


def ensure_plan_wizard_state() -> Dict[str, Any]:
    """経営計画ウィザード用のセッション情報を返す。"""

    if "plan_wizard" not in st.session_state:
        reset_plan_wizard_state()

    state: Dict[str, Any] = st.session_state["plan_wizard"]
    state.setdefault("current_step", 0)
    state.setdefault("completed", False)
    state.setdefault("basic_info", {})

    basic_info = state["basic_info"]
    if not isinstance(basic_info.get("fiscal_year_start"), date):
        basic_info["fiscal_year_start"] = date.today().replace(day=1)
    basic_info.setdefault("company_name", "")
    basic_info.setdefault("preparer", "")
    basic_info.setdefault("plan_period_months", 12)
    basic_info.setdefault("target_margin", 15.0)
    basic_info.setdefault("strategic_focus", "")

    state["sales_table"] = prepare_plan_table(
        state.get("sales_table"), SALES_PLAN_COLUMNS, ["月次売上"]
    )
    state["expense_table"] = prepare_plan_table(
        state.get("expense_table"), EXPENSE_PLAN_COLUMNS, ["月次金額"]
    )
    state.setdefault("sales_import_hash", None)
    state.setdefault("expense_import_hash", None)
    state.setdefault("sales_import_feedback", None)
    state.setdefault("expense_import_feedback", None)
    state.setdefault("metrics", {})
    return state


def append_plan_rows(
    df: pd.DataFrame,
    label_column: str,
    numeric_column: str,
    default_values: Optional[Dict[str, Any]],
    items: List[str],
) -> Tuple[pd.DataFrame, int]:
    """プルダウンで選択した項目を既存の表に追加する。"""

    if df is None or df.empty:
        df = pd.DataFrame(columns=[label_column, numeric_column])

    existing = set(df[label_column].astype(str).str.strip())
    new_rows: List[Dict[str, Any]] = []
    for item in items:
        normalized = str(item).strip()
        if normalized and normalized not in existing:
            row = {label_column: normalized, numeric_column: 0.0}
            if default_values:
                for key, value in default_values.items():
                    row[key] = value
            new_rows.append(row)
            existing.add(normalized)

    if new_rows:
        df = pd.concat([df, pd.DataFrame(new_rows)], ignore_index=True)

    return df, len(new_rows)


def normalize_plan_import(
    df: pd.DataFrame,
    column_candidates: Dict[str, List[str]],
    required_columns: List[str],
    numeric_columns: List[str],
) -> pd.DataFrame:
    """CSV取り込み時に列名を標準化し、必要列を抽出する。"""

    if df is None or df.empty:
        raise ValueError("CSVにデータがありません。")

    working = df.copy()
    working.columns = [str(col).strip() for col in working.columns]
    rename_map: Dict[str, str] = {}
    for target, candidates in column_candidates.items():
        found = next((col for col in candidates if col in working.columns), None)
        if found:
            rename_map[found] = target

    missing = [col for col in required_columns if col not in rename_map.values()]
    if missing:
        raise ValueError(
            f"必要な列({', '.join(missing)})がCSV内に見つかりませんでした。列名を確認してください。"
        )

    normalized = working[list(rename_map.keys())].rename(columns=rename_map)
    label_column = required_columns[0]
    normalized = normalized.dropna(subset=[label_column])
    normalized[label_column] = normalized[label_column].astype(str).str.strip()
    normalized = normalized[normalized[label_column] != ""]
    for column in numeric_columns:
        normalized[column] = pd.to_numeric(normalized[column], errors="coerce").fillna(0.0)
    for target in column_candidates.keys():
        if target not in normalized.columns:
            normalized[target] = "" if target not in numeric_columns else 0.0
    return normalized


def import_plan_csv(
    file_bytes: bytes,
    column_candidates: Dict[str, List[str]],
    required_columns: List[str],
    numeric_columns: List[str],
) -> Tuple[pd.DataFrame, Optional[str]]:
    """会計ソフトからエクスポートしたCSVを標準形式に変換する。"""

    if not file_bytes:
        return pd.DataFrame(columns=required_columns), "CSVファイルが空です。"

    last_error: Optional[str] = None
    for encoding in ["utf-8-sig", "utf-8", "cp932", "shift_jis"]:
        try:
            text = file_bytes.decode(encoding)
            buffer = io.StringIO(text)
            raw_df = pd.read_csv(buffer)
            normalized = normalize_plan_import(
                raw_df, column_candidates, required_columns, numeric_columns
            )
            return normalized, None
        except UnicodeDecodeError:
            last_error = f"文字コード{encoding}での読み込みに失敗しました。"
            continue
        except pd.errors.ParserError:
            last_error = "CSVの解析に失敗しました。フォーマットを確認してください。"
            continue
        except ValueError as exc:
            return pd.DataFrame(columns=required_columns), str(exc)

    return pd.DataFrame(columns=required_columns), last_error or "CSVの読み込みに失敗しました。"


def calculate_plan_metrics_from_state(state: Dict[str, Any]) -> Dict[str, Any]:
    """売上・経費入力から主要な財務指標を算出する。"""

    sales_df = prepare_plan_table(state.get("sales_table"), SALES_PLAN_COLUMNS, ["月次売上"])
    expense_df = prepare_plan_table(
        state.get("expense_table"), EXPENSE_PLAN_COLUMNS, ["月次金額"]
    )
    state["sales_table"] = sales_df
    state["expense_table"] = expense_df

    info = state.get("basic_info", {})
    period_months = int(info.get("plan_period_months") or 0)
    monthly_sales = float(sales_df["月次売上"].sum()) if not sales_df.empty else 0.0
    monthly_expenses = float(expense_df["月次金額"].sum()) if not expense_df.empty else 0.0
    monthly_profit = monthly_sales - monthly_expenses
    margin = monthly_profit / monthly_sales if monthly_sales else np.nan
    target_margin_pct = float(info.get("target_margin") or 0.0)
    margin_gap_pct = (margin * 100 - target_margin_pct) if monthly_sales else np.nan

    metrics = {
        "monthly_sales": monthly_sales,
        "monthly_expenses": monthly_expenses,
        "monthly_profit": monthly_profit,
        "monthly_margin": margin,
        "annual_sales": monthly_sales * period_months,
        "annual_expenses": monthly_expenses * period_months,
        "annual_profit": monthly_profit * period_months,
        "target_margin_pct": target_margin_pct,
        "margin_gap_pct": margin_gap_pct,
        "period_months": period_months,
        "burn_rate": monthly_expenses - monthly_sales,
    }
    state["metrics"] = metrics
    return metrics


def build_plan_summary_df(metrics: Dict[str, Any]) -> pd.DataFrame:
    """計画の要約表を作成する。"""

    rows: List[Dict[str, Any]] = [
        {
            "指標": "売上",
            "月次計画額": metrics.get("monthly_sales", 0.0),
            "年間計画額": metrics.get("annual_sales", 0.0),
            "指標値": np.nan,
        },
        {
            "指標": "経費",
            "月次計画額": metrics.get("monthly_expenses", 0.0),
            "年間計画額": metrics.get("annual_expenses", 0.0),
            "指標値": np.nan,
        },
        {
            "指標": "営業利益",
            "月次計画額": metrics.get("monthly_profit", 0.0),
            "年間計画額": metrics.get("annual_profit", 0.0),
            "指標値": np.nan,
        },
        {
            "指標": "月次バーンレート (費用-売上)",
            "月次計画額": metrics.get("burn_rate", 0.0),
            "年間計画額": metrics.get("burn_rate", 0.0)
            * metrics.get("period_months", 0),
            "指標値": np.nan,
        },
    ]

    margin = metrics.get("monthly_margin")
    if margin is not None and np.isfinite(margin):
        rows.append(
            {
                "指標": "営業利益率",
                "月次計画額": np.nan,
                "年間計画額": np.nan,
                "指標値": margin * 100,
            }
        )

    margin_gap = metrics.get("margin_gap_pct")
    if margin_gap is not None and np.isfinite(margin_gap):
        rows.append(
            {
                "指標": "目標比差分 (pt)",
                "月次計画額": np.nan,
                "年間計画額": np.nan,
                "指標値": margin_gap,
            }
        )

    return pd.DataFrame(rows)


def compute_actual_reference(actual_sales: Optional[pd.DataFrame]) -> Dict[str, float]:
    """実績データから平均売上・利益などを算出して比較指標を返す。"""

    if actual_sales is None or actual_sales.empty:
        return {}
    if "order_date" not in actual_sales.columns or "sales_amount" not in actual_sales.columns:
        return {}

    working = actual_sales.copy()
    working["order_month"] = working["order_date"].dt.to_period("M")
    monthly_sales = working.groupby("order_month")["sales_amount"].sum()
    reference: Dict[str, float] = {}
    if not monthly_sales.empty:
        reference["monthly_sales_avg"] = float(monthly_sales.mean())

    profit_column = None
    if "net_gross_profit" in working.columns:
        profit_column = "net_gross_profit"
    elif "gross_profit" in working.columns:
        profit_column = "gross_profit"

    if profit_column:
        monthly_profit = working.groupby("order_month")[profit_column].sum()
        if not monthly_profit.empty:
            reference["monthly_profit_avg"] = float(monthly_profit.mean())
            sales_avg = reference.get("monthly_sales_avg")
            if sales_avg:
                reference["margin_avg"] = reference["monthly_profit_avg"] / sales_avg

    return reference


def validate_plan_basic_info(info: Dict[str, Any]) -> Tuple[bool, List[str], List[str]]:
    """基本情報入力の妥当性を確認する。"""

    errors: List[str] = []
    warnings: List[str] = []

    if not info.get("company_name", "").strip():
        errors.append("事業所名を入力してください。")
    if not isinstance(info.get("fiscal_year_start"), date):
        errors.append("計画開始月を選択してください。")

    period = int(info.get("plan_period_months") or 0)
    if period <= 0:
        errors.append("計画期間は1ヶ月以上を指定してください。")

    if not info.get("preparer", "").strip():
        warnings.append("作成担当者を入力すると共有がスムーズになります。")

    target_margin = float(info.get("target_margin") or 0.0)
    if target_margin < 0:
        errors.append("目標利益率は0%以上で設定してください。")
    elif target_margin > 80:
        warnings.append("目標利益率が高すぎる可能性があります。")

    return len(errors) == 0, errors, warnings


def validate_plan_sales(df: pd.DataFrame) -> Tuple[bool, List[str], List[str]]:
    """売上予測入力の妥当性を確認する。"""

    errors: List[str] = []
    warnings: List[str] = []

    if df is None or df.empty:
        errors.append("売上予測を1件以上入力してください。")
        return False, errors, warnings

    if "項目" not in df.columns or "月次売上" not in df.columns:
        errors.append("売上予測の列構成が不正です。")
        return False, errors, warnings

    empty_label = df["項目"].astype(str).str.strip() == ""
    if empty_label.any():
        errors.append("空欄の売上項目があります。名称を入力してください。")

    negative = df["月次売上"] < 0
    if negative.any():
        errors.append("売上金額は0以上で入力してください。")

    zero_rows = df["月次売上"] == 0
    if zero_rows.any():
        warnings.append("0円の売上項目があります。必要でなければ削除してください。")

    duplicates = df["項目"].astype(str).str.strip().duplicated()
    if duplicates.any():
        warnings.append("同名の売上項目が複数あります。集計が重複する可能性があります。")

    return len(errors) == 0, errors, warnings


def validate_plan_expenses(df: pd.DataFrame) -> Tuple[bool, List[str], List[str]]:
    """経費計画入力の妥当性を確認する。"""

    errors: List[str] = []
    warnings: List[str] = []

    if df is None or df.empty:
        errors.append("経費計画を1件以上入力してください。")
        return False, errors, warnings

    if "費目" not in df.columns or "月次金額" not in df.columns:
        errors.append("経費計画の列構成が不正です。")
        return False, errors, warnings

    empty_label = df["費目"].astype(str).str.strip() == ""
    if empty_label.any():
        errors.append("空欄の経費科目があります。名称を入力してください。")

    negative = df["月次金額"] < 0
    if negative.any():
        errors.append("経費金額は0以上で入力してください。")

    zero_rows = df["月次金額"] == 0
    if zero_rows.any():
        warnings.append("0円の経費項目があります。必要でなければ削除してください。")

    if "区分" in df.columns and (df["区分"].astype(str).str.strip() == "").any():
        warnings.append("区分が未選択の経費があります。")

    return len(errors) == 0, errors, warnings


def validate_plan_metrics(metrics: Dict[str, Any]) -> Tuple[bool, List[str], List[str]]:
    """財務指標計算ステップの妥当性を確認する。"""

    errors: List[str] = []
    warnings: List[str] = []

    if not metrics:
        errors.append("売上と経費の入力を完了してください。")
        return False, errors, warnings

    if metrics.get("monthly_sales", 0.0) <= 0:
        errors.append("売上予測が未入力または0円のため、指標を計算できません。")

    if metrics.get("monthly_expenses", 0.0) < 0:
        errors.append("経費金額が不正です。")

    if metrics.get("period_months", 0) <= 0:
        errors.append("計画期間を見直してください。")

    if (
        metrics.get("monthly_sales", 0.0) > 0
        and metrics.get("monthly_profit", 0.0) < 0
    ):
        warnings.append("月次営業利益がマイナスです。コスト構成を確認してください。")

    margin_gap = metrics.get("margin_gap_pct")
    if margin_gap is not None and np.isfinite(margin_gap) and margin_gap < 0:
        warnings.append("計画上の利益率が目標を下回っています。")

    return len(errors) == 0, errors, warnings


def render_instruction_popover(label: str, content: str) -> None:
    """ポップオーバーまたはエクスパンダーで操作ガイドを表示する。"""

    popover_fn = getattr(st, "popover", None)
    if callable(popover_fn):
        with popover_fn(label):
            st.markdown(content)
    else:
        with st.expander(label):
            st.markdown(content)


@contextmanager
def form_section(
    title: Optional[str],
    description: Optional[str] = None,
    *,
    tone: str = "primary",
) -> None:
    """フォーム入力をカード化し、余白と階層を整える。"""

    classes = ["form-section"]
    if tone and tone != "primary":
        classes.append(f"form-section--{tone}")

    with st.container():
        st.markdown(
            f"<div class='{ ' '.join(classes) }'>", unsafe_allow_html=True
        )
        if title:
            st.markdown(
                f"<div class='form-section__title'>{html.escape(title)}</div>",
                unsafe_allow_html=True,
            )
        if description:
            st.markdown(
                f"<p class='form-section__description'>{html.escape(description)}</p>",
                unsafe_allow_html=True,
            )
        try:
            yield
        finally:
            st.markdown("</div>", unsafe_allow_html=True)


def render_plan_stepper(current_step: int) -> None:
    """ウィザードの進行状況を視覚的なタイムラインで表示する。"""

    items: List[str] = []
    total_steps = len(PLAN_WIZARD_STEPS)
    for idx, step in enumerate(PLAN_WIZARD_STEPS):
        if idx < current_step:
            state_class = "stepper__item stepper__item--done"
            status = "完了"
        elif idx == current_step:
            state_class = "stepper__item stepper__item--active"
            status = "進行中"
        else:
            state_class = "stepper__item"
            status = "未着手"

        items.append(
            """
            <div class="{state_class}">
                <div class="stepper__index">{index}</div>
                <div class="stepper__body">
                    <div class="stepper__title">{title}</div>
                    <div class="stepper__desc">{description}</div>
                </div>
                <div class="stepper__status">{status}</div>
            </div>
            """.format(
                state_class=state_class,
                index=idx + 1,
                title=html.escape(step["title"]),
                description=html.escape(step.get("description", "")),
                status=status,
            )
        )

    st.markdown(
        f"<div class='stepper'>{''.join(items)}</div>",
        unsafe_allow_html=True,
    )


def render_plan_step_basic_info(state: Dict[str, Any]) -> Tuple[bool, Dict[str, Any]]:
    """ウィザードの基本情報入力ステップを描画する。"""

    state.setdefault("basic_info", {})
    info = state["basic_info"]
    render_instruction_popover(
        "基本情報の入力ガイド",
        """
- 会社名や担当者などの基本情報を入力します。
- 計画開始月と期間は年間換算の計算に利用されます。
- 目標利益率を設定すると達成状況のチェックが自動化されます。
""",
    )

    with st.form("plan_basic_info"):
        form_info: Dict[str, Any] = {}

        with form_section(
            "事業所と担当者",
            "共有時に識別される基本情報を先に押さえておきます。",
        ):
            form_info["company_name"] = st.text_input(
                "事業所名",
                value=info.get("company_name", ""),
                key="plan_company_name",
                help="経営計画書に記載する正式な社名または店舗名を入力してください。",
            )
            form_info["preparer"] = st.text_input(
                "作成担当者",
                value=info.get("preparer", ""),
                key="plan_preparer",
                help="計画の作成者または責任者を入力すると共有がスムーズになります。",
            )

        with form_section(
            "計画期間と利益目標",
            "期間と目標値は後続のシミュレーションに自動反映されます。",
        ):
            col1, col2 = st.columns(2)
            default_start = info.get("fiscal_year_start")
            if not isinstance(default_start, date):
                default_start = date.today().replace(day=1)
            form_info["fiscal_year_start"] = col1.date_input(
                "計画開始月",
                value=default_start,
                key="plan_fiscal_start",
                help="事業計画の初月を選択します。月次予測の起点として使用されます。",
            )

            period_default = int(info.get("plan_period_months") or 12)
            form_info["plan_period_months"] = col2.slider(
                "計画期間（月）",
                min_value=3,
                max_value=36,
                value=period_default,
                step=1,
                key="plan_period_months",
                help="3〜36ヶ月の範囲で計画期間を指定します。",
            )

            target_margin_default = float(info.get("target_margin") or 15.0)
            form_info["target_margin"] = col1.slider(
                "目標営業利益率(%)",
                min_value=0.0,
                max_value=50.0,
                value=target_margin_default,
                step=0.5,
                key="plan_target_margin",
                help="経営チームが目指す営業利益率を設定します。",
            )

        with form_section(
            "重点施策メモ",
            "将来の振り返りで意図を再確認できるよう、戦略メモを残せます。",
            tone="secondary",
        ):
            st.markdown(
                "<span class='form-section__status'>任意入力</span>",
                unsafe_allow_html=True,
            )
            form_info["strategic_focus"] = st.text_area(
                "重点施策メモ",
                value=info.get("strategic_focus", ""),
                key="plan_strategic_focus",
                help="成長戦略や重点施策をメモできます。後続ステップの指標と合わせて検討してください。",
            )

        st.caption(
            "段階的なウィザードと統一されたツールチップを用いたインターフェースは、Nielsen Norman Groupの調査 (moldstud.com) によればユーザー満足度を約20%向上させます。"
        )

        submitted = st.form_submit_button("基本情報を保存")

    return submitted, form_info


def render_plan_step_sales(state: Dict[str, Any], context: Dict[str, Any]) -> Tuple[bool, Dict[str, Any]]:
    """売上予測入力ステップを描画する。"""

    existing_table = prepare_plan_table(
        state.get("sales_table"), SALES_PLAN_COLUMNS, ["月次売上"]
    )
    if not isinstance(state.get("sales_table"), pd.DataFrame):
        state["sales_table"] = existing_table
    manual_categories = state.get("sales_manual_categories")
    if not isinstance(manual_categories, list):
        manual_categories = []
        state["sales_manual_categories"] = manual_categories
    else:
        manual_categories = list(manual_categories)

    render_instruction_popover(
        "売上入力のヒント",
        """
- 会計ソフトから出力したCSVを取り込むと科目と金額を自動で整形します。
- テンプレートを読み込めば、よくあるチャネル構成を一度で入力できます。
- プルダウンから追加した科目は0円で挿入されるため、数値を上書きするだけで済みます。
""",
    )

    import_feedback = state.get("sales_import_feedback")
    import_hash = state.get("sales_import_hash")

    with st.form("plan_sales"):
        with form_section(
            "売上データの取り込み",
            "CSVやAPI連携からエクスポートしたデータを一括で整形します。",
        ):
            uploaded = st.file_uploader(
                "会計ソフトの売上CSVを取り込む",
                type=["csv"],
                key="plan_sales_upload",
                help="勘定奉行やfreeeなどの会計ソフトから出力したCSVをアップロードすると自動でマッピングされます。",
            )
            if uploaded is not None:
                persist_uploaded_artifact(
                    uploaded,
                    category="plan_sales",
                    label=getattr(uploaded, "name", "plan_sales"),
                )
            download_button_from_df(
                "売上計画テンプレートをダウンロード",
                get_plan_sales_template(),
                _build_sample_filename("plan_sales", "wizard"),
            )
            st.caption("CSVの列構成を確認できるテンプレートファイルです。")

            feedback = import_feedback
            if feedback:
                level, message = feedback
                if level == "error":
                    st.error(message)
                elif level == "success":
                    st.success(message)

        with form_section(
            "テンプレートと科目の追加",
            "よく使うチャネル構成を呼び出し、入力の手戻りを防ぎます。",
            tone="secondary",
        ):
            template_cols = st.columns([3, 1])
            template_options = ["テンプレートを選択"] + list(SALES_PLAN_TEMPLATES.keys())
            selected_template = template_cols[0].selectbox(
                "売上テンプレートを適用",
                options=template_options,
                key="plan_sales_template",
                help="売上の典型的な構成をテンプレートとして呼び出せます。保存ボタンで適用されます。",
            )
            template_cols[1].markdown("&nbsp;")

            common_candidates = list(
                dict.fromkeys(
                    COMMON_SALES_ITEMS
                    + context.get("category_options", [])
                    + manual_categories
                )
            )
            selected_common = st.multiselect(
                "よく使う売上科目を追加",
                options=common_candidates,
                key="plan_sales_common_select",
                help="複数選択すると、0円の行として追加され数値だけ入力すれば完了です。候補にない科目は下の入力欄から手入力できます。",
            )
            manual_input = st.text_input(
                "候補にない売上科目を追加",
                key="plan_sales_manual_input",
                help="一覧に表示されない科目名を入力し、保存ボタンを押して候補に加えます。",
            )

        with form_section(
            "売上計画の編集",
            "取り込んだ行はここで月次金額とチャネルを整えます。",
        ):
            channel_options = list(
                dict.fromkeys(context.get("channel_options", PLAN_CHANNEL_OPTIONS_BASE))
            )
            channel_select_options = [""] + channel_options
            column_module = getattr(st, "column_config", None)
            column_config = {}
            if column_module:
                column_config["項目"] = column_module.TextColumn(
                    "項目",
                    help="売上項目の名称を入力します。",
                )
                column_config["月次売上"] = column_module.NumberColumn(
                    "月次売上 (円)",
                    min_value=0.0,
                    step=50_000.0,
                    help="各項目の月次売上計画を入力します。",
                )
                if hasattr(column_module, "SelectboxColumn"):
                    column_config["チャネル"] = column_module.SelectboxColumn(
                        "チャネル/メモ",
                        options=channel_select_options,
                        help="主要チャネルやメモを選択・入力します。",
                    )
                else:
                    column_config["チャネル"] = column_module.TextColumn(
                        "チャネル/メモ",
                        help="主要チャネルやメモを入力します。",
                    )
            else:
                column_config = None

            editor_kwargs: Dict[str, Any] = {
                "num_rows": "dynamic",
                "use_container_width": True,
                "hide_index": True,
            }
            if column_config:
                editor_kwargs["column_config"] = column_config

            sales_editor_value = st.data_editor(
                existing_table,
                key="plan_sales_editor",
                **editor_kwargs,
            )
            working_table = prepare_plan_table(
                sales_editor_value, SALES_PLAN_COLUMNS, ["月次売上"]
            )
            monthly_total = (
                float(working_table["月次売上"].sum())
                if not working_table.empty
                else 0.0
            )
            st.metric("月次売上計画合計", f"{monthly_total:,.0f} 円")
            st.caption("CSV取り込みとテンプレートで手入力を軽減し、小規模企業でも負荷を抑えられます。")

        submitted = st.form_submit_button("売上計画を保存")

    updates: Dict[str, Any] = {}
    if submitted:
        updated_table = prepare_plan_table(
            sales_editor_value, SALES_PLAN_COLUMNS, ["月次売上"]
        )
        updated_feedback = import_feedback
        updated_hash = import_hash
        updated_manual_categories = list(dict.fromkeys(manual_categories))

        if uploaded is not None:
            file_bytes = uploaded.getvalue()
            file_hash = hashlib.md5(file_bytes).hexdigest()
            if file_hash and file_hash != import_hash:
                imported_df, error = import_plan_csv(
                    file_bytes,
                    SALES_IMPORT_CANDIDATES,
                    ["項目", "月次売上"],
                    ["月次売上"],
                )
                if error:
                    updated_feedback = ("error", error)
                else:
                    updated_table = prepare_plan_table(
                        imported_df, SALES_PLAN_COLUMNS, ["月次売上"]
                    )
                    updated_feedback = (
                        "success",
                        f"CSVから{len(updated_table)}件の売上科目を読み込みました。",
                    )
                updated_hash = file_hash
        elif selected_template != "テンプレートを選択":
            template_df = pd.DataFrame(SALES_PLAN_TEMPLATES[selected_template])
            updated_table = prepare_plan_table(
                template_df, SALES_PLAN_COLUMNS, ["月次売上"]
            )
            updated_feedback = (
                "success",
                f"テンプレート『{selected_template}』を適用しました。",
            )
            st.session_state["plan_sales_template"] = "テンプレートを選択"

        manual_value = (manual_input or "").strip()
        if manual_value:
            if manual_value not in updated_manual_categories:
                updated_manual_categories.append(manual_value)
            updated_manual_categories = list(dict.fromkeys(updated_manual_categories))
            st.session_state["plan_sales_manual_input"] = ""
            st.success(f"売上科目『{manual_value}』を候補に追加しました。")

        if selected_common:
            updated_table, added = append_plan_rows(
                updated_table,
                "項目",
                "月次売上",
                {"チャネル": ""},
                selected_common,
            )
            if added:
                st.success(f"{added}件の売上科目を追加しました。")
            else:
                st.info("新しく追加できる科目がありませんでした。")
            st.session_state["plan_sales_common_select"] = []

        updates = {
            "sales_table": updated_table,
            "sales_manual_categories": updated_manual_categories,
            "sales_import_feedback": updated_feedback,
            "sales_import_hash": updated_hash,
        }

    return submitted, updates

def render_plan_step_expenses(state: Dict[str, Any], context: Dict[str, Any]) -> Tuple[bool, Dict[str, Any]]:
    """経費入力ステップを描画する。"""

    existing_table = prepare_plan_table(
        state.get("expense_table"), EXPENSE_PLAN_COLUMNS, ["月次金額"]
    )
    if not isinstance(state.get("expense_table"), pd.DataFrame):
        state["expense_table"] = existing_table

    render_instruction_popover(
        "経費入力のヒント",
        """
- 会計ソフトから出力した支出CSVを読み込むと費目と金額を自動で整形します。
- テンプレートは小規模ECでよく使う固定費と変動費の構成を含んでいます。
- プルダウンから費目を追加して月次金額を入力すれば経費計画が完成します。
""",
    )

    import_feedback = state.get("expense_import_feedback")
    import_hash = state.get("expense_import_hash")

    with st.form("plan_expenses"):
        with form_section(
            "経費データの取り込み",
            "支出CSVを読み込むと費目と金額を自動整形します。",
        ):
            uploaded = st.file_uploader(
                "会計ソフトの経費CSVを取り込む",
                type=["csv"],
                key="plan_expense_upload",
                help="freeeや弥生会計などから出力した経費CSVをアップロードすると自動でマッピングします。",
            )
            if uploaded is not None:
                persist_uploaded_artifact(
                    uploaded,
                    category="plan_expense",
                    label=getattr(uploaded, "name", "plan_expense"),
                )
            download_button_from_df(
                "経費計画テンプレートをダウンロード",
                get_plan_expense_template(),
                _build_sample_filename("plan_expense", "wizard"),
            )
            st.caption("CSVの列構成を確認できるテンプレートファイルです。")

            feedback = import_feedback
            if feedback:
                level, message = feedback
                if level == "error":
                    st.error(message)
                elif level == "success":
                    st.success(message)

        with form_section(
            "テンプレートと費目の追加",
            "固定費・変動費のひな形を呼び出し、抜け漏れを防ぎます。",
            tone="secondary",
        ):
            template_cols = st.columns([3, 1])
            template_options = ["テンプレートを選択"] + list(EXPENSE_PLAN_TEMPLATES.keys())
            selected_template = template_cols[0].selectbox(
                "経費テンプレートを適用",
                options=template_options,
                key="plan_expense_template",
                help="固定費・変動費の代表的な構成をテンプレートから読み込めます。保存ボタンで適用されます。",
            )
            template_cols[1].markdown("&nbsp;")

            selected_common = st.multiselect(
                "よく使う経費科目を追加",
                options=COMMON_EXPENSE_ITEMS,
                key="plan_expense_common_select",
                help="複数選択で0円の行を追加し、金額だけ入力できるようにします。",
            )

        with form_section(
            "経費計画の編集",
            "費目ごとの月次金額と区分を整えます。",
        ):
            column_module = getattr(st, "column_config", None)
            column_config = {}
            if column_module:
                column_config["費目"] = column_module.TextColumn(
                    "費目",
                    help="経費の科目名を入力します。",
                )
                column_config["月次金額"] = column_module.NumberColumn(
                    "月次金額 (円)",
                    min_value=0.0,
                    step=20_000.0,
                    help="各費目の月次金額を入力します。",
                )
                if hasattr(column_module, "SelectboxColumn"):
                    column_config["区分"] = column_module.SelectboxColumn(
                        "区分",
                        options=PLAN_EXPENSE_CLASSIFICATIONS,
                        help="固定費/変動費/投資などの区分を選択します。",
                    )
                else:
                    column_config["区分"] = column_module.TextColumn(
                        "区分",
                        help="固定費や変動費などの区分を入力します。",
                    )
            else:
                column_config = None

            editor_kwargs: Dict[str, Any] = {
                "num_rows": "dynamic",
                "use_container_width": True,
                "hide_index": True,
            }
            if column_config:
                editor_kwargs["column_config"] = column_config

            expense_editor_value = st.data_editor(
                existing_table,
                key="plan_expense_editor",
                **editor_kwargs,
            )
            working_table = prepare_plan_table(
                expense_editor_value, EXPENSE_PLAN_COLUMNS, ["月次金額"]
            )
            monthly_total = (
                float(working_table["月次金額"].sum())
                if not working_table.empty
                else 0.0
            )
            st.metric("月次経費計画合計", f"{monthly_total:,.0f} 円")
            st.caption("テンプレートと自動補完で経費入力も数クリックで完了します。")

        submitted = st.form_submit_button("経費計画を保存")

    updates: Dict[str, Any] = {}
    if submitted:
        updated_table = prepare_plan_table(
            expense_editor_value, EXPENSE_PLAN_COLUMNS, ["月次金額"]
        )
        updated_feedback = import_feedback
        updated_hash = import_hash

        if uploaded is not None:
            file_bytes = uploaded.getvalue()
            file_hash = hashlib.md5(file_bytes).hexdigest()
            if file_hash and file_hash != import_hash:
                imported_df, error = import_plan_csv(
                    file_bytes,
                    EXPENSE_IMPORT_CANDIDATES,
                    ["費目", "月次金額"],
                    ["月次金額"],
                )
                if error:
                    updated_feedback = ("error", error)
                else:
                    updated_table = prepare_plan_table(
                        imported_df, EXPENSE_PLAN_COLUMNS, ["月次金額"]
                    )
                    updated_feedback = (
                        "success",
                        f"CSVから{len(updated_table)}件の経費科目を読み込みました。",
                    )
                updated_hash = file_hash
        elif selected_template != "テンプレートを選択":
            template_df = pd.DataFrame(EXPENSE_PLAN_TEMPLATES[selected_template])
            updated_table = prepare_plan_table(
                template_df, EXPENSE_PLAN_COLUMNS, ["月次金額"]
            )
            updated_feedback = (
                "success",
                f"テンプレート『{selected_template}』を適用しました。",
            )
            st.session_state["plan_expense_template"] = "テンプレートを選択"

        if selected_common:
            updated_table, added = append_plan_rows(
                updated_table,
                "費目",
                "月次金額",
                {"区分": "固定費"},
                selected_common,
            )
            if added:
                st.success(f"{added}件の経費科目を追加しました。")
            else:
                st.info("新しく追加できる科目がありませんでした。")
            st.session_state["plan_expense_common_select"] = []

        updates = {
            "expense_table": updated_table,
            "expense_import_feedback": updated_feedback,
            "expense_import_hash": updated_hash,
        }

    return submitted, updates

def render_plan_step_metrics(state: Dict[str, Any], context: Dict[str, Any]) -> None:
    """財務指標計算ステップを描画する。"""

    metrics = calculate_plan_metrics_from_state(state)
    actual_reference = context.get("actual_reference", {})

    monthly_sales_delta = None
    if actual_reference.get("monthly_sales_avg") is not None:
        diff = metrics["monthly_sales"] - actual_reference["monthly_sales_avg"]
        monthly_sales_delta = f"{diff:,.0f} 円 vs 過去平均"

    monthly_profit_delta = None
    if actual_reference.get("monthly_profit_avg") is not None:
        diff_profit = metrics["monthly_profit"] - actual_reference["monthly_profit_avg"]
        monthly_profit_delta = f"{diff_profit:,.0f} 円 vs 過去平均"

    margin_value = metrics.get("monthly_margin")
    margin_display = (
        f"{margin_value * 100:.1f} %"
        if margin_value is not None and np.isfinite(margin_value)
        else "計算不可"
    )
    margin_delta = None
    if metrics.get("target_margin_pct") is not None and np.isfinite(metrics.get("margin_gap_pct")):
        margin_delta = f"{metrics['margin_gap_pct']:.1f} pt vs 目標"

    with form_section(
        "主要指標とアラート",
        "過去平均と比較して計画値の妥当性を確認します。",
    ):
        col1, col2, col3 = st.columns(3)
        col1.metric(
            "月次売上計画",
            f"{metrics['monthly_sales']:,.0f} 円",
            delta=monthly_sales_delta,
        )
        col2.metric(
            "月次営業利益",
            f"{metrics['monthly_profit']:,.0f} 円",
            delta=monthly_profit_delta,
        )
        col3.metric("営業利益率", margin_display, delta=margin_delta)

        if metrics.get("monthly_profit", 0.0) < 0:
            st.error("月次営業利益がマイナスです。コスト配分や売上計画を見直してください。")
        elif metrics.get("monthly_profit", 0.0) == 0:
            st.warning("月次営業利益が0円です。余裕を持たせるために売上・経費を再検討しましょう。")

    summary_df = build_plan_summary_df(metrics)
    with form_section(
        "計画サマリー表",
        "月次・年間の計画額を一覧で確認し、そのままCSVに出力できます。",
    ):
        formatters: Dict[str, str] = {}
        if "月次計画額" in summary_df.columns:
            formatters["月次計画額"] = "{:,.0f}"
        if "年間計画額" in summary_df.columns:
            formatters["年間計画額"] = "{:,.0f}"
        if "指標値" in summary_df.columns:
            formatters["指標値"] = "{:,.1f}"
        st.dataframe(summary_df.style.format(formatters), use_container_width=True)

        if actual_reference.get("margin_avg") is not None:
            st.caption(
                f"参考: 過去平均の営業利益率は{actual_reference['margin_avg'] * 100:.1f}%です。"
            )


def render_plan_step_review(state: Dict[str, Any], context: Dict[str, Any]) -> None:
    """ウィザード最終ステップの結果確認を描画する。"""

    metrics = state.get("metrics") or calculate_plan_metrics_from_state(state)
    info = state.get("basic_info", {})

    st.success("入力内容を確認し、必要に応じて修正してください。")

    with form_section(
        "基本情報サマリー",
        "共有前に必須項目を再確認します。",
    ):
        st.markdown(
            "<span class='form-section__status'>入力完了</span>",
            unsafe_allow_html=True,
        )
        st.markdown(
            f"**事業所名**: {info.get('company_name') or '-'} / **担当者**: {info.get('preparer') or '-'} / "
            f"**計画開始月**: {info.get('fiscal_year_start')} / **期間**: {info.get('plan_period_months')}ヶ月"
        )

    with form_section(
        "売上予測一覧",
        "CSVエクスポート前に最新の売上予測を確認します。",
    ):
        if state["sales_table"].empty:
            st.info("売上予測が未入力です。前のステップで追加してください。")
        else:
            st.dataframe(
                state["sales_table"].style.format({"月次売上": "{:,.0f}"}),
                use_container_width=True,
            )

    with form_section(
        "経費計画一覧",
        "費目別の月次コストを確認し、共有前の抜け漏れを防ぎます。",
    ):
        if state["expense_table"].empty:
            st.info("経費計画が未入力です。前のステップで追加してください。")
        else:
            st.dataframe(
                state["expense_table"].style.format({"月次金額": "{:,.0f}"}),
                use_container_width=True,
            )

    with form_section(
        "財務指標サマリー",
        "年間換算を含む主要指標を一覧で確認できます。",
    ):
        summary_df = build_plan_summary_df(metrics)
        formatters: Dict[str, str] = {}
        if "月次計画額" in summary_df.columns:
            formatters["月次計画額"] = "{:,.0f}"
        if "年間計画額" in summary_df.columns:
            formatters["年間計画額"] = "{:,.0f}"
        if "指標値" in summary_df.columns:
            formatters["指標値"] = "{:,.1f}"
        st.dataframe(summary_df.style.format(formatters), use_container_width=True)

        download_button_from_df(
            "計画サマリーをCSVでダウンロード",
            summary_df,
            "business_plan_summary.csv",
        )

        actual_reference = context.get("actual_reference", {})
        actual_caption: List[str] = []
        if actual_reference.get("monthly_sales_avg") is not None:
            actual_caption.append(f"平均売上 {actual_reference['monthly_sales_avg']:,.0f}円/月")
        if actual_reference.get("monthly_profit_avg") is not None:
            actual_caption.append(f"平均営業利益 {actual_reference['monthly_profit_avg']:,.0f}円/月")
        if actual_reference.get("margin_avg") is not None:
            actual_caption.append(f"平均利益率 {actual_reference['margin_avg'] * 100:.1f}%")
        if actual_caption:
            st.caption("過去実績: " + " / ".join(actual_caption))

        st.caption("入力内容はブラウザセッションに一時保存されます。CSVをダウンロードして関係者と共有してください。")


def render_business_plan_wizard(actual_sales: Optional[pd.DataFrame]) -> None:
    """経営計画ウィザードの全体を描画する。"""

    state = ensure_plan_wizard_state()
    if state.get("current_step", 0) < len(PLAN_WIZARD_STEPS) - 1:
        state["completed"] = False

    channel_options = list(PLAN_CHANNEL_OPTIONS_BASE)
    category_options: List[str] = []
    if actual_sales is not None and not actual_sales.empty:
        if "channel" in actual_sales.columns:
            for channel in actual_sales["channel"].dropna().unique():
                channel_str = str(channel).strip()
                if channel_str and channel_str not in channel_options:
                    channel_options.append(channel_str)
        if "category" in actual_sales.columns:
            raw_categories = (
                actual_sales["category"].dropna().map(lambda cat: str(cat).strip())
            )
            filtered_categories = raw_categories[raw_categories != ""]
            if not filtered_categories.empty:
                category_counts = filtered_categories.value_counts()
                category_options = list(category_counts.index[:CATEGORY_SUGGESTION_LIMIT])

    channel_options = list(dict.fromkeys(channel_options))
    context = {
        "channel_options": channel_options,
        "category_options": category_options,
        "actual_reference": compute_actual_reference(actual_sales),
    }

    header_cols = st.columns([3, 1])
    with header_cols[0]:
        st.markdown("### 経営計画ウィザード")
    with header_cols[1]:
        if st.button("リセット", key="plan_reset_button"):
            reset_plan_wizard_state()
            trigger_rerun()

    step_index = int(state.get("current_step", 0))
    total_steps = len(PLAN_WIZARD_STEPS)
    progress_fraction = (step_index + 1) / total_steps
    progress_label = (
        f"ステップ {step_index + 1} / {total_steps}: {PLAN_WIZARD_STEPS[step_index]['title']}"
    )
    try:
        st.progress(progress_fraction, text=progress_label)
    except TypeError:
        st.progress(progress_fraction)
        st.caption(progress_label)

    render_plan_stepper(step_index)

    st.markdown(f"#### {PLAN_WIZARD_STEPS[step_index]['title']}")
    st.write(PLAN_WIZARD_STEPS[step_index]["description"])

    def _validation_default() -> Dict[str, Any]:
        return {"is_valid": False, "errors": [], "warnings": []}

    for key in ("basic_info_validation", "sales_validation", "expense_validation"):
        state.setdefault(key, _validation_default())

    validation_result: Dict[str, Any]
    errors: List[str] = []
    warnings: List[str] = []
    is_valid = False

    if step_index == 0:
        submitted, form_info = render_plan_step_basic_info(state)
        if submitted:
            state["basic_info"] = form_info
            is_valid, errors, warnings = validate_plan_basic_info(form_info)
            state["basic_info_validation"] = {
                "is_valid": is_valid,
                "errors": errors,
                "warnings": warnings,
            }
        else:
            validation_result = state.get("basic_info_validation", _validation_default())
            is_valid = bool(validation_result.get("is_valid"))
            errors = list(validation_result.get("errors", []))
            warnings = list(validation_result.get("warnings", []))
    elif step_index == 1:
        submitted, updates = render_plan_step_sales(state, context)
        if submitted:
            state.update(updates)
            is_valid, errors, warnings = validate_plan_sales(state["sales_table"])
            state["sales_validation"] = {
                "is_valid": is_valid,
                "errors": errors,
                "warnings": warnings,
            }
        else:
            validation_result = state.get("sales_validation", _validation_default())
            is_valid = bool(validation_result.get("is_valid"))
            errors = list(validation_result.get("errors", []))
            warnings = list(validation_result.get("warnings", []))
    elif step_index == 2:
        submitted, updates = render_plan_step_expenses(state, context)
        if submitted:
            state.update(updates)
            is_valid, errors, warnings = validate_plan_expenses(state["expense_table"])
            state["expense_validation"] = {
                "is_valid": is_valid,
                "errors": errors,
                "warnings": warnings,
            }
        else:
            validation_result = state.get("expense_validation", _validation_default())
            is_valid = bool(validation_result.get("is_valid"))
            errors = list(validation_result.get("errors", []))
            warnings = list(validation_result.get("warnings", []))
    elif step_index == 3:
        render_plan_step_metrics(state, context)
        is_valid, errors, warnings = validate_plan_metrics(state.get("metrics", {}))
    else:
        render_plan_step_review(state, context)
        is_valid, errors, warnings = True, [], []

    for message in errors:
        st.error(message)
    for message in warnings:
        st.warning(message)

    nav_cols = st.columns([1, 1, 1])
    if nav_cols[0].button("戻る", disabled=step_index == 0, key=f"plan_prev_{step_index}"):
        state["current_step"] = max(step_index - 1, 0)
        trigger_rerun()

    next_label = "完了" if step_index == total_steps - 1 else "次へ進む"
    next_disabled = step_index < total_steps - 1 and not is_valid
    if nav_cols[2].button(next_label, disabled=next_disabled, key=f"plan_next_{step_index}"):
        if step_index < total_steps - 1:
            state["current_step"] = min(step_index + 1, total_steps - 1)
        else:
            state["completed"] = True
        trigger_rerun()

    if step_index == total_steps - 1 and state.get("completed"):
        st.success("経営計画ウィザードの入力が完了しました。CSV出力で関係者と共有できます。")


def _nanmean(series: pd.Series) -> float:
    """np.nanmeanの警告を避けつつ平均値を計算する。"""

    if series is None:
        return float("nan")
    clean = pd.to_numeric(series, errors="coerce").dropna()
    if clean.empty:
        return float("nan")
    return float(clean.mean())


def format_period_label(period: pd.Period, freq: str) -> str:
    """表示用の期間ラベルを生成する。"""

    if freq in {"M", "Q", "Y"}:
        return str(period)
    start = period.start_time
    end = period.end_time
    if freq.startswith("W"):
        return f"{start.strftime('%Y-%m-%d')}週 ({start.strftime('%m/%d')}〜{end.strftime('%m/%d')})"
    return f"{start.strftime('%Y-%m-%d')}〜{end.strftime('%Y-%m-%d')}"


def summarize_sales_by_period(df: pd.DataFrame, freq: str) -> pd.DataFrame:
    """売上と粗利を指定粒度で集計する。"""

    columns = [
        "period",
        "period_start",
        "period_end",
        "period_label",
        "sales_amount",
        "gross_profit",
        "net_gross_profit",
        "gross_margin_rate",
        "prev_period_sales",
        "sales_mom",
        "prev_year_sales",
        "sales_yoy",
        "prev_period_gross",
        "gross_mom",
        "prev_year_gross",
        "gross_yoy",
    ]
    if df.empty:
        return pd.DataFrame(columns=columns)

    working = df.copy()
    working["period"] = working["order_date"].dt.to_period(freq)
    summary = (
        working.groupby("period")[["sales_amount", "gross_profit", "net_gross_profit"]]
        .sum()
        .reset_index()
        .sort_values("period")
    )
    summary["period_start"] = summary["period"].dt.to_timestamp()
    summary["period_end"] = summary["period"].dt.to_timestamp(how="end")
    summary["period_label"] = summary["period"].apply(lambda p: format_period_label(p, freq))

    summary["gross_margin_rate"] = np.where(
        summary["sales_amount"] != 0,
        summary["net_gross_profit"] / summary["sales_amount"],
        np.nan,
    )

    summary["prev_period_sales"] = summary["sales_amount"].shift(1)
    summary["sales_mom"] = np.where(
        (summary["prev_period_sales"].notna()) & (summary["prev_period_sales"] != 0),
        (summary["sales_amount"] - summary["prev_period_sales"]) / summary["prev_period_sales"],
        np.nan,
    )

    yoy_lag = PERIOD_YOY_LAG.get(freq, 0)
    if yoy_lag:
        summary["prev_year_sales"] = summary["sales_amount"].shift(yoy_lag)
        summary["sales_yoy"] = np.where(
            (summary["prev_year_sales"].notna()) & (summary["prev_year_sales"] != 0),
            (summary["sales_amount"] - summary["prev_year_sales"]) / summary["prev_year_sales"],
            np.nan,
        )
    else:
        summary["prev_year_sales"] = np.nan
        summary["sales_yoy"] = np.nan

    summary["prev_period_gross"] = summary["net_gross_profit"].shift(1)
    summary["gross_mom"] = np.where(
        (summary["prev_period_gross"].notna()) & (summary["prev_period_gross"] != 0),
        (summary["net_gross_profit"] - summary["prev_period_gross"]) / summary["prev_period_gross"],
        np.nan,
    )

    if yoy_lag:
        summary["prev_year_gross"] = summary["net_gross_profit"].shift(yoy_lag)
        summary["gross_yoy"] = np.where(
            (summary["prev_year_gross"].notna()) & (summary["prev_year_gross"] != 0),
            (summary["net_gross_profit"] - summary["prev_year_gross"]) / summary["prev_year_gross"],
            np.nan,
        )
    else:
        summary["prev_year_gross"] = np.nan
        summary["gross_yoy"] = np.nan

    return summary[columns]


def build_kpi_history_df(
    merged_df: pd.DataFrame,
    subscription_df: Optional[pd.DataFrame],
    overrides: Optional[Dict[str, float]],
) -> pd.DataFrame:
    """月次KPI履歴を作成する。"""

    if merged_df.empty:
        return pd.DataFrame()

    months = (
        merged_df["order_month"].dropna().sort_values().unique()
        if "order_month" in merged_df.columns
        else []
    )
    history: List[Dict[str, Any]] = []
    for month in months:
        kpi_row = calculate_kpis(merged_df, subscription_df, month=month, overrides=overrides)
        if kpi_row:
            history.append(kpi_row)

    if not history:
        return pd.DataFrame()

    history_df = pd.DataFrame(history)
    if "month" in history_df.columns:
        history_df["month"] = pd.PeriodIndex(history_df["month"], freq="M")
    return history_df


def aggregate_kpi_history(history_df: pd.DataFrame, freq: str) -> pd.DataFrame:
    """KPI履歴を指定した粒度で集計する。"""

    columns = [
        "period",
        "period_start",
        "period_end",
        "period_label",
        "sales",
        "gross_profit",
        "marketing_cost",
        "active_customers_avg",
        "new_customers",
        "repeat_customers",
        "cancelled_subscriptions",
        "previous_active_customers",
        "ltv",
        "arpu",
        "churn_rate",
        "repeat_rate",
        "gross_margin_rate",
        "inventory_turnover_days",
        "stockout_rate",
        "training_sessions",
        "new_product_count",
        "ltv_prev",
        "ltv_delta",
        "arpu_prev",
        "arpu_delta",
        "churn_prev",
        "churn_delta",
        "gross_margin_prev",
        "gross_margin_delta",
        "repeat_prev",
        "repeat_delta",
        "inventory_turnover_prev",
        "inventory_turnover_delta",
        "stockout_prev",
        "stockout_delta",
        "training_prev",
        "training_delta",
        "new_product_prev",
        "new_product_delta",
    ]
    if history_df.empty:
        return pd.DataFrame(columns=columns)

    working = history_df.dropna(subset=["month"]).copy()
    if working.empty:
        return pd.DataFrame(columns=columns)

    working["timestamp"] = working["month"].dt.to_timestamp()
    working["period"] = working["timestamp"].dt.to_period(freq)
    aggregated = (
        working.groupby("period").agg(
            sales=("sales", "sum"),
            gross_profit=("gross_profit", "sum"),
            marketing_cost=("marketing_cost", "sum"),
            active_customers=("active_customers", _nanmean),
            new_customers=("new_customers", "sum"),
            repeat_customers=("repeat_customers", "sum"),
            cancelled_subscriptions=("cancelled_subscriptions", "sum"),
            previous_active_customers=("previous_active_customers", "sum"),
            ltv=("ltv", _nanmean),
            inventory_turnover_days=("inventory_turnover_days", _nanmean),
            stockout_rate=("stockout_rate", _nanmean),
            training_sessions=("training_sessions", "sum"),
            new_product_count=("new_product_count", "sum"),
        )
    ).reset_index()

    if aggregated.empty:
        return pd.DataFrame(columns=columns)

    aggregated.rename(columns={"active_customers": "active_customers_avg"}, inplace=True)
    aggregated["period_start"] = aggregated["period"].dt.to_timestamp()
    aggregated["period_end"] = aggregated["period"].dt.to_timestamp(how="end")
    aggregated["period_label"] = aggregated["period"].apply(lambda p: format_period_label(p, freq))

    aggregated["arpu"] = aggregated.apply(
        lambda row: row["sales"] / row["active_customers_avg"]
        if row["active_customers_avg"]
        else np.nan,
        axis=1,
    )
    aggregated["churn_rate"] = aggregated.apply(
        lambda row: row["cancelled_subscriptions"] / row["previous_active_customers"]
        if row["previous_active_customers"]
        else np.nan,
        axis=1,
    )
    aggregated["repeat_rate"] = aggregated.apply(
        lambda row: row["repeat_customers"] / row["active_customers_avg"]
        if row["active_customers_avg"]
        else np.nan,
        axis=1,
    )
    aggregated["gross_margin_rate"] = aggregated.apply(
        lambda row: row["gross_profit"] / row["sales"] if row["sales"] else np.nan,
        axis=1,
    )

    aggregated.sort_values("period", inplace=True)
    aggregated["ltv_prev"] = aggregated["ltv"].shift(1)
    aggregated["ltv_delta"] = aggregated["ltv"] - aggregated["ltv_prev"]
    aggregated["arpu_prev"] = aggregated["arpu"].shift(1)
    aggregated["arpu_delta"] = aggregated["arpu"] - aggregated["arpu_prev"]
    aggregated["churn_prev"] = aggregated["churn_rate"].shift(1)
    aggregated["churn_delta"] = aggregated["churn_rate"] - aggregated["churn_prev"]
    aggregated["gross_margin_prev"] = aggregated["gross_margin_rate"].shift(1)
    aggregated["gross_margin_delta"] = aggregated["gross_margin_rate"] - aggregated["gross_margin_prev"]
    aggregated["repeat_prev"] = aggregated["repeat_rate"].shift(1)
    aggregated["repeat_delta"] = aggregated["repeat_rate"] - aggregated["repeat_prev"]
    aggregated["inventory_turnover_prev"] = aggregated["inventory_turnover_days"].shift(1)
    aggregated["inventory_turnover_delta"] = (
        aggregated["inventory_turnover_days"] - aggregated["inventory_turnover_prev"]
    )
    aggregated["stockout_prev"] = aggregated["stockout_rate"].shift(1)
    aggregated["stockout_delta"] = aggregated["stockout_rate"] - aggregated["stockout_prev"]
    aggregated["training_prev"] = aggregated["training_sessions"].shift(1)
    aggregated["training_delta"] = aggregated["training_sessions"] - aggregated["training_prev"]
    aggregated["new_product_prev"] = aggregated["new_product_count"].shift(1)
    aggregated["new_product_delta"] = (
        aggregated["new_product_count"] - aggregated["new_product_prev"]
    )

    return aggregated[columns]


def format_currency(value: Optional[float]) -> str:
    """通貨表記で値を整形する。"""

    if value is None or pd.isna(value):
        return "-"
    return f"{value:,.0f} 円"


def format_percent(value: Optional[float], digits: int = 1) -> str:
    """割合値を%表示に変換する。"""

    if value is None or pd.isna(value):
        return "-"
    return f"{value * 100:.{digits}f}%"


def format_number(value: Optional[float], *, digits: int = 1, unit: str = "") -> str:
    """一般的な数値を文字列化する。"""

    if value is None or pd.isna(value):
        return "-"
    formatted = f"{value:,.{digits}f}" if digits > 0 else f"{value:,.0f}"
    return f"{formatted}{unit}"


def format_delta(
    value: Optional[float], *, digits: int = 1, unit: str = "", percentage: bool = False
) -> Optional[str]:
    """指標変化量の表示を整える。"""

    if value is None or pd.isna(value):
        return None
    if abs(float(value)) < 1e-9:
        return None
    if percentage:
        return f"{value * 100:+.{digits}f} pt"
    formatted = f"{value:+.{digits}f}"
    if unit:
        formatted = f"{formatted}{unit}"
    return formatted


def render_bsc_card(
    *,
    title: str,
    icon: str,
    subtitle: Optional[str],
    metrics: List[Dict[str, Any]],
    variant: str = "neutral",
) -> None:
    """バランスト・スコアカードのカードUIを描画する。"""

    card_classes = ["bsc-card"]
    if variant and variant != "neutral":
        card_classes.append(f"bsc-card--{variant}")
    st.markdown(
        f"<div class='{' '.join(card_classes)}'>",
        unsafe_allow_html=True,
    )
    st.markdown(
        f"<div class='bsc-card__title'>{icon} {html.escape(title)}</div>", unsafe_allow_html=True
    )
    if subtitle:
        st.markdown(
            f"<div class='bsc-card__subtitle'>{html.escape(subtitle)}</div>",
            unsafe_allow_html=True,
        )
    for metric in metrics:
        st.metric(metric["label"], metric["value"], delta=metric.get("delta"))
        target_text = metric.get("target_text")
        if target_text:
            st.caption(target_text)
    st.markdown("</div>", unsafe_allow_html=True)


def build_bsc_quadrants(selected_kpi_row: pd.Series) -> List[Dict[str, Any]]:
    """KPIサマリーからBSC四象限のデータ構造を生成する。"""

    if selected_kpi_row is None or selected_kpi_row.empty:
        return []

    configs: List[Dict[str, Any]] = [
        {
            "key": "financial",
            "title": "財務",
            "title_icon": "💹",
            "subtitle": "長期的な収益性を測る指標",
            "metric_key": "ltv",
            "metric_label": "LTV",
            "formatter": lambda v: format_currency(v),
            "delta_key": "ltv_delta",
            "delta_formatter": lambda v: format_delta(v, digits=0, unit=" 円"),
            "direction": "above",
            "icons": {"success": "🌟", "warning": "🎯", "neutral": "•"},
            "target_text": lambda target: f"目標: {target:,.0f} 円以上"
            if target is not None
            else None,
        },
        {
            "key": "customer",
            "title": "顧客",
            "title_icon": "🤝",
            "subtitle": "リピート構造とロイヤルティ",
            "metric_key": "repeat_rate",
            "metric_label": "リピート率",
            "formatter": lambda v: format_percent(v, digits=1),
            "delta_key": "repeat_delta",
            "delta_formatter": lambda v: format_delta(v, digits=1, percentage=True),
            "direction": "above",
            "icons": {"success": "✨", "warning": "🚧", "neutral": "•"},
            "target_text": lambda target: f"目標: {target * 100:.1f}%以上"
            if target is not None
            else None,
        },
        {
            "key": "internal",
            "title": "内部プロセス",
            "title_icon": "🛠",
            "subtitle": "在庫運用と業務効率",
            "metric_key": "inventory_turnover_days",
            "metric_label": "在庫回転日数",
            "formatter": lambda v: format_number(v, digits=1, unit=" 日"),
            "delta_key": "inventory_turnover_delta",
            "delta_formatter": lambda v: format_delta(v, digits=1, unit=" 日"),
            "direction": "below",
            "icons": {"success": "✅", "warning": "⚠️", "neutral": "•"},
            "target_text": lambda target: f"目標: {target:.0f} 日以内"
            if target is not None
            else None,
        },
        {
            "key": "learning",
            "title": "学習と成長",
            "title_icon": "🎓",
            "subtitle": "人材育成と改善サイクル",
            "metric_key": "training_sessions",
            "metric_label": "研修実施回数",
            "formatter": lambda v: format_number(v, digits=0, unit=" 回"),
            "delta_key": "training_delta",
            "delta_formatter": lambda v: format_delta(v, digits=0, unit=" 回"),
            "direction": "above",
            "icons": {"success": "📈", "warning": "💡", "neutral": "•"},
            "target_text": lambda target: f"目標: {target:.0f} 回以上"
            if target is not None
            else None,
        },
    ]

    quadrants: List[Dict[str, Any]] = []
    for config in configs:
        metric_key = config["metric_key"]
        delta_key = config["delta_key"]
        raw_value = selected_kpi_row.get(metric_key)
        delta_value = selected_kpi_row.get(delta_key)
        if pd.isna(raw_value):
            raw_value = None
        if pd.isna(delta_value):
            delta_value = None

        formatter: Callable[[Optional[float]], str] = config["formatter"]
        formatted_value = formatter(raw_value)
        delta_formatter: Callable[[Optional[float]], Optional[str]] = config["delta_formatter"]
        delta_text = delta_formatter(delta_value) if delta_value is not None else None

        target_value = BSC_TARGETS.get(metric_key)
        target_formatter: Callable[[Optional[float]], Optional[str]] = config["target_text"]
        target_text = target_formatter(target_value)

        icons = config["icons"]
        direction = config.get("direction", "above")
        variant = "neutral"
        status_icon = icons.get("neutral", "")
        if raw_value is not None and target_value is not None:
            if direction == "below":
                meets_target = float(raw_value) <= float(target_value)
            else:
                meets_target = float(raw_value) >= float(target_value)
            variant = "success" if meets_target else "warning"
            status_icon = icons["success"] if meets_target else icons["warning"]

        label_prefix = f"{status_icon} " if status_icon else ""
        metric_entry = {
            "label": f"{label_prefix}{config['metric_label']}",
            "value": formatted_value,
            "delta": delta_text,
            "target_text": target_text,
            "raw_value": raw_value,
            "target_value": target_value,
            "delta_value": delta_value,
        }

        quadrants.append(
            {
                "key": config["key"],
                "title": config["title"],
                "icon": config["title_icon"],
                "subtitle": config.get("subtitle"),
                "metrics": [metric_entry],
                "variant": variant,
                "metric_label": config["metric_label"],
            }
        )

    has_values = any(
        metric.get("raw_value") is not None
        for quadrant in quadrants
        for metric in quadrant.get("metrics", [])
    )
    return quadrants if has_values else []


def render_bsc_cards_grid(quadrants: Sequence[Dict[str, Any]]) -> None:
    """BSCカード群を2x2レイアウトで表示する。"""

    if not quadrants:
        st.info("BSCを表示するための指標が不足しています。")
        return

    columns_per_row = 2
    for row_start in range(0, len(quadrants), columns_per_row):
        row_quadrants = quadrants[row_start : row_start + columns_per_row]
        cols = st.columns(len(row_quadrants))
        for col, quadrant in zip(cols, row_quadrants):
            with col:
                render_bsc_card(
                    title=quadrant.get("title", ""),
                    icon=quadrant.get("icon", ""),
                    subtitle=quadrant.get("subtitle"),
                    metrics=quadrant.get("metrics", []),
                    variant=quadrant.get("variant", "neutral"),
                )


def render_bsc_quadrant_chart(quadrants: Sequence[Dict[str, Any]]) -> None:
    """BSC四象限をPlotlyチャートとして描画する。"""

    if not quadrants:
        st.info("BSCを表示するための指標が不足しています。")
        return

    ensure_theme_state_defaults()
    tokens = st.session_state.get("ui_active_tokens", LIGHT_THEME_TOKENS)

    variant_color_map = {
        "success": SUCCESS_COLOR,
        "warning": WARNING_COLOR,
        "neutral": ACCENT_COLOR,
    }
    variant_surface_map = {
        "success": SUCCESS_SURFACE_COLOR,
        "warning": WARNING_SURFACE_COLOR,
        "neutral": tokens.get("surface_tint", tokens.get("surface", "#ffffff")),
    }

    quadrant_bounds = {
        "financial": (0.0, 1.0, 1.0, 2.0),
        "customer": (1.0, 2.0, 1.0, 2.0),
        "internal": (0.0, 1.0, 0.0, 1.0),
        "learning": (1.0, 2.0, 0.0, 1.0),
    }

    fig = go.Figure()

    for quadrant in quadrants:
        key = quadrant.get("key")
        bounds = quadrant_bounds.get(key)
        if not bounds:
            continue
        x0, x1, y0, y1 = bounds
        variant = quadrant.get("variant", "neutral")
        border_color = variant_color_map.get(variant, ACCENT_COLOR)
        fill_color = variant_surface_map.get(variant, tokens.get("surface_tint"))

        fig.add_shape(
            type="rect",
            x0=x0,
            x1=x1,
            y0=y0,
            y1=y1,
            line=dict(color=border_color, width=2),
            fillcolor=fill_color,
            layer="below",
        )

        center_x = (x0 + x1) / 2
        center_y = (y0 + y1) / 2
        title_y = y1 - 0.18
        value_y = center_y + 0.15
        target_y = center_y - 0.2

        metric_info = quadrant.get("metrics", [{}])[0]
        value_text = metric_info.get("value", "-")
        target_text = metric_info.get("target_text")
        delta_text = metric_info.get("delta")

        fig.add_annotation(
            x=center_x,
            y=title_y,
            text=f"{quadrant.get('icon', '')} {quadrant.get('title', '')}",
            showarrow=False,
            font=dict(size=16, color=tokens.get("text"), family=MCKINSEY_FONT_STACK),
        )
        fig.add_annotation(
            x=center_x,
            y=value_y,
            text=value_text,
            showarrow=False,
            font=dict(size=18, color=tokens.get("text"), family=NUMERIC_FONT_STACK),
        )
        if target_text:
            fig.add_annotation(
                x=center_x,
                y=target_y,
                text=target_text,
                showarrow=False,
                font=dict(size=12, color=tokens.get("muted"), family=MCKINSEY_FONT_STACK),
            )

        hover_lines = [
            f"{quadrant.get('title', '')} ({quadrant.get('metric_label', '')})",
            f"実績: {value_text}",
        ]
        if delta_text:
            hover_lines.append(f"前期差: {delta_text}")
        if target_text:
            hover_lines.append(target_text)

        fig.add_trace(
            go.Scatter(
                x=[center_x],
                y=[center_y],
                mode="markers",
                marker=dict(size=20, opacity=0),
                showlegend=False,
                hovertemplate="<br>".join(hover_lines) + "<extra></extra>",
            )
        )

    apply_chart_theme(fig)
    fig.update_layout(
        xaxis=dict(range=[0, 2], showgrid=False, zeroline=False, showticklabels=False),
        yaxis=dict(range=[0, 2], showgrid=False, zeroline=False, showticklabels=False),
        height=480,
        margin=dict(l=40, r=40, t=40, b=40),
        hovermode="closest",
    )

    st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})


def persistent_segmented_control(
    key: str,
    options: Sequence[T],
    *,
    default: Optional[T] = None,
    label: str = "表示切替",
    help_text: Optional[str] = None,
    label_visibility: str = "collapsed",
    format_func: Optional[Callable[[T], str]] = None,
) -> T:
    """セッションに選択状態を保持するセグメントコントロールを描画する。"""

    if not options:
        raise ValueError("options must not be empty")

    default_value = default if default is not None else options[0]
    if key not in st.session_state or st.session_state[key] not in options:
        st.session_state[key] = default_value

    widget_value = st.segmented_control(
        label,
        options,
        selection_mode="single",
        default=st.session_state[key],
        format_func=format_func,
        key=f"{key}_segmented",
        help=help_text,
        label_visibility=label_visibility,
    )

    if isinstance(widget_value, list):
        selected_value: T = widget_value[0] if widget_value else default_value
    else:
        selected_value = widget_value if widget_value in options else default_value

    if selected_value not in options:
        selected_value = default_value

    st.session_state[key] = selected_value
    return selected_value


def _format_numeric_text(
    value: Union[int, float, None],
    *,
    is_integer: bool,
    text_format: Optional[str] = None,
) -> str:
    """数値入力用のテキスト表示を生成する。"""

    if value is None:
        return ""
    try:
        numeric_value = float(value)
    except (TypeError, ValueError):
        return ""
    if np.isnan(numeric_value):
        return ""
    if text_format:
        try:
            return format(numeric_value, text_format)
        except ValueError:
            pass
    if is_integer:
        return f"{int(round(numeric_value)):,}"
    if abs(numeric_value) >= 1000:
        return f"{numeric_value:,.1f}"
    if abs(numeric_value) >= 1:
        return f"{numeric_value:,.2f}"
    return f"{numeric_value:.4f}"


def _normalize_numeric_text(value: str) -> str:
    """全角文字や記号を考慮して数値文字列を標準化する。"""

    normalized = (value or "").strip().translate(NUMERIC_TEXT_TRANSLATION)
    normalized = normalized.replace("%", "")
    normalized = normalized.replace(",", "")
    return normalized


def enhanced_number_input(
    label: str,
    *,
    key: str,
    default_value: Union[int, float],
    min_value: Optional[Union[int, float]] = None,
    max_value: Optional[Union[int, float]] = None,
    step: Optional[Union[int, float]] = None,
    number_format: Optional[str] = None,
    help_text: Optional[str] = None,
    text_format: Optional[str] = None,
    placeholder: Optional[str] = None,
    allow_float: bool = True,
) -> Union[int, float]:
    """ボタン操作と直接入力を併用できる数値入力コンポーネント。"""

    is_integer = not allow_float
    caster: Callable[[Union[int, float]], Union[int, float]] = int if is_integer else float

    min_cast = caster(min_value) if min_value is not None else None
    max_cast = caster(max_value) if max_value is not None else None
    if step is None:
        step_cast: Union[int, float] = 1 if is_integer else 1.0
    else:
        step_cast = caster(step)

    if number_format is None:
        number_format = "%d" if is_integer else "%.2f"

    container = st.container()
    container.markdown(f"**{label}**")
    if help_text:
        container.caption(help_text)

    spinner_key = f"{key}_spinner"
    text_key = f"{key}_text"
    pending_spinner_key = f"{spinner_key}_pending"
    pending_text_key = f"{text_key}_pending"

    current_value = caster(st.session_state.get(key, default_value))
    st.session_state.setdefault(key, current_value)
    st.session_state.setdefault(spinner_key, current_value)

    if pending_spinner_key in st.session_state:
        st.session_state[spinner_key] = caster(st.session_state[pending_spinner_key])
        del st.session_state[pending_spinner_key]

    columns = container.columns([0.55, 0.45])
    previous_spinner_value = caster(st.session_state.get(spinner_key, current_value))
    spinner_value = columns[0].number_input(
        "数値入力",
        min_value=min_cast,
        max_value=max_cast,
        value=previous_spinner_value,
        step=step_cast,
        format=number_format,
        key=spinner_key,
        label_visibility="collapsed",
    )

    spinner_value = caster(spinner_value)
    st.session_state[key] = spinner_value

    formatted_spinner = _format_numeric_text(
        spinner_value, is_integer=is_integer, text_format=text_format
    )
    previous_formatted = _format_numeric_text(
        previous_spinner_value, is_integer=is_integer, text_format=text_format
    )
    spinner_changed = spinner_value != previous_spinner_value
    if text_key not in st.session_state:
        st.session_state[text_key] = formatted_spinner
    elif spinner_changed and st.session_state[text_key] == previous_formatted:
        st.session_state[text_key] = formatted_spinner

    if pending_text_key in st.session_state:
        st.session_state[text_key] = st.session_state[pending_text_key]
        del st.session_state[pending_text_key]

    text_value = columns[1].text_input(
        "直接入力",
        value=st.session_state[text_key],
        key=text_key,
        placeholder=placeholder or "例: 1,000",
        label_visibility="visible",
    )

    parsed_value: Union[int, float] = spinner_value
    error_message: Optional[str] = None
    normalized = _normalize_numeric_text(text_value)
    if normalized:
        try:
            parsed_float = float(normalized)
        except ValueError:
            error_message = "数値として認識できません。"
        else:
            if is_integer and not parsed_float.is_integer():
                error_message = "整数で入力してください。"
            else:
                parsed_value = caster(parsed_float)
                if min_cast is not None and parsed_value < min_cast:
                    error_message = f"{min_cast:,}以上を入力してください。"
                if max_cast is not None and parsed_value > max_cast:
                    error_message = f"{max_cast:,}以下を入力してください。"

    if error_message:
        container.error(error_message)
        final_value = spinner_value
    else:
        final_value = parsed_value
        st.session_state[key] = final_value
        st.session_state[pending_spinner_key] = final_value
        st.session_state[pending_text_key] = _format_numeric_text(
            final_value, is_integer=is_integer, text_format=text_format
        )

    return final_value


def render_navigation() -> Tuple[str, str]:
    """トップレベルのナビゲーションを描画し、選択されたキーと表示ラベルを返す。"""

    label_options = list(NAV_OPTION_LOOKUP.values())
    label_to_key = {value: key for key, value in NAV_OPTION_LOOKUP.items()}

    current_key = st.session_state.get("main_nav", PRIMARY_NAV_ITEMS[0]["key"])
    if current_key not in NAV_OPTION_LOOKUP:
        current_key = PRIMARY_NAV_ITEMS[0]["key"]
    selected_key = current_key
    selected_label = NAV_OPTION_LOOKUP[selected_key]

    st.markdown("<div class='main-nav-tabs'>", unsafe_allow_html=True)
    nav_columns = st.columns(len(label_options))
    for idx, label in enumerate(label_options):
        key = label_to_key[label]
        is_active = key == selected_key
        button_pressed = nav_columns[idx].button(
            label,
            key=f"main_nav_tab_{key}",
            use_container_width=True,
            type="primary" if is_active else "secondary",
        )
        if button_pressed:
            selected_key = key
            selected_label = label
    st.markdown("</div>", unsafe_allow_html=True)

    st.session_state["main_nav"] = selected_key
    st.session_state["main_nav_display"] = selected_label
    return selected_key, NAV_LABEL_LOOKUP[selected_key]


def render_breadcrumb(current_label: str) -> None:
    """現在地がわかるパンくずリストを表示する。"""

    root_label = NAV_LABEL_LOOKUP.get("dashboard", "Dashboard")
    if current_label == root_label:
        parts = [current_label]
    else:
        parts = [root_label, current_label]
    breadcrumb = " / ".join(parts)
    st.markdown(
        f"<div class='breadcrumb-trail'>{html.escape(breadcrumb)}</div>",
        unsafe_allow_html=True,
    )


def render_hero_section(
    latest_label: str, period_label: str, record_count: int, alert_count: int
) -> None:
    """ヒーローエリアをマッキンゼー風に表示する。"""

    if alert_count > 0:
        status_label = f"要確認: {alert_count}件"
        status_class = "hero-badge hero-badge--alert with-icon"
        status_icon = build_ui_icon("!", tone="warning")
    else:
        status_label = "主要指標は安定しています"
        status_class = "hero-badge hero-badge--accent with-icon"
        status_icon = build_ui_icon("OK", tone="success")

    latest_icon = build_ui_icon("L", tone="accent")
    period_icon = build_ui_icon("P", tone="accent")
    records_icon = build_ui_icon("R", tone="accent")
    executive_icon = build_ui_icon("EX", tone="accent")
    store_icon = build_ui_icon("ST", tone="accent")
    finance_icon = build_ui_icon("FA", tone="accent")

    st.markdown(
        """
        <div class="hero-panel">
            <div class="hero-title">くらしいきいき社 計数管理ダッシュボード</div>
            <p class="hero-subtitle">高粗利商材のパフォーマンスを即座に把握し、迅速な意思決定を支援します。</p>
            <div class="hero-meta">
                <span class="hero-badge with-icon">{latest_icon}<span class="hero-badge__label">最新データ: {latest}</span></span>
                <span class="hero-badge with-icon">{period_icon}<span class="hero-badge__label">表示期間: {period}</span></span>
                <span class="hero-badge with-icon">{records_icon}<span class="hero-badge__label">対象レコード: {records}</span></span>
                <span class="{status_class}">{status_icon}<span class="hero-badge__label">{status}</span></span>
            </div>
            <div class="hero-persona">
                <span class="hero-chip">{executive_icon}<span>社長: 売上・粗利を5秒確認</span></span>
                <span class="hero-chip">{store_icon}<span>店長: リピーターと在庫</span></span>
                <span class="hero-chip">{finance_icon}<span>経理: 資金繰りと育成</span></span>
            </div>
        </div>
        """.format(
            latest=html.escape(latest_label or "-"),
            period=html.escape(period_label or "-"),
            records=f"{record_count:,} 件",
            status_class=status_class,
            status=html.escape(status_label),
            status_icon=status_icon,
            latest_icon=latest_icon,
            period_icon=period_icon,
            records_icon=records_icon,
            executive_icon=executive_icon,
            store_icon=store_icon,
            finance_icon=finance_icon,
        ),
        unsafe_allow_html=True,
    )


def render_status_banner(alerts: Optional[List[str]]) -> None:
    """アラート状況をアクセントカラーで表示する。"""

    if alerts:
        items = "".join(f"<li>{html.escape(msg)}</li>" for msg in alerts)
        alert_icon = build_ui_icon("!", tone="warning")
        st.markdown(
            f"""
            <div class="alert-banner alert-banner--warning">
                <div class="alert-banner__title">{alert_icon}<span>警告が検知されました</span></div>
                <ul>{items}</ul>
            </div>
            """,
            unsafe_allow_html=True,
        )
    else:
        ok_icon = build_ui_icon("OK", tone="success")
        st.markdown(
            """
            <div class="alert-banner alert-banner--ok">
                <div class="alert-banner__title">{ok_icon}<span>主要指標は設定した閾値内に収まっています。</span></div>
            </div>
            """.format(ok_icon=ok_icon),
            unsafe_allow_html=True,
        )


def render_search_bar() -> str:
    """ヒーロー直下のクイック検索をカードスタイルで表示する。"""

    with st.container():
        st.markdown(
            "<div class='surface-card search-card'>", unsafe_allow_html=True
        )
        st.markdown(
            "<div class='search-title'>{title}</div>".format(
                title=html.escape(
                    translate("quick_search_title", default="クイック検索")
                )
            ),
            unsafe_allow_html=True,
        )
        query = st.text_input(
            translate("quick_search_title", default="クイック検索"),
            placeholder=translate(
                "quick_search_placeholder",
                default="商品名、チャネル、チュートリアルを検索",
            ),
            key="global_search",
            label_visibility="collapsed",
        )
        st.markdown("</div>", unsafe_allow_html=True)
    return query


def render_global_search_results(query: str, merged_df: pd.DataFrame) -> None:
    """検索クエリに一致するデータやチュートリアルをまとめて表示する。"""

    query = (query or "").strip()
    if not query:
        return

    query_lower = query.lower()
    with st.container():
        st.markdown("<div class='surface-card search-results-card'>", unsafe_allow_html=True)
        st.markdown(
            "### {heading}".format(
                heading=html.escape(
                    translate("quick_search_heading", default="クイック検索結果")
                )
            )
        )

        tab_entries: List[Tuple[str, str]] = [
            ("sales", translate("quick_search_tab_sales", default="売上データ")),
            ("tutorials", translate("quick_search_tab_tutorials", default="チュートリアル")),
        ]
        tab_lookup = {key: label for key, label in tab_entries}
        selected_tab = persistent_segmented_control(
            "quick_search_active_tab",
            [key for key, _ in tab_entries],
            default="sales",
            label="Quick search tab",
            help_text=translate(
                "quick_search_tab_help",
                default="検索結果の種類を切り替えます。",
            ),
            label_visibility="collapsed",
            format_func=lambda key: tab_lookup.get(str(key), str(key)),
        )

        sales_summary: Optional[pd.DataFrame] = None
        sales_message: Optional[str] = None
        if merged_df is not None and not merged_df.empty:
            searchable = merged_df.copy()
            for column in ["product_name", "channel", "category"]:
                if column in searchable.columns:
                    searchable[column] = searchable[column].astype(str)
            fallback = pd.Series([False] * len(searchable), index=searchable.index)
            product_series = (
                searchable["product_name"].str.contains(query, case=False, na=False)
                if "product_name" in searchable.columns
                else fallback
            )
            channel_series = (
                searchable["channel"].str.contains(query, case=False, na=False)
                if "channel" in searchable.columns
                else fallback
            )
            category_series = (
                searchable["category"].str.contains(query, case=False, na=False)
                if "category" in searchable.columns
                else fallback
            )
            mask = product_series | channel_series | category_series
            matched_sales = searchable[mask].copy()
            if not matched_sales.empty and "order_date" in matched_sales.columns:
                matched_sales.sort_values("order_date", ascending=False, inplace=True)
            if not matched_sales.empty:
                display_cols = []
                if "order_date" in matched_sales.columns:
                    matched_sales["order_date"] = pd.to_datetime(matched_sales["order_date"])
                    matched_sales["order_date_str"] = matched_sales["order_date"].dt.strftime(
                        "%Y-%m-%d"
                    )
                    display_cols.append("order_date_str")
                if "channel" in matched_sales.columns:
                    display_cols.append("channel")
                if "product_name" in matched_sales.columns:
                    display_cols.append("product_name")
                if "sales_amount" in matched_sales.columns:
                    display_cols.append("sales_amount")
                summary_table = matched_sales.head(10)[display_cols].rename(
                    columns={
                        "order_date_str": "受注日",
                        "channel": "チャネル",
                        "product_name": "商品名",
                        "sales_amount": "売上高",
                    }
                )
                if "売上高" in summary_table.columns:
                    summary_table["売上高"] = summary_table["売上高"].map(
                        lambda v: f"{v:,.0f}"
                    )
                sales_summary = summary_table
            else:
                sales_message = translate(
                    "quick_search_no_match",
                    default="売上データに一致する項目は見つかりませんでした。",
                )
        else:
            sales_message = translate(
                "quick_search_no_data",
                default="売上データが読み込まれていないため検索できません。",
            )

        matches = [
            tutorial
            for tutorial in TUTORIAL_INDEX
            if query_lower in tutorial["title"].lower()
            or any(query_lower in keyword.lower() for keyword in tutorial.get("keywords", []))
        ]
        if selected_tab == "sales":
            if sales_summary is not None and not sales_summary.empty:
                st.caption(
                    translate(
                        "quick_search_caption", default="直近の一致10件を表示しています。"
                    )
                )
                st.dataframe(sales_summary, hide_index=True, use_container_width=True)
            elif sales_message:
                st.caption(sales_message)
        else:
            st.markdown(
                "**{heading}**".format(
                    heading=html.escape(
                        translate(
                            "quick_search_tutorial_heading",
                            default="関連チュートリアル",
                        )
                    )
                )
            )
            if matches:
                for tutorial in matches:
                    st.markdown(f"- [{tutorial['title']}]({tutorial['path']})")
            else:
                st.caption(
                    translate(
                        "quick_search_no_tutorial",
                        default="該当するチュートリアルは見つかりませんでした。",
                    )
                )
        st.markdown("</div>", unsafe_allow_html=True)


def _format_currency_compact(value: Optional[float]) -> str:
    """通貨をスペースなしの円表示に整形する。"""

    if value is None or pd.isna(value):
        return "-"
    return f"{float(value):,.0f}円"


def format_percentage_delta(value: Optional[float], *, digits: int = 1) -> Optional[str]:
    """百分率の変化量を%表記で返す。"""

    if value is None or pd.isna(value):
        return None
    return f"{float(value) * 100:+.{digits}f}%"


def format_target_gap(
    value: Optional[float],
    target: Optional[float],
    *,
    percentage: bool = False,
    digits: int = 1,
) -> Tuple[str, Optional[float]]:
    """値と目標値の差分をテキストと数値で返す。"""

    if value is None or pd.isna(value) or target is None or pd.isna(target):
        return "-", None
    gap = float(value) - float(target)
    if percentage:
        text = f"{gap * 100:+.{digits}f} pt"
    else:
        text = f"{gap:+,.0f} 円"
    return text, gap


def delta_class_from_value(value: Optional[float]) -> str:
    """KGIカード用のデルタクラスを決定する。"""

    if value is None or pd.isna(value):
        return ""
    numeric = float(value)
    if numeric > 0:
        return "kgi-card__delta--up"
    if numeric < 0:
        return "kgi-card__delta--down"
    return ""


def kpi_delta_class(value: Optional[float]) -> str:
    """KPIストリップ用のデルタクラスを返す。"""

    if value is None or pd.isna(value):
        return ""
    return "kpi-strip__delta--up" if float(value) >= 0 else "kpi-strip__delta--down"


def build_delta_label(prefix: str, formatted: Optional[str], raw_value: Optional[float]) -> str:
    """矢印付きのデルタ表示を生成する。"""

    if not formatted:
        return f"{prefix} -"
    arrow = "―"
    if raw_value is not None and not pd.isna(raw_value):
        numeric = float(raw_value)
        if numeric > 0:
            arrow = "▲"
        elif numeric < 0:
            arrow = "▼"
    return f"{prefix} {arrow} {formatted}"


def show_kpi_card(
    label: str,
    current: Optional[float],
    previous: Optional[float],
    *,
    unit: str = "",
    value_format: str = "number",
    digits: int = 0,
    inverse: bool = False,
) -> None:
    """st.metric を用いてKPIカードを描画する。"""

    display_value = "-"
    if current is not None and not pd.isna(current):
        numeric = float(current)
        if value_format == "percent":
            display_value = f"{numeric:.{digits}f}{unit}"
        elif digits > 0:
            display_value = f"{numeric:,.{digits}f}{unit}"
        else:
            display_value = f"{numeric:,.0f}{unit}"

    delta_text: Optional[str] = None
    if (
        current is not None
        and not pd.isna(current)
        and previous is not None
        and not pd.isna(previous)
        and float(previous) != 0
    ):
        change_ratio = (float(current) - float(previous)) / float(previous)
        delta_text = f"{change_ratio * 100:+.1f}%"

    st.metric(
        label=label,
        value=display_value,
        delta=delta_text,
        delta_color="inverse" if inverse else "normal",
    )


KPI_LANGUAGE_OPTIONS = LANGUAGE_OPTIONS


def get_kpi_help(metric_key: str, language: str) -> Optional[Dict[str, str]]:
    """指定したKPIの定義情報を取得する。"""

    entry = KPI_HELP_CONTENT.get(metric_key)
    if not entry:
        return None
    return entry.get(language) or entry.get("ja")


def render_kpi_help_panel(language: str) -> None:
    """選択した言語でKPIヘルプテーブルを表示する。"""

    rows: List[Dict[str, str]] = []
    for key, translations in KPI_HELP_CONTENT.items():
        content = translations.get(language) or translations.get("ja")
        if not content:
            continue
        rows.append(
            {
                "KPI": content.get("title", key),
                "Formula": content.get("formula", "-"),
                "Interpretation": content.get("interpretation", "-"),
            }
        )

    if not rows:
        st.info("定義情報が登録されていません。")
        return

    df = pd.DataFrame(rows)
    if language == "ja":
        df.rename(columns={"KPI": "指標", "Formula": "計算式", "Interpretation": "見方"}, inplace=True)
    st.dataframe(df, hide_index=True, use_container_width=True)


def render_kpi_help_controls() -> str:
    """KPIヘルプ用の言語切替UIを描画し、選択言語コードを返す。"""

    default_lang = st.session_state.get("kpi_help_language", get_ui_language())
    lang_index = next(
        (idx for idx, (code, _) in enumerate(KPI_LANGUAGE_OPTIONS) if code == default_lang),
        0,
    )

    controls = st.columns([3, 1])
    with controls[0]:
        selected_label = st.selectbox(
            translate("kpi_help_language_label", default="KPIヘルプ言語"),
            options=[label for _, label in KPI_LANGUAGE_OPTIONS],
            index=lang_index,
            help=translate(
                "theme_language_help",
                default="日本語/英語の表記とヘルプを切り替えます。",
            ),
        )
    label_to_code = {label: code for code, label in KPI_LANGUAGE_OPTIONS}
    language_code = label_to_code.get(selected_label, "ja")
    st.session_state["kpi_help_language"] = language_code

    with controls[1]:
        if st.button(
            translate("kpi_help_button_label", default="KPIヘルプ"),
            type="secondary",
            help=translate(
                "kpi_help_button_help", default="指標の定義一覧を表示します。"
            ),
        ):
            st.session_state["kpi_help_visible"] = not st.session_state.get(
                "kpi_help_visible", False
            )

    if st.session_state.get("kpi_help_visible"):
        with st.expander(
            translate("kpi_help_panel_title", default="KPI定義ヘルプ"), expanded=True
        ):
            render_kpi_help_panel(language_code)
            st.button(
                translate("kpi_help_close", default="ヘルプを閉じる"),
                key="close_kpi_help_panel",
                on_click=lambda: st.session_state.update({"kpi_help_visible": False}),
            )

    return language_code


def render_kgi_cards(
    selected_kpi_row: pd.Series,
    period_row: Optional[pd.DataFrame],
    cash_forecast: pd.DataFrame,
    starting_cash: float,
) -> None:
    """KGI3指標のカードを描画する。"""

    if selected_kpi_row is None or selected_kpi_row.empty:
        return

    sales_value = selected_kpi_row.get("sales")
    sales_previous: Optional[float] = None
    if (
        period_row is not None
        and not period_row.empty
        and "prev_period_sales" in period_row.columns
    ):
        prev_value = period_row.iloc[0].get("prev_period_sales")
        if prev_value is not None and not pd.isna(prev_value):
            sales_previous = float(prev_value)
    sales_gap_text, sales_gap_val = format_target_gap(
        sales_value,
        KGI_TARGETS.get("sales"),
    )

    gross_margin_rate = selected_kpi_row.get("gross_margin_rate")
    gross_prev_rate = selected_kpi_row.get("gross_margin_prev")
    gross_current_pct: Optional[float] = None
    if gross_margin_rate is not None and not pd.isna(gross_margin_rate):
        gross_current_pct = float(gross_margin_rate) * 100
    gross_previous_pct: Optional[float] = None
    if gross_prev_rate is not None and not pd.isna(gross_prev_rate):
        gross_previous_pct = float(gross_prev_rate) * 100
    gross_gap_text, gross_gap_val = format_target_gap(
        gross_margin_rate,
        KGI_TARGETS.get("gross_margin_rate"),
        percentage=True,
    )

    cash_balance = starting_cash
    previous_cash_balance: Optional[float] = None
    if cash_forecast is not None and not cash_forecast.empty:
        first_row = cash_forecast.iloc[0]
        cash_balance = float(first_row.get("cash_balance", starting_cash))
        net_cf_val = first_row.get("net_cf")
        if net_cf_val is not None and not pd.isna(net_cf_val):
            previous_cash_balance = cash_balance - float(net_cf_val)
    cash_gap_text, cash_gap_val = format_target_gap(
        cash_balance,
        KGI_TARGETS.get("cash_balance"),
        digits=0,
    )

    cards_info: List[Dict[str, Any]] = [
        {
            "label": "月次売上高",
            "current": sales_value,
            "previous": sales_previous,
            "unit": "円",
            "value_format": "number",
            "digits": 0,
            "target_text": sales_gap_text,
            "gap_value": sales_gap_val,
        },
        {
            "label": "粗利率",
            "current": gross_current_pct,
            "previous": gross_previous_pct,
            "unit": "%",
            "value_format": "percent",
            "digits": 1,
            "target_text": gross_gap_text,
            "gap_value": gross_gap_val,
        },
        {
            "label": "資金残高",
            "current": cash_balance,
            "previous": previous_cash_balance,
            "unit": "円",
            "value_format": "number",
            "digits": 0,
            "target_text": cash_gap_text,
            "gap_value": cash_gap_val,
        },
    ]

    columns = st.columns(len(cards_info))
    for column, info in zip(columns, cards_info):
        with column:
            show_kpi_card(
                info["label"],
                info.get("current"),
                info.get("previous"),
                unit=info.get("unit", ""),
                value_format=info.get("value_format", "number"),
                digits=int(info.get("digits", 0)),
                inverse=info.get("inverse", False),
            )
            target_text = info.get("target_text")
            gap_value = info.get("gap_value")
            if target_text and target_text != "-":
                icon_tone = "warning" if gap_value is not None and gap_value < 0 else "success"
                icon_code = "!" if icon_tone == "warning" else "T"
                icon_html = build_ui_icon(icon_code, tone=icon_tone)
                st.markdown(
                    "<div class='caption-with-icon'>{icon}<span>目標差 {text}</span></div>".format(
                        icon=icon_html,
                        text=html.escape(str(target_text)),
                    ),
                    unsafe_allow_html=True,
                )
            else:
                st.caption("目標差 -")


def clear_filter_selection(filter_name: str) -> None:
    """指定したフィルタの選択状態をクリアしてリロードする。"""

    state_key = FILTER_STATE_KEYS.get(filter_name)
    if not state_key:
        return

    if filter_name == "store":
        set_state_and_widget(state_key, None)
    else:
        set_state_and_widget(state_key, [])
    trigger_rerun()


def render_dashboard_meta(
    latest_label: str,
    period_label: str,
    record_count: int,
    alert_count: int,
    store_selection: Optional[str] = None,
    channel_selection: Optional[Sequence[str]] = None,
    category_selection: Optional[Sequence[str]] = None,
) -> None:
    """データのメタ情報とフィルタ状態をチップ状に表示する。"""

    chips: List[Tuple[str, str, str, str]] = [
        ("最新データ", latest_label or "-", "L", "accent"),
        ("表示期間", period_label or "-", "P", "accent"),
        ("対象レコード", f"{record_count:,} 件", "R", "accent"),
    ]
    if alert_count:
        chips.append(("アラート", f"{alert_count} 件", "!", "warning"))

    chips_html = "".join(
        "<span class='dashboard-meta__chip'>{icon}<span>{label}: {value}</span></span>".format(
            icon=build_ui_icon(icon_code, tone=tone),
            label=html.escape(label),
            value=html.escape(value),
        )
        for label, value, icon_code, tone in chips
    )
    st.markdown(f"<div class='dashboard-meta'>{chips_html}</div>", unsafe_allow_html=True)

    filter_entries: List[Dict[str, Any]] = []

    def _format_list(values: Sequence[str]) -> str:
        display_values = [str(value) for value in values if value]
        if not display_values:
            return "-"
        if len(display_values) <= 2:
            return "、".join(display_values)
        return "、".join(display_values[:2]) + f" ほか{len(display_values) - 2}件"

    if store_selection:
        filter_entries.append(
            {
                "label": "店舗",
                "value": str(store_selection),
                "filter": "store",
                "help": "店舗フィルタをクリア",
            }
        )

    if channel_selection:
        formatted = _format_list(channel_selection)
        filter_entries.append(
            {
                "label": "チャネル",
                "value": formatted,
                "filter": "channels",
                "help": "チャネルの選択をリセット",
            }
        )

    if category_selection:
        formatted = _format_list(category_selection)
        filter_entries.append(
            {
                "label": "カテゴリ",
                "value": formatted,
                "filter": "categories",
                "help": "カテゴリの選択をリセット",
            }
        )

    chip_container = st.container()
    if filter_entries:
        chip_container.markdown(
            "<div class='dashboard-filter-chips-anchor'></div>",
            unsafe_allow_html=True,
        )
        chip_columns = chip_container.columns(len(filter_entries))
        for column, entry in zip(chip_columns, filter_entries):
            button_label = f"{entry['label']}: {entry['value']} ✕"
            with column:
                st.button(
                    button_label,
                    key=f"filter_chip_{entry['filter']}",
                    on_click=lambda name=entry["filter"]: clear_filter_selection(name),
                    type="secondary",
                    help=entry.get("help"),
                )
    else:
        chip_container.markdown(
            "<div class='dashboard-meta dashboard-meta--empty'>適用中のフィルタはありません</div>",
            unsafe_allow_html=True,
        )


def _build_first_level_kpi_metrics(
    kpi_period_summary: Optional[pd.DataFrame],
    selected_kpi_row: Optional[pd.Series],
) -> List[Dict[str, Any]]:
    """第1階層KPIカードに必要な値をまとめて返す。"""

    if selected_kpi_row is None or selected_kpi_row.empty:
        return []

    prev_row: Optional[pd.Series] = None
    if (
        kpi_period_summary is not None
        and not kpi_period_summary.empty
        and "period" in kpi_period_summary.columns
    ):
        current_period = selected_kpi_row.get("period")
        if current_period is not None:
            candidates = kpi_period_summary[kpi_period_summary["period"] < current_period]
            if not candidates.empty:
                prev_row = candidates.iloc[-1]

    metrics: List[Dict[str, Any]] = []

    if "active_customers_avg" in selected_kpi_row.index:
        active_value = selected_kpi_row.get("active_customers_avg")
        prev_active = prev_row.get("active_customers_avg") if prev_row is not None else np.nan
        active_delta: Optional[float] = None
        if pd.notna(active_value) and pd.notna(prev_active):
            active_delta = float(active_value) - float(prev_active)
        metrics.append(
            {
                "key": "active_customers",
                "label": "月次顧客数",
                "value": format_number(active_value, digits=0, unit=" 人"),
                "raw_value": active_value,
                "previous_raw_value": prev_active,
                "delta_value": active_delta,
                "delta_text": format_delta(active_delta, digits=0, unit=" 人")
                if active_delta is not None
                else None,
                "value_column": "active_customers_avg",
                "format_func": lambda v, unit=" 人": format_number(v, digits=0, unit=unit),
                "chart_axis_format": ",.0f",
                "is_percentage": False,
                "definition_key": "active_customers",
            }
        )

    if "ltv" in selected_kpi_row.index:
        ltv_value = selected_kpi_row.get("ltv")
        ltv_delta = selected_kpi_row.get("ltv_delta")
        if pd.isna(ltv_delta):
            ltv_delta = None
        metrics.append(
            {
                "key": "ltv",
                "label": "LTV",
                "value": _format_currency_compact(ltv_value),
                "raw_value": ltv_value,
                "previous_raw_value": None,
                "delta_value": ltv_delta,
                "delta_text": format_delta(ltv_delta, digits=0, unit=" 円")
                if ltv_delta is not None
                else None,
                "value_column": "ltv",
                "format_func": _format_currency_compact,
                "chart_axis_format": ",.0f",
                "is_percentage": False,
                "definition_key": "ltv",
            }
        )

    if "arpu" in selected_kpi_row.index:
        arpu_value = selected_kpi_row.get("arpu")
        arpu_delta = selected_kpi_row.get("arpu_delta")
        if pd.isna(arpu_delta):
            arpu_delta = None
        metrics.append(
            {
                "key": "arpu",
                "label": "ARPU",
                "value": _format_currency_compact(arpu_value),
                "raw_value": arpu_value,
                "previous_raw_value": None,
                "delta_value": arpu_delta,
                "delta_text": format_delta(arpu_delta, digits=0, unit=" 円")
                if arpu_delta is not None
                else None,
                "value_column": "arpu",
                "format_func": _format_currency_compact,
                "chart_axis_format": ",.0f",
                "is_percentage": False,
                "definition_key": "arpu",
            }
        )

    if "churn_rate" in selected_kpi_row.index:
        churn_value = selected_kpi_row.get("churn_rate")
        churn_delta = selected_kpi_row.get("churn_delta")
        if pd.isna(churn_delta):
            churn_delta = None
        metrics.append(
            {
                "key": "churn_rate",
                "label": "解約率",
                "value": format_percent(churn_value),
                "raw_value": churn_value,
                "previous_raw_value": None,
                "delta_value": churn_delta,
                "delta_text": format_delta(churn_delta, percentage=True)
                if churn_delta is not None
                else None,
                "value_column": "churn_rate",
                "format_func": format_percent,
                "chart_axis_format": ".1%",
                "is_percentage": True,
                "definition_key": "churn_rate",
            }
        )

    if "gross_margin_rate" in selected_kpi_row.index:
        gross_margin_value = selected_kpi_row.get("gross_margin_rate")
        gross_delta = selected_kpi_row.get("gross_margin_delta")
        if pd.isna(gross_delta):
            gross_delta = None
        metrics.append(
            {
                "key": "gross_margin_rate",
                "label": "粗利率",
                "value": format_percent(gross_margin_value),
                "raw_value": gross_margin_value,
                "previous_raw_value": None,
                "delta_value": gross_delta,
                "delta_text": format_delta(gross_delta, percentage=True)
                if gross_delta is not None
                else None,
                "value_column": "gross_margin_rate",
                "format_func": format_percent,
                "chart_axis_format": ".1%",
                "is_percentage": True,
                "definition_key": "gross_margin_rate",
            }
        )

    if "repeat_rate" in selected_kpi_row.index:
        repeat_value = selected_kpi_row.get("repeat_rate")
        repeat_delta = selected_kpi_row.get("repeat_delta")
        if pd.isna(repeat_delta):
            repeat_delta = None
        metrics.append(
            {
                "key": "repeat_rate",
                "label": "リピート率",
                "value": format_percent(repeat_value),
                "raw_value": repeat_value,
                "previous_raw_value": None,
                "delta_value": repeat_delta,
                "delta_text": format_delta(repeat_delta, percentage=True)
                if repeat_delta is not None
                else None,
                "value_column": "repeat_rate",
                "format_func": format_percent,
                "chart_axis_format": ".1%",
                "is_percentage": True,
                "definition_key": "repeat_rate",
            }
        )

    if "inventory_turnover_days" in selected_kpi_row.index:
        inventory_value = selected_kpi_row.get("inventory_turnover_days")
        inventory_delta = selected_kpi_row.get("inventory_turnover_delta")
        if pd.isna(inventory_delta):
            inventory_delta = None
        metrics.append(
            {
                "key": "inventory_turnover_days",
                "label": "在庫回転日数",
                "value": format_number(inventory_value, digits=1, unit=" 日"),
                "raw_value": inventory_value,
                "previous_raw_value": None,
                "delta_value": inventory_delta,
                "delta_text": format_delta(inventory_delta, digits=1, unit=" 日")
                if inventory_delta is not None
                else None,
                "value_column": "inventory_turnover_days",
                "format_func": lambda v: format_number(v, digits=1, unit=" 日"),
                "chart_axis_format": ",.0f",
                "is_percentage": False,
                "definition_key": "inventory_turnover_days",
            }
        )

    return metrics


def render_first_level_kpi_strip(
    kpi_period_summary: pd.DataFrame,
    selected_kpi_row: pd.Series,
    *,
    help_language: str = "ja",
) -> List[Dict[str, Any]]:
    """第1階層KPIをStreamlitコンポーネントで表示する。"""

    metrics = _build_first_level_kpi_metrics(kpi_period_summary, selected_kpi_row)
    if not metrics:
        st.info("表示可能なKPIデータが不足しています。")
        st.session_state.pop("active_kpi_drilldown", None)
        return []

    st.session_state.setdefault("active_kpi_drilldown", None)
    active_key = st.session_state.get("active_kpi_drilldown")
    metric_keys = {metric["key"] for metric in metrics}
    if active_key and active_key not in metric_keys:
        st.session_state["active_kpi_drilldown"] = None
        active_key = None

    columns = st.columns(len(metrics))
    for column, metric in zip(columns, metrics):
        with column:
            st.metric(
                metric["label"],
                metric["value"],
                delta=metric.get("delta_text"),
            )
            definition = get_kpi_help(metric.get("definition_key", ""), help_language)
            if definition:
                formula_text = definition.get("formula")
                interpretation_text = definition.get("interpretation")
                caption_parts = []
                if formula_text:
                    caption_parts.append(f"🔢 {formula_text}")
                if interpretation_text:
                    caption_parts.append(f"🎯 {interpretation_text}")
                if caption_parts:
                    st.caption(" | ".join(caption_parts))
            is_active = active_key == metric["key"]
            button_label = "詳細を閉じる" if is_active else "詳細を表示"
            button_type = "primary" if is_active else "secondary"
            if st.button(
                button_label,
                key=f"kpi_card_button_{metric['key']}",
                type=button_type,
                use_container_width=True,
            ):
                if is_active:
                    st.session_state["active_kpi_drilldown"] = None
                else:
                    st.session_state["active_kpi_drilldown"] = metric["key"]

    return metrics


def render_active_kpi_details(
    kpi_period_summary: Optional[pd.DataFrame],
    metrics: Sequence[Dict[str, Any]],
    *,
    help_language: str = "ja",
) -> None:
    """選択中のKPIに応じた詳細ビューを表示する。"""

    if not metrics or kpi_period_summary is None or kpi_period_summary.empty:
        return

    active_key = st.session_state.get("active_kpi_drilldown")
    if not active_key:
        return

    metric_lookup = {metric["key"]: metric for metric in metrics}
    metric = metric_lookup.get(active_key)
    if metric is None:
        st.session_state["active_kpi_drilldown"] = None
        return

    value_column = metric.get("value_column")
    if not value_column or value_column not in kpi_period_summary.columns:
        st.info(f"{metric['label']}の詳細を表示するためのデータが不足しています。")
        return

    required_columns = {"period_start", "period_label", value_column}
    if not required_columns.issubset(kpi_period_summary.columns):
        st.info("期間情報が不足しているため詳細を表示できません。")
        return

    detail_container = st.container()
    with detail_container:
        header_col, close_col = st.columns([6, 1])
        close_clicked = close_col.button("閉じる", key="close_kpi_drilldown_button")
        header_col.subheader(f"{metric['label']}の詳細")
        if close_clicked:
            st.session_state["active_kpi_drilldown"] = None
            trigger_rerun()
            return

        st.metric(
            metric["label"],
            metric["value"],
            delta=metric.get("delta_text"),
        )
        definition = get_kpi_help(metric.get("definition_key", ""), help_language)
        if definition:
            explanation_parts = [definition.get("formula"), definition.get("interpretation")]
            explanation = " | ".join(filter(None, explanation_parts))
            if explanation:
                st.caption(f"ℹ️ {explanation}")

        history = kpi_period_summary[["period_start", "period_label", value_column]].copy()
        history["period_start"] = pd.to_datetime(history["period_start"], errors="coerce")
        history.dropna(subset=["period_start"], inplace=True)
        history.sort_values("period_start", inplace=True)
        history = history.tail(12)
        history.dropna(subset=[value_column], inplace=True)

        if history.empty:
            st.info(f"{metric['label']}の履歴データが不足しています。")
            return

        y_format = metric.get("chart_axis_format", ",.0f")
        tooltip_format = y_format if not metric.get("is_percentage") else ".1%"
        chart = (
            alt.Chart(history)
            .mark_line(point=alt.OverlayMarkDef(size=60, filled=True))
            .encode(
                x=alt.X(
                    "period_start:T",
                    title="期間",
                    axis=alt.Axis(format="%Y-%m", labelOverlap=True),
                ),
                y=alt.Y(
                    f"{value_column}:Q",
                    title=metric["label"],
                    axis=alt.Axis(format=y_format),
                ),
                tooltip=[
                    alt.Tooltip("period_label:N", title="期間"),
                    alt.Tooltip(
                        f"{value_column}:Q",
                        title=metric["label"],
                        format=tooltip_format,
                    ),
                ],
            )
            .properties(height=280, title=f"{metric['label']}の推移")
        )
        st.altair_chart(apply_altair_theme(chart), use_container_width=True)

        table_df = history[["period_label", value_column]].copy()
        format_func: Callable[[Any], str] = metric.get("format_func", lambda v: "-")
        table_df[value_column] = table_df[value_column].map(format_func)
        table_df.rename(columns={"period_label": "期間", value_column: metric["label"]}, inplace=True)
        st.dataframe(table_df, use_container_width=True)


def render_kpi_overview_tab(kpi_period_summary: pd.DataFrame) -> None:
    """KPIタブ向けに主要指標のトレンドとテーブルを表示する。"""

    if kpi_period_summary is None or kpi_period_summary.empty:
        st.info("KPI履歴が読み込まれていません。")
        return

    history = kpi_period_summary.tail(12).copy()
    history["period_start"] = pd.to_datetime(history["period_start"])
    history["period_label"] = history["period_label"].astype(str)

    metric_configs = [
        ("ltv", "LTV", "円", ACCENT_COLOR, False),
        ("arpu", "ARPU", "円", GROSS_SERIES_COLOR, False),
        ("repeat_rate", "リピート率", "％", ACCENT_COLOR, True),
        ("churn_rate", "チャーン率", "％", ERROR_COLOR, True),
    ]
    chart_columns = st.columns(2)
    for (metric, label, unit, color, is_percent), column in zip(metric_configs, chart_columns * 2):
        if metric not in history.columns:
            continue
        series = history[["period_start", "period_label", metric]].dropna()
        if series.empty:
            column.info(f"{label}の履歴データが不足しています。")
            continue
        encoding = alt.Y(
            f"{metric}:Q",
            title=f"{label} ({unit})",
            axis=alt.Axis(format=".1%" if is_percent else ",.0f"),
        )
        chart = (
            alt.Chart(series)
            .mark_line(color=color, point=alt.OverlayMarkDef(size=60, filled=True))
            .encode(
                x=alt.X("period_start:T", title="期間", axis=alt.Axis(format="%Y-%m", labelOverlap=True)),
                y=encoding,
                tooltip=[
                    alt.Tooltip("period_label:N", title="期間"),
                    alt.Tooltip(
                        f"{metric}:Q",
                        title=label,
                        format=".1%" if is_percent else ",.0f",
                    ),
                ],
            )
            .properties(title=f"{label}の推移", height=260)
        )
        column.altair_chart(apply_altair_theme(chart), use_container_width=True)

    table_columns = [
        "period_label",
        "sales",
        "gross_profit",
        "ltv",
        "arpu",
        "repeat_rate",
        "churn_rate",
    ]
    available_columns = [col for col in table_columns if col in history.columns]
    if available_columns:
        display_df = history[available_columns].rename(columns={"period_label": "期間"}).copy()
        for currency_col in ["sales", "gross_profit", "ltv", "arpu"]:
            if currency_col in display_df.columns:
                display_df[currency_col] = display_df[currency_col].map(
                    lambda v: f"{v:,.0f}" if pd.notna(v) else "-"
                )
        for pct_col in ["repeat_rate", "churn_rate"]:
            if pct_col in display_df.columns:
                display_df[pct_col] = display_df[pct_col].map(
                    lambda v: f"{v * 100:.1f}%" if pd.notna(v) else "-"
                )
        st.dataframe(display_df, use_container_width=True)
    else:
        st.info("KPIサマリーを表示する列が不足しています。")


def render_sales_tab(
    merged_df: pd.DataFrame,
    period_summary: pd.DataFrame,
    channel_share_df: pd.DataFrame,
    category_share_df: pd.DataFrame,
    selected_granularity_label: str,
) -> None:
    """売上タブの可視化と明細を描画する。"""

    if period_summary is not None and not period_summary.empty:
        latest_row = period_summary.iloc[-1]
        prev_row = period_summary.iloc[-2] if len(period_summary) > 1 else None
        card_cols = st.columns(3)

        latest_sales = float(latest_row.get("sales_amount", 0.0))
        sales_delta = latest_row.get("sales_mom")
        card_cols[0].metric(
            "当期売上高",
            f"{latest_sales:,.0f} 円",
            delta=f"{sales_delta * 100:+.1f}%" if pd.notna(sales_delta) else "-",
        )

        latest_gross = float(latest_row.get("net_gross_profit", 0.0))
        gross_delta = latest_row.get("gross_mom")
        card_cols[1].metric(
            "当期粗利",
            f"{latest_gross:,.0f} 円",
            delta=f"{gross_delta * 100:+.1f}%" if pd.notna(gross_delta) else "-",
        )

        latest_margin = latest_row.get("gross_margin_rate")
        prev_margin = prev_row.get("gross_margin_rate") if prev_row is not None else np.nan
        margin_delta = (
            (latest_margin - prev_margin) if pd.notna(latest_margin) and pd.notna(prev_margin) else np.nan
        )
        card_cols[2].metric(
            "粗利率",
            f"{latest_margin:.1%}" if pd.notna(latest_margin) else "-",
            delta=f"{margin_delta * 100:+.1f}pt" if pd.notna(margin_delta) else "-",
        )

        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>売上推移</div></div>",
            unsafe_allow_html=True,
        )
        latest_periods = period_summary.tail(12).copy()
        latest_periods["period_start"] = pd.to_datetime(latest_periods["period_start"])
        sales_chart_source = latest_periods.rename(
            columns={
                "period_start": "期間開始",
                "period_label": "期間",
                "sales_amount": "現状売上",
                "prev_year_sales": "前年同期間売上",
            }
        )
        value_columns = [
            col for col in ["現状売上", "前年同期間売上"] if col in sales_chart_source.columns
        ]
        if value_columns:
            melted = (
                sales_chart_source.melt(
                    id_vars=["期間開始", "期間"],
                    value_vars=value_columns,
                    var_name="指標",
                    value_name="金額",
                )
                .dropna(subset=["金額"])
                .sort_values("期間開始")
            )
            color_domain: List[str] = []
            color_range: List[str] = []
            for column in value_columns:
                color_domain.append(column)
                if column == "現状売上":
                    color_range.append(SALES_SERIES_COLOR)
                elif column == "前年同期間売上":
                    color_range.append(YOY_SERIES_COLOR)
                else:
                    color_range.append(SALES_SERIES_COLOR)

            sales_line = alt.Chart(melted).mark_line(
                point=alt.OverlayMarkDef(size=70, filled=True)
            ).encode(
                x=alt.X(
                    "期間開始:T",
                    title=f"{selected_granularity_label}開始日",
                    axis=alt.Axis(format="%Y-%m", labelOverlap=True),
                ),
                y=alt.Y(
                    "金額:Q",
                    title="売上高 (円)",
                    axis=alt.Axis(format=",.0f"),
                ),
                color=alt.Color(
                    "指標:N",
                    scale=alt.Scale(domain=color_domain, range=color_range),
                    legend=alt.Legend(title="系列"),
                ),
                tooltip=[
                    alt.Tooltip("期間:T", title="期間"),
                    alt.Tooltip("指標:N", title="系列"),
                    alt.Tooltip("金額:Q", title="金額", format=",.0f"),
                ],
            )

            chart_layers: List[alt.Chart] = [sales_line]
            sales_target = KGI_TARGETS.get("sales")
            if sales_target is not None and not pd.isna(sales_target):
                target_df = pd.DataFrame({"基準": ["売上目標"], "金額": [float(sales_target)]})
                target_rule = alt.Chart(target_df).mark_rule(strokeDash=[6, 4]).encode(
                    y="金額:Q",
                    color=alt.Color(
                        "基準:N",
                        scale=alt.Scale(domain=["売上目標"], range=[BASELINE_SERIES_COLOR]),
                        legend=alt.Legend(title="基準"),
                    ),
                    tooltip=[alt.Tooltip("金額:Q", title="売上目標", format=",.0f")],
                )
                chart_layers.append(target_rule)

            sales_chart = alt.layer(*chart_layers).resolve_scale(color="independent").properties(
                height=320,
            )
            sales_chart = apply_altair_theme(sales_chart)
            st.altair_chart(sales_chart, use_container_width=True)
        else:
            st.caption("売上推移を表示するための指標が不足しています。")

        latest_row = latest_periods.iloc[-1]
        peak_idx = latest_periods["sales_amount"].idxmax()
        peak_row = latest_periods.loc[peak_idx]
        latest_sales = float(latest_row.get("sales_amount", 0.0))
        yoy_value = latest_row.get("sales_yoy")
        yoy_text = f"{float(yoy_value) * 100:+.1f}%" if pd.notna(yoy_value) else "前年比データなし"
        sales_target = KGI_TARGETS.get("sales")
        target_gap_text, _ = format_target_gap(latest_sales, sales_target)
        summary_parts = [
            f"売上は{latest_row['period_label']}に{latest_sales:,.0f}円で、前年同期間比 {yoy_text}。",
            f"ピークは{peak_row['period_label']}の{float(peak_row['sales_amount']):,.0f}円です。",
        ]
        if target_gap_text != "-":
            summary_parts.append(f"目標値との差は{target_gap_text}です。")
        st.caption(" ".join(summary_parts))
        st.markdown("</div>", unsafe_allow_html=True)
    else:
        st.info("売上推移を表示するデータが不足しています。")

    if (channel_share_df is not None and not channel_share_df.empty) or (
        category_share_df is not None and not category_share_df.empty
    ):
        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>チャネル・カテゴリ内訳</div></div>",
            unsafe_allow_html=True,
        )
        chart_cols = st.columns(2)
        if channel_share_df is not None and not channel_share_df.empty:
            channel_rank = channel_share_df.sort_values("sales_amount", ascending=False).copy()
            channel_rank["構成比"] = channel_rank["sales_amount"] / channel_rank["sales_amount"].sum()
            channel_rank.rename(
                columns={"channel": "チャネル", "sales_amount": "売上高"}, inplace=True
            )
            bar = alt.Chart(channel_rank.head(10)).mark_bar(
                cornerRadiusTopLeft=3,
                cornerRadiusTopRight=3,
            ).encode(
                y=alt.Y("チャネル:N", sort="-x", title=None),
                x=alt.X("売上高:Q", title="売上高 (円)", axis=alt.Axis(format=",.0f")),
                color=alt.value(SALES_SERIES_COLOR),
                tooltip=[
                    alt.Tooltip("チャネル:N", title="チャネル"),
                    alt.Tooltip("売上高:Q", title="売上高", format=",.0f"),
                    alt.Tooltip("構成比:Q", title="構成比", format=".1%"),
                ],
            )
            labels = alt.Chart(channel_rank.head(10)).mark_text(
                align="left",
                baseline="middle",
                dx=6,
                color=TEXT_COLOR,
                fontWeight="bold",
            ).encode(
                y=alt.Y("チャネル:N", sort="-x"),
                x=alt.X("売上高:Q"),
                text=alt.Text("構成比:Q", format=".1%"),
            )
            channel_chart = apply_altair_theme((bar + labels).properties(height=260))
            chart_cols[0].altair_chart(channel_chart, use_container_width=True)

            top_channel = channel_rank.iloc[0]
            if len(channel_rank) >= 5:
                fifth_channel = channel_rank.iloc[4]
                diff_value = float(top_channel["売上高"]) - float(fifth_channel["売上高"])
                chart_cols[0].caption(
                    f"売上上位チャネルは{top_channel['チャネル']}で構成比{top_channel['構成比']:.1%}。5位との差は{diff_value:,.0f}円です。"
                )
            else:
                chart_cols[0].caption(
                    f"売上上位チャネルは{top_channel['チャネル']}で構成比{top_channel['構成比']:.1%}です。"
                )
        else:
            chart_cols[0].info("チャネル別の集計データがありません。")

        if category_share_df is not None and not category_share_df.empty:
            category_rank = category_share_df.sort_values("sales_amount", ascending=False).copy()
            category_rank["構成比"] = (
                category_rank["sales_amount"] / category_rank["sales_amount"].sum()
            )
            category_rank.rename(
                columns={"category": "カテゴリ", "sales_amount": "売上高"}, inplace=True
            )
            bar = alt.Chart(category_rank.head(10)).mark_bar(
                cornerRadiusTopLeft=3,
                cornerRadiusTopRight=3,
                color=GROSS_SERIES_COLOR,
            ).encode(
                y=alt.Y("カテゴリ:N", sort="-x", title=None),
                x=alt.X("売上高:Q", title="売上高 (円)", axis=alt.Axis(format=",.0f")),
                tooltip=[
                    alt.Tooltip("カテゴリ:N", title="カテゴリ"),
                    alt.Tooltip("売上高:Q", title="売上高", format=",.0f"),
                    alt.Tooltip("構成比:Q", title="構成比", format=".1%"),
                ],
            )
            labels = alt.Chart(category_rank.head(10)).mark_text(
                align="left",
                baseline="middle",
                dx=6,
                color=TEXT_COLOR,
                fontWeight="bold",
            ).encode(
                y=alt.Y("カテゴリ:N", sort="-x"),
                x=alt.X("売上高:Q"),
                text=alt.Text("構成比:Q", format=".1%"),
            )
            category_chart = apply_altair_theme((bar + labels).properties(height=260))
            chart_cols[1].altair_chart(category_chart, use_container_width=True)

            top_category = category_rank.iloc[0]
            chart_cols[1].caption(
                f"売上トップカテゴリは{top_category['カテゴリ']}で、構成比は{top_category['構成比']:.1%}です。"
            )
        else:
            chart_cols[1].info("カテゴリ別の集計データがありません。")
        st.markdown("</div>", unsafe_allow_html=True)

    with st.expander("売上明細（商品別・上位50件）", expanded=False):
        if merged_df is None or merged_df.empty:
            st.info("売上データがありません。")
        else:
            detail_df = (
                merged_df.groupby(["product_code", "product_name", "category"])
                .agg(
                    売上高=("sales_amount", "sum"),
                    粗利=("net_gross_profit", "sum"),
                    販売数量=("quantity", "sum"),
                )
                .reset_index()
                .sort_values("売上高", ascending=False)
                .head(50)
            )
            if detail_df.empty:
                st.info("表示できる明細がありません。")
            else:
                detail_df["粗利率"] = np.where(
                    detail_df["売上高"] != 0,
                    detail_df["粗利"] / detail_df["売上高"],
                    np.nan,
                )
                display_df = detail_df.rename(
                    columns={
                        "product_code": "商品コード",
                        "product_name": "商品名",
                        "category": "カテゴリ",
                    }
                )
                column_order = [
                    "商品コード",
                    "商品名",
                    "カテゴリ",
                    "売上高",
                    "粗利",
                    "粗利率",
                    "販売数量",
                ]
                display_df = display_df[column_order]
                column_config = {
                    "売上高": st.column_config.NumberColumn("売上高 (円)", format=",.0f"),
                    "粗利": st.column_config.NumberColumn("粗利 (円)", format=",.0f"),
                    "販売数量": st.column_config.NumberColumn("販売数量", format=",.0f"),
                    "粗利率": st.column_config.NumberColumn("粗利率 (%)", format="0.0%"),
                }
                st.dataframe(
                    display_df,
                    hide_index=True,
                    use_container_width=True,
                    column_config=column_config,
                )
                toolbar = st.columns(4)
                export_filename_base = "sales_detail"
                with toolbar[0]:
                    download_button_from_df("CSV出力", detail_df, f"{export_filename_base}.csv")
                with toolbar[1]:
                    download_excel_from_df("Excel出力", detail_df, f"{export_filename_base}.xlsx")
                with toolbar[2]:
                    download_pdf_from_df(
                        "PDF出力",
                        detail_df,
                        f"{export_filename_base}.pdf",
                        title="商品別売上明細",
                    )
                with toolbar[3]:
                    download_ppt_from_df(
                        "PowerPoint",
                        detail_df.head(20),
                        f"{export_filename_base}.pptx",
                        title="商品別売上トップ20",
                    )


def render_gross_tab(
    merged_df: pd.DataFrame,
    period_summary: pd.DataFrame,
    selected_granularity_label: str,
) -> None:
    """粗利タブのグラフと明細を描画する。"""

    if period_summary is not None and not period_summary.empty:
        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>粗利と粗利率の推移</div></div>",
            unsafe_allow_html=True,
        )
        latest_periods = period_summary.tail(12).copy()
        latest_periods["period_start"] = pd.to_datetime(latest_periods["period_start"])

        if "gross_margin_rate" not in latest_periods.columns:
            if {"net_gross_profit", "sales_amount"}.issubset(latest_periods.columns):
                latest_periods["gross_margin_rate"] = np.where(
                    latest_periods["sales_amount"] != 0,
                    latest_periods["net_gross_profit"] / latest_periods["sales_amount"],
                    np.nan,
                )
            else:
                latest_periods["gross_margin_rate"] = np.nan

        latest_periods["gross_margin_pct"] = latest_periods["gross_margin_rate"] * 100

        gross_bar = alt.Chart(latest_periods).mark_bar(color=GROSS_SERIES_COLOR).encode(
            x=alt.X(
                "period_start:T",
                title=f"{selected_granularity_label}開始日",
                axis=alt.Axis(format="%Y-%m", labelOverlap=True),
            ),
            y=alt.Y(
                "net_gross_profit:Q",
                title="粗利 (円)",
                axis=alt.Axis(format=",.0f"),
            ),
            tooltip=[
                alt.Tooltip("period_label:N", title="期間"),
                alt.Tooltip("net_gross_profit:Q", title="粗利", format=",.0f"),
            ],
        )

        gross_line = alt.Chart(latest_periods).mark_line(
            color=YOY_SERIES_COLOR, point=alt.OverlayMarkDef(size=60, filled=True)
        ).encode(
            x=alt.X("period_start:T"),
            y=alt.Y(
                "gross_margin_pct:Q",
                title="粗利率 (%)",
                axis=alt.Axis(format=".1f", orient="right"),
            ),
            tooltip=[
                alt.Tooltip("period_label:N", title="期間"),
                alt.Tooltip("gross_margin_pct:Q", title="粗利率", format=".1f"),
            ],
        )

        gross_layers: List[alt.Chart] = [gross_bar, gross_line]
        gross_target = KGI_TARGETS.get("gross_margin_rate")
        if gross_target is not None and not pd.isna(gross_target):
            gross_target_df = pd.DataFrame(
                {"基準": ["粗利率目標"], "粗利率": [float(gross_target) * 100]}
            )
            gross_target_rule = alt.Chart(gross_target_df).mark_rule(strokeDash=[6, 4]).encode(
                y=alt.Y(
                    "粗利率:Q",
                    title="粗利率 (%)",
                ),
                color=alt.Color(
                    "基準:N",
                    scale=alt.Scale(domain=["粗利率目標"], range=[BASELINE_SERIES_COLOR]),
                    legend=alt.Legend(title="基準"),
                ),
                tooltip=[alt.Tooltip("粗利率:Q", title="粗利率目標", format=".1f")],
            )
            gross_layers.append(gross_target_rule)

        gross_chart = (
            alt.layer(*gross_layers)
            .resolve_scale(y="independent", color="independent")
            .properties(height=320)
        )
        st.altair_chart(apply_altair_theme(gross_chart), use_container_width=True)

        latest_row = latest_periods.iloc[-1]
        latest_gross = float(latest_row.get("net_gross_profit", 0.0))
        gross_yoy = latest_row.get("gross_yoy")
        gross_margin = latest_row.get("gross_margin_rate")
        gross_margin_text = format_percent(gross_margin)
        gross_yoy_text = (
            f"{float(gross_yoy) * 100:+.1f}%" if pd.notna(gross_yoy) else "前年比データなし"
        )
        peak_idx = latest_periods["net_gross_profit"].idxmax()
        peak_row = latest_periods.loc[peak_idx]
        gross_target_gap_text, _ = format_target_gap(
            gross_margin, KGI_TARGETS.get("gross_margin_rate"), percentage=True
        )
        summary_parts = [
            f"最新の粗利は{latest_row['period_label']}で{latest_gross:,.0f}円、粗利率は{gross_margin_text}です。",
            f"前年同期間比は{gross_yoy_text}、粗利のピークは{peak_row['period_label']}の{float(peak_row['net_gross_profit']):,.0f}円です。",
        ]
        if gross_target_gap_text != "-":
            summary_parts.append(f"粗利率目標との差は{gross_target_gap_text}です。")
        st.caption(" ".join(summary_parts))
        st.markdown("</div>", unsafe_allow_html=True)
    else:
        st.info("粗利推移を表示するデータが不足しています。")

    if merged_df is not None and not merged_df.empty:
        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>粗利構成</div></div>",
            unsafe_allow_html=True,
        )
        chart_cols = st.columns(2)
        category_gross = (
            merged_df.groupby("category")["net_gross_profit"].sum().reset_index().sort_values("net_gross_profit", ascending=False).head(10)
        )
        if not category_gross.empty:
            category_gross.rename(
                columns={"category": "カテゴリ", "net_gross_profit": "粗利"}, inplace=True
            )
            total = category_gross["粗利"].sum()
            if total:
                category_gross["構成比"] = category_gross["粗利"] / total
            else:
                category_gross["構成比"] = 0
            bar = alt.Chart(category_gross).mark_bar(
                cornerRadiusTopLeft=3,
                cornerRadiusTopRight=3,
                color=GROSS_SERIES_COLOR,
            ).encode(
                y=alt.Y("カテゴリ:N", sort="-x", title=None),
                x=alt.X("粗利:Q", title="粗利 (円)", axis=alt.Axis(format=",.0f")),
                tooltip=[
                    alt.Tooltip("カテゴリ:N", title="カテゴリ"),
                    alt.Tooltip("粗利:Q", title="粗利", format=",.0f"),
                    alt.Tooltip("構成比:Q", title="構成比", format=".1%"),
                ],
            )
            labels = alt.Chart(category_gross).mark_text(
                align="left",
                baseline="middle",
                dx=6,
                color=TEXT_COLOR,
                fontWeight="bold",
            ).encode(
                y=alt.Y("カテゴリ:N", sort="-x"),
                x=alt.X("粗利:Q"),
                text=alt.Text("構成比:Q", format=".1%"),
            )
            chart_cols[0].altair_chart(
                apply_altair_theme((bar + labels).properties(height=260)),
                use_container_width=True,
            )
            top_category = category_gross.iloc[0]
            chart_cols[0].caption(
                f"粗利が最も高いカテゴリは{top_category['カテゴリ']}で、構成比は{top_category['構成比']:.1%}です。"
            )
        else:
            chart_cols[0].info("カテゴリ別の粗利データがありません。")

        product_gross = (
            merged_df.groupby("product_name")["net_gross_profit"].sum().reset_index().sort_values("net_gross_profit", ascending=False).head(10)
        )
        if not product_gross.empty:
            product_gross.rename(
                columns={"product_name": "商品", "net_gross_profit": "粗利"}, inplace=True
            )
            total = product_gross["粗利"].sum()
            if total:
                product_gross["構成比"] = product_gross["粗利"] / total
            else:
                product_gross["構成比"] = 0
            bar = alt.Chart(product_gross).mark_bar(
                cornerRadiusTopLeft=3,
                cornerRadiusTopRight=3,
                color=GROSS_SERIES_COLOR,
            ).encode(
                y=alt.Y("商品:N", sort="-x", title=None),
                x=alt.X("粗利:Q", title="粗利 (円)", axis=alt.Axis(format=",.0f")),
                tooltip=[
                    alt.Tooltip("商品:N", title="商品"),
                    alt.Tooltip("粗利:Q", title="粗利", format=",.0f"),
                    alt.Tooltip("構成比:Q", title="構成比", format=".1%"),
                ],
            )
            labels = alt.Chart(product_gross).mark_text(
                align="left",
                baseline="middle",
                dx=6,
                color=TEXT_COLOR,
                fontWeight="bold",
            ).encode(
                y=alt.Y("商品:N", sort="-x"),
                x=alt.X("粗利:Q"),
                text=alt.Text("構成比:Q", format=".1%"),
            )
            chart_cols[1].altair_chart(
                apply_altair_theme((bar + labels).properties(height=260)),
                use_container_width=True,
            )
            top_product = product_gross.iloc[0]
            chart_cols[1].caption(
                f"粗利トップ商品は{top_product['商品']}で、構成比は{top_product['構成比']:.1%}です。"
            )
        else:
            chart_cols[1].info("商品別の粗利データがありません。")
        st.markdown("</div>", unsafe_allow_html=True)

    with st.expander("原価率・粗利テーブル", expanded=False):
        if merged_df is None or merged_df.empty:
            st.info("データがありません。")
        else:
            detail_df = (
                merged_df.groupby(["product_code", "product_name", "category"])
                .agg(
                    売上高=("sales_amount", "sum"),
                    粗利=("net_gross_profit", "sum"),
                    推定原価=("estimated_cost", "sum"),
                    原価率=("cost_rate", "mean"),
                )
                .reset_index()
            )
            if detail_df.empty:
                st.info("表示できる明細がありません。")
            else:
                detail_df["粗利率"] = np.where(
                    detail_df["売上高"] != 0,
                    detail_df["粗利"] / detail_df["売上高"],
                    np.nan,
                )
                detail_df.sort_values("粗利", ascending=False, inplace=True)
                display_df = detail_df.copy()
                for column in ["売上高", "粗利", "推定原価"]:
                    display_df[column] = display_df[column].map(lambda v: f"{v:,.0f}")
                display_df["原価率"] = display_df["原価率"].map(
                    lambda v: f"{v * 100:.1f}%" if pd.notna(v) else "-"
                )
                display_df["粗利率"] = display_df["粗利率"].map(
                    lambda v: f"{v * 100:.1f}%" if pd.notna(v) else "-"
                )
                st.dataframe(display_df.head(50), hide_index=True, use_container_width=True)
                toolbar = st.columns(4)
                export_base = "gross_profit_detail"
                with toolbar[0]:
                    download_button_from_df("CSV出力", detail_df, f"{export_base}.csv")
                with toolbar[1]:
                    download_excel_from_df("Excel出力", detail_df, f"{export_base}.xlsx")
                with toolbar[2]:
                    download_pdf_from_df(
                        "PDF出力",
                        detail_df,
                        f"{export_base}.pdf",
                        title="粗利・原価テーブル",
                    )
                with toolbar[3]:
                    download_ppt_from_df(
                        "PowerPoint",
                        detail_df.head(20),
                        f"{export_base}.pptx",
                        title="粗利トップ20",
                    )


def render_store_comparison_chart(
    analysis_df: pd.DataFrame,
    fixed_cost: float,
    *,
    selected_store: Optional[str] = None,
) -> None:
    """店舗別の売上・粗利・営業利益(推計)を横棒で比較表示する。"""

    if analysis_df is None or analysis_df.empty:
        st.info("店舗別の比較に利用できるデータがありません。")
        return
    if "store" not in analysis_df.columns or analysis_df["store"].nunique(dropna=True) <= 1:
        st.caption("※ 店舗情報が不足しているため全社集計のみを表示しています。")
        return

    store_summary = (
        analysis_df.groupby("store")[["sales_amount", "net_gross_profit"]]
        .sum()
        .reset_index()
    )
    if store_summary.empty:
        st.info("店舗別に集計できる売上データがありません。")
        return

    total_sales = float(store_summary["sales_amount"].sum())
    if total_sales <= 0:
        st.info("売上高が0のため比較グラフを表示できません。")
        return

    fixed_cost_value = float(fixed_cost or 0.0)
    allocation_ratio = store_summary["sales_amount"] / total_sales
    store_summary["estimated_operating_profit"] = (
        store_summary["net_gross_profit"] - allocation_ratio * fixed_cost_value
    )

    metric_map = {
        "sales_amount": "売上高",
        "net_gross_profit": "粗利",
        "estimated_operating_profit": "営業利益(推計)",
    }
    melted = store_summary.melt(
        id_vars="store",
        value_vars=list(metric_map.keys()),
        var_name="metric",
        value_name="value",
    )
    if melted.empty:
        st.info("店舗別の比較に利用できる指標がありません。")
        return

    melted["metric_label"] = melted["metric"].map(metric_map)
    color_sequence = [SALES_SERIES_COLOR, GROSS_SERIES_COLOR, OPERATING_SERIES_COLOR]
    chart_title = (
        f"{selected_store}の売上・利益比較"
        if selected_store
        else "店舗別売上・利益比較"
    )
    comparison_chart = px.bar(
        melted,
        x="value",
        y="store",
        color="metric_label",
        orientation="h",
        barmode="group",
        labels={"value": "金額（円）", "store": "店舗", "metric_label": "指標"},
        color_discrete_sequence=color_sequence,
    )
    comparison_chart = apply_chart_theme(comparison_chart)
    comparison_chart.update_layout(
        title=chart_title,
        legend=dict(title="指標", orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1.0),
        xaxis_title="金額（円）",
        yaxis_title="店舗",
    )
    comparison_chart.update_traces(hovertemplate="店舗=%{y}<br>%{legendgroup}=%{x:,.0f}円<extra></extra>")
    st.plotly_chart(comparison_chart, use_container_width=True)

    top_store = store_summary.sort_values("sales_amount", ascending=False).iloc[0]
    st.caption(
        f"売上トップ店舗は{top_store['store']}で{top_store['sales_amount']:,.0f}円、推計営業利益は{top_store['estimated_operating_profit']:,.0f}円です。"
    )


def render_abc_analysis(df: pd.DataFrame) -> None:
    """ABC分析を縦棒と累積折れ線の組み合わせで描画する。"""

    if df is None or df.empty or "product_name" not in df.columns:
        st.info("ABC分析に利用できる商品データがありません。")
        return

    product_sales = (
        df.groupby(["product_code", "product_name"])["sales_amount"]
        .sum()
        .reset_index()
        .sort_values("sales_amount", ascending=False)
    )
    if product_sales.empty:
        st.info("ABC分析に利用できる売上データがありません。")
        return

    product_sales["累積売上"] = product_sales["sales_amount"].cumsum()
    total_sales = float(product_sales["sales_amount"].sum())
    if total_sales <= 0:
        st.info("売上総額が0のためABC分析を表示できません。")
        return

    product_sales["累積構成比"] = product_sales["累積売上"] / total_sales
    product_sales["ランク"] = np.where(
        product_sales["累積構成比"] <= 0.8,
        "A",
        np.where(product_sales["累積構成比"] <= 0.95, "B", "C"),
    )
    product_sales = product_sales.head(30)

    rank_colors = {"A": SALES_SERIES_COLOR, "B": GROSS_SERIES_COLOR, "C": YOY_SERIES_COLOR}
    bar_colors = [rank_colors.get(rank, SALES_SERIES_COLOR) for rank in product_sales["ランク"]]

    fig = go.Figure()
    fig.add_bar(
        x=product_sales["product_name"],
        y=product_sales["sales_amount"],
        name="売上高",
        marker_color=bar_colors,
        hovertemplate="商品=%{x}<br>売上高=%{y:,.0f}円<extra></extra>",
    )
    fig.add_scatter(
        x=product_sales["product_name"],
        y=product_sales["累積構成比"] * 100,
        mode="lines+markers",
        name="累積構成比",
        yaxis="y2",
        line=dict(color=GROSS_SERIES_COLOR, width=3),
        marker=dict(size=8),
        hovertemplate="商品=%{x}<br>累積構成比=%{y:.1f}%<extra></extra>",
    )
    fig.update_layout(
        xaxis_title="商品",
        yaxis=dict(title="売上高（円）", showgrid=True, gridcolor="rgba(11,31,51,0.08)"),
        yaxis2=dict(
            title="累積構成比（％）",
            overlaying="y",
            side="right",
            range=[0, 110],
            tickformat=".0f",
            showgrid=False,
        ),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1.0),
        margin=dict(l=40, r=60, t=60, b=80),
    )
    fig.add_shape(
        type="line",
        x0=-0.5,
        x1=len(product_sales) - 0.5,
        y0=80,
        y1=80,
        yref="y2",
        line=dict(color=SUCCESS_COLOR, dash="dash"),
    )
    fig.update_layout(plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)")
    st.plotly_chart(fig, use_container_width=True)

    boundary_index = product_sales[product_sales["累積構成比"] > 0.8].index.min()
    if boundary_index is not None and not np.isnan(boundary_index):
        boundary_product = product_sales.iloc[int(boundary_index)]
        st.caption(
            f"累積構成比80%の境界は{boundary_product['product_name']}で、売上高は{boundary_product['sales_amount']:,.0f}円です。"
        )


def render_inventory_heatmap(
    merged_df: pd.DataFrame, selected_kpi_row: Optional[pd.Series]
) -> None:
    """店舗×カテゴリの在庫状況をヒートマップで表示する。"""

    if merged_df is None or merged_df.empty:
        st.info("在庫ヒートマップを準備中です。元データのアップロードをお待ちください。")
        return
    required_columns = {"store", "category", "estimated_cost"}
    if not required_columns.issubset(merged_df.columns):
        st.info("在庫推計に必要な列が不足しています。データ整備が完了次第、自動で表示されます。")
        return

    turnover_days = None
    if selected_kpi_row is not None:
        turnover_days = selected_kpi_row.get("inventory_turnover_days")
    if turnover_days is None or pd.isna(turnover_days) or float(turnover_days) <= 0:
        turnover_days = 45.0

    inventory_value = (
        merged_df.groupby(["store", "category"])["estimated_cost"].sum().reset_index()
    )
    if inventory_value.empty:
        st.info("カテゴリ別の在庫データを準備中です。もうしばらくお待ちください。")
        return

    inventory_value["推定在庫金額"] = (
        inventory_value["estimated_cost"] / 30.0 * float(turnover_days)
    )
    heatmap_source = inventory_value.pivot(
        index="store", columns="category", values="推定在庫金額"
    ).fillna(0.0)
    if heatmap_source.empty:
        st.info("在庫ヒートマップのデータが揃っていません。準備完了までしばらくお待ちください。")
        return

    fig = go.Figure(
        data=
        [
            go.Heatmap(
                z=heatmap_source.values,
                x=heatmap_source.columns.astype(str),
                y=heatmap_source.index.astype(str),
                colorscale=HEATMAP_BLUE_SCALE,
                colorbar=dict(title="推定在庫金額（円）", tickformat=",.0f"),
                hovertemplate="店舗=%{y}<br>カテゴリ=%{x}<br>推定在庫=%{z:,.0f}円<extra></extra>",
            )
        ]
    )
    fig.update_layout(
        height=420,
        xaxis_title="カテゴリ",
        yaxis_title="店舗",
        margin=dict(l=60, r=60, t=50, b=60),
        plot_bgcolor="rgba(0,0,0,0)",
        paper_bgcolor="rgba(0,0,0,0)",
    )
    st.plotly_chart(fig, use_container_width=True)
    st.caption(
        f"在庫回転日数{float(turnover_days):.0f}日を基準に推定した金額です。濃い青は安全在庫を上回る余剰在庫を示唆します。"
    )


def render_inventory_tab(
    merged_df: pd.DataFrame,
    kpi_period_summary: pd.DataFrame,
    selected_kpi_row: pd.Series,
) -> None:
    """在庫タブの主要指標と推計表を表示する。"""

    turnover_days_value: float = 0.0
    if selected_kpi_row is not None:
        candidate = selected_kpi_row.get("inventory_turnover_days")
        if candidate is not None and not pd.isna(candidate) and float(candidate) > 0:
            turnover_days_value = float(candidate)

    if turnover_days_value <= 0:
        turnover_days_value = 45.0

    if kpi_period_summary is not None and not kpi_period_summary.empty:
        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>在庫KPIの推移</div></div>",
            unsafe_allow_html=True,
        )
        history = kpi_period_summary.tail(12).copy()
        history["period_start"] = pd.to_datetime(history["period_start"], errors="coerce")
        history["period_end"] = pd.to_datetime(history.get("period_end"), errors="coerce")
        history["period_days"] = (
            (history["period_end"] - history["period_start"]).dt.days.fillna(0).astype(int) + 1
        )
        if {"sales", "gross_profit"}.issubset(history.columns):
            history["cogs"] = history["sales"] - history["gross_profit"]
        else:
            history["cogs"] = np.nan
        history["inventory_balance"] = np.where(
            (history["cogs"].notna())
            & (history["cogs"] > 0)
            & (history["inventory_turnover_days"].notna())
            & (history["inventory_turnover_days"] > 0)
            & (history["period_days"] > 0),
            history["cogs"] / history["period_days"] * history["inventory_turnover_days"],
            np.nan,
        )
        chart_cols = st.columns(2)
        balance_source = history.dropna(subset=["inventory_balance"]).copy()
        if not balance_source.empty:
            balance_chart = (
                alt.Chart(balance_source)
                .mark_area(line={"color": INVENTORY_SERIES_COLOR}, color="rgba(37,99,235,0.18)")
                .encode(
                    x=alt.X(
                        "period_start:T",
                        title="期間開始",
                        axis=alt.Axis(format="%Y-%m", labelOverlap=True),
                    ),
                    y=alt.Y("inventory_balance:Q", title="推定在庫残高", axis=alt.Axis(format=",.0f")),
                    tooltip=[
                        alt.Tooltip("period_label:N", title="期間"),
                        alt.Tooltip("inventory_balance:Q", title="推定在庫残高", format=",.0f"),
                        alt.Tooltip("inventory_turnover_days:Q", title="在庫回転日数", format=",.1f"),
                    ],
                )
                .properties(height=260)
            )
            st.altair_chart(apply_altair_theme(balance_chart), use_container_width=True)
            st.caption(
                "在庫残高は当期の売上原価と在庫回転日数から推計しています。濃い部分は在庫金額の増加局面です。"
            )
        else:
            st.info("在庫残高の推移を算出するための指標が揃っていません。データのアップロード状況をご確認ください。")

        turnover_line = alt.Chart(history).mark_line(
            color=INVENTORY_SERIES_COLOR, point=alt.OverlayMarkDef(size=60, filled=True)
        ).encode(
            x=alt.X("period_start:T", title="期間開始", axis=alt.Axis(format="%Y-%m", labelOverlap=True)),
            y=alt.Y("inventory_turnover_days:Q", title="在庫回転日数", axis=alt.Axis(format=",.0f")),
            tooltip=[
                alt.Tooltip("period_label:N", title="期間"),
                alt.Tooltip("inventory_turnover_days:Q", title="在庫回転日数", format=",.1f"),
            ],
        )
        turnover_target = pd.DataFrame({"target": [45.0]})
        turnover_rule = (
            alt.Chart(turnover_target)
            .mark_rule(color=SUCCESS_COLOR, strokeDash=[6, 4])
            .encode(y="target:Q")
        )
        chart_cols[0].altair_chart(
            apply_altair_theme((turnover_line + turnover_rule).properties(height=260)),
            use_container_width=True,
        )

        stockout_chart = alt.Chart(history).mark_line(
            color=YOY_SERIES_COLOR, point=alt.OverlayMarkDef(size=60, filled=True)
        ).encode(
            x=alt.X("period_start:T", title="期間開始", axis=alt.Axis(format="%Y-%m", labelOverlap=True)),
            y=alt.Y(
                "stockout_rate:Q",
                title="欠品率",
                axis=alt.Axis(format=".1%"),
            ),
            tooltip=[
                alt.Tooltip("period_label:N", title="期間"),
                alt.Tooltip("stockout_rate:Q", title="欠品率", format=".1%"),
            ],
        )
        stockout_threshold = pd.DataFrame({"target": [0.05]})
        stockout_rule = (
            alt.Chart(stockout_threshold)
            .mark_rule(color=WARNING_COLOR, strokeDash=[6, 4])
            .encode(y="target:Q")
        )
        chart_cols[1].altair_chart(
            apply_altair_theme((stockout_chart + stockout_rule).properties(height=260)),
            use_container_width=True,
        )
        latest_inventory_row = history.iloc[-1]
        turnover_value = latest_inventory_row.get("inventory_turnover_days")
        stockout_value = latest_inventory_row.get("stockout_rate")
        chart_cols[0].caption(
            f"最新の在庫回転日数は{turnover_value:,.1f}日で、直近最大値は{history['inventory_turnover_days'].max():,.1f}日です。"
            if pd.notna(turnover_value)
            else "在庫回転日数の最新値が取得できません。"
        )
        chart_cols[1].caption(
            f"最新の欠品率は{stockout_value:.1%}で、最小値は{history['stockout_rate'].min():.1%}です。"
            if pd.notna(stockout_value)
            else "欠品率の最新値が取得できません。"
        )
        st.markdown("</div>", unsafe_allow_html=True)
    else:
        st.info("在庫指標を準備中です。必要なデータが揃い次第こちらに表示されます。")

    if merged_df is not None and not merged_df.empty:
        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>在庫構成の推計</div></div>",
            unsafe_allow_html=True,
        )
        chart_cols = st.columns(2)
        category_qty = (
            merged_df.groupby("category")["quantity"].sum().reset_index().sort_values("quantity", ascending=False).head(10)
        )
        if not category_qty.empty:
            category_qty.rename(columns={"quantity": "販売数量"}, inplace=True)
            total_qty = category_qty["販売数量"].sum()
            if total_qty:
                category_qty["構成比"] = category_qty["販売数量"] / total_qty
            else:
                category_qty["構成比"] = 0
            bar = alt.Chart(category_qty).mark_bar(
                cornerRadiusTopLeft=3,
                cornerRadiusTopRight=3,
                color=INVENTORY_SERIES_COLOR,
            ).encode(
                y=alt.Y("category:N", sort="-x", title="カテゴリ"),
                x=alt.X("販売数量:Q", title="販売数量", axis=alt.Axis(format=",.0f")),
                tooltip=[
                    alt.Tooltip("category:N", title="カテゴリ"),
                    alt.Tooltip("販売数量:Q", title="販売数量", format=",.0f"),
                    alt.Tooltip("構成比:Q", title="構成比", format=".1%"),
                ],
            )
            labels = alt.Chart(category_qty).mark_text(
                align="left",
                baseline="middle",
                dx=6,
                color=TEXT_COLOR,
                fontWeight="bold",
            ).encode(
                y=alt.Y("category:N", sort="-x"),
                x=alt.X("販売数量:Q"),
                text=alt.Text("構成比:Q", format=".1%"),
            )
            chart_cols[0].altair_chart(
                apply_altair_theme((bar + labels).properties(height=260)),
                use_container_width=True,
            )
            top_category = category_qty.iloc[0]
            chart_cols[0].caption(
                f"在庫数量が最も多いカテゴリは{top_category['category']}で、構成比は{top_category['構成比']:.1%}です。"
            )
        else:
            chart_cols[0].info("カテゴリ別の在庫情報を準備中です。元データを確認してください。")

        product_qty = (
            merged_df.groupby("product_name")["quantity"].sum().reset_index().sort_values("quantity", ascending=False).head(10)
        )
        if not product_qty.empty:
            product_qty.rename(columns={"quantity": "販売数量"}, inplace=True)
            total_qty = product_qty["販売数量"].sum()
            if total_qty:
                product_qty["構成比"] = product_qty["販売数量"] / total_qty
            else:
                product_qty["構成比"] = 0
            bar = alt.Chart(product_qty).mark_bar(
                cornerRadiusTopLeft=3,
                cornerRadiusTopRight=3,
                color=INVENTORY_SERIES_COLOR,
            ).encode(
                y=alt.Y("product_name:N", sort="-x", title="商品"),
                x=alt.X("販売数量:Q", title="販売数量", axis=alt.Axis(format=",.0f")),
                tooltip=[
                    alt.Tooltip("product_name:N", title="商品"),
                    alt.Tooltip("販売数量:Q", title="販売数量", format=",.0f"),
                    alt.Tooltip("構成比:Q", title="構成比", format=".1%"),
                ],
            )
            labels = alt.Chart(product_qty).mark_text(
                align="left",
                baseline="middle",
                dx=6,
                color=TEXT_COLOR,
                fontWeight="bold",
            ).encode(
                y=alt.Y("product_name:N", sort="-x"),
                x=alt.X("販売数量:Q"),
                text=alt.Text("構成比:Q", format=".1%"),
            )
            chart_cols[1].altair_chart(
                apply_altair_theme((bar + labels).properties(height=260)),
                use_container_width=True,
            )
            top_product = product_qty.iloc[0]
            chart_cols[1].caption(
                f"在庫数量が最も多い商品は{top_product['product_name']}で、構成比は{top_product['構成比']:.1%}です。"
            )
        else:
            chart_cols[1].info("商品別の在庫情報を準備中です。データが揃い次第更新されます。")
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>在庫ヒートマップ</div></div>",
            unsafe_allow_html=True,
        )
        render_inventory_heatmap(merged_df, selected_kpi_row)

        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>売れ残りリスク商品</div></div>",
            unsafe_allow_html=True,
        )
        slow_mover_df = (
            merged_df.groupby(["product_code", "product_name", "category"], as_index=False)[
                ["sales_amount", "quantity", "estimated_cost"]
            ]
            .sum()
        )
        if slow_mover_df.empty:
            st.info("売れ残りリスクを評価するための商品別データが不足しています。")
        else:
            slow_mover_df.rename(columns={"quantity": "販売数量", "estimated_cost": "推定原価"}, inplace=True)
            slow_mover_df["推定在庫金額"] = np.where(
                slow_mover_df["推定原価"] > 0,
                slow_mover_df["推定原価"] / 30.0 * turnover_days_value,
                np.nan,
            )
            slow_mover_df["回転率(推計)"] = np.where(
                slow_mover_df["推定在庫金額"] > 0,
                slow_mover_df["sales_amount"] / slow_mover_df["推定在庫金額"],
                np.nan,
            )
            candidates = slow_mover_df.sort_values(
                ["回転率(推計)", "販売数量", "推定在庫金額"], ascending=[True, True, False]
            ).head(5)
            display_cols = candidates[
                [
                    "product_code",
                    "product_name",
                    "category",
                    "販売数量",
                    "sales_amount",
                    "推定在庫金額",
                    "回転率(推計)",
                ]
            ].copy()
            display_cols.rename(
                columns={
                    "product_code": "商品コード",
                    "product_name": "商品名",
                    "category": "カテゴリ",
                    "sales_amount": "売上高",
                },
                inplace=True,
            )
            for column in ["販売数量", "売上高", "推定在庫金額"]:
                display_cols[column] = display_cols[column].map(
                    lambda v: f"{v:,.0f}" if pd.notna(v) else "-"
                )
            display_cols["回転率(推計)"] = display_cols["回転率(推計)"].map(
                lambda v: f"{v:,.2f}" if pd.notna(v) else "-"
            )
            st.dataframe(display_cols, use_container_width=True, hide_index=True)
            st.caption(
                "回転率が1未満の商品は売れ残りリスクが高いため、販促や在庫調整の検討をおすすめします。"
            )
        st.markdown("</div>", unsafe_allow_html=True)

    with st.expander("在庫推計テーブル", expanded=False):
        if merged_df is None or merged_df.empty:
            st.info("在庫推計データを準備中です。アップロード状況をご確認ください。")
        else:
            detail_df = (
                merged_df.groupby(["product_code", "product_name", "category"])
                .agg(
                    販売数量=("quantity", "sum"),
                    売上高=("sales_amount", "sum"),
                    推定原価=("estimated_cost", "sum"),
                )
                .reset_index()
            )
            if detail_df.empty:
                st.info("表示できる明細がありません。")
            else:
                if turnover_days_value > 0:
                    detail_df["推定在庫金額"] = detail_df["推定原価"] / 30.0 * float(turnover_days_value)
                else:
                    detail_df["推定在庫金額"] = np.nan
                detail_df.sort_values("推定在庫金額", ascending=False, inplace=True)
                display_df = detail_df.copy()
                display_df["販売数量"] = display_df["販売数量"].map(lambda v: f"{v:,.0f}")
                for column in ["売上高", "推定原価", "推定在庫金額"]:
                    display_df[column] = display_df[column].map(lambda v: f"{v:,.0f}" if pd.notna(v) else "-")
                st.dataframe(display_df.head(50), hide_index=True, use_container_width=True)
                toolbar = st.columns(4)
                export_base = "inventory_overview"
                with toolbar[0]:
                    download_button_from_df("CSV出力", detail_df, f"{export_base}.csv")
                with toolbar[1]:
                    download_excel_from_df("Excel出力", detail_df, f"{export_base}.xlsx")
                with toolbar[2]:
                    download_pdf_from_df(
                        "PDF出力",
                        detail_df,
                        f"{export_base}.pdf",
                        title="在庫推計テーブル",
                    )
                with toolbar[3]:
                    download_ppt_from_df(
                        "PowerPoint",
                        detail_df.head(20),
                        f"{export_base}.pptx",
                        title="在庫上位商品",
                    )


def render_cash_tab(
    cash_plan: pd.DataFrame,
    cash_forecast: pd.DataFrame,
    starting_cash: float,
    monthly_summary: pd.DataFrame,
) -> None:
    """資金タブのグラフと明細を描画する。"""

    base_cash_plan = cash_plan.copy() if isinstance(cash_plan, pd.DataFrame) else pd.DataFrame()
    base_cash_forecast = cash_forecast.copy() if isinstance(cash_forecast, pd.DataFrame) else pd.DataFrame()
    if (base_cash_forecast is None or base_cash_forecast.empty) and not base_cash_plan.empty:
        base_cash_forecast = forecast_cashflow(base_cash_plan, starting_cash)

    active_cash_plan = base_cash_plan.copy()
    active_cash_forecast = base_cash_forecast.copy()
    forecast_df = pd.DataFrame()
    forecast_summary_df = pd.DataFrame()
    savings_summary_df = pd.DataFrame()
    updated_cash_plan: Optional[pd.DataFrame] = None
    updated_cash_forecast: Optional[pd.DataFrame] = None
    artifact = None

    st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
    st.markdown(
        "<div class='chart-section__header'><div class='chart-section__title'>AI売上予測と資金インパクト</div></div>",
        unsafe_allow_html=True,
    )

    if monthly_summary is None or monthly_summary.empty:
        st.info("売上データが不足しているため機械学習予測を実行できません。直近数ヶ月分の売上データをアップロードしてください。")
    else:
        model_options = [("arima", "ARIMA (統計モデル)"), ("lstm", "LSTM (ニューラルネット)")]
        default_model = st.session_state.get("cash_forecast_model_type", "arima")
        default_index = next((idx for idx, opt in enumerate(model_options) if opt[0] == default_model), 0)
        controls = st.columns([2, 1, 1])
        selected_model = controls[0].selectbox(
            "モデルタイプ",
            options=model_options,
            index=default_index,
            format_func=lambda opt: opt[1],
            key="cash_forecast_model_selector",
            help="統計的なARIMAモデルと深層学習LSTMモデルから選択できます。",
        )
        model_type = selected_model[0]
        st.session_state["cash_forecast_model_type"] = model_type

        default_horizon = int(st.session_state.get("cash_forecast_horizon_value", min(6, max(3, len(base_cash_plan) or 6))))
        horizon = controls[1].slider(
            "予測期間 (月)",
            min_value=3,
            max_value=12,
            value=default_horizon,
            step=1,
            key="cash_forecast_horizon",
        )
        st.session_state["cash_forecast_horizon_value"] = horizon

        default_confidence = float(st.session_state.get("cash_forecast_confidence_value", 0.90))
        confidence = controls[2].slider(
            "信頼水準",
            min_value=0.80,
            max_value=0.99,
            value=default_confidence,
            step=0.01,
            key="cash_forecast_confidence",
        )
        st.session_state["cash_forecast_confidence_value"] = confidence

        window_size = min(12, max(6, len(monthly_summary) - 1))
        if model_type == "lstm":
            max_window = max(6, min(24, len(monthly_summary) - 1))
            window_size = controls[0].slider(
                "学習ウィンドウ (月)",
                min_value=6,
                max_value=max_window,
                value=min(window_size, max_window),
                step=1,
                key="cash_forecast_window",
                help="LSTMが過去何ヶ月分のデータを学習に使うかを指定します。",
            )

        try:
            with st.spinner("売上予測を更新しています..."):
                artifact = train_forecast_model(
                    monthly_summary,
                    model_type=model_type,
                    freq="M",
                    window_size=window_size,
                    lstm_params={"units": 32, "epochs": 120, "batch_size": 8},
                )
                forecast_df = predict_sales_forecast(
                    artifact,
                    periods=horizon,
                    confidence_level=confidence,
                )
        except Exception as exc:  # pragma: no cover - UI通知用
            st.warning(f"売上予測の計算に失敗しました: {exc}")
        else:
            history = monthly_summary[["order_month", "sales_amount"]].copy()
            history = history.tail(max(24, horizon + 3))
            history["period_start"] = history["order_month"].dt.to_timestamp()

            forecast_chart_df = forecast_df.copy()
            forecast_chart_df["period_start"] = forecast_chart_df["month"].dt.to_timestamp()
            forecast_chart_df["lower_ci"] = forecast_chart_df["lower_ci"].clip(lower=0.0)
            forecast_chart_df["upper_ci"] = forecast_chart_df["upper_ci"].clip(lower=0.0)

            history_line = alt.Chart(history).mark_line(
                color=SALES_SERIES_COLOR, point=alt.OverlayMarkDef(size=60, filled=True)
            ).encode(
                x=alt.X("period_start:T", title="期間開始", axis=alt.Axis(format="%Y-%m", labelOverlap=True)),
                y=alt.Y("sales_amount:Q", title="売上高 (円)", axis=alt.Axis(format=",.0f")),
                tooltip=[
                    alt.Tooltip("order_month:N", title="月"),
                    alt.Tooltip("sales_amount:Q", title="実績売上", format=",.0f"),
                ],
            )

            forecast_band = alt.Chart(forecast_chart_df).mark_area(
                opacity=0.2, color=FORECAST_SERIES_COLOR
            ).encode(
                x=alt.X("period_start:T", title="期間開始", axis=alt.Axis(format="%Y-%m", labelOverlap=True)),
                y=alt.Y("lower_ci:Q", title=""),
                y2="upper_ci:Q",
                tooltip=[
                    alt.Tooltip("month:N", title="月"),
                    alt.Tooltip("lower_ci:Q", title="下限", format=",.0f"),
                    alt.Tooltip("upper_ci:Q", title="上限", format=",.0f"),
                ],
            )

            forecast_line = alt.Chart(forecast_chart_df).mark_line(
                color=FORECAST_SERIES_COLOR, point=alt.OverlayMarkDef(size=60, filled=True)
            ).encode(
                x=alt.X("period_start:T", title="期間開始", axis=alt.Axis(format="%Y-%m", labelOverlap=True)),
                y=alt.Y("forecast:Q", title="売上予測 (円)", axis=alt.Axis(format=",.0f")),
                tooltip=[
                    alt.Tooltip("month:N", title="月"),
                    alt.Tooltip("forecast:Q", title="予測売上", format=",.0f"),
                    alt.Tooltip("lower_ci:Q", title="下限", format=",.0f"),
                    alt.Tooltip("upper_ci:Q", title="上限", format=",.0f"),
                ],
            )

            forecast_chart = (
                alt.layer(history_line, forecast_band, forecast_line)
                .resolve_scale(color="independent")
                .properties(height=320)
            )
            st.altair_chart(apply_altair_theme(forecast_chart), use_container_width=True)

            if not forecast_df.empty:
                last_row = forecast_df.iloc[-1]
                st.caption(
                    f"{len(forecast_df)}ヶ月先の予測売上は{last_row['forecast']:,.0f}円 (信頼区間 {last_row['lower_ci']:,.0f}〜{last_row['upper_ci']:,.0f}円) を想定しています。"
                )

            metric_cols = st.columns(3)
            baseline_mape = artifact.metrics.get("baseline_mape") if artifact else None
            model_mape = artifact.metrics.get("model_mape") if artifact else None
            improvement = artifact.metrics.get("mape_improvement") if artifact else None
            if baseline_mape is not None and not np.isnan(baseline_mape):
                metric_cols[0].metric("ベースラインMAPE", f"{baseline_mape:.1%}")
            if model_mape is not None and not np.isnan(model_mape):
                metric_cols[1].metric("モデルMAPE", f"{model_mape:.1%}")
            if improvement is not None and not np.isnan(improvement):
                metric_cols[2].metric("誤差改善", f"{improvement:.1%}")

            if not base_cash_plan.empty:
                updated_cash_plan, forecast_summary_df = apply_forecast_to_cashflow(
                    base_cash_plan, forecast_df, monthly_summary
                )
                updated_cash_forecast = forecast_cashflow(updated_cash_plan, starting_cash)
                savings_summary_df = estimate_forecast_savings(artifact, monthly_summary)
                apply_toggle = st.toggle(
                    "売上予測を資金計画に反映する",
                    value=True,
                    key="cash_forecast_apply_toggle",
                )
                if apply_toggle:
                    active_cash_plan = updated_cash_plan
                    active_cash_forecast = updated_cash_forecast
                else:
                    active_cash_plan = base_cash_plan
                    active_cash_forecast = base_cash_forecast

            detail_cols = st.columns(2)
            forecast_display = forecast_df.copy()
            if not forecast_display.empty:
                forecast_display["予測月"] = forecast_display["month"].astype(str)
                for column in ["forecast", "lower_ci", "upper_ci"]:
                    forecast_display[column] = forecast_display[column].map(
                        lambda v: f"{v:,.0f}" if pd.notna(v) else "-"
                    )
                detail_cols[0].dataframe(
                    forecast_display.rename(
                        columns={"forecast": "予測売上", "lower_ci": "下限", "upper_ci": "上限"}
                    )[["予測月", "予測売上", "下限", "上限"]],
                    hide_index=True,
                    use_container_width=True,
                )

            def _format_summary_value(label: str, value: Any) -> str:
                if value is None or pd.isna(value):
                    return "-"
                if "MAPE" in label or "率" in label:
                    return f"{float(value):.1%}"
                return f"{float(value):,.0f}"

            summary_container = detail_cols[1]
            if not forecast_summary_df.empty:
                summary_container.markdown("**キャッシュフロー影響**")
                formatted_summary = forecast_summary_df.copy()
                for column in ["現状", "予測適用", "差分"]:
                    formatted_summary[column] = formatted_summary.apply(
                        lambda row: _format_summary_value(row["指標"], row[column]), axis=1
                    )
                summary_container.dataframe(
                    formatted_summary,
                    hide_index=True,
                    use_container_width=True,
                )
            if not savings_summary_df.empty:
                summary_container.markdown("**投資効果シミュレーション**")
                formatted_savings = savings_summary_df.copy()
                for column in ["現状", "予測適用", "差分"]:
                    formatted_savings[column] = formatted_savings.apply(
                        lambda row: _format_summary_value(row["指標"], row[column]), axis=1
                    )
                summary_container.dataframe(
                    formatted_savings,
                    hide_index=True,
                    use_container_width=True,
                )

    st.markdown("</div>", unsafe_allow_html=True)

    if active_cash_forecast is not None and not active_cash_forecast.empty:
        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>キャッシュ残高推移</div></div>",
            unsafe_allow_html=True,
        )
        forecast_df_plot = active_cash_forecast.copy()
        forecast_df_plot["period_start"] = forecast_df_plot["month"].dt.to_timestamp()
        forecast_df_plot["period_label"] = forecast_df_plot["month"].astype(str)
        cash_line = alt.Chart(forecast_df_plot).mark_line(
            color=CASH_SERIES_COLOR, point=alt.OverlayMarkDef(size=60, filled=True)
        ).encode(
            x=alt.X("period_start:T", title="期間開始", axis=alt.Axis(format="%Y-%m", labelOverlap=True)),
            y=alt.Y("cash_balance:Q", title="期末現金残高 (円)", axis=alt.Axis(format=",.0f")),
            tooltip=[
                alt.Tooltip("period_label:N", title="期間"),
                alt.Tooltip("cash_balance:Q", title="期末現金残高", format=",.0f"),
                alt.Tooltip("net_cf:Q", title="純キャッシュフロー", format=",.0f"),
            ],
        )

        cash_layers: List[alt.Chart] = [cash_line]
        cash_target = KGI_TARGETS.get("cash_balance")
        if cash_target is not None and not pd.isna(cash_target):
            cash_target_df = pd.DataFrame({"基準": ["目標残高"], "金額": [float(cash_target)]})
            target_rule = alt.Chart(cash_target_df).mark_rule(strokeDash=[6, 4]).encode(
                y="金額:Q",
                color=alt.Color(
                    "基準:N",
                    scale=alt.Scale(domain=["目標残高"], range=[BASELINE_SERIES_COLOR]),
                    legend=alt.Legend(title="基準"),
                ),
                tooltip=[alt.Tooltip("金額:Q", title="目標残高", format=",.0f")],
            )
            cash_layers.append(target_rule)

        cash_chart = alt.layer(*cash_layers).resolve_scale(color="independent").properties(
            height=320,
        )
        st.altair_chart(apply_altair_theme(cash_chart), use_container_width=True)

        latest_row = forecast_df_plot.iloc[-1]
        latest_cash = float(latest_row.get("cash_balance", starting_cash))
        net_cf = latest_row.get("net_cf")
        net_cf_text = f"{float(net_cf):,.0f}円" if pd.notna(net_cf) else "-"
        target_gap_text, _ = format_target_gap(latest_cash, cash_target)
        summary_parts = [
            f"最新の期末現金残高は{latest_cash:,.0f}円、純キャッシュフローは{net_cf_text}です。",
        ]
        if target_gap_text != "-":
            summary_parts.append(f"目標残高との差は{target_gap_text}です。")
        st.caption(" ".join(summary_parts))
        st.markdown("</div>", unsafe_allow_html=True)
    else:
        st.info("資金繰り予測を表示するデータが不足しています。")

    if active_cash_plan is not None and not active_cash_plan.empty:
        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>キャッシュフロー内訳</div></div>",
            unsafe_allow_html=True,
        )
        plan_df = active_cash_plan.copy()
        plan_df["period_start"] = plan_df["month"].dt.to_timestamp()
        melted = plan_df.melt(
            id_vars=["period_start"],
            value_vars=["operating_cf", "investment_cf", "financing_cf", "loan_repayment"],
            var_name="type",
            value_name="金額",
        )
        label_map = {
            "operating_cf": "営業CF",
            "investment_cf": "投資CF",
            "financing_cf": "財務CF",
            "loan_repayment": "返済",
        }
        melted["区分"] = melted["type"].map(label_map)
        melted = melted.dropna(subset=["区分"])
        domain = [label_map[key] for key in label_map]
        range_colors = [CF_COLOR_MAPPING[label] for label in domain]
        cf_chart = alt.Chart(melted).mark_bar().encode(
            x=alt.X("period_start:T", title="期間開始", axis=alt.Axis(format="%Y-%m", labelOverlap=True)),
            y=alt.Y("金額:Q", title="キャッシュフロー (円)", axis=alt.Axis(format=",.0f")),
            color=alt.Color("区分:N", scale=alt.Scale(domain=domain, range=range_colors), legend=alt.Legend(title="区分")),
            tooltip=[
                alt.Tooltip("period_start:T", title="期間"),
                alt.Tooltip("区分:N", title="区分"),
                alt.Tooltip("金額:Q", title="金額", format=",.0f"),
            ],
        )
        st.altair_chart(apply_altair_theme(cf_chart.properties(height=320)), use_container_width=True)

        latest_plan = plan_df.iloc[-1]
        dominant_key = max(label_map, key=lambda key: abs(float(latest_plan.get(key, 0.0))))
        dominant_label = label_map[dominant_key]
        dominant_value = float(latest_plan.get(dominant_key, 0.0))
        st.caption(
            f"直近の主要キャッシュフローは{dominant_label}で{dominant_value:,.0f}円です。"
        )
        st.markdown("</div>", unsafe_allow_html=True)

    with st.expander("キャッシュフロー明細", expanded=False):
        if active_cash_plan is None or active_cash_plan.empty:
            st.info("キャッシュフロー計画データがありません。")
        else:
            table_df = active_cash_plan.copy()
            table_df["month_label"] = table_df["month"].astype(str)
            export_df = table_df[[
                "month_label",
                "operating_cf",
                "investment_cf",
                "financing_cf",
                "loan_repayment",
            ]].copy()
            if active_cash_forecast is not None and not active_cash_forecast.empty:
                forecast_export = active_cash_forecast.copy()
                forecast_export["month_label"] = forecast_export["month"].astype(str)
                export_df = export_df.merge(
                    forecast_export[["month_label", "net_cf", "cash_balance"]],
                    on="month_label",
                    how="left",
                )
            else:
                export_df["net_cf"] = (
                    export_df["operating_cf"]
                    + export_df["financing_cf"]
                    - export_df["investment_cf"]
                    - export_df["loan_repayment"]
                )
                export_df["cash_balance"] = (
                    export_df["net_cf"].cumsum() + float(starting_cash)
                )

            display_df = export_df.rename(
                columns={
                    "month_label": "月",
                    "operating_cf": "営業CF",
                    "investment_cf": "投資CF",
                    "financing_cf": "財務CF",
                    "loan_repayment": "返済",
                    "net_cf": "純キャッシュフロー",
                    "cash_balance": "期末現金残高",
                }
            )
            format_columns = ["営業CF", "投資CF", "財務CF", "返済", "純キャッシュフロー", "期末現金残高"]
            formatted_df = display_df.copy()
            for column in format_columns:
                formatted_df[column] = formatted_df[column].map(lambda v: f"{v:,.0f}" if pd.notna(v) else "-")
            st.dataframe(formatted_df, hide_index=True, use_container_width=True)
            toolbar = st.columns(4)
            export_base = "cash_flow_plan"
            with toolbar[0]:
                download_button_from_df("CSV出力", display_df, f"{export_base}.csv")
            with toolbar[1]:
                download_excel_from_df("Excel出力", display_df, f"{export_base}.xlsx")
            with toolbar[2]:
                download_pdf_from_df(
                    "PDF出力",
                    display_df,
                    f"{export_base}.pdf",
                    title="キャッシュフロー明細",
                )
            with toolbar[3]:
                download_ppt_from_df(
                    "PowerPoint",
                    display_df.head(20),
                    f"{export_base}.pptx",
                    title="キャッシュフロー概要",
                )


def render_fixed_cost_breakdown(
    expense_df: Optional[pd.DataFrame], fixed_cost: float
) -> None:
    """固定費の内訳を積み上げ棒グラフで表示する。"""

    if expense_df is not None and isinstance(expense_df, pd.DataFrame) and not expense_df.empty:
        working = expense_df.copy()
    else:
        working = pd.DataFrame(EXPENSE_PLAN_TEMPLATES.get("スリム型コスト構成", []))

    if working.empty:
        st.info("固定費内訳を表示するデータがありません。")
        return

    rename_map = {col: col for col in ["費目", "月次金額", "区分"] if col in working.columns}
    working = working.rename(columns=rename_map)
    if "区分" in working.columns:
        working = working[working["区分"].isin(["固定費", "固定費用", "固定費用計", "固定"])]
    if working.empty:
        st.info("固定費区分のデータがありません。")
        return

    breakdown = working.groupby("費目")["月次金額"].sum().reset_index()
    total_current = float(breakdown["月次金額"].sum())
    target_total = float(fixed_cost or 0.0)
    if total_current > 0 and target_total > 0:
        breakdown["月次金額"] = breakdown["月次金額"] * target_total / total_current

    breakdown["店舗"] = "全社"
    palette = get_active_chart_colorway() + [ACCENT_COLOR, SECONDARY_COLOR]
    fig = go.Figure()
    for idx, row in enumerate(breakdown.itertuples()):
        fig.add_bar(
            name=str(row.費目),
            x=[row.店舗],
            y=[row.月次金額],
            marker_color=palette[idx % len(palette)],
            hovertemplate="費目=%{fullData.name}<br>金額=%{y:,.0f}円<extra></extra>",
        )

    if target_total > 0:
        fig.add_scatter(
            x=["全社"],
            y=[target_total],
            name="固定費目標",
            mode="lines+markers",
            line=dict(color=BASELINE_SERIES_COLOR, dash="dash"),
            marker=dict(size=10, color=BASELINE_SERIES_COLOR),
            hovertemplate="固定費目標=%{y:,.0f}円<extra></extra>",
        )

    fig.update_layout(
        barmode="stack",
        xaxis_title="店舗",
        yaxis_title="金額（円）",
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1.0),
        plot_bgcolor="rgba(0,0,0,0)",
        paper_bgcolor="rgba(0,0,0,0)",
        margin=dict(l=40, r=40, t=60, b=60),
    )
    st.plotly_chart(fig, use_container_width=True)

    top_item = breakdown.sort_values("月次金額", ascending=False).iloc[0]
    st.caption(
        f"主要固定費は{top_item['費目']}で{top_item['月次金額']:,.0f}円です。目標固定費は{target_total:,.0f}円に調整しています。"
    )


def render_profit_meter(pl_result: pd.DataFrame, base_pl: Dict[str, float]) -> None:
    """シナリオ売上の進捗をゲージ表示し、損益状況を補足する。"""

    if pl_result is None or pl_result.empty:
        st.info("シミュレーション結果がまだ計算されていません。")
        return

    try:
        scenario_sales = float(
            pl_result.loc[pl_result["項目"] == "売上高", "シナリオ"].iloc[0]
        )
        scenario_profit = float(
            pl_result.loc[pl_result["項目"] == "営業利益", "シナリオ"].iloc[0]
        )
    except IndexError:
        st.info("シミュレーション指標が不足しています。")
        return

    base_sales = float(base_pl.get("sales", 0.0))
    base_cogs = float(base_pl.get("cogs", 0.0))
    base_sga = float(base_pl.get("sga", 0.0))
    contribution = 0.0
    if base_sales > 0:
        contribution = 1.0 - (base_cogs / base_sales if base_sales else 0.0)
    break_even = None
    if contribution > 0:
        break_even = base_sga / contribution

    gauge_upper = max(scenario_sales, break_even or 0.0, base_sales) * 1.2
    if gauge_upper <= 0:
        gauge_upper = max(scenario_sales, 1.0)

    steps = []
    if break_even and gauge_upper > break_even:
        steps = [
            {"range": [0, break_even], "color": ERROR_GAUGE_RANGE_COLOR},
            {"range": [break_even, gauge_upper], "color": SUCCESS_GAUGE_RANGE_COLOR},
        ]

    indicator = go.Figure(
        go.Indicator(
            mode="gauge+number",
            value=scenario_sales,
            number=dict(valueformat=",.0f", suffix=" 円"),
            gauge=dict(
                axis=dict(range=[0, gauge_upper], tickformat=",.0f"),
                bar=dict(color=SALES_SERIES_COLOR),
                steps=steps,
                threshold=dict(
                    line=dict(
                        color=SUCCESS_COLOR if scenario_sales >= (break_even or 0) else ERROR_COLOR,
                        width=4,
                    ),
                    value=break_even if break_even is not None else scenario_sales,
                ),
            ),
        )
    )
    indicator.update_layout(height=340, margin=dict(t=40, b=20, l=20, r=20))
    st.plotly_chart(indicator, use_container_width=True)

    profit_text = (
        f"営業利益は{scenario_profit:,.0f}円" if pd.notna(scenario_profit) else "営業利益は算出できません"
    )
    if break_even is not None:
        st.caption(
            f"損益分岐点売上は約{break_even:,.0f}円です。現在のシナリオ売上{scenario_sales:,.0f}円では{profit_text}となります。"
        )
    else:
        st.caption(
            f"現状の原価率では損益分岐点を計算できませんが、シナリオ売上{scenario_sales:,.0f}円で{profit_text}です。"
        )


def render_data_upload_guide(
    merged_df: pd.DataFrame,
    cost_df: pd.DataFrame,
    subscription_df: pd.DataFrame,
) -> None:
    """アップロード手順とサンプルファイルを案内し、プレビューを表示する。"""

    st.markdown("### データアップロードガイド")
    st.caption("自社データを取り込むための手順とテンプレートをまとめています。")

    with st.expander("Step-by-step ガイド", expanded=False):
        st.markdown(
            """
            1. サイドバー上部の **「はじめに」** からサンプルデータを読み込むか、下のテンプレートをダウンロードします。
            2. 自社の売上・原価率・サブスクKPIの列名をテンプレートに合わせて整形し、CSVまたはExcelで保存します。
            3. サイドバーのアップロードセクションにファイルをドラッグ&ドロップすると、フォーマットチェックとプレビューが自動表示されます。
            """
        )
        st.caption("必要項目が不足している場合は読み込み時に警告が表示されます。")

    download_cols = st.columns(3)
    with download_cols[0]:
        st.markdown("**売上データテンプレート**")
        sample_sales = generate_sample_sales_data().head(200)
        download_button_from_df("CSV", sample_sales, "sample_sales.csv")
        download_excel_from_df("Excel", sample_sales, "sample_sales.xlsx")
    with download_cols[1]:
        st.markdown("**原価率テンプレート**")
        sample_cost = generate_sample_cost_data()
        download_button_from_df("CSV", sample_cost, "sample_cost.csv")
        download_excel_from_df("Excel", sample_cost, "sample_cost.xlsx")
    with download_cols[2]:
        st.markdown("**KPIテンプレート**")
        sample_kpi = generate_sample_subscription_data()
        download_button_from_df("CSV", sample_kpi, "sample_kpi.csv")
        download_excel_from_df("Excel", sample_kpi, "sample_kpi.xlsx")

    st.markdown("#### アップロード済みデータのプレビュー")
    preview_tabs = st.tabs(["売上", "原価率", "KPI"])
    with preview_tabs[0]:
        if merged_df is None or merged_df.empty:
            st.info("売上データがアップロードされていません。")
        else:
            st.dataframe(merged_df.head(20), hide_index=True, use_container_width=True)
    with preview_tabs[1]:
        if cost_df is None or cost_df.empty:
            st.info("原価率データがアップロードされていません。")
        else:
            st.dataframe(cost_df.head(20), hide_index=True, use_container_width=True)
    with preview_tabs[2]:
        if subscription_df is None or subscription_df.empty:
            st.info("サブスクKPIデータがアップロードされていません。")
        else:
            st.dataframe(subscription_df.head(20), hide_index=True, use_container_width=True)

    st.markdown("---")


def render_data_status_section(
    merged_df: pd.DataFrame,
    cost_df: pd.DataFrame,
    subscription_df: pd.DataFrame,
    *,
    use_sample_data: bool,
    automated_sales_data: Dict[str, Any],
) -> None:
    """データアップロード状況をカード形式で表示する。"""

    st.markdown("### データアップロード状況")
    st.caption("チャネルや補助データの最新状態を確認できます。")

    cards: List[str] = []

    def _status_markup(label: str, tone: str, code: str) -> str:
        return f"{build_ui_icon(code, tone=tone)}<span>{html.escape(label)}</span>"

    if merged_df is not None and not merged_df.empty:
        channel_summary = (
            merged_df.groupby("channel")
            .agg(
                records=("sales_amount", "size"),
                amount=("sales_amount", "sum"),
                latest=("order_date", "max"),
                earliest=("order_date", "min"),
            )
            .reset_index()
            .sort_values("records", ascending=False)
        )
        for _, row in channel_summary.iterrows():
            latest = pd.to_datetime(row["latest"]).strftime("%Y-%m-%d") if pd.notna(row["latest"]) else "-"
            earliest = pd.to_datetime(row["earliest"]).strftime("%Y-%m-%d") if pd.notna(row["earliest"]) else "-"
            meta = f"{earliest} 〜 {latest}"
            body = f"件数: {int(row['records']):,} / 売上高: {row['amount']:,.0f}円"
            cards.append(
                """
                <div class="data-status-card">
                    <div class="data-status-card__title">{title}</div>
                    <div class="data-status-card__meta">{meta}</div>
                    <div class="data-status-card__body">{body}</div>
                    <div class="data-status-card__status data-status-card__status--ok">{status}</div>
                </div>
                """.format(
                    title=html.escape(str(row["channel"])),
                    meta=html.escape(meta),
                    body=html.escape(body),
                    status=_status_markup("正常", "success", "OK"),
                )
            )
    else:
        cards.append(
            """
            <div class="data-status-card">
                <div class="data-status-card__title">売上データ</div>
                <div class="data-status-card__meta">-</div>
                <div class="data-status-card__body">売上ファイルが未読み込みです。</div>
                <div class="data-status-card__status data-status-card__status--warning">{status}</div>
            </div>
            """.format(status=_status_markup("未取込", "warning", "!"))
        )

    cost_loaded = cost_df is not None and not cost_df.empty
    cost_status_class = (
        "data-status-card__status data-status-card__status--ok"
        if cost_loaded
        else "data-status-card__status data-status-card__status--warning"
    )
    cost_status_label = (
        _status_markup("正常", "success", "OK")
        if cost_loaded
        else _status_markup("未登録", "warning", "!")
    )
    cost_body = (
        f"登録済みアイテム: {len(cost_df):,}件" if cost_loaded else "原価率データが未設定です。"
    )
    cards.append(
        """
            <div class="data-status-card">
                <div class="data-status-card__title">原価率マスタ</div>
                <div class="data-status-card__meta">-</div>
                <div class="data-status-card__body">{body}</div>
                <div class="{status_class}">{status}</div>
            </div>
            """.format(
                body=html.escape(cost_body),
                status_class=cost_status_class,
                status=cost_status_label,
            )
        )

    sub_loaded = subscription_df is not None and not subscription_df.empty
    sub_status_class = (
        "data-status-card__status data-status-card__status--ok"
        if sub_loaded
        else "data-status-card__status data-status-card__status--warning"
    )
    sub_status_label = (
        _status_markup("正常", "success", "OK")
        if sub_loaded
        else _status_markup("未登録", "warning", "!")
    )
    sub_body = (
        f"月次レコード: {len(subscription_df):,}件" if sub_loaded else "サブスクKPIが未入力です。"
    )
    cards.append(
        """
            <div class="data-status-card">
                <div class="data-status-card__title">定期購買 / KPIデータ</div>
                <div class="data-status-card__meta">-</div>
                <div class="data-status-card__body">{body}</div>
                <div class="{status_class}">{status}</div>
            </div>
            """.format(
                body=html.escape(sub_body),
                status_class=sub_status_class,
                status=sub_status_label,
            )
        )

    if automated_sales_data:
        api_last_fetched = st.session_state.get("api_last_fetched", {})
        api_reports = st.session_state.get("api_sales_validation", {})
        api_lines: List[str] = []
        error_count = 0
        warning_count = 0
        ok_count = 0
        for channel, df in automated_sales_data.items():
            last_fetch = api_last_fetched.get(channel)
            report = api_reports.get(channel)
            status_label = "正常"
            if report and getattr(report, "has_errors", lambda: False)():
                status_label = "エラー"
                error_count += 1
            elif report and getattr(report, "has_warnings", lambda: False)():
                status_label = "警告あり"
                warning_count += 1
            else:
                ok_count += 1
            timestamp = last_fetch.strftime("%Y-%m-%d %H:%M") if last_fetch else "-"
            api_lines.append(f"{channel}: {status_label} / 取得 {timestamp}")
        if error_count:
            api_status_class = "data-status-card__status data-status-card__status--error"
            api_status_label = _status_markup(
                f"エラー {error_count}件", "error", "×"
            )
        elif warning_count:
            api_status_class = "data-status-card__status data-status-card__status--warning"
            api_status_label = _status_markup(
                f"警告 {warning_count}件", "warning", "!"
            )
        else:
            api_status_class = "data-status-card__status data-status-card__status--ok"
            api_status_label = _status_markup(
                f"正常 {ok_count}件", "success", "OK"
            )

        footnote_html = ""
        if api_lines:
            footnote_html = "<div class='data-status-card__footnote'>{}</div>".format(
                "<br />".join(html.escape(line) for line in api_lines)
            )

        cards.append(
            """
            <div class="data-status-card">
                <div class="data-status-card__title">API連携</div>
                <div class="data-status-card__meta">接続チャネル: {count}件</div>
                <div class="data-status-card__body">自動取得の最終実行状況を表示します。</div>
                <div class="{status_class}">{status}</div>
                {footnote}
            </div>
            """.format(
                count=len(automated_sales_data),
                status_class=api_status_class,
                status=html.escape(api_status_label),
                footnote=footnote_html,
            )
        )

    st.markdown(
        "<div class='data-status-grid'>{}</div>".format("".join(cards)),
        unsafe_allow_html=True,
    )

    if use_sample_data:
        st.caption("※ 現在はサンプルデータを表示しています。実データをアップロードすると自動的に置き換わります。")


def normalize_scenario_input(df: Optional[pd.DataFrame]) -> pd.DataFrame:
    """Phase2シナリオモジュール向けのベースデータを正規化する。"""

    if df is None or df.empty:
        return pd.DataFrame(columns=["order_date", "order_month", "sales_amount", "net_gross_profit"])

    working = df.copy()
    column_lookup = {col.lower(): col for col in working.columns}

    def _match_column(candidates: Sequence[str]) -> Optional[str]:
        for candidate in candidates:
            if candidate in working.columns:
                return candidate
            lowered = candidate.lower()
            if lowered in column_lookup:
                return column_lookup[lowered]
        return None

    sales_col = _match_column(["sales_amount", "sales", "売上", "売上高", "revenue", "total_sales"])
    if sales_col and sales_col != "sales_amount":
        working.rename(columns={sales_col: "sales_amount"}, inplace=True)
    if "sales_amount" not in working.columns:
        working["sales_amount"] = 0.0
    working["sales_amount"] = pd.to_numeric(working["sales_amount"], errors="coerce").fillna(0.0)

    date_col = _match_column(["order_date", "date", "日付", "年月日", "month"])
    year_col = _match_column(["year", "年度", "会計年度"])
    if date_col and date_col != "order_date":
        working["order_date"] = pd.to_datetime(working[date_col], errors="coerce")
    elif "order_date" in working.columns:
        working["order_date"] = pd.to_datetime(working["order_date"], errors="coerce")
    elif year_col:
        working["order_date"] = pd.to_datetime(working[year_col].astype(str) + "-01", errors="coerce")
    else:
        working["order_date"] = pd.date_range(
            end=pd.Timestamp.today(), periods=len(working), freq="M"
        )

    working.dropna(subset=["order_date"], inplace=True)
    working.sort_values("order_date", inplace=True)
    working["order_month"] = pd.PeriodIndex(working["order_date"], freq="M")

    profit_col = _match_column(["net_gross_profit", "gross_profit", "profit", "粗利", "営業利益"])
    if profit_col and profit_col != "net_gross_profit":
        working.rename(columns={profit_col: "net_gross_profit"}, inplace=True)
    if "net_gross_profit" in working.columns:
        working["net_gross_profit"] = pd.to_numeric(
            working["net_gross_profit"], errors="coerce"
        ).fillna(working["sales_amount"] * 0.45)
    else:
        working["net_gross_profit"] = working["sales_amount"] * 0.45

    base_columns = ["order_date", "order_month", "sales_amount", "net_gross_profit"]
    remaining_columns = [col for col in working.columns if col not in base_columns]
    return working[base_columns + remaining_columns]


def calculate_recent_growth(series: Optional[pd.Series]) -> Optional[float]:
    """直近2期間の成長率（割合）を返す。"""

    if series is None or series.empty or len(series) < 2:
        return None
    latest = float(series.iloc[-1])
    previous = float(series.iloc[-2])
    if previous == 0:
        return None
    growth = (latest - previous) / previous
    if not np.isfinite(growth):
        return None
    return growth


def build_swot_insights(
    kpi: Dict[str, Optional[float]], growth_rate: Optional[float]
) -> Dict[str, List[str]]:
    """KPIからSWOT分析コメントを生成する。"""

    strengths: List[str] = []
    weaknesses: List[str] = []
    opportunities: List[str] = []
    threats: List[str] = []

    gross_margin = kpi.get("gross_margin_rate")
    repeat_rate = kpi.get("repeat_rate")
    roas = kpi.get("roas")
    churn_rate = kpi.get("churn_rate")
    ltv = kpi.get("ltv")
    cac = kpi.get("cac")

    if gross_margin is not None and np.isfinite(gross_margin):
        if gross_margin >= 0.55:
            strengths.append("粗利率が業界平均を上回り、利益創出力が高い状態です。")
        else:
            weaknesses.append("粗利率が業界水準を下回っています。原価と販管費の最適化が必要です。")

    if repeat_rate is not None and np.isfinite(repeat_rate):
        if repeat_rate >= 0.40:
            strengths.append("リピート率が40%超で、ファン顧客が育っています。")
        else:
            weaknesses.append("リピート率が伸び悩んでいます。CRM施策の強化を検討しましょう。")

    if roas is not None and np.isfinite(roas):
        if roas >= 4.0:
            strengths.append("広告投資の回収効率が高く、成長投資を加速できます。")
        elif roas < 2.5:
            threats.append("ROASが低位で広告費の回収が遅れています。チャネルポートフォリオの見直しが必要です。")

    if churn_rate is not None and np.isfinite(churn_rate) and churn_rate > 0.05:
        threats.append("解約率が5%を超過しています。オンボーディングやロイヤルティ施策を強化しましょう。")

    if (
        ltv is not None
        and cac is not None
        and np.isfinite(ltv)
        and np.isfinite(cac)
        and cac > 0
    ):
        ratio = ltv / cac
        if ratio >= 3.0:
            strengths.append("LTV/CACが3倍以上で投資リターンが十分です。")
        elif ratio < 2.0:
            threats.append("LTV/CACが2倍未満のため、顧客獲得コストの圧縮が課題です。")

    if growth_rate is not None:
        if growth_rate > 0.05:
            opportunities.append("直近売上が5%以上成長しており、攻めの投資タイミングです。")
        elif growth_rate < 0:
            threats.append("売上が減速傾向にあります。販促や価格政策の再設計が必要です。")

    opportunities.append("新チャネル開拓や外部調達により成長余地を拡大できます。")

    return {
        "strengths": strengths or ["強みを特定するデータが不足しています。"],
        "weaknesses": weaknesses or ["大きな弱みは検出されませんでした。"],
        "opportunities": opportunities or ["追加の市場調査により機会を探索できます。"],
        "threats": threats or ["重大な脅威は検出されませんでした。"],
    }


def build_industry_benchmark_table(kpi: Dict[str, Optional[float]]) -> pd.DataFrame:
    """業界平均と比較するベンチマーク表を作成する。"""

    benchmarks = {
        "gross_margin_rate": 0.52,
        "repeat_rate": 0.38,
        "roas": 3.5,
        "ltv": 32000.0,
        "cac": 12000.0,
    }
    labels = {
        "gross_margin_rate": "売上総利益率 (%)",
        "repeat_rate": "リピート率 (%)",
        "roas": "ROAS",
        "ltv": "LTV (円)",
        "cac": "CAC (円)",
    }

    rows: List[Dict[str, Any]] = []
    for key, label in labels.items():
        company_value = kpi.get(key)
        industry_value = benchmarks.get(key)
        if key in {"gross_margin_rate", "repeat_rate"}:
            company_value = (company_value or 0.0) * 100 if company_value is not None else np.nan
            industry_value = (industry_value or 0.0) * 100 if industry_value is not None else np.nan
        rows.append(
            {
                "指標": label,
                "自社": float(company_value) if company_value is not None else np.nan,
                "業界平均": float(industry_value) if industry_value is not None else np.nan,
                "差分": (
                    float(company_value) - float(industry_value)
                    if company_value is not None
                    and industry_value is not None
                    and np.isfinite(company_value)
                    and np.isfinite(industry_value)
                    else np.nan
                ),
            }
        )

    return pd.DataFrame(rows)


def _generate_adjustment_mesh(adjust_pct_map: Dict[str, float]) -> Tuple[List[str], np.ndarray]:
    """調整幅を元にした%-表現の組み合わせを生成する。"""

    keys = list(adjust_pct_map.keys())
    if not keys:
        return keys, np.zeros((1, 0))

    factors: List[np.ndarray] = []
    for key in keys:
        width = max(0.0, float(adjust_pct_map.get(key, 0.0))) / 100.0
        if width == 0.0:
            factors.append(np.array([0.0]))
        else:
            factors.append(np.array([-width, 0.0, width], dtype=float))

    mesh = np.array(np.meshgrid(*factors, indexing="ij"))
    mesh = mesh.reshape(len(keys), -1).T if mesh.size else np.zeros((1, len(keys)))
    return keys, mesh


def _build_profit_range_cases(
    scenario_name: str,
    annual_sales: float,
    margin: float,
    monthly_fixed_cost: float,
    horizon: int,
    adjust_pct_map: Dict[str, float],
) -> pd.DataFrame:
    """単価・数量・固定費を揺らした場合の利益レンジを計算する。"""

    keys, mesh = _generate_adjustment_mesh(adjust_pct_map)
    if mesh.size == 0:
        mesh = np.zeros((1, len(keys)))

    label_map = {
        "unit_price": "単価",
        "quantity": "数量",
        "fixed_cost": "固定費",
    }

    rows: List[Dict[str, Any]] = []
    annual_fixed_cost = monthly_fixed_cost * max(1, horizon)
    for combo in mesh:
        multipliers = {key: 1.0 + delta for key, delta in zip(keys, combo)}
        adjusted_sales = annual_sales * multipliers.get("unit_price", 1.0) * multipliers.get("quantity", 1.0)
        adjusted_fixed_cost = annual_fixed_cost * multipliers.get("fixed_cost", 1.0)
        adjusted_profit = (adjusted_sales * margin) - adjusted_fixed_cost
        label_parts = []
        for key, delta in zip(keys, combo):
            label = label_map.get(key, key)
            label_parts.append(f"{label} {delta * 100:+.1f}%")
        adjustment_label = " / ".join(label_parts) if label_parts else "基準"
        rows.append(
            {
                "scenario": scenario_name,
                "unit_price_pct": multipliers.get("unit_price", 1.0) - 1.0,
                "quantity_pct": multipliers.get("quantity", 1.0) - 1.0,
                "fixed_cost_pct": multipliers.get("fixed_cost", 1.0) - 1.0,
                "sales": adjusted_sales,
                "profit": adjusted_profit,
                "adjustment_label": adjustment_label,
            }
        )

    range_df = pd.DataFrame(rows)
    if not range_df.empty:
        range_df["unit_price_pct"] = range_df["unit_price_pct"] * 100.0
        range_df["quantity_pct"] = range_df["quantity_pct"] * 100.0
        range_df["fixed_cost_pct"] = range_df["fixed_cost_pct"] * 100.0

    return range_df


def build_tornado_dataframe(profit_range_df: pd.DataFrame) -> pd.DataFrame:
    """単変量ごとの利益インパクトを算出しトルネード図用に整形する。"""

    if profit_range_df is None or profit_range_df.empty:
        return pd.DataFrame()

    factors = [
        ("unit_price_pct", "単価"),
        ("quantity_pct", "数量"),
        ("fixed_cost_pct", "固定費"),
    ]

    rows: List[Dict[str, Any]] = []
    tolerance = 1e-6
    for scenario, scenario_df in profit_range_df.groupby("scenario"):
        base_row = scenario_df[
            (scenario_df["unit_price_pct"].abs() < tolerance)
            & (scenario_df["quantity_pct"].abs() < tolerance)
            & (scenario_df["fixed_cost_pct"].abs() < tolerance)
        ]
        if base_row.empty:
            continue
        base_profit = float(base_row["profit"].iloc[0])

        for column, label in factors:
            others = [col for col, _ in factors if col != column]

            positive_df = scenario_df[
                (scenario_df[column] > tolerance)
                & np.logical_and.reduce([scenario_df[other].abs() < tolerance for other in others])
            ]
            negative_df = scenario_df[
                (scenario_df[column] < -tolerance)
                & np.logical_and.reduce([scenario_df[other].abs() < tolerance for other in others])
            ]

            if not positive_df.empty:
                impact = float(positive_df["profit"].max()) - base_profit
                rows.append(
                    {
                        "scenario": scenario,
                        "factor": label,
                        "direction": "増加",
                        "impact": impact,
                    }
                )
            if not negative_df.empty:
                impact = float(negative_df["profit"].min()) - base_profit
                rows.append(
                    {
                        "scenario": scenario,
                        "factor": label,
                        "direction": "減少",
                        "impact": impact,
                    }
                )

    tornado_df = pd.DataFrame(rows)
    if tornado_df.empty:
        return tornado_df

    tornado_df["impact_abs"] = tornado_df["impact"].abs()
    return tornado_df


def render_scenario_share_controls(
    summary_df: Optional[pd.DataFrame],
    profit_range_summary: Optional[pd.DataFrame] = None,
) -> None:
    """シナリオ比較結果の共有リンク生成UIを表示する。"""

    if summary_df is None or summary_df.empty:
        return

    share_cols = st.columns([3, 1])
    with share_cols[0]:
        st.caption("共有リンクを作成すると同じ結果をチームに展開できます。")
    with share_cols[1]:
        if st.button("共有リンクを生成", key="generate_scenario_share"):
            payload = {
                "generated_at": datetime.utcnow().isoformat(),
                "summary": summary_df.to_dict(orient="records"),
            }
            if profit_range_summary is not None and not profit_range_summary.empty:
                payload["range_summary"] = profit_range_summary.to_dict(orient="records")
            token = encode_share_payload(payload)
            st.session_state["scenario_share_token"] = token
            st.experimental_set_query_params(share=token)
            st.success("共有用のクエリパラメータを生成しました。ブラウザのURLをコピーしてください。")

    token = st.session_state.get("scenario_share_token")
    if token:
        share_suffix = f"?share={token}"
        st.text_input("共有クエリ", share_suffix, key="scenario_share_suffix")
        mail_body = quote(
            "このダッシュボードにアクセスし、URL末尾に" + share_suffix + "を付与して確認してください。"
        )
        st.markdown(f"[メールで共有する](mailto:?subject=Scenario%20Share&body={mail_body})")


def render_scenario_radar_chart(summary_df: Optional[pd.DataFrame]) -> None:
    """KPIのレーダーチャートで複数シナリオを比較表示する。"""

    if summary_df is None or summary_df.empty:
        return
    metric_candidates = ["年間売上", "年間利益", "期末キャッシュ", "成長率(%)", "利益率(%)"]
    metrics = [metric for metric in metric_candidates if metric in summary_df.columns]
    if len(metrics) < 3 or "シナリオ" not in summary_df.columns:
        return

    normalized = summary_df[["シナリオ"] + metrics].copy()
    percentage_metrics = {metric for metric in metrics if "(%)" in metric}
    for metric in metrics:
        series = pd.to_numeric(normalized[metric], errors="coerce")
        max_val = series.max()
        min_val = series.min()
        if pd.isna(max_val):
            normalized[metric] = 0.0
        elif max_val == min_val:
            normalized[metric] = 100.0
        else:
            normalized[metric] = ((series - min_val) / (max_val - min_val) * 100).fillna(0.0)

    categories = metrics + [metrics[0]]
    fig = go.Figure()

    for idx, row in normalized.iterrows():
        scaled_values = [float(row[metric]) for metric in metrics]
        scaled_values.append(scaled_values[0])
        actual_row = summary_df.loc[idx, metrics]
        hover_lines = []
        for metric in metrics:
            value = actual_row[metric]
            if metric in percentage_metrics:
                hover_lines.append(f"{metric}: {value:,.1f}%")
            else:
                hover_lines.append(f"{metric}: {value:,.0f}")
        fig.add_trace(
            go.Scatterpolar(
                r=scaled_values,
                theta=categories,
                name=str(summary_df.loc[idx, "シナリオ"]),
                fill="toself",
                hovertemplate="<br>".join(hover_lines) + "<extra></extra>",
            )
        )

    fig.update_layout(
        polar=dict(
            radialaxis=dict(range=[0, 100], showticklabels=False, gridcolor="rgba(15,23,42,0.08)")
        ),
        legend=dict(orientation="h", yanchor="bottom", y=1.05, xanchor="center", x=0.5),
        margin=dict(l=40, r=40, t=60, b=30),
    )

    st.markdown("### KPIレーダーチャート")
    st.plotly_chart(fig, use_container_width=True)


def render_shared_scenario_preview() -> None:
    """共有リンクから読み込んだシナリオ比較結果を表示する。"""

    summary_df = st.session_state.get("shared_scenario_summary_df")
    if summary_df is None or summary_df.empty:
        return

    st.markdown("### 共有されたシナリオ比較")
    st.info("共有リンクから読み込んだ結果を表示しています。サンプルデータがなくても比較概要を確認できます。")
    format_map = {
        column: fmt
        for column, fmt in {
            "年間売上": "{:.0f}",
            "年間利益": "{:.0f}",
            "平均月次売上": "{:.0f}",
            "期末キャッシュ": "{:.0f}",
            "成長率(%)": "{:.1f}",
            "利益率(%)": "{:.1f}",
        }.items()
        if column in summary_df.columns
    }
    st.dataframe(summary_df.style.format(format_map), use_container_width=True)
    range_df = st.session_state.get("shared_scenario_range_df")
    if isinstance(range_df, pd.DataFrame) and not range_df.empty:
        range_format = {
            column: "{:.0f}"
            for column in ["最小利益", "最大利益", "利益幅"]
            if column in range_df.columns
        }
        st.dataframe(range_df.style.format(range_format), use_container_width=True)


def run_scenario_projection(
    monthly_sales: pd.Series, scenario: Dict[str, Any]
) -> Tuple[str, pd.DataFrame, Dict[str, Any], pd.DataFrame, pd.DataFrame]:
    """シナリオ設定に基づき将来12〜36ヶ月の推移を試算する。"""

    scenario_name = scenario.get("name") or "新規シナリオ"
    growth = float(scenario.get("growth", 0.0)) / 100.0
    margin = max(0.0, float(scenario.get("margin", 0.0))) / 100.0
    funding = float(scenario.get("funding", 0.0))
    horizon = int(scenario.get("horizon", 12) or 12)
    if horizon <= 0:
        horizon = 12

    monthly_fixed_cost = float(scenario.get("base_fixed_cost", DEFAULT_FIXED_COST))
    unit_price_adjust_pct = float(scenario.get("unit_price_adjust_pct", 10.0))
    quantity_adjust_pct = float(scenario.get("quantity_adjust_pct", 10.0))
    fixed_cost_adjust_pct = float(scenario.get("fixed_cost_adjust_pct", 10.0))

    base_series = monthly_sales.copy() if monthly_sales is not None else pd.Series(dtype=float)
    if base_series.empty:
        base_series = pd.Series([1_000_000.0], index=pd.PeriodIndex([pd.Period(date.today(), freq="M")]))
    base_series = base_series.sort_index()

    seasonal_values = base_series.tail(min(len(base_series), horizon)).to_list()
    if not seasonal_values:
        seasonal_values = [float(base_series.iloc[-1])]
    while len(seasonal_values) < horizon:
        seasonal_values.extend(seasonal_values[: horizon - len(seasonal_values)])
    base_array = np.array(seasonal_values[:horizon], dtype=float)

    if isinstance(base_series.index, pd.PeriodIndex):
        start_period = base_series.index[-1]
    else:
        try:
            start_period = pd.Period(base_series.index.max(), freq="M")
        except Exception:
            start_period = pd.Period(date.today(), freq="M")
    projection_periods = pd.period_range(start=start_period + 1, periods=horizon, freq="M")

    cumulative_cash = funding
    rows: List[Dict[str, Any]] = []
    for idx, (base_value, period) in enumerate(zip(base_array, projection_periods), start=1):
        projected_sales = base_value * ((1 + growth) ** idx)
        projected_profit = (projected_sales * margin) - monthly_fixed_cost
        cumulative_cash += projected_profit
        rows.append(
            {
                "scenario": scenario_name,
                "month_index": idx,
                "period": period.to_timestamp(),
                "period_label": period.strftime("%Y-%m"),
                "sales": projected_sales,
                "profit": projected_profit,
                "cash": cumulative_cash,
                "funding": funding,
                "growth_pct": growth * 100,
                "margin_pct": margin * 100,
            }
        )

    projection_df = pd.DataFrame(rows)
    annual_sales = float(projection_df["sales"].sum()) if not projection_df.empty else 0.0
    annual_profit = float(projection_df["profit"].sum()) if not projection_df.empty else 0.0
    summary_row = {
        "シナリオ": scenario_name,
        "年間売上": annual_sales,
        "年間利益": annual_profit,
        "平均月次売上": float(projection_df["sales"].mean()) if not projection_df.empty else 0.0,
        "期末キャッシュ": float(projection_df["cash"].iloc[-1]) if not projection_df.empty else funding,
        "成長率(%)": growth * 100,
        "利益率(%)": margin * 100,
        "調整固定費(月額)": monthly_fixed_cost,
        "調達額": funding,
    }

    margin_center = margin * 100
    margin_points = np.linspace(max(0.0, margin_center - 10), min(100.0, margin_center + 10), 5)
    sensitivity_rows = [
        {
            "scenario": scenario_name,
            "margin": point,
            "annual_profit": annual_sales * (point / 100.0) - (monthly_fixed_cost * horizon),
        }
        for point in margin_points
    ]
    sensitivity_df = pd.DataFrame(sensitivity_rows)

    adjust_pct_map = {
        "unit_price": unit_price_adjust_pct,
        "quantity": quantity_adjust_pct,
        "fixed_cost": fixed_cost_adjust_pct,
    }
    profit_range_df = _build_profit_range_cases(
        scenario_name,
        annual_sales,
        margin,
        monthly_fixed_cost,
        horizon,
        adjust_pct_map,
    )

    return scenario_name, projection_df, summary_row, sensitivity_df, profit_range_df


def generate_phase2_report(
    summary_df: Optional[pd.DataFrame],
    swot: Optional[Dict[str, List[str]]],
    benchmark_df: Optional[pd.DataFrame],
    profit_range_summary: Optional[pd.DataFrame] = None,
) -> str:
    """シナリオ比較のサマリーを含むテキストレポートを生成する。"""

    buffer = io.StringIO()
    buffer.write("=== シナリオ分析レポート ===\n")
    buffer.write(f"生成日時: {datetime.now():%Y-%m-%d %H:%M}\n\n")

    if summary_df is not None and not summary_df.empty:
        buffer.write("[シナリオ比較サマリー]\n")
        for _, row in summary_df.iterrows():
            buffer.write(
                "- {name}: 年間売上 {sales:,.0f} 円 / 年間利益 {profit:,.0f} 円 / 期末キャッシュ {cash:,.0f} 円\n".format(
                    name=row.get("シナリオ", "-"),
                    sales=row.get("年間売上", 0.0),
                    profit=row.get("年間利益", 0.0),
                    cash=row.get("期末キャッシュ", 0.0),
                )
            )

    if profit_range_summary is not None and not profit_range_summary.empty:
        buffer.write("\n[単価・数量・固定費の利益レンジ]\n")
        for _, row in profit_range_summary.iterrows():
            buffer.write(
                "- {scenario}: 最小利益 {min_profit:,.0f} 円 / 最大利益 {max_profit:,.0f} 円 / 利益幅 {range_width:,.0f} 円\n".format(
                    scenario=row.get("scenario", "-"),
                    min_profit=row.get("最小利益", 0.0),
                    max_profit=row.get("最大利益", 0.0),
                    range_width=row.get("利益幅", 0.0),
                )
            )

    if swot:
        buffer.write("\n[SWOT分析]\n")
        for title, key in [
            ("Strengths", "strengths"),
            ("Weaknesses", "weaknesses"),
            ("Opportunities", "opportunities"),
            ("Threats", "threats"),
        ]:
            buffer.write(f"{title}:\n")
            for item in swot.get(key, []):
                buffer.write(f"  - {item}\n")

    if benchmark_df is not None and not benchmark_df.empty:
        buffer.write("\n[業界ベンチマーク]\n")
        for _, row in benchmark_df.iterrows():
            indicator = row.get("指標", "-")
            company_value = row.get("自社")
            industry_value = row.get("業界平均")
            diff_value = row.get("差分")

            def _fmt(value: Any, suffix: str = "") -> str:
                if value is None:
                    return "-"
                try:
                    if np.isnan(value):
                        return "-"
                except TypeError:
                    pass
                return f"{value:,.2f}{suffix}"

            suffix = "%" if "率" in str(indicator) else ""
            buffer.write(
                f"- {indicator}: 自社 {_fmt(company_value, suffix)} / 業界 {_fmt(industry_value, suffix)} / 差分 {_fmt(diff_value, suffix)}\n"
            )

    return buffer.getvalue()


def render_scenario_analysis_section(
    merged_df: pd.DataFrame,
    subscription_df: Optional[pd.DataFrame],
) -> None:
    """Phase2で追加するシナリオ分析ハブを描画する。"""

    st.markdown(
        """
        <div class="surface-card" style="display:flex;justify-content:space-between;align-items:center;gap:1rem;">
            <div>
                <div style="font-size:1.1rem;font-weight:700;" class="with-icon with-icon--tight">{icon}<span>戦略意思決定センター</span></div>
                <div style="color:var(--muted-text-color);font-size:0.9rem;">Scenario Intelligence Hub</div>
            </div>
            <div style="display:flex;gap:0.4rem;">
                <a href="https://www.linkedin.com" style="text-decoration:none;border-radius:999px;padding:0.35rem 0.75rem;background:var(--accent-color);color:#ffffff;font-weight:600;">in</a>
                <a href="https://twitter.com" style="text-decoration:none;border-radius:999px;padding:0.35rem 0.75rem;border:1px solid rgba(255,255,255,0.25);color:var(--text-color);">X</a>
            </div>
        </div>
        """.format(icon=build_ui_icon("SC", tone="accent")),
        unsafe_allow_html=True,
    )

    tabs = st.tabs(["データ入力", "分析結果", "シナリオ比較", "レポート出力"])

    stored_base = st.session_state.get("scenario_uploaded_df")
    base_df = (
        stored_base.copy()
        if isinstance(stored_base, pd.DataFrame) and not stored_base.empty
        else normalize_scenario_input(merged_df)
    )

    monthly_sales = None
    if isinstance(base_df, pd.DataFrame) and not base_df.empty and "sales_amount" in base_df.columns:
        if "order_month" not in base_df.columns:
            base_df["order_month"] = pd.PeriodIndex(pd.to_datetime(base_df["order_date"]), freq="M")
        monthly_sales = base_df.groupby("order_month")["sales_amount"].sum().sort_index()

    base_sales_total = float(base_df["sales_amount"].sum()) if isinstance(base_df, pd.DataFrame) and "sales_amount" in base_df.columns else 0.0
    base_quantity_total = 0.0
    if isinstance(base_df, pd.DataFrame) and "quantity" in base_df.columns:
        base_quantity_total = float(pd.to_numeric(base_df["quantity"], errors="coerce").fillna(0.0).sum())
    base_unit_price = 0.0
    if base_quantity_total > 0:
        base_unit_price = base_sales_total / base_quantity_total
    elif isinstance(base_df, pd.DataFrame) and not base_df.empty:
        base_unit_price = float(base_df["sales_amount"].mean()) if "sales_amount" in base_df.columns else 0.0
    st.session_state["scenario_base_sales_total"] = base_sales_total
    st.session_state["scenario_base_quantity_total"] = base_quantity_total
    st.session_state["scenario_base_unit_price"] = base_unit_price

    with tabs[0]:
        st.header("入力データ")
        uploaded_file = st.file_uploader(
            "シナリオ用の売上データをアップロード", type=["csv", "xlsx"], key="scenario_file_uploader"
        )
        if uploaded_file is not None:
            persist_uploaded_artifact(
                uploaded_file,
                category="scenario",
                label=getattr(uploaded_file, "name", "scenario"),
            )
        download_button_from_df(
            "シナリオ用サンプルCSVをダウンロード",
            get_sample_sales_template(limit=200),
            _build_sample_filename("scenario", "sales"),
        )
        st.caption("フォーマット例を確認したい場合はサンプルCSVをご利用ください。")
        if uploaded_file is not None:
            try:
                if uploaded_file.name.lower().endswith(".csv"):
                    raw_df = pd.read_csv(uploaded_file)
                else:
                    raw_df = pd.read_excel(uploaded_file)
                normalized_df = normalize_scenario_input(raw_df)
                st.session_state["scenario_uploaded_df"] = normalized_df
                st.success("アップロードしたデータをシナリオ基礎データとして設定しました。")
            except Exception as exc:  # pragma: no cover - runtime保護
                st.error(f"データの読み込みに失敗しました: {exc}")

        saved_scenarios = list_saved_uploads("scenario")
        if saved_scenarios:
            option_map = {
                f"{entry['label']} ({entry['saved_at']})": entry for entry in saved_scenarios
            }
            saved_labels = ["保存済みデータを読み込まない"] + list(option_map.keys())
            selected_saved = st.selectbox(
                "保存済みのシナリオデータ",
                saved_labels,
                key="scenario_saved_select",
            )
            if selected_saved in option_map:
                entry = option_map[selected_saved]
                saved_df = load_saved_upload(entry)
                if saved_df is not None and not saved_df.empty:
                    st.session_state["scenario_uploaded_df"] = normalize_scenario_input(saved_df)
                    st.success(f"保存済みファイル『{entry['name']}』を読み込みました。")
                st.session_state["scenario_saved_select"] = saved_labels[0]
                trigger_rerun()

        if st.button("現在のダッシュボードデータを基準にする", key="scenario_use_dashboard"):
            normalized_df = normalize_scenario_input(merged_df)
            if normalized_df is not None and not normalized_df.empty:
                st.session_state["scenario_uploaded_df"] = normalized_df
                st.success("ダッシュボードの集計結果を基にシナリオ分析を行います。")
            else:
                st.warning("利用可能なダッシュボードデータがありません。")

        preview_df = st.session_state.get("scenario_uploaded_df")
        if isinstance(preview_df, pd.DataFrame) and not preview_df.empty:
            st.caption("現在アクティブなシナリオ基礎データ")
            st.dataframe(preview_df.tail(10))
        else:
            st.info("シナリオ用のデータをアップロードするか、既存データを読み込んでください。")

        scenarios = st.session_state.setdefault("scenario_inputs", [])
        form_defaults = st.session_state.setdefault(
            "scenario_form_defaults",
            {
                "unit_price_adjust_pct": 10.0,
                "quantity_adjust_pct": 10.0,
                "fixed_cost_adjust_pct": 10.0,
            },
        )
        with st.form("scenario_entry_form", clear_on_submit=True):
            st.subheader("シナリオパラメータ")
            default_name = f"シナリオ {len(scenarios) + 1}"
            name_col, horizon_col = st.columns([2, 1])
            scenario_name = name_col.text_input("シナリオ名", value=default_name)
            horizon = horizon_col.slider("分析期間 (ヶ月)", min_value=3, max_value=36, value=12)

            st.markdown("#### 主要指標")
            metric_cols = st.columns(3)
            with metric_cols[0]:
                growth = enhanced_number_input(
                    "売上成長率 (%)",
                    key="scenario_growth_input",
                    default_value=5.0,
                    min_value=-50.0,
                    max_value=150.0,
                    step=0.5,
                    number_format="%.1f",
                    text_format=",.1f",
                    placeholder="例: 5",
                )
            with metric_cols[1]:
                margin = enhanced_number_input(
                    "営業利益率 (%)",
                    key="scenario_margin_input",
                    default_value=12.0,
                    min_value=0.0,
                    max_value=100.0,
                    step=0.5,
                    number_format="%.1f",
                    text_format=",.1f",
                    placeholder="例: 12.5",
                )
            with metric_cols[2]:
                funding = enhanced_number_input(
                    "資金調達額 (円)",
                    key="scenario_funding_input",
                    default_value=0.0,
                    min_value=0.0,
                    step=100_000.0,
                    number_format="%.0f",
                    text_format=",.0f",
                    placeholder="例: 1,500,000",
                )

            with st.expander("詳細な感度設定", expanded=False):
                adjust_cols = st.columns(3)
                with adjust_cols[0]:
                    unit_price_adjust_pct = enhanced_number_input(
                        "単価調整幅 (%)",
                        key="scenario_unit_price_adjust",
                        default_value=float(form_defaults.get("unit_price_adjust_pct", 10.0)),
                        min_value=0.0,
                        max_value=100.0,
                        step=0.5,
                        number_format="%.1f",
                        text_format=",.1f",
                        placeholder="例: 10",
                        help_text="基準単価に対して感度分析で揺らす範囲を指定します。",
                    )
                with adjust_cols[1]:
                    quantity_adjust_pct = enhanced_number_input(
                        "数量調整幅 (%)",
                        key="scenario_quantity_adjust",
                        default_value=float(form_defaults.get("quantity_adjust_pct", 10.0)),
                        min_value=0.0,
                        max_value=100.0,
                        step=0.5,
                        number_format="%.1f",
                        text_format=",.1f",
                        placeholder="例: 8",
                        help_text="基準数量に対して感度分析で揺らす範囲を指定します。",
                    )
                with adjust_cols[2]:
                    fixed_cost_adjust_pct = enhanced_number_input(
                        "固定費調整幅 (%)",
                        key="scenario_fixed_cost_adjust",
                        default_value=float(form_defaults.get("fixed_cost_adjust_pct", 10.0)),
                        min_value=0.0,
                        max_value=100.0,
                        step=0.5,
                        number_format="%.1f",
                        text_format=",.1f",
                        placeholder="例: 12",
                        help_text="基準固定費に対する増減幅です。",
                    )
            submitted = st.form_submit_button("シナリオを追加")
            if submitted:
                form_defaults.update(
                    {
                        "unit_price_adjust_pct": float(unit_price_adjust_pct),
                        "quantity_adjust_pct": float(quantity_adjust_pct),
                        "fixed_cost_adjust_pct": float(fixed_cost_adjust_pct),
                    }
                )
                st.session_state["scenario_form_defaults"] = form_defaults
                scenarios.append(
                    {
                        "name": scenario_name,
                        "growth": growth,
                        "margin": margin,
                        "funding": funding,
                        "horizon": horizon,
                        "unit_price_adjust_pct": float(unit_price_adjust_pct),
                        "quantity_adjust_pct": float(quantity_adjust_pct),
                        "fixed_cost_adjust_pct": float(fixed_cost_adjust_pct),
                        "base_sales_total": base_sales_total,
                        "base_quantity_total": base_quantity_total,
                        "base_unit_price": base_unit_price,
                        "base_fixed_cost": float(st.session_state.get("sidebar_fixed_cost", DEFAULT_FIXED_COST)),
                    }
                )
                st.session_state["scenario_inputs"] = scenarios
                st.success(f"シナリオ『{scenario_name}』を登録しました。")

        if scenarios:
            st.markdown("### 登録済みシナリオ")
            for idx, scenario in enumerate(list(scenarios)):
                info_col, remove_col = st.columns([5, 1])
                info_col.markdown(
                    f"**{scenario.get('name', 'シナリオ')}** — 成長率 {scenario.get('growth', 0.0):.1f}% / 利益率 {scenario.get('margin', 0.0):.1f}% / "
                    f"調達額 {scenario.get('funding', 0.0):,.0f} 円 / 期間 {scenario.get('horizon', 0)}ヶ月 / "
                    f"調整幅(単価 {scenario.get('unit_price_adjust_pct', 0.0):.1f}%, 数量 {scenario.get('quantity_adjust_pct', 0.0):.1f}%, 固定費 {scenario.get('fixed_cost_adjust_pct', 0.0):.1f}%)"
                )
                if remove_col.button("削除", key=f"remove_scenario_{idx}"):
                    scenarios.pop(idx)
                    st.session_state["scenario_inputs"] = scenarios
                    trigger_rerun()

    with tabs[1]:
        st.header("分析結果")
        if base_df is None or base_df.empty:
            st.info("データが読み込まれていません。先にデータ入力タブで設定してください。")
            st.session_state["phase2_swot"] = None
            st.session_state["phase2_benchmark"] = None
        else:
            growth_rate = calculate_recent_growth(monthly_sales)
            if "order_month" not in base_df.columns:
                base_df["order_month"] = pd.PeriodIndex(pd.to_datetime(base_df["order_date"]), freq="M")
            try:
                kpi_dict = calculate_kpis(base_df, subscription_df)
            except Exception:  # pragma: no cover - 安全策
                kpi_dict = {}

            latest_sales = (
                float(monthly_sales.iloc[-1]) if monthly_sales is not None and not monthly_sales.empty else 0.0
            )
            delta_label = f"{growth_rate * 100:.1f}%" if growth_rate is not None else None
            metric_cols = st.columns(3)
            metric_cols[0].metric("最新月売上", f"{latest_sales:,.0f} 円", delta=delta_label)
            gross_margin_pct = kpi_dict.get("gross_margin_rate")
            if gross_margin_pct is not None and np.isfinite(gross_margin_pct):
                metric_cols[1].metric("粗利率", f"{gross_margin_pct * 100:,.1f}%")
            active_customers = kpi_dict.get("active_customers")
            if active_customers is not None and np.isfinite(active_customers):
                metric_cols[2].metric("アクティブ顧客", f"{active_customers:,.0f} 人")

            if monthly_sales is not None and not monthly_sales.empty:
                trend_df = monthly_sales.reset_index()
                trend_df["period_start"] = trend_df["order_month"].dt.to_timestamp()
                chart = alt.Chart(trend_df).mark_line(point=True).encode(
                    x=alt.X("period_start:T", title="期間"),
                    y=alt.Y("sales_amount:Q", title="売上高"),
                    tooltip=["period_start:T", alt.Tooltip("sales_amount", title="売上高", format=",")],
                ).properties(height=320)
                st.altair_chart(apply_altair_theme(chart), use_container_width=True)

            swot = build_swot_insights(kpi_dict, growth_rate)
            st.session_state["phase2_swot"] = swot
            swot_cols = st.columns(4)
            swot_titles = [
                ("Strengths", "strengths", "S", "success"),
                ("Weaknesses", "weaknesses", "W", "warning"),
                ("Opportunities", "opportunities", "O", "accent"),
                ("Threats", "threats", "T", "error"),
            ]
            for col, (title, key, code, tone) in zip(swot_cols, swot_titles):
                icon_html = build_ui_icon(code, tone=tone)
                col.markdown(
                    "#### <span class='with-icon with-icon--tight'>{icon}<span>{title}</span></span>".format(
                        icon=icon_html,
                        title=html.escape(title),
                    ),
                    unsafe_allow_html=True,
                )
                for item in swot.get(key, []):
                    col.markdown(f"- {item}")

            benchmark_df = build_industry_benchmark_table(kpi_dict)
            st.session_state["phase2_benchmark"] = benchmark_df
            st.markdown("### 業界ベンチマーク比較")
            if benchmark_df.empty:
                st.info("比較可能なKPIが不足しています。")
            else:
                st.dataframe(benchmark_df.style.format({"自社": "{:.2f}", "業界平均": "{:.2f}", "差分": "{:.2f}"}))

            if gross_margin_pct is not None and monthly_sales is not None and not monthly_sales.empty:
                base_sales = float(monthly_sales.iloc[-1])
                growth_points = np.linspace(-0.05, 0.20, 6)
                sensitivity_rows = []
                for point in growth_points:
                    annual_sales_projection = base_sales * ((1 + point) ** 12)
                    annual_profit_projection = annual_sales_projection * gross_margin_pct
                    sensitivity_rows.append(
                        {
                            "growth_pct": point * 100,
                            "annual_profit": annual_profit_projection,
                        }
                    )
                sensitivity_df = pd.DataFrame(sensitivity_rows)
                chart = alt.Chart(sensitivity_df).mark_line(point=True).encode(
                    x=alt.X("growth_pct:Q", title="成長率 (%)"),
                    y=alt.Y("annual_profit:Q", title="年間利益"),
                    tooltip=[
                        alt.Tooltip("growth_pct", title="成長率", format=".1f"),
                        alt.Tooltip("annual_profit", title="年間利益", format=",")
                    ],
                ).properties(height=280)
                st.markdown("### 感度分析: 成長率と年間利益")
                st.altair_chart(apply_altair_theme(chart), use_container_width=True)

    with tabs[2]:
        st.header("シナリオ比較")
        scenarios = st.session_state.get("scenario_inputs", [])
        if base_df is None or base_df.empty:
            st.info("先にデータ入力タブで基礎データを設定してください。")
            st.session_state["phase2_summary_df"] = None
            st.session_state["scenario_profit_ranges"] = None
            st.session_state["scenario_profit_range_summary"] = None
            render_shared_scenario_preview()
        elif not scenarios:
            st.info("比較するシナリオを追加してください。")
            st.session_state["phase2_summary_df"] = None
            st.session_state["scenario_profit_ranges"] = None
            st.session_state["scenario_profit_range_summary"] = None
            render_shared_scenario_preview()
        else:
            total = len(scenarios)
            progress = st.progress(0.0)
            results: List[pd.DataFrame] = []
            summaries: List[Dict[str, Any]] = []
            sensitivity_frames: List[pd.DataFrame] = []
            profit_range_frames: List[pd.DataFrame] = []
            with st.spinner("シナリオを計算しています..."):
                with ThreadPoolExecutor(max_workers=min(4, total)) as executor:
                    futures = {
                        executor.submit(run_scenario_projection, monthly_sales, scenario): scenario
                        for scenario in scenarios
                    }
                    for idx, future in enumerate(as_completed(futures), start=1):
                        (
                            scenario_name,
                            projection_df,
                            summary_row,
                            sensitivity_df,
                            profit_range_df,
                        ) = future.result()
                        results.append(projection_df)
                        summaries.append(summary_row)
                        if sensitivity_df is not None and not sensitivity_df.empty:
                            sensitivity_frames.append(sensitivity_df)
                        if profit_range_df is not None and not profit_range_df.empty:
                            profit_range_frames.append(profit_range_df)
                        progress.progress(idx / total)
            progress.empty()

            if results:
                combined_df = pd.concat(results, ignore_index=True)
                st.session_state["scenario_results"] = combined_df
                summary_df = pd.DataFrame(summaries)
                st.session_state["phase2_summary_df"] = summary_df
                range_summary: Optional[pd.DataFrame] = None
                if profit_range_frames:
                    combined_profit_range = pd.concat(profit_range_frames, ignore_index=True)
                    st.session_state["scenario_profit_ranges"] = combined_profit_range
                    range_summary = (
                        combined_profit_range.groupby("scenario")["profit"]
                        .agg(["min", "max"])
                        .rename(columns={"min": "最小利益", "max": "最大利益"})
                    )
                    range_summary["利益幅"] = range_summary["最大利益"] - range_summary["最小利益"]
                    range_summary = range_summary.reset_index()
                    st.session_state["scenario_profit_range_summary"] = range_summary
                else:
                    st.session_state["scenario_profit_ranges"] = None
                    st.session_state["scenario_profit_range_summary"] = None
                    range_summary = None

                st.markdown("### 年間売上・利益比較")
                st.dataframe(
                    summary_df.style.format(
                        {
                            "年間売上": "{:.0f}",
                            "年間利益": "{:.0f}",
                            "平均月次売上": "{:.0f}",
                            "期末キャッシュ": "{:.0f}",
                            "成長率(%)": "{:.1f}",
                            "利益率(%)": "{:.1f}",
                            "調整固定費(月額)": "{:.0f}",
                            "調達額": "{:.0f}",
                        }
                    )
                )

                render_scenario_radar_chart(summary_df)

                export_cols = st.columns(3)
                with export_cols[0]:
                    download_button_from_df("比較サマリーCSV", summary_df, "scenario_summary.csv")
                with export_cols[1]:
                    download_excel_from_df("比較サマリーExcel", summary_df, "scenario_summary.xlsx")
                with export_cols[2]:
                    if REPORTLAB_AVAILABLE:
                        download_pdf_from_df(
                            "比較サマリーPDF",
                            summary_df,
                            "scenario_summary.pdf",
                            title="シナリオ比較サマリー",
                        )
                    else:
                        st.button(
                            "比較サマリーPDF",
                            disabled=True,
                            help="reportlabパッケージをインストールするとPDF出力が利用できます。",
                        )
                        st.caption(
                            "PDF出力には外部パッケージreportlabが必要です。"
                            " サーバー環境にインストールできない場合はExcelエクスポートをご利用ください。"
                        )
                        st.caption("pip install reportlab で追加できます。")

                sales_chart = alt.Chart(combined_df).mark_line().encode(
                    x=alt.X("period:T", title="期間"),
                    y=alt.Y("sales:Q", title="売上高"),
                    color=alt.Color("scenario:N", title="シナリオ"),
                    tooltip=[
                        "scenario", "period_label", alt.Tooltip("sales", title="売上高", format=",")
                    ],
                ).properties(height=360)
                st.altair_chart(apply_altair_theme(sales_chart), use_container_width=True)

                if sensitivity_frames:
                    sensitivity_combined = pd.concat(sensitivity_frames, ignore_index=True)
                    sensitivity_chart = alt.Chart(sensitivity_combined).mark_line(point=True).encode(
                        x=alt.X("margin:Q", title="利益率 (%)"),
                        y=alt.Y("annual_profit:Q", title="年間利益"),
                        color=alt.Color("scenario:N", title="シナリオ"),
                        tooltip=[
                            "scenario",
                            alt.Tooltip("margin", title="利益率", format=".1f"),
                            alt.Tooltip("annual_profit", title="年間利益", format=",")
                        ],
                    ).properties(height=320)
                    st.markdown("### 感度分析: 利益率別の年間利益")
                    st.altair_chart(apply_altair_theme(sensitivity_chart), use_container_width=True)
                if profit_range_frames:
                    combined_profit_range = pd.concat(profit_range_frames, ignore_index=True)
                    range_chart = alt.Chart(combined_profit_range).mark_line(point=True).encode(
                        x=alt.X("adjustment_label:N", title="調整ケース"),
                        y=alt.Y("profit:Q", title="年間利益"),
                        color=alt.Color("scenario:N", title="シナリオ"),
                        tooltip=[
                            "scenario",
                            "adjustment_label",
                            alt.Tooltip("profit", title="年間利益", format=",")
                        ],
                    ).properties(height=320)
                    st.markdown("### 単価・数量・固定費の利益レンジ")
                    st.altair_chart(apply_altair_theme(range_chart), use_container_width=True)
                    tornado_df = build_tornado_dataframe(combined_profit_range)
                    if not tornado_df.empty:
                        st.markdown("### トルネードチャート: 利益感度")
                        tornado_chart = (
                            alt.Chart(tornado_df)
                            .mark_bar()
                            .encode(
                                y=alt.Y(
                                    "factor:N",
                                    sort=alt.SortField("impact_abs", order="descending"),
                                    title="要因",
                                ),
                                x=alt.X(
                                    "impact:Q",
                                    title="年間利益インパクト",
                                    axis=alt.Axis(format=",.0f"),
                                ),
                                color=alt.Color(
                                    "direction:N",
                                    title="方向",
                                    scale=alt.Scale(range=[ERROR_COLOR, SUCCESS_COLOR]),
                                ),
                                tooltip=[
                                    alt.Tooltip("scenario:N", title="シナリオ"),
                                    alt.Tooltip("factor:N", title="要因"),
                                    alt.Tooltip("direction:N", title="変化方向"),
                                    alt.Tooltip("impact:Q", title="利益影響", format=",.0f"),
                                ],
                            )
                            .properties(height=320)
                            .facet(column=alt.Column("scenario:N", title="シナリオ"))
                        )
                        st.altair_chart(apply_altair_theme(tornado_chart), use_container_width=True)
                    display_profit_range = combined_profit_range.rename(
                        columns={
                            "unit_price_pct": "単価変化(%)",
                            "quantity_pct": "数量変化(%)",
                            "fixed_cost_pct": "固定費変化(%)",
                            "sales": "年間売上",
                            "profit": "年間利益",
                        }
                    )
                    st.dataframe(
                        display_profit_range.style.format(
                            {
                                "単価変化(%)": "{:.1f}",
                                "数量変化(%)": "{:.1f}",
                                "固定費変化(%)": "{:.1f}",
                                "年間売上": "{:.0f}",
                                "年間利益": "{:.0f}",
                            }
                        )
                    )
                    range_summary = st.session_state.get("scenario_profit_range_summary")
                    if isinstance(range_summary, pd.DataFrame) and not range_summary.empty:
                        st.dataframe(
                            range_summary.style.format(
                                {
                                    "最小利益": "{:.0f}",
                                    "最大利益": "{:.0f}",
                                    "利益幅": "{:.0f}",
                                }
                            )
                        )
                render_scenario_share_controls(summary_df, range_summary)
            else:
                st.info("シナリオ計算結果を取得できませんでした。")
                st.session_state["phase2_summary_df"] = None
                st.session_state["scenario_profit_ranges"] = None
                st.session_state["scenario_profit_range_summary"] = None

    with tabs[3]:
        st.header("レポート出力")
        scenario_results = st.session_state.get("scenario_results")
        summary_df = st.session_state.get("phase2_summary_df")
        swot = st.session_state.get("phase2_swot")
        benchmark_df = st.session_state.get("phase2_benchmark")

        if scenario_results is None or scenario_results.empty:
            st.info("シナリオ比較を実行するとレポートを出力できます。")
        else:
            profit_range_summary = st.session_state.get("scenario_profit_range_summary")
            report_text = generate_phase2_report(summary_df, swot, benchmark_df, profit_range_summary)
            st.session_state["phase2_report_summary"] = report_text
            st.download_button(
                "PDF出力 (テキスト)",
                report_text,
                file_name="scenario_report.txt",
                mime="text/plain",
            )

            excel_buffer = io.BytesIO()
            with pd.ExcelWriter(excel_buffer, engine="xlsxwriter") as writer:
                scenario_results.to_excel(writer, sheet_name="scenarios", index=False)
                if summary_df is not None:
                    summary_df.to_excel(writer, sheet_name="summary", index=False)
                if benchmark_df is not None:
                    benchmark_df.to_excel(writer, sheet_name="benchmark", index=False)
                profit_range_df = st.session_state.get("scenario_profit_ranges")
                if isinstance(profit_range_df, pd.DataFrame) and not profit_range_df.empty:
                    profit_range_df.to_excel(writer, sheet_name="profit_ranges", index=False)
                if isinstance(profit_range_summary, pd.DataFrame) and not profit_range_summary.empty:
                    profit_range_summary.to_excel(writer, sheet_name="profit_range_summary", index=False)
            st.download_button(
                "Excel出力",
                excel_buffer.getvalue(),
                file_name="scenario_report.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

            st.markdown("#### レポートサマリー")
            st.text(report_text)


def render_sidebar_upload_expander(
    label: str,
    *,
    uploader_key: str,
    description: str,
    multiple: bool,
    meta_text: str,
    help_text: str,
    file_types: Optional[List[str]] = None,
    sample_label: Optional[str] = None,
    sample_generator: Optional[Callable[[], pd.DataFrame]] = None,
    sample_filename: Optional[str] = None,
    sample_note: Optional[str] = None,
) -> Any:
    """サイドバーにアイコン付きのアップロード用アコーディオンを描画する。"""

    file_types = file_types or ["xlsx", "xls", "csv"]
    with st.sidebar.expander(label, expanded=False):
        st.markdown(
            f"""
            <div class="sidebar-upload-card">
                <div class="sidebar-upload-card__icons">
                    <span class="sidebar-upload-card__icon sidebar-upload-card__icon--csv">CSV</span>
                    <span class="sidebar-upload-card__icon sidebar-upload-card__icon--excel">XLSX</span>
                </div>
                <div class="sidebar-upload-card__body">
                    <div class="sidebar-upload-card__title">CSV / Excelファイルに対応</div>
                    <div class="sidebar-upload-card__meta">{meta_text}</div>
                    <p class="sidebar-upload-card__desc">{description}</p>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        uploaded = st.file_uploader(
            "ファイルを選択",
            type=file_types,
            accept_multiple_files=multiple,
            key=f"{uploader_key}_uploader",
            label_visibility="collapsed",
            help=help_text,
        )
        if uploaded:
            if multiple:
                for item in uploaded:
                    persist_uploaded_artifact(
                        item,
                        category=uploader_key,
                        label=getattr(item, "name", uploader_key),
                    )
            else:
                persist_uploaded_artifact(
                    uploaded,
                    category=uploader_key,
                    label=getattr(uploaded, "name", uploader_key),
                )
        if sample_generator is not None:
            try:
                sample_df = sample_generator()
            except Exception as exc:  # pragma: no cover - 安全装置
                st.caption(f"サンプルデータの生成に失敗しました: {exc}")
            else:
                if isinstance(sample_df, pd.DataFrame) and not sample_df.empty:
                    if sample_note:
                        st.caption(sample_note)
                    download_button_from_df(
                        sample_label or "サンプルフォーマットをダウンロード",
                        sample_df,
                        sample_filename or "sample.csv",
                    )
    return uploaded


def _assignment_widget_key(file_name: str) -> str:
    """ファイル割当セレクトボックス用のセッションキーを生成する。"""

    digest = hashlib.md5(file_name.encode("utf-8")).hexdigest()
    return f"sales_wizard_assign_{digest}"


def infer_channel_from_name(
    file_name: str, configs: Sequence[Dict[str, str]]
) -> Optional[str]:
    """ファイル名に含まれるキーワードからチャネルを推測する。"""

    if not file_name:
        return None

    normalized = file_name.lower().replace(" ", "")
    for config in configs:
        channel = config.get("channel", "")
        label = config.get("label", "")
        candidates = {channel, label}
        for candidate in list(candidates):
            if candidate:
                candidates.add(candidate.replace(" ", ""))

        for candidate in candidates:
            token = str(candidate).lower()
            if token and token in normalized:
                return channel

    return None


def format_file_size(num_bytes: Optional[int]) -> str:
    """ファイルサイズを人が読みやすい形式で返す。"""

    if not num_bytes or num_bytes < 0:
        return "-"

    units = ["B", "KB", "MB", "GB", "TB"]
    size = float(num_bytes)
    for unit in units:
        if size < 1024.0 or unit == units[-1]:
            return f"{size:.1f} {unit}"
        size /= 1024.0

    return f"{num_bytes} B"


def render_sales_upload_wizard(
    configs: Sequence[Dict[str, str]],
    parent_container: Optional[Any] = None,
) -> Dict[str, List[Any]]:
    """売上データ取り込みウィザードを描画し、チャネルごとのファイルを返す。"""

    channel_files: Dict[str, List[Any]] = {config["channel"]: [] for config in configs}
    assignments: Dict[str, str] = st.session_state.setdefault("sales_wizard_assignments", {})
    uploaded_files: List[Any] = []
    preview_rows: List[Dict[str, str]] = []
    unassigned_count = 0

    container = (
        parent_container.container() if parent_container is not None else st.sidebar.container()
    )

    with container:
        st.markdown("<div class='sidebar-wizard-title'>売上データ取り込みウィザード</div>", unsafe_allow_html=True)
        st.caption("複数チャネルの売上ファイルをまとめてアップロードし、チャネルへ一括割当できます。")

        st.markdown("**ステップ1. ファイルをまとめてアップロード**")
        uploaded = st.file_uploader(
            "売上データファイルを追加",
            type=["xlsx", "xls", "csv"],
            accept_multiple_files=True,
            key="sales_wizard_uploader",
            help=UPLOAD_HELP_MULTIPLE,
            label_visibility="collapsed",
        )
        if uploaded:
            uploaded_files = list(uploaded)
        else:
            uploaded_files = []

        signature_index: Dict[str, Optional[str]] = {}
        for uploaded_file in uploaded_files:
            _, _, signature = persist_uploaded_artifact(
                uploaded_file,
                category="sales",
                label=getattr(uploaded_file, "name", "sales"),
            )
            if signature:
                signature_index[getattr(uploaded_file, "name", "")] = signature

        st.caption(UPLOAD_META_MULTIPLE)

        saved_entries = list_saved_uploads("sales")
        if saved_entries:
            option_map = {
                f"{entry['label']} ({entry['saved_at']})": entry for entry in saved_entries
            }
            saved_labels = ["保存済みファイルを追加しない"] + list(option_map.keys())
            selected_saved = st.selectbox(
                "保存済みファイルを利用",
                saved_labels,
                key="sales_saved_select",
            )
            if selected_saved in option_map:
                entry = option_map[selected_saved]
                active_saved = set(st.session_state.get("active_saved_sales", []))
                if entry["path"] not in active_saved:
                    active_saved.add(entry["path"])
                    st.session_state["active_saved_sales"] = list(active_saved)
                    st.success(
                        f"保存済みファイル『{entry['name']}』を取り込み対象に追加しました。"
                    )
                else:
                    st.info("選択した保存済みファイルは既に追加されています。")
                st.session_state["sales_saved_select"] = saved_labels[0]
                trigger_rerun()

            active_saved = st.session_state.get("active_saved_sales", [])
            if active_saved:
                summary = ", ".join(Path(path).name for path in active_saved[:3])
                if len(active_saved) > 3:
                    summary += f" 他{len(active_saved) - 3}件"
                st.caption(f"保存済みファイル: {summary}")
                if st.button("保存済みファイルの割当をクリア", key="clear_saved_sales"):
                    st.session_state["active_saved_sales"] = []
                    st.success("保存済みファイルの選択をクリアしました。")
                    trigger_rerun()

        with st.expander("チャネル別サンプルフォーマットを確認"):
            for config in configs:
                channel = config["channel"]
                try:
                    sample_df = get_sample_sales_template(channel)
                except Exception as exc:  # pragma: no cover - サンプル生成失敗時の保護
                    st.caption(f"{channel}サンプルの生成に失敗しました: {exc}")
                    continue
                download_button_from_df(
                    f"{channel}サンプルCSVをダウンロード",
                    sample_df,
                    _build_sample_filename("sales", channel),
                )

        st.caption("主なチャネルと想定ファイルの例:")
        for config in configs:
            st.caption(f"・{config['channel']}: {config['description']}")

        st.markdown("**ステップ2. ファイルごとにチャネルを割当**")
        options = [CHANNEL_ASSIGNMENT_PLACEHOLDER] + list(channel_files.keys())
        current_names = {getattr(file, "name", "") for file in uploaded_files}

        for stale in list(assignments.keys()):
            if stale not in current_names:
                assignments.pop(stale, None)

        if not uploaded_files:
            st.caption("ファイルをアップロードするとチャネル割当の設定が表示されます。")
        else:
            for uploaded_file in uploaded_files:
                file_name = getattr(uploaded_file, "name", "")
                detected = assignments.get(file_name) or infer_channel_from_name(file_name, configs)
                default_option = detected if detected in channel_files else CHANNEL_ASSIGNMENT_PLACEHOLDER
                widget_key = _assignment_widget_key(file_name)
                if widget_key not in st.session_state or st.session_state[widget_key] not in options:
                    st.session_state[widget_key] = default_option

                st.markdown(
                    "<div class='wizard-file-item'>"
                    f"<span class='wizard-file-item__name'>{html.escape(file_name)}</span>"
                    f"<span class='wizard-file-item__size'>{format_file_size(getattr(uploaded_file, 'size', None))}</span>"
                    "</div>",
                    unsafe_allow_html=True,
                )

                selection = st.selectbox(
                    f"{file_name}のチャネル",
                    options=options,
                    key=widget_key,
                    label_visibility="collapsed",
                )

                if selection in channel_files:
                    channel_files[selection].append(uploaded_file)
                    assignments[file_name] = selection
                    assigned_label = selection
                    update_upload_metadata(
                        signature_index.get(file_name), channel=selection
                    )
                else:
                    assignments.pop(file_name, None)
                    assigned_label = "未割当"
                    unassigned_count += 1
                    update_upload_metadata(signature_index.get(file_name), channel=None)

                preview_rows.append(
                    {
                        "ファイル名": file_name,
                        "割当チャネル": assigned_label,
                        "サイズ": format_file_size(getattr(uploaded_file, "size", None)),
                    }
                )

        st.markdown("**ステップ3. プレビュー & 検証メッセージ**")
        if preview_rows:
            preview_df = pd.DataFrame(preview_rows)
            st.table(preview_df)
            if unassigned_count:
                st.warning(f"{unassigned_count}件のファイルが未割当です。チャネルを選択してください。")
            else:
                st.success("すべてのファイルにチャネルが割り当てられました。取り込みを実行できます。")
        else:
            st.caption("チャネル割当結果のプレビューがここに表示されます。")

    st.session_state["sales_wizard_assignments"] = assignments

    return channel_files


def main() -> None:
    init_phase2_session_state()
    ensure_theme_state_defaults()
    ensure_language_state_defaults()
    load_share_payload_from_query()

    st.markdown("<div id='page-top'></div>", unsafe_allow_html=True)

    st.session_state.setdefault("hotkey_last_processed", None)

    if st.session_state.pop("pending_enable_sample_data", False):
        set_state_and_widget("use_sample_data", True)

    if "sidebar_onboarding_visible" not in st.session_state:
        st.session_state["sidebar_onboarding_visible"] = True

    show_onboarding = bool(st.session_state.get("sidebar_onboarding_visible", True))
    toggle_label = ("▼ " if show_onboarding else "▶ ") + "はじめに"

    def _toggle_sidebar_onboarding() -> None:
        current = bool(st.session_state.get("sidebar_onboarding_visible", True))
        st.session_state["sidebar_onboarding_visible"] = not current

    st.sidebar.button(
        toggle_label,
        key="toggle_sidebar_onboarding",
        use_container_width=True,
        type="primary",
        help="クリックして『はじめに』セクションの表示/非表示を切り替えます。",
        on_click=_toggle_sidebar_onboarding,
    )

    show_onboarding = bool(st.session_state.get("sidebar_onboarding_visible", True))

    onboarding_container = st.sidebar.container()

    if "use_sample_data" not in st.session_state:
        set_state_and_widget("use_sample_data", True)
    else:
        ensure_widget_mirror("use_sample_data")
    use_sample_data = bool(st.session_state.get("use_sample_data", True))

    theme_expander = st.sidebar.expander("テーマと表示設定", expanded=False)
    with theme_expander:
        st.caption(translate("theme_caption", default="配色やフォントサイズを調整して見やすいダッシュボードにカスタマイズできます。"))

        current_language = get_ui_language()
        language_index = next(
            (idx for idx, (code, _) in enumerate(LANGUAGE_OPTIONS) if code == current_language),
            0,
        )
        selected_language_label = st.selectbox(
            translate("theme_language_label", default="表示言語"),
            options=[label for _, label in LANGUAGE_OPTIONS],
            index=language_index,
            help=translate(
                "theme_language_help",
                default="日本語/英語の表記とヘルプを切り替えます。",
            ),
        )
        label_to_code = {label: code for code, label in LANGUAGE_OPTIONS}
        st.session_state["ui_language"] = label_to_code.get(
            selected_language_label, DEFAULT_LANGUAGE
        )

        default_theme_mode = st.session_state.get("ui_theme_mode", "light")
        dark_mode = st.toggle(
            translate("theme_toggle_label", default="ダークテーマ"),
            value=(default_theme_mode == "dark"),
            key="ui_theme_toggle",
            help=translate(
                "theme_toggle_help",
                default="ライトテーマに切り替えると背景が明るい配色になります。",
            ),
        )
        st.session_state["ui_theme_mode"] = "dark" if dark_mode else "light"

        font_scale_default = int(round(st.session_state.get("ui_font_scale", 1.0) * 100))
        font_scale_default = max(85, min(120, font_scale_default))
        font_scale_percent = st.slider(
            translate("font_size_label", default="本文フォントサイズ"),
            min_value=85,
            max_value=120,
            value=font_scale_default,
            step=5,
            help=translate(
                "font_size_help",
                default="本文や表の文字サイズを調整します (基準値=100)。",
            ),
        )
        font_scale = font_scale_percent / 100.0
        st.session_state["ui_font_scale"] = font_scale

        if dark_mode:
            variant_options = list(DARK_THEME_VARIANTS.keys())
            current_variant = st.session_state.get("ui_dark_variant_saved", DEFAULT_DARK_THEME_VARIANT)
            try:
                variant_index = variant_options.index(current_variant)
            except ValueError:
                variant_index = 0
            selected_variant = st.selectbox(
                translate("dark_variant_label", default="暗色テーマのコントラスト"),
                variant_options,
                index=variant_index,
                format_func=lambda key: DARK_THEME_VARIANT_LABELS.get(key, key),
                help=translate(
                    "dark_variant_help",
                    default="暗色テーマ時の背景/境界のコントラストを調整します。",
                ),
            )
            st.session_state["ui_dark_variant_saved"] = selected_variant
            st.session_state["ui_dark_variant"] = selected_variant

            palette_options = list(COLOR_PALETTE_PRESETS.keys())
            current_palette = st.session_state.get("ui_dark_palette_saved", DEFAULT_CHART_PALETTE_KEY)
            try:
                palette_index = palette_options.index(current_palette)
            except ValueError:
                palette_index = 0
            selected_palette = st.selectbox(
                translate("palette_label", default="チャートカラーパレット"),
                palette_options,
                index=palette_index,
                format_func=lambda key: str(COLOR_PALETTE_PRESETS[key]["label"]),
                help=translate(
                    "palette_help",
                    default="色覚多様性に配慮した配色に切り替えられます。",
                ),
            )
            st.session_state["ui_dark_palette_saved"] = selected_palette
            st.session_state["ui_color_palette"] = selected_palette
        else:
            st.session_state["ui_dark_variant"] = st.session_state.get(
                "ui_dark_variant_saved", DEFAULT_DARK_THEME_VARIANT
            )
            st.session_state["ui_color_palette"] = DEFAULT_CHART_PALETTE_KEY

    inject_mckinsey_style(
        dark_mode=dark_mode,
        theme_variant=st.session_state.get("ui_dark_variant", DEFAULT_DARK_THEME_VARIANT),
        font_scale=font_scale,
    )

    render_intro_section()

    st.sidebar.toggle(
        "管理者モード",
        value=bool(st.session_state.get("admin_mode_toggle", False)),
        key="admin_mode_toggle",
        help="管理者向けの詳細なログ表示を有効化します。",
    )

    st.sidebar.image(COMPANY_LOGO_URL, width=140)
    st.sidebar.caption("McKinsey inspired analytics suite")
    st.sidebar.header("データ設定")

    if use_sample_data:
        ensure_sample_data_cached()

    wizard_panel = st.sidebar.expander("売上データ設定ウィザード", expanded=False)
    with wizard_panel:
        st.markdown(
            "<div class='sidebar-subheading'>売上データアップロード</div>",
            unsafe_allow_html=True,
        )
        st.caption("複数チャネルのデータをまとめて読み込む場合はこちらを利用してください。")
        channel_files = render_sales_upload_wizard(
            SALES_UPLOAD_CONFIGS, parent_container=wizard_panel
        )

    ancillary_results: Dict[str, Any] = {}
    ancillary_panel = st.sidebar.expander("補助データのアップロード", expanded=False)
    with ancillary_panel:
        st.caption("原価率やサブスクリプションなどの補助データを必要に応じて登録できます。")
        for config in ANCILLARY_UPLOAD_CONFIGS:
            ancillary_results[config["key"]] = render_sidebar_upload_expander(
                config["label"],
                uploader_key=config["key"],
                description=config["description"],
                multiple=config.get("multiple", False),
                meta_text=config.get("meta_text", UPLOAD_META_SINGLE),
                help_text=config.get("help_text", UPLOAD_HELP_SINGLE),
                sample_label=f"{config['label']}サンプルCSVをダウンロード",
                sample_generator=(
                    get_sample_cost_template
                    if config["key"] == "cost"
                    else get_sample_subscription_template
                ),
                sample_filename=_build_sample_filename("ancillary", config["key"]),
                sample_note="期待される列構成を確認できるサンプルです。",
            )

    cost_file = ancillary_results.get("cost")
    subscription_file = ancillary_results.get("subscription")

    remember_last_uploaded_files(channel_files, cost_file, subscription_file)

    last_uploaded = st.session_state.get("last_uploaded")
    if last_uploaded:
        preview = ", ".join(last_uploaded[:3])
        if len(last_uploaded) > 3:
            preview += f" 他{len(last_uploaded) - 3}件"
        st.sidebar.caption(f"前回アップロード: {preview}")

    with st.sidebar.expander("設定の保存・読み込み", expanded=False):
        st.caption("現在のフィルタや手入力KPI、シナリオをJSONとして保存・復元できます。")
        settings_payload = json.dumps(collect_dashboard_settings(), ensure_ascii=False, indent=2)
        st.download_button(
            "現在の設定をダウンロード",
            settings_payload.encode("utf-8"),
            file_name="dashboard_settings.json",
            mime="application/json",
        )
        uploaded_settings = st.file_uploader(
            "保存した設定ファイルを読み込む",
            type="json",
            key="dashboard_settings_loader",
        )
        if uploaded_settings is not None:
            try:
                loaded_settings = json.load(uploaded_settings)
            except json.JSONDecodeError as exc:
                st.error(f"設定ファイルの読み込みに失敗しました: {exc}")
            else:
                apply_dashboard_settings(loaded_settings)
                st.success("保存された設定を反映しました。")
                trigger_rerun()

    if "api_sales_data" not in st.session_state:
        st.session_state["api_sales_data"] = {}
    if "api_sales_validation" not in st.session_state:
        st.session_state["api_sales_validation"] = {}
    if "api_last_fetched" not in st.session_state:
        st.session_state["api_last_fetched"] = {}

    st.sidebar.markdown("---")
    with st.sidebar.expander("API/RPA自動連携設定", expanded=False):
        st.caption("各モールのAPIやRPAが出力したURLを登録すると、手動アップロードなしで売上データを取得できます。")
        for channel in channel_files.keys():
            endpoint = st.text_input(f"{channel} APIエンドポイント", key=f"api_endpoint_{channel}")
            token = st.text_input(
                f"{channel} APIトークン/キー",
                key=f"api_token_{channel}",
                type="password",
                help="必要に応じてBasic認証やBearerトークンを設定してください。",
            )
            params_raw = st.text_input(
                f"{channel} クエリパラメータ (key=value&...)",
                key=f"api_params_{channel}",
                help="日付範囲などの条件が必要な場合に指定します。",
            )

            params_dict = _parse_api_params(params_raw)

            fetch_now = st.button(f"{channel}の最新データを取得", key=f"fetch_api_{channel}")
            if fetch_now:
                if not endpoint:
                    st.warning("エンドポイントURLを入力してください。")
                else:
                    with st.spinner(f"{channel}のデータを取得中..."):
                        fetched_df, fetch_report = fetch_sales_from_endpoint(
                            endpoint,
                            token=token or None,
                            params=params_dict,
                            channel_hint=channel,
                        )
                    st.session_state["api_sales_data"][channel] = fetched_df
                    st.session_state["api_sales_validation"][channel] = fetch_report
                    st.session_state["api_last_fetched"][channel] = datetime.now()
                    if fetch_report.has_errors():
                        st.error(f"{channel}のAPI連携でエラーが発生しました。詳細はデータ管理タブをご確認ください。")
                    elif fetch_report.has_warnings():
                        st.warning(f"{channel}のデータは取得しましたが警告があります。データ管理タブで確認してください。")
                    else:
                        st.success(f"{channel}のデータ取得が完了しました。")

            last_fetch = st.session_state["api_last_fetched"].get(channel)
            if last_fetch:
                status_report: Optional[ValidationReport] = st.session_state["api_sales_validation"].get(channel)
                latest_df = st.session_state["api_sales_data"].get(channel)
                record_count = len(latest_df) if isinstance(latest_df, pd.DataFrame) else 0
                if status_report and status_report.has_errors():
                    status_level = "error"
                elif status_report and status_report.has_warnings():
                    status_level = "warning"
                else:
                    status_level = "ok"
                icon_code, tone, status_label = STATUS_PILL_DETAILS.get(
                    status_level, ("i", "muted", "情報")
                )
                status_icon = build_ui_icon(icon_code, tone=tone)
                st.markdown(
                    "<div class='status-pill status-pill--{cls}'>{icon}<span>状態: {label}</span></div>".format(
                        cls=status_level,
                        icon=status_icon,
                        label=html.escape(status_label),
                    ),
                    unsafe_allow_html=True,
                )
                st.markdown(
                    f"<div class='sidebar-meta'>最終取得: {last_fetch.strftime('%Y-%m-%d %H:%M')} / {record_count:,} 件</div>",
                    unsafe_allow_html=True,
                )

        if st.button("自動取得データをクリア", key="clear_api_sales"):
            st.session_state["api_sales_data"].clear()
            st.session_state["api_sales_validation"].clear()
            st.session_state["api_last_fetched"].clear()
            st.success("保存されていたAPI取得データをクリアしました。")

    finance_controls = st.sidebar.container()
    with finance_controls:
        fixed_cost = enhanced_number_input(
            "月間固定費（販管費のうち人件費・地代等）",
            key="sidebar_fixed_cost_input",
            default_value=float(st.session_state.get("sidebar_fixed_cost", DEFAULT_FIXED_COST)),
            min_value=0.0,
            step=50_000.0,
            number_format="%.0f",
            text_format=",.0f",
            placeholder="例: 450,000",
            help_text="固定費に該当する販管費の合計額です。人件費・地代家賃・システム利用料などを含めて設定します。",
        )
        st.session_state["sidebar_fixed_cost"] = float(fixed_cost)
        starting_cash = enhanced_number_input(
            "現在の現金残高（円）",
            key="sidebar_starting_cash_input",
            default_value=float(st.session_state.get("sidebar_starting_cash", 3_000_000.0)),
            min_value=0.0,
            step=100_000.0,
            number_format="%.0f",
            text_format=",.0f",
            placeholder="例: 3,000,000",
            help_text="ダッシュボード表示時点の現預金残高です。資金繰りの初期値として利用されます。",
        )
        st.session_state["sidebar_starting_cash"] = float(starting_cash)

    with st.sidebar.expander("KPIの手入力（任意）"):
        manual_active = enhanced_number_input(
            "当月アクティブ顧客数",
            key="manual_active_customers",
            default_value=0,
            min_value=0,
            step=50,
            allow_float=False,
            text_format=",d",
            placeholder="例: 1,200",
        )
        manual_new = enhanced_number_input(
            "当月新規顧客数",
            key="manual_new_customers",
            default_value=0,
            min_value=0,
            step=10,
            allow_float=False,
            text_format=",d",
            placeholder="例: 180",
        )
        manual_repeat = enhanced_number_input(
            "当月リピート顧客数",
            key="manual_repeat_customers",
            default_value=0,
            min_value=0,
            step=10,
            allow_float=False,
            text_format=",d",
            placeholder="例: 320",
        )
        manual_cancel = enhanced_number_input(
            "当月解約件数",
            key="manual_cancel_customers",
            default_value=0,
            min_value=0,
            step=5,
            allow_float=False,
            text_format=",d",
            placeholder="例: 25",
        )
        manual_prev_active = enhanced_number_input(
            "前月契約数",
            key="manual_previous_customers",
            default_value=0,
            min_value=0,
            step=50,
            allow_float=False,
            text_format=",d",
            placeholder="例: 1,150",
        )
        manual_marketing = enhanced_number_input(
            "当月広告費",
            key="manual_marketing_cost",
            default_value=0.0,
            min_value=0.0,
            step=50_000.0,
            number_format="%.0f",
            text_format=",.0f",
            placeholder="例: 250,000",
        )
        manual_ltv = enhanced_number_input(
            "LTV試算値",
            key="manual_ltv_value",
            default_value=0.0,
            min_value=0.0,
            step=1_000.0,
            number_format="%.0f",
            text_format=",.0f",
            placeholder="例: 120,000",
        )

        st.markdown("#### バランスト・スコアカード指標")
        manual_inventory_days = enhanced_number_input(
            "在庫回転日数（目標: 45日以下）",
            key="manual_inventory_days",
            default_value=45.0,
            min_value=0.0,
            step=1.0,
            number_format="%.0f",
            text_format=",.0f",
            placeholder="例: 40",
            help_text="内部プロセス視点: 在庫を現金化するまでの日数を把握します。",
        )
        manual_stockout_pct = enhanced_number_input(
            "欠品率（%）",
            key="manual_stockout_pct",
            default_value=4.0,
            min_value=0.0,
            max_value=100.0,
            step=0.5,
            number_format="%.1f",
            text_format=",.1f",
            placeholder="例: 4",
            help_text="内部プロセス視点: 欠品による販売機会損失を監視します。",
        )
        manual_training_sessions = enhanced_number_input(
            "従業員研修実施数（月内）",
            key="manual_training_sessions",
            default_value=2,
            min_value=0,
            step=1,
            allow_float=False,
            text_format=",d",
            placeholder="例: 2",
            help_text="学習・成長視点: 店長や経理がスキルを磨いた回数です。",
        )
        manual_new_products = enhanced_number_input(
            "新商品リリース数（月内）",
            key="manual_new_products",
            default_value=1,
            min_value=0,
            step=1,
            allow_float=False,
            text_format=",d",
            placeholder="例: 1",
            help_text="学習・成長視点: 新しい価値提案の数を追跡します。",
        )

    automated_sales_data = st.session_state.get("api_sales_data", {})
    automated_reports = list(st.session_state.get("api_sales_validation", {}).values())

    try:
        with st.spinner(STATE_MESSAGES["loading"]["text"]):
            data_dict = load_data(
                use_sample_data,
                channel_files,
                cost_file,
                subscription_file,
                automated_sales=automated_sales_data,
                automated_reports=automated_reports,
            )
    except Exception as exc:
        logger.exception("Failed to load dashboard data")
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        error_entry = {
            "timestamp": timestamp,
            "message": str(exc),
            "traceback": "".join(
                traceback.format_exception(type(exc), exc, exc.__traceback__)
            ),
        }
        st.session_state.setdefault("admin_error_log", []).append(error_entry)
        st.session_state["has_loaded_data"] = False
        render_onboarding_wizard(
            onboarding_container,
            data_loaded=False,
            filters_ready=False,
            analysis_ready=False,
            sample_checked=use_sample_data,
            visible=show_onboarding,
        )
        display_state_message(
            "server_error",
            action=trigger_rerun,
            action_key="reload_after_error",
        )
        return

    sales_df = data_dict["sales"].copy()
    cost_df = data_dict["cost"].copy()
    subscription_df = data_dict["subscription"].copy()
    sales_validation: ValidationReport = data_dict.get("sales_validation", ValidationReport())

    if sales_df.empty:
        st.session_state["has_loaded_data"] = False
        render_sidebar_disabled_placeholder()
        render_empty_dashboard_placeholder()
        render_onboarding_wizard(
            onboarding_container,
            data_loaded=False,
            filters_ready=False,
            analysis_ready=False,
            sample_checked=use_sample_data,
            visible=show_onboarding,
        )

        def _enable_sample_data() -> None:
            st.session_state["pending_enable_sample_data"] = True
            trigger_rerun()

        def _navigate_to_upload() -> None:
            st.session_state["main_nav"] = "data"
            st.session_state["preferred_dashboard_tab"] = "データ管理"
            trigger_rerun()

        display_state_message(
            "data_unloaded",
            action=_enable_sample_data,
            secondary_action=_navigate_to_upload,
            secondary_action_key="data_unloaded_upload",
        )
        return

    merged_full = merge_sales_and_costs(sales_df, cost_df)
    sales_validation.extend(validate_channel_fees(merged_full))

    freq_lookup = {label: freq for label, freq in PERIOD_FREQ_OPTIONS}
    freq_labels = list(freq_lookup.keys())
    default_freq_label = next(
        (label for label, freq in PERIOD_FREQ_OPTIONS if freq == "M"),
        freq_labels[0],
    )

    unique_channels = sorted(sales_df["channel"].dropna().unique().tolist())
    unique_categories = sorted(sales_df["category"].dropna().unique().tolist())
    global_min_date = normalize_date_input(sales_df["order_date"].min()) or date.today()
    global_max_date = normalize_date_input(sales_df["order_date"].max()) or date.today()
    if global_min_date > global_max_date:
        global_min_date, global_max_date = global_max_date, global_min_date
    global_default_period = suggest_default_period(global_min_date, global_max_date)

    store_candidates = ["全社"]
    candidate_values: List[str] = []
    if "store" in sales_df.columns:
        candidate_values = [str(value) for value in sales_df["store"].dropna().unique()]
        store_candidates.extend(candidate_values)
    store_candidates.extend(option for option in DEFAULT_STORE_OPTIONS if option not in store_candidates)
    store_options = list(dict.fromkeys(store_candidates)) or ["全社"]
    if "全社" in store_options:
        preferred_store = "全社"
    elif candidate_values:
        preferred_store = candidate_values[0]
    else:
        preferred_store = store_options[0]
    default_store = preferred_store

    default_filters = {
        FILTER_STATE_KEYS["store"]: default_store,
        FILTER_STATE_KEYS["channels"]: unique_channels,
        FILTER_STATE_KEYS["categories"]: unique_categories,
        FILTER_STATE_KEYS["period"]: global_default_period,
        FILTER_STATE_KEYS["freq"]: default_freq_label,
    }

    store_state_key = FILTER_STATE_KEYS["store"]
    current_store_state = st.session_state.get(store_state_key)
    if current_store_state not in store_options:
        current_store_state = default_store
    set_state_and_widget(store_state_key, current_store_state)
    store_widget_key = widget_key_for(store_state_key)
    current_store = st.session_state[store_state_key]
    store_index = store_options.index(current_store) if current_store in store_options else 0

    if current_store and current_store != "全社" and "store" in sales_df.columns:
        store_sales_df = sales_df[sales_df["store"] == current_store].copy()
    else:
        store_sales_df = sales_df.copy()

    store_min_candidate = normalize_date_input(store_sales_df["order_date"].min()) if not store_sales_df.empty else None
    store_max_candidate = normalize_date_input(store_sales_df["order_date"].max()) if not store_sales_df.empty else None
    min_date = store_min_candidate or global_min_date
    max_date = store_max_candidate or global_max_date
    if min_date > max_date:
        min_date, max_date = max_date, min_date

    period_state_key = FILTER_STATE_KEYS["period"]
    stored_period = st.session_state.get(period_state_key, global_default_period)
    default_period = suggest_default_period(min_date, max_date)
    normalized_period_state = normalize_period_state_value(
        stored_period,
        min_date,
        max_date,
        default_period,
    )
    set_state_and_widget(period_state_key, normalized_period_state)
    period_widget_key = widget_key_for(period_state_key)

    channel_series = store_sales_df["channel"].dropna()
    if not channel_series.empty:
        channel_counts = channel_series.value_counts()
        top_channel_options = channel_counts.head(CATEGORY_SUGGESTION_LIMIT).index.tolist()
    else:
        top_channel_options = []
    channel_state_key = FILTER_STATE_KEYS["channels"]
    manual_channel_state_key = f"{channel_state_key}_manual_entries"
    manual_channel_entries: List[str] = []
    for value in st.session_state.get(manual_channel_state_key, []):
        normalized = str(value).strip()
        if not normalized or normalized in manual_channel_entries:
            continue
        manual_channel_entries.append(normalized)
    st.session_state[manual_channel_state_key] = manual_channel_entries
    previous_channel_selection = list(st.session_state.get(channel_state_key, []))
    available_channels = list(
        dict.fromkeys(top_channel_options + manual_channel_entries + previous_channel_selection)
    )
    preserved_channels = [
        ch for ch in previous_channel_selection if ch in available_channels
    ]
    if available_channels and not preserved_channels:
        preserved_channels = available_channels
    set_state_and_widget(channel_state_key, preserved_channels)
    channel_widget_key = widget_key_for(channel_state_key)
    channel_manual_input_key = widget_key_for(f"{channel_state_key}_manual_input")

    category_series = store_sales_df["category"].dropna()
    if not category_series.empty:
        category_counts = category_series.value_counts()
        top_category_options = category_counts.head(CATEGORY_SUGGESTION_LIMIT).index.tolist()
    else:
        top_category_options = []
    category_state_key = FILTER_STATE_KEYS["categories"]
    manual_category_state_key = f"{category_state_key}_manual_entries"
    manual_category_entries: List[str] = []
    for value in st.session_state.get(manual_category_state_key, []):
        normalized = str(value).strip()
        if not normalized or normalized in manual_category_entries:
            continue
        manual_category_entries.append(normalized)
    st.session_state[manual_category_state_key] = manual_category_entries
    previous_category_selection = list(st.session_state.get(category_state_key, []))
    available_categories = list(
        dict.fromkeys(top_category_options + manual_category_entries + previous_category_selection)
    )
    preserved_categories = [
        cat for cat in previous_category_selection if cat in available_categories
    ]
    if available_categories and not preserved_categories:
        preserved_categories = available_categories
    set_state_and_widget(category_state_key, preserved_categories)
    category_widget_key = widget_key_for(category_state_key)
    category_manual_input_key = widget_key_for(f"{category_state_key}_manual_input")

    freq_state_key = FILTER_STATE_KEYS["freq"]
    current_freq_label = st.session_state.get(freq_state_key, default_freq_label)
    if current_freq_label not in freq_lookup:
        current_freq_label = default_freq_label
    set_state_and_widget(freq_state_key, current_freq_label)
    freq_widget_key = widget_key_for(freq_state_key)
    freq_index = (
        freq_labels.index(st.session_state[freq_state_key])
        if st.session_state[freq_state_key] in freq_labels
        else 0
    )

    def _apply_filter_form(state_key: str) -> None:
        update_state_from_widget(state_key)
        trigger_rerun()

    st.sidebar.selectbox(
        "店舗選択",
        options=store_options,
        index=store_index,
        key=store_widget_key,
        help="最後に選択した店舗は次回アクセス時も自動で設定されます。",
        on_change=_apply_filter_form,
        args=(store_state_key,),
    )
    st.sidebar.date_input(
        "表示期間（開始日 / 終了日）",
        value=st.session_state[period_state_key],
        min_value=min_date,
        max_value=max_date,
        key=period_widget_key,
        help="ダッシュボードに表示する対象期間です。開始日と終了日を指定してください。",
        on_change=_apply_filter_form,
        args=(period_state_key,),
    )
    st.sidebar.multiselect(
        "表示するチャネル",
        options=available_channels,
        default=st.session_state[channel_state_key] if available_channels else [],
        key=channel_widget_key,
        help="チャネル選択は関連レポートでも共有されます。",
        on_change=_apply_filter_form,
        args=(channel_state_key,),
    )
    st.sidebar.text_input(
        "その他チャネルを追加",
        key=channel_manual_input_key,
        placeholder="チャネル名を入力",
        help="候補一覧にないチャネル名を入力して追加できます。",
        on_change=add_manual_filter_value,
        args=(channel_state_key, manual_channel_state_key, channel_manual_input_key),
    )
    st.sidebar.multiselect(
        "表示するカテゴリ",
        options=available_categories,
        default=st.session_state[category_state_key] if available_categories else [],
        key=category_widget_key,
        help="カテゴリ選択は粗利・在庫の分析タブにも共有されます。",
        on_change=_apply_filter_form,
        args=(category_state_key,),
    )
    st.sidebar.text_input(
        "その他カテゴリを追加",
        key=category_manual_input_key,
        placeholder="カテゴリ名を入力",
        help="候補一覧にないカテゴリ名を入力して追加できます。",
        on_change=add_manual_filter_value,
        args=(category_state_key, manual_category_state_key, category_manual_input_key),
    )
    st.sidebar.selectbox(
        "ダッシュボード表示粒度",
        options=freq_labels,
        index=freq_index,
        key=freq_widget_key,
        help="売上やKPIの集計粒度を選べます。月次・週次・四半期などの粒度に対応しています。",
        on_change=_apply_filter_form,
        args=(freq_state_key,),
    )
    st.sidebar.caption("選択内容は変更と同時にダッシュボードへ反映されます。")

    current_period = st.session_state[period_state_key]
    selected_granularity_label = st.session_state[freq_state_key]
    selected_freq = freq_lookup[selected_granularity_label]

    st.sidebar.markdown("---")
    if st.sidebar.button("設定をリセット", key="reset_filter_button"):
        reset_filters(default_filters)
    if st.sidebar.button("セッション状態を初期化", key="clear_session_button"):
        st.session_state.clear()
        trigger_rerun()

    selected_store = st.session_state[store_state_key]
    selected_channels = st.session_state[channel_state_key]
    selected_categories = st.session_state[category_state_key]
    date_range = current_period
    filters_ready = bool(selected_channels) and bool(selected_categories)

    filter_signature = build_filter_signature(
        selected_store,
        selected_channels,
        selected_categories,
        date_range,
        selected_granularity_label,
    )
    signature_key = FILTER_STATE_KEYS["signature"]
    if signature_key not in st.session_state:
        st.session_state[signature_key] = filter_signature
    elif st.session_state[signature_key] != filter_signature:
        display_state_message("success", action_key="filters_success_message")
        st.session_state[signature_key] = filter_signature

    store_filter: Optional[List[str]] = None
    if selected_store and selected_store not in ("全社", None):
        store_filter = [selected_store]
    if isinstance(date_range, (tuple, list)) and len(date_range) == 2:
        date_range_list = [date_range[0], date_range[1]]
    else:
        date_range_list = [date_range, date_range]

    filtered_sales = apply_filters(
        sales_df,
        selected_channels,
        date_range_list,
        selected_categories,
        stores=store_filter,
    )
    analysis_ready = not filtered_sales.empty

    st.session_state["has_loaded_data"] = True
    st.session_state["filters_ready"] = filters_ready
    st.session_state["analysis_ready"] = analysis_ready

    render_onboarding_wizard(
        onboarding_container,
        data_loaded=True,
        filters_ready=filters_ready,
        analysis_ready=analysis_ready,
        sample_checked=use_sample_data,
        visible=show_onboarding,
    )

    if filtered_sales.empty:
        display_state_message(
            "filter_no_result",
            action=lambda: reset_filters(default_filters),
            action_key="reset_after_empty",
        )
    merged_df = merge_sales_and_costs(filtered_sales, cost_df)
    segmented_sales_df = annotate_customer_segments(merged_df)
    monthly_summary = monthly_sales_summary(merged_df)
    period_summary = summarize_sales_by_period(merged_df, selected_freq)

    kpi_overrides = {}
    if manual_active > 0:
        kpi_overrides["active_customers"] = manual_active
    if manual_new > 0:
        kpi_overrides["new_customers"] = manual_new
    if manual_repeat > 0:
        kpi_overrides["repeat_customers"] = manual_repeat
    if manual_cancel > 0:
        kpi_overrides["cancelled_subscriptions"] = manual_cancel
    if manual_prev_active > 0:
        kpi_overrides["previous_active_customers"] = manual_prev_active
    if manual_marketing > 0:
        kpi_overrides["marketing_cost"] = manual_marketing
    if manual_ltv > 0:
        kpi_overrides["ltv"] = manual_ltv

    kpi_overrides["inventory_turnover_days"] = manual_inventory_days
    kpi_overrides["stockout_rate"] = manual_stockout_pct / 100 if manual_stockout_pct >= 0 else np.nan
    kpi_overrides["training_sessions"] = manual_training_sessions
    kpi_overrides["new_product_count"] = manual_new_products

    kpis = calculate_kpis(merged_df, subscription_df, overrides=kpi_overrides)
    kpi_history_df = build_kpi_history_df(merged_df, subscription_df, kpi_overrides)
    kpi_period_summary = aggregate_kpi_history(kpi_history_df, selected_freq)
    latest_kpi_row: Optional[pd.Series]
    if kpi_period_summary.empty:
        latest_kpi_row = None
    else:
        latest_kpi_row = kpi_period_summary.iloc[-1]

    base_pl = create_current_pl(merged_df, subscription_df, fixed_cost=fixed_cost)
    default_cash_plan = create_default_cashflow_plan(merged_df)
    default_cash_forecast = forecast_cashflow(default_cash_plan, starting_cash)

    alerts = build_alerts(monthly_summary, kpis, default_cash_forecast)

    channel_share_df = compute_channel_share(merged_df)
    category_share_df = compute_category_share(merged_df)

    latest_timestamp = None
    if not merged_df.empty and "order_date" in merged_df.columns:
        latest_timestamp = merged_df["order_date"].max()
    if latest_timestamp is not None and pd.notna(latest_timestamp):
        if isinstance(latest_timestamp, pd.Timestamp):
            latest_label = latest_timestamp.strftime("%Y-%m-%d")
        else:
            latest_label = str(latest_timestamp)
    else:
        latest_label = "-"

    range_label = "-"
    if isinstance(date_range, (tuple, list)) and len(date_range) == 2:
        start_value, end_value = date_range
        start_label = start_value.strftime("%Y-%m-%d") if hasattr(start_value, "strftime") else str(start_value)
        end_label = end_value.strftime("%Y-%m-%d") if hasattr(end_value, "strftime") else str(end_value)
        range_label = f"{start_label} 〜 {end_label}"

    total_records = int(len(merged_df)) if not merged_df.empty else 0
    alert_count = len(alerts) if alerts else 0

    search_query = render_search_bar()

    with st.container():
        st.markdown("<div class='surface-card main-nav-block'>", unsafe_allow_html=True)
        selected_nav_key, selected_nav_label = render_navigation()
        st.markdown("</div>", unsafe_allow_html=True)

    quick_action_feedback = st.session_state.pop("quick_action_feedback", None)

    with st.container():
        st.markdown("<div class='surface-card quick-actions-block'>", unsafe_allow_html=True)
        action_cols = st.columns([1, 1])
        action_cols[0].button(
            "最新データを再取得",
            key="quick_action_refetch_data",
            use_container_width=True,
            type="secondary",
            on_click=lambda: _queue_hotkey_action(HOTKEY_REFETCH_DATA),
        )
        action_cols[1].markdown(
            """
            <div class='quick-action-info'>
                <strong>自動更新</strong><br>
                フィルタやタブの切り替えは即座に再計算されます。
            </div>
            """,
            unsafe_allow_html=True,
        )

        st.caption("ショートカット: Shift+Uキーで最新データを再取得")

        hotkey_event = components.html(
            f"""
<script>
(function() {{
  const frameId = window.frameElement ? window.frameElement.id : null;
  const sendPayload = (payload) => {{
    if (!frameId) {{ return; }}
    window.parent.postMessage({{
      isStreamlitMessage: true,
      type: "streamlit:setComponentValue",
      id: frameId,
      value: payload,
    }}, "*");
    window.parent.postMessage({{
      type: "kuraiki-hotkey",
      payload: payload,
    }}, "*");
  }};
  const shouldIgnoreTarget = (target) => {{
    if (!target) {{ return false; }}
    const tag = target.tagName;
    if (tag === "INPUT" || tag === "TEXTAREA" || target.isContentEditable) {{
      return true;
    }}
    return target.closest && !!target.closest('[contenteditable="true"]');
  }};
  const resolveAction = (event) => {{
    if (!event.shiftKey) {{ return null; }}
    const key = event.key ? event.key.toUpperCase() : "";
    if (key === "U") {{ return "{HOTKEY_REFETCH_DATA}"; }}
    return null;
  }};
  document.addEventListener('keydown', (event) => {{
    const action = resolveAction(event);
    if (!action) {{ return; }}
    if (shouldIgnoreTarget(event.target)) {{ return; }}
    event.preventDefault();
    const payload = {{ action: action, timestamp: Date.now() }};
    sendPayload(payload);
  }}, {{ capture: true }});
}})();
</script>
""",
            height=0,
            width=0,
        )
        if hotkey_event and isinstance(hotkey_event, dict):
            st.session_state[HOTKEY_ACTION_STATE] = hotkey_event

        if quick_action_feedback:
            if quick_action_feedback.get("errors"):
                st.error(
                    "再取得でエラーが発生しました: "
                    + " / ".join(quick_action_feedback["errors"])
                )
            if quick_action_feedback.get("warnings"):
                st.warning(
                    "警告が検出されました: "
                    + " / ".join(quick_action_feedback["warnings"])
                )
            if quick_action_feedback.get("success"):
                st.success(
                    "最新データの再取得が完了しました: "
                    + "、".join(quick_action_feedback["success"])
                )
            if (
                quick_action_feedback.get("info")
                and not (
                    quick_action_feedback.get("success")
                    or quick_action_feedback.get("warnings")
                    or quick_action_feedback.get("errors")
                )
            ):
                st.info(quick_action_feedback["info"])
        st.markdown("</div>", unsafe_allow_html=True)

    def _handle_pending_hotkey_action() -> None:
        payload = st.session_state.get(HOTKEY_ACTION_STATE)
        if not isinstance(payload, dict):
            return
        action = payload.get("action")
        timestamp = payload.get("timestamp")
        if not action or timestamp is None:
            return
        last_processed = st.session_state.get("hotkey_last_processed")
        if timestamp == last_processed:
            return
        st.session_state["hotkey_last_processed"] = timestamp

        if action == HOTKEY_REFETCH_DATA:
            feedback = refresh_configured_automations(channel_files)
            st.session_state["quick_action_feedback"] = feedback
            trigger_rerun()

    _handle_pending_hotkey_action()

    render_breadcrumb(selected_nav_label)

    if search_query:
        render_global_search_results(search_query, merged_df)
        st.divider()

    if selected_nav_key == "dashboard":
        st.subheader("経営ダッシュボード")
        if kpi_period_summary.empty:
            st.info(
                "KPI情報が不足しています。KPIデータをアップロードするか、サイドバーで数値を入力してください。"
            )
        else:
            period_options = kpi_period_summary["period_label"].tolist()
            default_period_idx = len(period_options) - 1 if period_options else 0
            selected_dashboard_period = st.selectbox(
                f"{selected_granularity_label}の表示期間",
                options=period_options,
                index=default_period_idx,
                key="dashboard_period_select",
            )
            selected_kpi_row = kpi_period_summary[
                kpi_period_summary["period_label"] == selected_dashboard_period
            ].iloc[0]
            selected_period = selected_kpi_row["period"]
            period_row = period_summary[period_summary["period"] == selected_period]
            period_start = pd.to_datetime(selected_kpi_row["period_start"]).date()
            period_end = pd.to_datetime(selected_kpi_row["period_end"]).date()

            gross_rate_value = selected_kpi_row.get("gross_margin_rate")
            gross_target = KGI_TARGETS.get("gross_margin_rate")
            if (
                gross_rate_value is not None
                and not pd.isna(gross_rate_value)
                and gross_target is not None
                and gross_rate_value < gross_target
            ):
                display_state_message(
                    "warning_gross_margin",
                    action=lambda: jump_to_section("gross"),
                    action_label="粗利タブを開く",
                    action_key="warning_gross_margin_button",
                )

            render_kgi_cards(selected_kpi_row, period_row, default_cash_forecast, starting_cash)
            render_dashboard_meta(
                latest_label,
                range_label,
                total_records,
                alert_count,
                store_selection=selected_store,
                channel_selection=selected_channels,
                category_selection=selected_categories,
            )
            render_status_banner(alerts)
            st.caption(f"対象期間: {period_start} 〜 {period_end}")

            help_language = render_kpi_help_controls()
            kpi_metrics = render_first_level_kpi_strip(
                kpi_period_summary,
                selected_kpi_row,
                help_language=help_language,
            )
            bsc_quadrants = build_bsc_quadrants(selected_kpi_row)
            if bsc_quadrants:
                st.markdown("### バランスト・スコアカード")
                bsc_view_mode = persistent_segmented_control(
                    "bsc_view_mode",
                    ["カード", "チャート"],
                    default="カード",
                    help_text="カード表示と四象限チャートを切り替えます。",
                    label="BSC表示切替",
                    label_visibility="collapsed",
                )
                if bsc_view_mode == "カード":
                    render_bsc_cards_grid(bsc_quadrants)
                else:
                    render_bsc_quadrant_chart(bsc_quadrants)
                st.caption("ターゲットを達成した象限はハイライト表示されます。")

            render_active_kpi_details(
                kpi_period_summary,
                kpi_metrics,
                help_language=help_language,
            )

            preferred_tab = st.session_state.pop("preferred_dashboard_tab", None)
            if preferred_tab:
                st.info(f"「{preferred_tab}」タブで詳細設定を確認できます。")

            sales_tab, finance_tab, data_tab = st.tabs(
                ["売上・粗利", "在庫・資金", "KPI・データ管理"]
            )

            with sales_tab:
                st.markdown("#### 売上ダッシュボード")
                render_sales_tab(
                    merged_df,
                    period_summary,
                    channel_share_df,
                    category_share_df,
                    selected_granularity_label,
                )
                st.divider()
                st.markdown("#### 粗利分析")
                render_gross_tab(merged_df, period_summary, selected_granularity_label)

            with finance_tab:
                st.markdown("#### 在庫・オペレーション")
                render_inventory_tab(merged_df, kpi_period_summary, selected_kpi_row)
                st.divider()
                st.markdown("#### 資金繰りとキャッシュフロー")
                render_cash_tab(
                    default_cash_plan,
                    default_cash_forecast,
                    starting_cash,
                    monthly_summary,
                )

            with data_tab:
                st.markdown("#### KPI詳細指標")
                render_kpi_overview_tab(kpi_period_summary)
                st.divider()
                st.markdown("#### データ連携状況")
                render_data_status_section(
                    merged_df,
                    cost_df,
                    subscription_df,
                    use_sample_data=use_sample_data,
                    automated_sales_data=automated_sales_data,
                )
            st.divider()

    elif selected_nav_key == "sales":
        st.subheader("売上分析")
        if merged_df.empty:
            st.info("売上データがありません。")
        else:
            st.caption("グラフをクリックすると他の可視化も同じ条件で絞り込まれます。")
            sales_cross_filters = st.session_state.setdefault(
                "sales_cross_filters", {"channel": None, "category": None}
            )

            available_analysis_channels = sorted(merged_df["channel"].unique())
            available_analysis_categories = sorted(merged_df["category"].unique())
            if (
                sales_cross_filters.get("channel")
                and sales_cross_filters["channel"] not in available_analysis_channels
            ):
                sales_cross_filters["channel"] = None
            if (
                sales_cross_filters.get("category")
                and sales_cross_filters["category"] not in available_analysis_categories
            ):
                sales_cross_filters["category"] = None

            analysis_df = merged_df.copy()
            active_highlights: List[str] = []
            if sales_cross_filters.get("channel"):
                analysis_df = analysis_df[analysis_df["channel"] == sales_cross_filters["channel"]]
                active_highlights.append(f"チャネル: {sales_cross_filters['channel']}")
            if sales_cross_filters.get("category"):
                analysis_df = analysis_df[analysis_df["category"] == sales_cross_filters["category"]]
                active_highlights.append(f"カテゴリ: {sales_cross_filters['category']}")

            if active_highlights:
                info_col, clear_col = st.columns([5, 1])
                info_col.info("ハイライト適用中: " + " / ".join(active_highlights))
                if clear_col.button("ハイライトをクリア", key="clear_sales_highlight"):
                    st.session_state["sales_cross_filters"] = {"channel": None, "category": None}
                    analysis_df = merged_df.copy()
                    active_highlights = []

            channel_trend_full = merged_df.copy()
            channel_trend_full["period"] = channel_trend_full["order_date"].dt.to_period(selected_freq)
            channel_trend_full = (
                channel_trend_full.groupby(["period", "channel"])["sales_amount"].sum().reset_index()
            )
            channel_trend_full["period_start"] = channel_trend_full["period"].dt.to_timestamp()
            channel_trend_full["period_label"] = channel_trend_full["period"].apply(
                lambda p: format_period_label(p, selected_freq)
            )
            channel_trend_full.sort_values(["channel", "period_start"], inplace=True)

            channel_chart = px.line(
                channel_trend_full,
                x="period_start",
                y="sales_amount",
                color="channel",
                markers=True,
                labels={
                    "sales_amount": "売上高",
                    "period_start": f"{selected_granularity_label}開始日",
                },
                custom_data=["channel", "period_label"],
                color_discrete_sequence=get_active_chart_colorway(),
            )
            channel_chart = apply_chart_theme(channel_chart)
            channel_chart.update_layout(
                clickmode="event+select",
                legend=dict(title="", itemclick="toggleothers", itemdoubleclick="toggle"),
            )
            for trace in channel_chart.data:
                trace.update(
                    hovertemplate="期間=%{customdata[1]}<br>チャネル=%{customdata[0]}<br>売上高=%{y:,.0f}円<extra></extra>"
                )
                if sales_cross_filters.get("channel") and trace.name != sales_cross_filters["channel"]:
                    trace.update(opacity=0.25, line={"width": 1})
                else:
                    trace.update(line={"width": 3})
            channel_events = plotly_events(
                channel_chart,
                click_event=True,
                override_width="100%",
                override_height=420,
                key="channel_trend_events",
            )
            if channel_events:
                clicked_channel = channel_events[0]["customdata"][0]
                current = st.session_state["sales_cross_filters"].get("channel")
                if current == clicked_channel:
                    st.session_state["sales_cross_filters"]["channel"] = None
                else:
                    st.session_state["sales_cross_filters"]["channel"] = clicked_channel

            category_sales_full = merged_df.copy()
            category_sales_full["period"] = category_sales_full["order_date"].dt.to_period(selected_freq)
            category_sales_full = (
                category_sales_full.groupby(["period", "category"])["sales_amount"].sum().reset_index()
            )
            category_sales_full["period_start"] = category_sales_full["period"].dt.to_timestamp()
            category_sales_full["period_label"] = category_sales_full["period"].apply(
                lambda p: format_period_label(p, selected_freq)
            )
            category_sales_full.sort_values(["category", "period_start"], inplace=True)

            category_bar = px.bar(
                category_sales_full,
                x="period_start",
                y="sales_amount",
                color="category",
                labels={
                    "sales_amount": "売上高",
                    "period_start": f"{selected_granularity_label}開始日",
                },
                custom_data=["category", "period_label"],
                color_discrete_sequence=get_active_chart_colorway(),
            )
            category_bar = apply_chart_theme(category_bar)
            category_bar.update_layout(
                barmode="stack",
                clickmode="event+select",
                legend=dict(title="", itemclick="toggleothers", itemdoubleclick="toggle"),
            )
            for trace in category_bar.data:
                trace.update(
                    hovertemplate="期間=%{customdata[1]}<br>カテゴリ=%{customdata[0]}<br>売上高=%{y:,.0f}円<extra></extra>"
                )
                if sales_cross_filters.get("category") and trace.name != sales_cross_filters["category"]:
                    trace.update(opacity=0.35)
                else:
                    trace.update(opacity=0.9)
            category_events = plotly_events(
                category_bar,
                click_event=True,
                override_width="100%",
                override_height=420,
                key="category_sales_events",
            )
            if category_events:
                clicked_category = category_events[0]["customdata"][0]
                current_category = st.session_state["sales_cross_filters"].get("category")
                if current_category == clicked_category:
                    st.session_state["sales_cross_filters"]["category"] = None
                else:
                    st.session_state["sales_cross_filters"]["category"] = clicked_category

            analysis_summary = summarize_sales_by_period(analysis_df, selected_freq)
            if analysis_df.empty:
                st.warning("選択された条件に該当するデータがありません。")
            elif analysis_summary.empty:
                st.info("指定した粒度で集計できる期間データがありません。")
            else:
                yoy_table = analysis_summary.tail(12)[
                    ["period_label", "sales_amount", "sales_yoy", "sales_mom"]
                ]
                yoy_table = yoy_table.rename(
                    columns={
                        "period_label": "期間",
                        "sales_amount": "売上高",
                        "sales_yoy": "前年同期比",
                        "sales_mom": "前期比",
                    }
                )
                st.dataframe(yoy_table)

            st.markdown("### 店舗別売上・利益比較")
            render_store_comparison_chart(
                analysis_df,
                fixed_cost,
                selected_store=st.session_state.get(FILTER_STATE_KEYS["store"]),
            )

            st.markdown("### ABC分析（売上上位30商品）")
            render_abc_analysis(analysis_df)

    elif selected_nav_key == "inventory":
        st.subheader("在庫分析")
        if merged_df.empty and kpi_period_summary.empty:
            st.info("在庫分析に必要なデータが読み込まれていません。サンプルデータを利用するかCSVをアップロードしてください。")
        else:
            selected_inventory_row: Optional[pd.Series] = latest_kpi_row
            if not kpi_period_summary.empty:
                period_options = kpi_period_summary["period_label"].tolist()
                default_idx = len(period_options) - 1 if period_options else 0
                selected_period_label = st.selectbox(
                    "表示する期間",
                    options=period_options,
                    index=default_idx,
                    key="inventory_period_select",
                )
                try:
                    selected_inventory_row = kpi_period_summary[
                        kpi_period_summary["period_label"] == selected_period_label
                    ].iloc[0]
                except IndexError:
                    selected_inventory_row = latest_kpi_row
            render_inventory_tab(merged_df, kpi_period_summary, selected_inventory_row)

    elif selected_nav_key == "gross":
        st.subheader("利益分析")
        if merged_df.empty:
            st.info("データがありません。")
        else:
            product_profit = (
                merged_df.groupby(["product_code", "product_name", "category"], as_index=False)[
                    [
                        "sales_amount",
                        "estimated_cost",
                        "net_gross_profit",
                        "quantity",
                        "channel_fee_amount",
                    ]
                ]
                .sum()
            )
            product_profit["gross_margin_rate"] = product_profit["net_gross_profit"] / product_profit["sales_amount"]
            product_profit["average_unit_price"] = np.where(
                product_profit["quantity"] > 0,
                product_profit["sales_amount"] / product_profit["quantity"],
                np.nan,
            )
            product_profit["ad_ratio"] = np.where(
                product_profit["sales_amount"] != 0,
                product_profit["channel_fee_amount"] / product_profit["sales_amount"],
                np.nan,
            )
            product_profit.sort_values("net_gross_profit", ascending=False, inplace=True)
            display_columns = {
                "product_code": "商品コード",
                "product_name": "商品名",
                "category": "カテゴリ",
                "sales_amount": "売上高",
                "net_gross_profit": "粗利",
                "gross_margin_rate": "粗利率",
                "average_unit_price": "平均単価",
                "quantity": "販売個数",
                "ad_ratio": "広告費比率",
            }
            st.dataframe(
                product_profit[list(display_columns.keys())]
                .rename(columns=display_columns)
                .style.format({
                    "売上高": "{:,.0f}",
                    "粗利": "{:,.0f}",
                    "粗利率": "{:.2%}",
                    "平均単価": "{:,.0f}",
                    "販売個数": "{:,.0f}",
                    "広告費比率": "{:.2%}",
                }),
                use_container_width=True,
            )

            channel_profit = (
                merged_df.groupby("channel")["net_gross_profit"].sum().reset_index()
            )
            channel_profit_chart = px.bar(
                channel_profit,
                x="channel",
                y="net_gross_profit",
                labels={"channel": "チャネル", "net_gross_profit": "粗利"},
                title="チャネル別粗利比較",
                color_discrete_sequence=[GROSS_SERIES_COLOR],
            )
            channel_profit_chart = apply_chart_theme(channel_profit_chart)
            channel_profit_chart.update_layout(
                legend=dict(title=""),
                xaxis_title="チャネル",
                yaxis_title="粗利",
            )
            st.plotly_chart(channel_profit_chart, use_container_width=True)

            top_products = product_profit.head(10).copy()
            st.subheader("高利益商材トップ10")
            selected_product_code = st.session_state.setdefault("profit_focus_product", None)
            if selected_product_code and selected_product_code not in top_products["product_code"].values:
                st.session_state["profit_focus_product"] = None
                selected_product_code = None

            top_products_sorted = top_products.sort_values("net_gross_profit")
            top_products_chart = px.bar(
                top_products_sorted,
                x="net_gross_profit",
                y="product_name",
                orientation="h",
                labels={"net_gross_profit": "粗利", "product_name": "商品名"},
                custom_data=["product_code", "product_name"],
                color_discrete_sequence=[GROSS_SERIES_COLOR],
            )
            top_products_chart = apply_chart_theme(top_products_chart)
            highlight_code = st.session_state.get("profit_focus_product")
            bar_colors = [
                SUCCESS_COLOR if code == highlight_code else GROSS_SERIES_COLOR
                for code in top_products_sorted["product_code"]
            ]
            top_products_chart.update_traces(
                marker_color=bar_colors,
                hovertemplate="%{customdata[1]}<br>粗利=%{x:,.0f}円<extra></extra>",
            )
            top_products_chart.update_layout(
                height=420,
                xaxis_title="粗利",
                yaxis_title="商品名",
                clickmode="event+select",
            )
            events_top_products = plotly_events(
                top_products_chart,
                click_event=True,
                override_width="100%",
                override_height=420,
                key="top_products_events",
            )
            if events_top_products:
                clicked_code = events_top_products[0]["customdata"][0]
                current_code = st.session_state.get("profit_focus_product")
                if current_code == clicked_code:
                    st.session_state["profit_focus_product"] = None
                else:
                    st.session_state["profit_focus_product"] = clicked_code

            focus_code = st.session_state.get("profit_focus_product")
            if focus_code is None and not product_profit.empty:
                focus_code = product_profit.iloc[0]["product_code"]
                st.session_state["profit_focus_product"] = focus_code

            if focus_code and focus_code in product_profit["product_code"].values:
                focus_row = product_profit[product_profit["product_code"] == focus_code].iloc[0]
                st.markdown(
                    f"### 選択した商品の詳細: {focus_row['product_name']} ({focus_code})"
                )
                detail_cols = st.columns(5)
                detail_cols[0].metric("売上高", f"{focus_row['sales_amount']:,.0f} 円")
                detail_cols[1].metric("粗利", f"{focus_row['net_gross_profit']:,.0f} 円")
                detail_cols[2].metric(
                    "平均単価",
                    f"{focus_row['average_unit_price']:,.0f} 円"
                    if pd.notna(focus_row["average_unit_price"])
                    else "-",
                )
                detail_cols[3].metric(
                    "販売個数",
                    f"{focus_row['quantity']:,.0f} 個"
                    if pd.notna(focus_row["quantity"])
                    else "-",
                )
                detail_cols[4].metric(
                    "広告費比率",
                    f"{focus_row['ad_ratio'] * 100:.2f}%"
                    if pd.notna(focus_row["ad_ratio"])
                    else "-",
                )

                product_detail = merged_df[merged_df["product_code"] == focus_code].copy()
                channel_breakdown = (
                    product_detail.groupby("channel")[
                        ["sales_amount", "net_gross_profit", "quantity", "channel_fee_amount"]
                    ]
                    .sum()
                    .reset_index()
                )
                channel_breakdown["広告費比率"] = np.where(
                    channel_breakdown["sales_amount"] != 0,
                    channel_breakdown["channel_fee_amount"] / channel_breakdown["sales_amount"],
                    np.nan,
                )
                if not channel_breakdown.empty:
                    breakdown_chart = px.bar(
                        channel_breakdown,
                        x="channel",
                        y="net_gross_profit",
                        labels={"channel": "チャネル", "net_gross_profit": "粗利"},
                        title="選択商品のチャネル別粗利",
                        color_discrete_sequence=[GROSS_SERIES_COLOR],
                    )
                    breakdown_chart = apply_chart_theme(breakdown_chart)
                    breakdown_chart.update_layout(
                        legend=dict(title=""),
                        xaxis_title="チャネル",
                        yaxis_title="粗利",
                    )
                    st.plotly_chart(breakdown_chart, use_container_width=True)
                    st.dataframe(
                        channel_breakdown.rename(
                            columns={
                                "channel": "チャネル",
                                "sales_amount": "売上高",
                                "net_gross_profit": "粗利",
                                "quantity": "販売個数",
                            }
                        ).style.format(
                            {
                                "売上高": "{:,.0f}",
                                "粗利": "{:,.0f}",
                                "販売個数": "{:,.0f}",
                                "広告費比率": "{:.2%}",
                            }
                        ),
                        use_container_width=True,
                        hide_index=True,
                    )

                product_trend = product_detail.copy()
                product_trend["period"] = product_trend["order_date"].dt.to_period(selected_freq)
                product_trend_summary = (
                    product_trend.groupby("period")[
                        ["sales_amount", "net_gross_profit", "quantity"]
                    ]
                    .sum()
                    .reset_index()
                )
                if not product_trend_summary.empty:
                    product_trend_summary["period_start"] = product_trend_summary["period"].dt.to_timestamp()
                    product_trend_summary["period_label"] = product_trend_summary["period"].apply(
                        lambda p: format_period_label(p, selected_freq)
                    )
                    profit_trend_chart = px.line(
                        product_trend_summary,
                        x="period_start",
                        y="net_gross_profit",
                        markers=True,
                        labels={
                            "period_start": f"{selected_granularity_label}開始日",
                            "net_gross_profit": "粗利",
                        },
                        hover_data={"period_label": True},
                        color_discrete_sequence=[GROSS_SERIES_COLOR],
                    )
                    profit_trend_chart = apply_chart_theme(profit_trend_chart)
                    profit_trend_chart.update_layout(title="選択商品の粗利推移")
                    st.plotly_chart(profit_trend_chart, use_container_width=True)
                    st.dataframe(
                        product_trend_summary.rename(
                            columns={
                                "period_label": "期間",
                                "sales_amount": "売上高",
                                "net_gross_profit": "粗利",
                                "quantity": "販売個数",
                            }
                        ).style.format(
                            {
                                "売上高": "{:,.0f}",
                                "粗利": "{:,.0f}",
                                "販売個数": "{:,.0f}",
                            }
                        ),
                        use_container_width=True,
                    )
            else:
                st.info("表示する高利益商材がありません。")

    elif selected_nav_key == "cash":
        st.subheader("財務モニタリング")
        plan_state = st.session_state.get("plan_wizard")
        expense_table_state = None
        if isinstance(plan_state, dict):
            expense_table_state = plan_state.get("expense_table")
        st.markdown("売上計画や広告費を調整してPL・キャッシュフローをシミュレートします。")

        col1, col2, col3, col4 = st.columns(4)
        sales_growth = col1.slider("売上成長率", min_value=-0.5, max_value=0.5, value=0.05, step=0.01)
        cost_adj = col2.slider(
            "原価率変動",
            min_value=-0.1,
            max_value=0.1,
            value=0.0,
            step=0.01,
            help="粗利がマイナスにならないよう、シミュレーションでは原価率を最大99%に制限します。",
        )
        sga_change = col3.slider("販管費変動率", min_value=-0.3, max_value=0.3, value=0.0, step=0.01)
        with col4:
            extra_ad = enhanced_number_input(
                "追加広告費",
                key="extra_ad_cost",
                default_value=0.0,
                min_value=0.0,
                step=50_000.0,
                number_format="%.0f",
                text_format=",.0f",
                placeholder="例: 300,000",
            )

        pl_result = simulate_pl(
            base_pl,
            sales_growth_rate=sales_growth,
            cost_rate_adjustment=cost_adj,
            sga_change_rate=sga_change,
            additional_ad_cost=extra_ad,
        )
        st.dataframe(pl_result.style.format({"現状": "{:,.0f}", "シナリオ": "{:,.0f}", "増減": "{:,.0f}"}))

        st.metric(
            "シナリオ営業利益",
            f"{pl_result.loc[pl_result['項目'] == '営業利益', 'シナリオ'].iloc[0]:,.0f} 円",
            delta=f"{pl_result.loc[pl_result['項目'] == '営業利益', '増減'].iloc[0]:,.0f} 円",
        )

        render_profit_meter(pl_result, base_pl)

        plan_edit = create_default_cashflow_plan(merged_df).copy()
        plan_edit["month"] = plan_edit["month"].astype(str)
        with st.expander("キャッシュフロープランを編集"):
            edited_plan = st.data_editor(
                plan_edit,
                num_rows="dynamic",
                use_container_width=True,
            )
        if isinstance(edited_plan, pd.DataFrame):
            plan_to_use = edited_plan.copy()
        else:
            plan_to_use = pd.DataFrame(edited_plan)
        if not plan_to_use.empty:
            plan_to_use["month"] = plan_to_use["month"].apply(lambda x: pd.Period(x, freq="M"))
        cash_forecast = forecast_cashflow(plan_to_use, starting_cash)
        if not cash_forecast.empty:
            cash_chart = px.line(
                cash_forecast.assign(month=cash_forecast["month"].astype(str)),
                x="month",
                y="cash_balance",
                markers=True,
                title="資金残高予測",
                color_discrete_sequence=[CASH_SERIES_COLOR],
            )
            cash_chart = apply_chart_theme(cash_chart)
            cash_chart.update_layout(yaxis_title="円", xaxis_title="月")
            st.plotly_chart(cash_chart, use_container_width=True)
            st.dataframe(cash_forecast)
        else:
            st.info("キャッシュフロープランが未設定です。")

        st.markdown("<div class='chart-section'>", unsafe_allow_html=True)
        st.markdown(
            "<div class='chart-section__header'><div class='chart-section__title'>固定費内訳</div></div>",
            unsafe_allow_html=True,
        )
        render_fixed_cost_breakdown(expense_table_state, fixed_cost)
        st.markdown("</div>", unsafe_allow_html=True)

    elif selected_nav_key == "kpi":
        st.subheader("KPIモニタリング")
        if kpi_history_df.empty:
            st.info("KPI履歴がありません。")
        else:
            kpi_history_display = kpi_history_df.sort_values("month").copy()
            kpi_history_display["month_str"] = kpi_history_display["month"].astype(str)
            kpi_tab_entries = [
                ("ltv", "LTV", px.line, {"color_discrete_sequence": [ACCENT_COLOR]}),
                ("cac", "CAC", px.line, {"color_discrete_sequence": [WARNING_COLOR]}),
                ("repeat_rate", "リピート率", px.bar, {"color_discrete_sequence": [ACCENT_COLOR]}),
                ("churn_rate", "チャーン率", px.bar, {"color_discrete_sequence": [ERROR_COLOR]}),
                ("roas", "ROAS", px.line, {"color_discrete_sequence": [SALES_SERIES_COLOR]}),
            ]
            value_to_label = {value: label for value, label, *_ in kpi_tab_entries}
            default_kpi_tab = st.session_state.get("kpi_chart_tab", kpi_tab_entries[0][0])
            selected_kpi_chart = persistent_segmented_control(
                "kpi_chart_tab",
                [value for value, *_ in kpi_tab_entries],
                default=default_kpi_tab,
                help_text="KPIタブの選択を保持し、再訪時に同じ指標から確認できます。",
                format_func=lambda value: value_to_label[value],
                label="KPI詳細切替",
                label_visibility="visible",
            )

            chart_factory = next(
                (
                    (chart_fn, extra_kwargs)
                    for value, _, chart_fn, extra_kwargs in kpi_tab_entries
                    if value == selected_kpi_chart
                ),
                (px.line, {}),
            )
            chart_fn, chart_kwargs = chart_factory
            if chart_fn is px.line:
                fig = chart_fn(
                    kpi_history_display,
                    x="month_str",
                    y=selected_kpi_chart,
                    title=f"{value_to_label[selected_kpi_chart]}推移",
                    markers=True,
                    **chart_kwargs,
                )
            else:
                fig = chart_fn(
                    kpi_history_display,
                    x="month_str",
                    y=selected_kpi_chart,
                    title=f"{value_to_label[selected_kpi_chart]}推移",
                    **chart_kwargs,
                )
            fig = apply_chart_theme(fig)
            st.plotly_chart(fig, use_container_width=True)

            st.dataframe(
                kpi_history_display[
                    [
                        "month_str",
                        "sales",
                        "gross_profit",
                        "ltv",
                        "arpu",
                        "repeat_rate",
                        "churn_rate",
                        "roas",
                        "cac",
                    ]
                ].rename(columns={"month_str": "month"})
            )

            st.markdown("### KPIセグメント分析")
            segment_months = (
                segmented_sales_df["order_month"].dropna().sort_values().unique()
                if not segmented_sales_df.empty and "order_month" in segmented_sales_df.columns
                else []
            )
            period_options = ["全期間"]
            period_map: Dict[str, Optional[pd.Period]] = {"全期間": None}
            for period_value in segment_months:
                label = str(period_value)
                period_options.append(label)
                period_map[label] = period_value
            default_period_index = len(period_options) - 1 if len(period_options) > 1 else 0
            period_select_key = "kpi_breakdown_period"
            if (
                period_select_key not in st.session_state
                or st.session_state[period_select_key] not in period_options
            ):
                st.session_state[period_select_key] = period_options[default_period_index]
            selected_period_label = st.selectbox(
                "分析対象期間",
                options=period_options,
                index=period_options.index(st.session_state[period_select_key]),
                key=period_select_key,
                help="チャネル別・カテゴリ別のKPI集計に用いる期間を選択します。",
            )
            selected_period_value = period_map.get(selected_period_label)
            if selected_period_value is None:
                segmented_target_df = segmented_sales_df.copy()
            else:
                segmented_target_df = segmented_sales_df[
                    segmented_sales_df["order_month"] == selected_period_value
                ]

            if segmented_target_df.empty:
                st.info("選択された期間に該当するデータがありません。")
            else:
                breakdown_configs = [
                    ("チャネル別", "channel", "チャネル"),
                    ("カテゴリ別", "category", "商品カテゴリ"),
                    ("顧客区分別", "customer_segment", "顧客区分"),
                ]
                breakdown_tables: List[Tuple[str, str, str, pd.DataFrame]] = []
                for title, column, label in breakdown_configs:
                    df_breakdown = compute_kpi_breakdown(
                        segmented_target_df, column, kpi_totals=kpis
                    )
                    breakdown_tables.append((title, column, label, df_breakdown))

                if "campaign" in segmented_target_df.columns:
                    campaign_breakdown = compute_kpi_breakdown(
                        segmented_target_df, "campaign", kpi_totals=kpis
                    )
                    breakdown_tables.append(
                        ("キャンペーン別", "campaign", "キャンペーン", campaign_breakdown)
                    )

                st.caption("広告費や解約率は最新KPI値をシェアに応じて按分した推計値です。")
                breakdown_titles = [title for title, *_ in breakdown_tables]
                selected_breakdown_title = persistent_segmented_control(
                    "kpi_breakdown_tab",
                    breakdown_titles,
                    default=st.session_state.get("kpi_breakdown_tab", breakdown_titles[0]),
                    help_text="前回表示していた切り口を記憶します。",
                )

                for title, column, label, df_breakdown in breakdown_tables:
                    if title != selected_breakdown_title:
                        continue
                    if df_breakdown is None or df_breakdown.empty:
                        st.info(f"{label}別のKPIを算出するためのデータが不足しています。")
                        break

                    chart_data = df_breakdown.nlargest(10, "sales_amount")
                    bar_chart = px.bar(
                        chart_data,
                        x=column,
                        y="sales_amount",
                        labels={column: label, "sales_amount": "売上高"},
                        title=f"{label}別売上高 (上位{min(len(chart_data), 10)}件)",
                        color_discrete_sequence=get_active_chart_colorway(),
                    )
                    bar_chart = apply_chart_theme(bar_chart)
                    bar_chart.update_layout(yaxis_title="円", xaxis_title=label)
                    st.plotly_chart(bar_chart, use_container_width=True)

                    display_df = df_breakdown.rename(
                        columns={
                            column: label,
                            "sales_amount": "売上高",
                            "gross_profit": "粗利",
                            "gross_margin_rate": "粗利率",
                            "sales_share": "売上構成比",
                            "active_customers": "顧客数",
                            "new_customers": "新規顧客数",
                            "repeat_customers": "リピート顧客数",
                            "reactivated_customers": "休眠復活顧客数",
                            "repeat_rate": "リピート率",
                            "churn_rate": "推定解約率",
                            "arpu": "ARPU",
                            "ltv": "推定LTV",
                            "cac": "CAC",
                            "roas": "ROAS",
                            "marketing_cost": "広告費配分",
                            "profit_contribution": "粗利貢献額",
                            "profit_per_customer": "顧客あたり利益",
                            "avg_order_value": "平均受注単価",
                            "orders": "注文件数",
                        }
                    )
                    ordered_columns = [
                            label,
                            "売上高",
                            "粗利",
                            "粗利率",
                            "売上構成比",
                            "顧客数",
                            "新規顧客数",
                            "リピート顧客数",
                            "休眠復活顧客数",
                            "リピート率",
                            "推定解約率",
                            "ARPU",
                            "推定LTV",
                            "CAC",
                            "ROAS",
                            "広告費配分",
                            "粗利貢献額",
                            "顧客あたり利益",
                            "平均受注単価",
                            "注文件数",
                        ]
                    existing_columns = [col for col in ordered_columns if col in display_df.columns]
                    formatters = {
                            "売上高": "{:,.0f}",
                            "粗利": "{:,.0f}",
                            "粗利率": "{:.1%}",
                            "売上構成比": "{:.1%}",
                            "顧客数": "{:,.0f}",
                            "新規顧客数": "{:,.0f}",
                            "リピート顧客数": "{:,.0f}",
                            "休眠復活顧客数": "{:,.0f}",
                            "リピート率": "{:.1%}",
                            "推定解約率": "{:.1%}",
                            "ARPU": "{:,.0f}",
                            "推定LTV": "{:,.0f}",
                            "CAC": "{:,.0f}",
                            "ROAS": "{:,.2f}倍",
                            "広告費配分": "{:,.0f}",
                            "粗利貢献額": "{:,.0f}",
                            "顧客あたり利益": "{:,.0f}",
                            "平均受注単価": "{:,.0f}",
                            "注文件数": "{:,.0f}",
                        }
                    st.dataframe(
                            display_df[existing_columns].style.format({k: v for k, v in formatters.items() if k in existing_columns}),
                            use_container_width=True,
                            hide_index=True,
                        )

            profit_column = (
                "net_gross_profit"
                if "net_gross_profit" in segmented_target_df.columns
                else "gross_profit"
                if "gross_profit" in segmented_target_df.columns
                else None
            )
            repeat_scope_df = (
                segmented_target_df[
                    segmented_target_df.get("customer_segment", "既存").ne("新規")
                ]
                if not segmented_target_df.empty
                else pd.DataFrame()
            )
            repeat_customer_count = (
                repeat_scope_df["customer_id"].nunique()
                if not repeat_scope_df.empty and "customer_id" in repeat_scope_df.columns
                else 0
            )
            avg_repeat_sales = (
                repeat_scope_df["sales_amount"].sum() / repeat_customer_count
                if repeat_customer_count
                else float("nan")
            )
            avg_repeat_profit = (
                repeat_scope_df[profit_column].sum() / repeat_customer_count
                if profit_column and repeat_customer_count
                else float("nan")
            )

            st.subheader("施策効果の簡易比較")
            with st.form("ab_test"):
                before_rate = enhanced_number_input(
                    "施策前リピート率(%)",
                    key="abtest_before_rate",
                    default_value=60.0,
                    min_value=0.0,
                    max_value=100.0,
                    step=1.0,
                    number_format="%.1f",
                    text_format=",.1f",
                    placeholder="例: 60",
                )
                after_rate = enhanced_number_input(
                    "施策後リピート率(%)",
                    key="abtest_after_rate",
                    default_value=68.0,
                    min_value=0.0,
                    max_value=100.0,
                    step=1.0,
                    number_format="%.1f",
                    text_format=",.1f",
                    placeholder="例: 68",
                )
                before_count = enhanced_number_input(
                    "施策前顧客数",
                    key="abtest_before_count",
                    default_value=100,
                    min_value=1,
                    step=1,
                    allow_float=False,
                    text_format=",d",
                    placeholder="例: 100",
                )
                after_count = enhanced_number_input(
                    "施策後顧客数",
                    key="abtest_after_count",
                    default_value=100,
                    min_value=1,
                    step=1,
                    allow_float=False,
                    text_format=",d",
                    placeholder="例: 120",
                )
                submitted = st.form_submit_button("改善効果を計算")
                if submitted:
                    improvement = after_rate - before_rate
                    st.write(f"リピート率改善幅: {improvement:.1f}ポイント")
                    lift = (after_rate / before_rate - 1) if before_rate else np.nan
                    st.write(f"相対改善率: {lift:.2%}" if before_rate else "施策前のリピート率が0のため計算できません。")

                    before_repeat_customers = before_count * (before_rate / 100.0)
                    after_repeat_customers = after_count * (after_rate / 100.0)
                    customer_delta = after_repeat_customers - before_repeat_customers

                    revenue_uplift = (
                        customer_delta * avg_repeat_sales
                        if np.isfinite(avg_repeat_sales)
                        else float("nan")
                    )
                    profit_uplift = (
                        customer_delta * avg_repeat_profit
                        if np.isfinite(avg_repeat_profit)
                        else float("nan")
                    )
                    uplift_cols = st.columns(2)
                    uplift_cols[0].metric(
                        "想定売上増加額",
                        f"{revenue_uplift:,.0f} 円" if np.isfinite(revenue_uplift) else "算出不可",
                    )
                    uplift_cols[1].metric(
                        "想定粗利増加額",
                        f"{profit_uplift:,.0f} 円" if np.isfinite(profit_uplift) else "算出不可",
                    )

            if np.isfinite(avg_repeat_sales):
                profit_note = (
                    f"、平均リピート粗利 {avg_repeat_profit:,.0f} 円"
                    if np.isfinite(avg_repeat_profit)
                    else ""
                )
                st.caption(
                    f"リピート顧客1人あたりの平均売上 {avg_repeat_sales:,.0f} 円{profit_note} を基準に試算しています。"
                )
            else:
                st.caption("リピート顧客の平均売上を算出できなかったため、金額の試算は参考値です。")

    elif selected_nav_key == "scenario":
        st.subheader("シナリオ分析")
        render_scenario_analysis_section(merged_df, subscription_df)

    elif selected_nav_key == "data":
        st.subheader("データアップロード/管理")
        st.markdown(
            """
            - サイドバーから各チャネルのExcel/CSVファイルをアップロードしてください。
            - データはローカルセッション内でのみ保持され、アプリ終了時に消去されます。
            - 列名が異なる場合でも代表的な項目は自動マッピングされます。
            """
        )

        render_data_upload_guide(merged_full, cost_df, subscription_df)
        render_business_plan_wizard(merged_full)
        st.markdown("---")

        if sales_validation:
            st.markdown("### 読み込みバリデーション結果")
            for idx, message in enumerate(sales_validation.messages):
                display_text = message.message
                if message.count is not None:
                    display_text += f" (対象: {message.count:,}件)"
                if message.level == "error":
                    st.error(display_text)
                else:
                    st.warning(display_text)
                if message.sample is not None and not message.sample.empty:
                    with st.expander(f"該当レコードの例 ({idx + 1})"):
                        st.dataframe(message.sample)
            if not sales_validation.duplicate_rows.empty:
                st.warning("重複している可能性があるレコード一覧 (先頭200件)")
                st.dataframe(sales_validation.duplicate_rows.head(200))
        else:
            st.success("データ読み込み時に重大な問題は検出されませんでした。")

        if st.session_state.get("admin_mode_toggle"):
            admin_logs = st.session_state.get("admin_error_log", [])
            with st.expander("管理者向けエラーログ", expanded=False):
                if admin_logs:
                    for log_entry in reversed(admin_logs):
                        timestamp = log_entry.get("timestamp", "")
                        message = log_entry.get("message", "")
                        trace = log_entry.get("traceback", "")
                        st.markdown(f"**{timestamp}** - {html.escape(message)}", unsafe_allow_html=False)
                        st.code(trace, language="text")
                else:
                    st.info("現在表示できるエラーログはありません。")

        if automated_sales_data:
            status_rows = []
            for channel, df in automated_sales_data.items():
                last_fetch = st.session_state["api_last_fetched"].get(channel)
                report: Optional[ValidationReport] = st.session_state["api_sales_validation"].get(channel)
                if last_fetch:
                    status = "エラー" if report and report.has_errors() else "警告あり" if report and report.has_warnings() else "正常"
                    status_rows.append(
                        {
                            "チャネル": channel,
                            "最終取得": last_fetch.strftime("%Y-%m-%d %H:%M"),
                            "取得件数": len(df) if isinstance(df, pd.DataFrame) else 0,
                            "ステータス": status,
                        }
                    )
            if status_rows:
                st.markdown("### API連携ステータス")
                st.dataframe(pd.DataFrame(status_rows))

        st.write("現在のデータ件数")
        summary_cols = st.columns(3)
        summary_cols[0].metric("売上明細件数", len(merged_full))
        summary_cols[1].metric("取り扱い商品数", merged_full["product_code"].nunique())
        summary_cols[2].metric("期間", f"{min_date} 〜 {max_date}")

        with st.expander("原価率データのプレビュー"):
            if cost_df.empty:
                st.info("原価率データが未設定です。")
            else:
                st.dataframe(cost_df)

        with st.expander("売上データのプレビュー"):
            st.dataframe(merged_full.head(100))

        st.markdown("---")
        st.markdown("アプリの使い方や改善要望があれば開発チームまでご連絡ください。")

    st.markdown(
        "<div class='back-to-top'><a href='#page-top' aria-label='{aria}'>{icon} {label}</a></div>".format(
            aria=html.escape(
                translate("back_to_top_aria", default="ページの先頭へ戻る")
            ),
            icon="⬆",
            label=html.escape(
                translate("back_to_top_label", default="トップへ戻る")
            ),
        ),
        unsafe_allow_html=True,
    )

    apply_accessibility_enhancements()


if __name__ == "__main__":
    main()
